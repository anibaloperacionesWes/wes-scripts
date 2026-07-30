# -*- coding: utf-8 -*-
"""
Genera una COPIA del deck "Informe WES __ Parque Arauco 7 Malls" con datos reales
del período 01/05/2026 – 07/07/2026, regenerando cada gráfico incrustado y
conservando la estructura/orden de slides del original.

Decisiones aplicadas (confirmadas con el usuario):
- Malls instalados/operativos.
- CUR (Curauma): sin slides de gráfico en el deck (solo foto). Puntos nuevos -37/-38 anotados.
- AEB (El Bosque): incluir "Matriz principal 1°piso" (000025-11); eliminar "Matriz A.A" (000025-30).
- MAM Falabella (000025-09): sin gráfico, con nota "a la espera de la OC para el cambio de equipo".
- MAQ slide 10: caption "San Ignacio 500" es error de copiado del deck original -> se corrige a
  "Matriz Principal" (000025-13) según el narrativo de la propia slide.
- Slide 21 (Resumen Administrativo): tabla contractual OC/HES -> se deja igual (marcar para revisión).

El texto de análisis (%/promedios) NO se recalcula: se marca para revisión manual.
"""

from __future__ import annotations

import re
import sys
import unicodedata
from datetime import datetime, timedelta
from pathlib import Path
from typing import Dict, List, Optional, Tuple

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Emu

sys.path.insert(0, str(Path(__file__).parent))

from generar_reporte_word import (
    get_hourly_measures_for_day,
    parse_date,
    format_number_chilean,
)
from generar_reportes_y_ppt_mall_maipu import obtener_datos_agregados

if sys.platform == "win32":
    try:
        sys.stdout.reconfigure(encoding="utf-8", errors="replace")
        sys.stderr.reconfigure(encoding="utf-8", errors="replace")
    except Exception:
        pass

# ------------------------------------------------------------------ #
# Configuración
# ------------------------------------------------------------------ #
SRC = Path(r"c:\Users\aniba\Downloads\Informe WES __ Parque Arauco 7 Malls  (01.04.2026).pptx")
OUT_DIR = Path(__file__).parent / "reports" / "_tmp_pa_7malls_charts"
OUT = OUT_DIR / "Informe WES __ Parque Arauco 7 Malls (07.07.2026).pptx"
TMP_CHARTS = OUT_DIR

DESDE = "01/05/2026"
HASTA = "07/07/2026"
DESDE_DT = parse_date(DESDE)
HASTA_DT = parse_date(HASTA, end_of_day=True)
PERIODO_TXT = "1/5/2026 a 7/7/2026"
FECHA_PORTADA = "7 Julio 2026"

AZUL = "#1f77b4"
AZUL2 = "#ff7f0e"

# Conjuntos de puntos por mall (instalados/operativos + ajustes del usuario)
MAE = ["000025-01", "000025-04", "000025-07", "000025-19"]          # Estación
MAM = ["000025-08", "000025-10", "000025-32", "000025-33"]          # Maipú (sin Falabella -09)
MAQ = ["000025-13", "000025-34"]                                    # Quilicura
BOM = ["000025-17", "000025-18"]                                    # Buenaventura (San Ignacio)
AEB = ["000025-11", "000025-12"]                                    # El Bosque (Matriz 1°piso + Anillo Plaza)

# Slides que llevan barra "Consumo total periodo" -> set de puntos del mall
SLIDE_TOTAL_POINTS: Dict[int, List[str]] = {
    4: MAE, 10: MAQ, 11: MAQ, 13: BOM, 14: BOM, 16: AEB, 17: AEB,
}

# Slides de análisis con gráficos
SLIDES_CON_GRAFICOS = [4, 5, 7, 8, 10, 11, 13, 14, 16, 17]

FALABELLA = "FALABELLA_NOTE"

WEEKDAYS = {
    "lunes": 0, "martes": 1, "miercoles": 2, "miércoles": 2, "jueves": 3,
    "viernes": 4, "sabado": 5, "sábado": 5, "domingo": 6,
}

_node_name_cache: Dict[str, str] = {}


def _norm(s: str) -> str:
    s = unicodedata.normalize("NFKD", s or "").encode("ascii", "ignore").decode().lower()
    return re.sub(r"[^a-z0-9]+", " ", s).strip()


# ------------------------------------------------------------------ #
# Resolución de nodo desde el texto del caption (con correcciones)
# ------------------------------------------------------------------ #
def resolver_nodo(slide_no: int, point_name: str) -> Optional[str]:
    k = _norm(point_name)
    # Correcciones específicas por slide
    if slide_no == 10 and "san ignacio 500" in k:
        return "000025-13"  # error de copiado -> Matriz Principal (Quilicura)
    if "matriz aa" in k or "matriz a a" in k:
        return "000025-11"  # AEB: usar Matriz principal 1°piso
    if "falabella" in k:
        return FALABELLA
    if "san ignacio 500" in k:
        return "000025-18"
    if "san ignacio 300" in k:
        return "000025-17"
    if "estanque sur" in k:
        return "000025-19"
    if "estanque norte" in k:
        return "000025-01"
    if "pizza hut" in k:
        return "000025-07"
    if "placa bancaria" in k or "placca bancaria" in k:
        return "000025-08"
    if "ripley" in k:
        return "000025-10"
    if "anillo plaza" in k:
        return "000025-12"
    if "matriz principal" in k:
        return "000025-13"
    if "banos" in k or "baños" in k:
        return "000025-04" if slide_no in (4, 5) else "000025-34"
    return None


# ------------------------------------------------------------------ #
# Datos
# ------------------------------------------------------------------ #
_datos_cache: Dict[str, dict] = {}


def _datos_nodo(node_id: str) -> dict:
    if node_id not in _datos_cache:
        print(f"    [data] descargando medidas {node_id} ...", flush=True)
        _datos_cache[node_id] = obtener_datos_agregados([node_id], DESDE, HASTA)
    return _datos_cache[node_id]


def _node_name(node_id: str) -> str:
    if node_id not in _node_name_cache:
        try:
            d = _datos_nodo(node_id)
            ns = d.get("nodes_summary") or []
            _node_name_cache[node_id] = ns[0]["node_name"] if ns else node_id
        except Exception:
            _node_name_cache[node_id] = node_id
    return _node_name_cache[node_id]


def _serie_diaria(node_id: str) -> Tuple[List, List]:
    d = _datos_nodo(node_id)
    ns = d.get("nodes_summary") or []
    measures = ns[0]["measures"] if ns else []
    daily: Dict = {}
    for m in measures:
        daily[m.date.date()] = daily.get(m.date.date(), 0.0) + m.total_m3
    dates = sorted(daily)
    return dates, [daily[x] for x in dates]


def _representative_days(caption: str) -> Tuple[datetime, datetime]:
    """Elige dos días del período según los días de semana mencionados en el caption."""
    found = [WEEKDAYS[w] for w in WEEKDAYS if w in caption.lower()]
    wd1 = found[0] if len(found) >= 1 else 0
    wd2 = found[1] if len(found) >= 2 else wd1
    # primer wd1 desde el inicio
    d1 = DESDE_DT
    for i in range(14):
        if (DESDE_DT + timedelta(days=i)).weekday() == wd1:
            d1 = DESDE_DT + timedelta(days=i)
            break
    # último wd2 antes del fin
    end = datetime(HASTA_DT.year, HASTA_DT.month, HASTA_DT.day)
    d2 = end
    for i in range(14):
        if (end - timedelta(days=i)).weekday() == wd2:
            d2 = end - timedelta(days=i)
            break
    return d1, d2


# ------------------------------------------------------------------ #
# Gráficos
# ------------------------------------------------------------------ #
def _figsize(w_in: float, h_in: float) -> Tuple[float, float]:
    base_w = 6.4
    ratio = max(0.35, min(1.4, h_in / w_in)) if w_in else 0.6
    return (base_w, base_w * ratio)


def chart_total(points: List[str], out: Path, w: float, h: float) -> Path:
    labels, values = [], []
    for nid in points:
        d = _datos_nodo(nid)
        tot = float((d.get("aggregate_summary") or {}).get("total") or 0.0)
        nm = _node_name(nid)
        labels.append(nm if len(nm) <= 22 else nm[:20] + "…")
        values.append(tot)
    fig, ax = plt.subplots(figsize=_figsize(w, h))
    bars = ax.bar(range(len(labels)), values, color=AZUL, alpha=0.85)
    ax.set_xticks(range(len(labels)))
    ax.set_xticklabels(labels, rotation=25, ha="right", fontsize=8)
    ax.set_ylabel("m³", fontsize=9)
    ax.grid(True, axis="y", alpha=0.3)
    mx = max(values) if values else 0
    for b, v in zip(bars, values):
        if v > 0:
            ax.text(b.get_x() + b.get_width() / 2, b.get_height() + mx * 0.01,
                    format_number_chilean(v, 0), ha="center", va="bottom", fontsize=8)
    fig.tight_layout()
    fig.savefig(out, dpi=200, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return out


def chart_diario(node_id: str, out: Path, w: float, h: float) -> Path:
    dates, values = _serie_diaria(node_id)
    fig, ax = plt.subplots(figsize=_figsize(w, h))
    if dates:
        ax.bar(dates, values, color=AZUL, alpha=0.8, width=0.9)
    ax.set_ylabel("m³/día", fontsize=9)
    ax.grid(True, axis="y", alpha=0.3)
    ax.tick_params(axis="x", labelsize=7, rotation=45)
    ax.tick_params(axis="y", labelsize=8)
    fig.tight_layout()
    fig.savefig(out, dpi=200, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return out


def chart_horario_dia(node_id: str, day: datetime, out: Path, w: float, h: float) -> Path:
    try:
        serie = get_hourly_measures_for_day(node_id, day)
    except Exception as e:
        print(f"    [warn] horario {node_id} {day.date()}: {e}")
        serie = []
    horas = list(range(24))
    valores = {h_: v for h_, v in serie}
    vals = [valores.get(h_, 0.0) for h_ in horas]
    fig, ax = plt.subplots(figsize=_figsize(w, h))
    ax.plot(horas, vals, color=AZUL, marker="o", markersize=2.5, linewidth=1.4)
    ax.fill_between(horas, vals, color=AZUL, alpha=0.12)
    ax.set_xlabel("Hora", fontsize=8)
    ax.set_ylabel("m³/h", fontsize=8)
    ax.set_xticks(range(0, 24, 3))
    ax.grid(True, alpha=0.3)
    ax.set_title(day.strftime("%a %d/%m"), fontsize=9)
    ax.tick_params(labelsize=7)
    fig.tight_layout()
    fig.savefig(out, dpi=200, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return out


def chart_nota_falabella(out: Path, w: float, h: float) -> Path:
    fig, ax = plt.subplots(figsize=_figsize(w, h))
    ax.axis("off")
    ax.text(0.5, 0.5,
            "Impulsión Falabella\n\nA la espera de la OC para\nel cambio de equipo.",
            ha="center", va="center", fontsize=11, color="#444444", wrap=True,
            bbox=dict(boxstyle="round,pad=0.6", fc="#f2f6fb", ec=AZUL, lw=1.2))
    fig.savefig(out, dpi=200, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return out


# ------------------------------------------------------------------ #
# PPTX helpers
# ------------------------------------------------------------------ #
def _center(sh) -> Tuple[float, float]:
    return (Emu(sh.left).inches + Emu(sh.width).inches / 2,
            Emu(sh.top).inches + Emu(sh.height).inches / 2)


def _replace_pic(pic, png: Path) -> None:
    slide_part = pic.part
    rId = pic._element.blip_rId
    img_part = slide_part.related_part(rId)
    with open(png, "rb") as f:
        img_part._blob = f.read()


def _clasificar(cap: str) -> str:
    c = cap.lower()
    if "total" in c:
        return "total"
    if "horario" in c:
        return "horario"
    return "diario"


def _extraer_punto(cap: str) -> str:
    c = cap.replace("\n", " ").replace("\x0b", " ")
    c = re.sub(r"(?i)^grafico:\s*", "", c)
    c = re.sub(r"(?i)consumos?\s+horario", "", c)
    c = re.sub(r"(?i)consumo\s+diario", "", c)
    c = re.sub(r"(?i)consumo\s+total", "", c)
    c = re.sub(r"(?i)consumo", "", c)
    c = re.sub(r"(?i)periodo.*$", "", c)
    c = re.sub(r"(?i)(lunes|martes|mi.rcoles|jueves|viernes|s.bado|domingo).*$", "", c)
    c = re.sub(r"\d.*$", "", c)
    return c.strip(" -·\u00a0")


def _actualizar_fechas_texto(prs) -> int:
    n = 0
    rango_re = re.compile(r"\d{1,2}\s*/\s*\d{1,2}\s*/\s*2026\s*a\s*\d{1,2}\s*/\s*\d{1,2}\s*/\s*2026")
    for slide in prs.slides:
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            for para in sh.text_frame.paragraphs:
                for run in para.runs:
                    t = run.text
                    if not t:
                        continue
                    nt = rango_re.sub(PERIODO_TXT, t)
                    nt = nt.replace("1 Abril 2026", FECHA_PORTADA)
                    if nt != t:
                        run.text = nt
                        n += 1
    # Segunda pasada a nivel de párrafo (la fecha de portada suele estar partida en runs)
    port_re = re.compile(r"\d{1,2}\s+[A-Za-zÁÉÍÓÚáéíóúñ]+\s+2026")
    for slide in prs.slides:
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            for para in sh.text_frame.paragraphs:
                if not para.runs:
                    continue
                full = "".join(r.text for r in para.runs)
                if "Abril 2026" in full or ("2026" in full and port_re.fullmatch(full.strip())):
                    nuevo = port_re.sub(FECHA_PORTADA, full)
                    nuevo = nuevo.replace("1 Abril 2026", FECHA_PORTADA)
                    if nuevo != full:
                        para.runs[0].text = nuevo
                        for extra in para.runs[1:]:
                            extra.text = ""
                        n += 1
    return n


# ------------------------------------------------------------------ #
# Main
# ------------------------------------------------------------------ #
def main() -> int:
    if not SRC.is_file():
        print(f"[ERROR] No existe el original: {SRC}")
        return 1
    TMP_CHARTS.mkdir(parents=True, exist_ok=True)

    prs = Presentation(str(SRC))
    changelog: List[str] = []

    for slide_no in SLIDES_CON_GRAFICOS:
        slide = prs.slides[slide_no - 1]
        caps = []
        pics = []
        for sh in slide.shapes:
            if sh.has_text_frame and "grafico:" in sh.text_frame.text.lower():
                caps.append(sh)
            elif sh.shape_type == 13:
                w = Emu(sh.width).inches
                if w and w > 0.6:  # excluir iconos de footer
                    pics.append(sh)
        if not caps or not pics:
            continue

        print(f"\n=== SLIDE {slide_no} : {len(caps)} captions, {len(pics)} gráficos ===")

        # Asignar pictures a captions en 2 fases para que ningún caption quede vacío:
        #  1) cada caption toma su picture más cercano aún libre (garantiza >=1 por caption)
        #  2) los pictures restantes van al caption más cercano
        grupos: Dict[int, List] = {i: [] for i in range(len(caps))}
        cap_centers = [_center(c) for c in caps]
        pic_centers = [_center(p) for p in pics]
        libres = set(range(len(pics)))

        def _dist(pi, ci):
            px, py = pic_centers[pi]
            cx, cy = cap_centers[ci]
            return ((px - cx) ** 2 + (py - cy) ** 2) ** 0.5

        for ci in range(len(caps)):
            if not libres:
                break
            pj = min(libres, key=lambda pi: _dist(pi, ci))
            grupos[ci].append(pics[pj])
            libres.discard(pj)
        for pj in list(libres):
            ci = min(range(len(caps)), key=lambda c: _dist(pj, c))
            grupos[ci].append(pics[pj])

        for i, cap in enumerate(caps):
            cap_txt = cap.text.replace("\n", " ").replace("\x0b", " ").strip()
            kind = _clasificar(cap_txt)
            group_pics = sorted(grupos[i], key=lambda s: (Emu(s.top).inches, Emu(s.left).inches))
            if not group_pics:
                continue
            w0 = Emu(group_pics[0].width).inches
            h0 = Emu(group_pics[0].height).inches

            if kind == "total":
                pts = SLIDE_TOTAL_POINTS.get(slide_no, [])
                if not pts:
                    continue
                png = TMP_CHARTS / f"s{slide_no}_total.png"
                chart_total(pts, png, w0, h0)
                for p in group_pics:
                    _replace_pic(p, png)
                changelog.append(f"S{slide_no} TOTAL -> {', '.join(pts)}")
                continue

            punto = _extraer_punto(cap_txt)
            nid = resolver_nodo(slide_no, cap_txt)  # caption completo: conserva "500/300"
            if nid == FALABELLA:
                png = TMP_CHARTS / f"s{slide_no}_falabella.png"
                chart_nota_falabella(png, w0, h0)
                for p in group_pics:
                    _replace_pic(p, png)
                changelog.append(f"S{slide_no} Falabella -> NOTA (espera OC)")
                continue
            if not nid:
                changelog.append(f"S{slide_no} '{punto}' -> SIN MATCH (no modificado)")
                continue

            if kind == "horario":
                d1, d2 = _representative_days(cap_txt)
                if len(group_pics) >= 2:
                    png1 = TMP_CHARTS / f"s{slide_no}_{nid}_h1.png"
                    png2 = TMP_CHARTS / f"s{slide_no}_{nid}_h2.png"
                    chart_horario_dia(nid, d1, png1, w0, h0)
                    chart_horario_dia(nid, d2, png2, w0, h0)
                    _replace_pic(group_pics[0], png1)
                    _replace_pic(group_pics[1], png2)
                    for extra in group_pics[2:]:
                        _replace_pic(extra, png1)
                    changelog.append(
                        f"S{slide_no} HORARIO {nid} ({_node_name(nid)}) -> {d1.date()} / {d2.date()}")
                else:
                    png1 = TMP_CHARTS / f"s{slide_no}_{nid}_h.png"
                    chart_horario_dia(nid, d1, png1, w0, h0)
                    _replace_pic(group_pics[0], png1)
                    changelog.append(f"S{slide_no} HORARIO {nid} -> {d1.date()}")
            else:  # diario
                png = TMP_CHARTS / f"s{slide_no}_{nid}_d.png"
                chart_diario(nid, png, w0, h0)
                for p in group_pics:
                    _replace_pic(p, png)
                changelog.append(f"S{slide_no} DIARIO {nid} ({_node_name(nid)})")

    n_fechas = _actualizar_fechas_texto(prs)
    print(f"\n[OK] Fechas/periodo actualizados en {n_fechas} runs.")

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))

    print("\n================ CHANGELOG ================")
    for c in changelog:
        print(" -", c)
    print("\n[OK] Copia generada:")
    print(f"     {OUT}")
    print("\n[REVISAR MANUALMENTE]")
    print(" - Texto de análisis (promedios/%): no recalculado.")
    print(" - Slide 21 (Resumen Administrativo): tabla contractual OC/HES sin cambios.")
    print(" - CUR/PAK: solo fotos de instalación (sin gráficos en el deck).")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
