# -*- coding: utf-8 -*-
"""
Edita UNA lámina del PPT de trabajo in-place (sin generar archivo nuevo).

Uso:
  python editar_lamina_pa.py --lamina 4
  python editar_lamina_pa.py --lamina 11   # MAQ análisis (estándar L04)
  python editar_lamina_pa.py --lamina 12   # MAQ horarios
  python editar_lamina_pa.py --lamina 13
  python editar_lamina_pa.py --lamina 14
  python editar_lamina_pa.py --lamina 15   # BOM análisis alza junio
  python editar_lamina_pa.py --lamina 16   # BOM patrón nocturno 0–6 h
  python editar_lamina_pa.py --lamina 17   # AEB análisis consumos (estándar L04)
  python editar_lamina_pa.py --lamina 18   # AEB comparativo nocturno
  python editar_lamina_pa.py --lamina 19   # CUR portada (actualizar puntos)
  python editar_lamina_pa.py --lamina 20   # CUR análisis 18/5–16/6 + vs cuenta
  python editar_lamina_pa.py --lamina 20 --factura-m3 680.5
  python editar_lamina_pa.py --lamina 22   # PAK portada (no modificar)
  python editar_lamina_pa.py --lamina 23   # PAK ranking dual (lámina completa)
  python editar_lamina_pa.py --lamina 24   # PAK cadena abastecimiento DL
  python editar_lamina_pa.py --lamina 25   # PAK consumo diario — 5 puntos (1/2)
  python editar_lamina_pa.py --lamina 26   # PAK consumo diario — 5 puntos (2/2)
  python editar_lamina_pa.py --lamina 27   # PAK patrón nocturno 0–8 h
"""

from __future__ import annotations

import argparse
import re
import sys
from copy import deepcopy
from datetime import date, datetime, timedelta
from pathlib import Path
from typing import Dict, List, Tuple

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.ticker import FormatStrFormatter, MultipleLocator

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Emu, Inches, Pt

sys.path.insert(0, str(Path(__file__).parent))

from generar_reporte_word import format_number_chilean, parse_date, get_hourly_measures_for_day
from generar_reportes_y_ppt_mall_maipu import (
    crear_grafico_ranking_nodos,
    obtener_datos_agregados,
)

if sys.platform == "win32":
    try:
        sys.stdout.reconfigure(encoding="utf-8", errors="replace")
        sys.stderr.reconfigure(encoding="utf-8", errors="replace")
    except Exception:
        pass

PPT = Path(__file__).parent / "reports" / "_tmp_pa_7malls_charts" / (
    "Informe WES __ Parque Arauco 7 Malls (16.08.2026).pptx"
)
CHARTS = PPT.parent

DESDE = "01/06/2026"
HASTA = "16/08/2026"
PERIODO_CAP = "1/6/2026 a 16/8/2026"
FECHA_PORTADA = "16 Agosto 2026"
PCT_TRIM_LABEL = "junio y agosto"
RANK_CAP = "Grafico: Ranking consumo Ene–Mar y Abr–Ago 2026"
RANK_Q2_LABEL = "Abr–Ago 2026"

# MAE — Estación
MAE_NODES = ["000025-01", "000025-04", "000025-07", "000025-19"]

# MAM — Maipú (Falabella activa desde 11/08/2026)
MAM_NODES = ["000025-08", "000025-09", "000025-10", "000025-32", "000025-33"]
FALABELLA_NODE = "000025-09"

# MAQ — Quilicura
MAQ_NODES = ["000025-13", "000025-34"]  # Matriz Principal, Baños

# BOM — Buenaventura (San Ignacio)
BOM_NODES = ["000025-17", "000025-18"]
BOM_NODE_300 = "000025-17"
BOM_NODE_500 = "000025-18"

# AEB — El Bosque
AEB_NODES = ["000025-11", "000025-12"]  # Matriz Principal, Anillo (Primer Piso)
AEB_NODE_MATRIZ = "000025-11"
AEB_NODE_ANILLO = "000025-12"
AEB_LABELS = {
    "000025-11": "Matriz Principal",
    "000025-12": "Anillo (Primer Piso)",
}

# Layout compartido L04/L07 (posiciones ajustadas manualmente en lámina 4)
LAYOUT_RANKING = (0.350, 1.010, 6.000, 2.580)
LAYOUT_CAP_TOTAL = (0.272, 3.698, 6.000, 0.280)
LAYOUT_LEFT_MID = (1.340, 3.970, 4.560, 1.520)
LAYOUT_CAP_LEFT_MID = (1.340, 5.540, 4.560, 0.280)
LAYOUT_RIGHT_MID = (7.832, 3.215, 4.600, 1.510)
LAYOUT_CAP_RIGHT_MID = (7.832, 4.790, 4.600, 0.280)
LAYOUT_LEFT_BOT = (1.436, 5.790, 4.600, 1.400)
LAYOUT_CAP_LEFT_BOT = (1.440, 7.000, 4.600, 0.280)
LAYOUT_RIGHT_BOT = (7.882, 5.134, 4.550, 1.521)
LAYOUT_CAP_RIGHT_BOT = (7.882, 6.805, 4.550, 0.280)
LAYOUT_ANALISIS = (7.170, 0.000, 6.160, 7.500)
CAPTION_H = 0.280

# L04 — gráficos (mismas posiciones que el deck editado)
L04_PICS: List[Tuple[str, str, float, float, float, float]] = [
    ("total", "ALL", *LAYOUT_RANKING),
    ("diario", "000025-19", *LAYOUT_LEFT_MID),   # Estanque Sur
    ("diario", "000025-04", *LAYOUT_RIGHT_MID),  # Baños Públicos
    ("diario", "000025-01", *LAYOUT_LEFT_BOT),   # Estanque Norte
    ("diario", "000025-07", *LAYOUT_RIGHT_BOT),  # Pizza Hut
]
L04_TEXTO_BOX = LAYOUT_ANALISIS
L04_RANK_ENE_MAR = ("01/01/2026", "31/03/2026")
L04_RANK_ABR_JUN = ("01/04/2026", "16/08/2026")
L04_PCT_MAY_JUN = ("01/06/2026", "16/08/2026")
L04_PCT_ANUAL = ("01/01/2026", "16/08/2026")
# Períodos compartidos ranking / panel % (estándar L04)
RANK_ENE_MAR = L04_RANK_ENE_MAR
RANK_ABR_JUN = L04_RANK_ABR_JUN
PCT_TRIM_REVIEW = L04_PCT_MAY_JUN
PCT_ANUAL = L04_PCT_ANUAL
L04_MAE_LABELS = {
    "000025-19": "Estanque Sur",
    "000025-01": "Estanque Norte",
    "000025-07": "Pizza Hut",
    "000025-04": "Baños Públicos",
}

# L07 — layout fijado manualmente en el deck (jul 2026); no sobrescribir panel de análisis
L07_RIPLEY_SLOT = (7.658, 3.475, 4.600, 2.290)
L07_CAP_RIPLEY = (7.658, 5.849, 4.600, 0.280)
L07_PLACA_SLOT = (0.350, 3.970, 6.000, 2.150)
L07_CAP_PLACA = (0.350, 6.170, 6.000, 0.280)
L07_FALABELLA_SLOT = (7.658, 6.295, 4.600, 0.888)
L07_CAP_FALABELLA = (7.392, 7.174, 4.600, 0.280)
L07_TEXTO_BOX = (7.174, 0.152, 6.160, 3.203)

L07_PICS: List[Tuple[str, str, float, float, float, float]] = [
    ("total", "ALL", *LAYOUT_RANKING),
    ("diario", "000025-10", *L07_RIPLEY_SLOT),
    ("diario", "000025-08", *L07_PLACA_SLOT),
    ("diario", FALABELLA_NODE, *L07_FALABELLA_SLOT),
]

# L08 — MAM análisis consumos (mismo estándar L04: ranking dual + % trimestre + anual)
L08_PICS = L07_PICS
L08_TEXTO_BOX = LAYOUT_ANALISIS

# L09 — MAM perfiles horarios (antes L08 en script)
L10_PICS: List[Tuple[str, str, float, float, float, float]] = [
    ("total", "ALL", *LAYOUT_RANKING),
    ("diario", "000025-13", *LAYOUT_LEFT_MID),   # Matriz Principal
    ("diario", "000025-34", *LAYOUT_RIGHT_MID),  # Baños
]
L10_TEXTO_BOX = LAYOUT_ANALISIS

# L11 — MAQ análisis consumos (estándar L04: ranking dual + % + 2 diarios)
L11_MAQ_PICS: List[Tuple[str, str, float, float, float, float]] = [
    ("total", "ALL", *LAYOUT_RANKING),
    ("diario", "000025-13", *LAYOUT_LEFT_MID),
    ("diario", "000025-34", *LAYOUT_RIGHT_MID),
]
L11_MAQ_TEXTO_BOX = LAYOUT_ANALISIS

# L13 — sección BOM (lámina divisoria; no regenerar con script)
# L14 — BOM análisis consumos (estándar L04: ranking dual + % trimestre + anual + 2 diarios)
L14_BOM_PICS: List[Tuple[str, str, float, float, float, float]] = [
    ("total", "ALL", *LAYOUT_RANKING),
    ("diario", BOM_NODE_300, *LAYOUT_LEFT_MID),
    ("diario", BOM_NODE_500, *LAYOUT_RIGHT_MID),
]
L14_BOM_TEXTO_BOX = LAYOUT_ANALISIS
L14_BOM_CHART_DESDE, L14_BOM_CHART_HASTA = L04_PCT_MAY_JUN
L14_BOM_PERIODO_CAP = PERIODO_CAP

# L13 legacy (obsoleto; usar --lamina 14)
L13_PICS: List[Tuple[str, str, float, float, float, float]] = [
    ("total", "ALL", *LAYOUT_RANKING),
    ("diario", BOM_NODE_300, *LAYOUT_LEFT_MID),   # San Ignacio 300
    ("diario", BOM_NODE_500, *LAYOUT_RIGHT_MID),  # San Ignacio 500
]
L13_TEXTO_BOX = LAYOUT_ANALISIS

# L09 — MAM horarios: Ripley (manual, no tocar) + Placa promedio 24 h (L08 fijada)
L09_NODE_PLACA = "000025-08"
L09_RIP_MAY = datetime(2026, 5, 4)
L09_RIP_JUN = datetime(2026, 6, 8)
L09_UMBRAL_ALERTA_M3H = 0.5  # m³/h — alerta Ripley

# Ripley — posiciones fijadas manualmente en el deck (solo referencia; no regenerar)
L09_SLOTS_RIPLEY: List[Tuple[str, datetime, float, float, float, float]] = [
    ("000025-10", L09_RIP_MAY, 0.161, 0.935, 6.026, 2.064),
    ("000025-10", L09_RIP_JUN, 0.241, 3.122, 5.946, 2.032),
]
L09_CAP_RIP_TOP = (1.334, 2.984, 3.902, 0.266)
L09_CAP_RIP_BOT = (1.334, 5.154, 3.902, 0.266)

# Placa — perfil promedio 24 h (columna derecha, bajo panel de texto)
L09_PLACA_PIC = (7.174, 3.380, 6.160, 2.150)
L09_CAP_PLACA = (7.658, 5.570, 4.600, 0.266)
L09_TEXTO_BOX = (7.174, 0.152, 6.160, 3.063)

# L10 — MAQ análisis consumos (mismo encabezado que L04/L07: ranking + % + 2 diarios)
# L12 — MAQ perfiles horarios (Matriz + Baños)
L11_MATRIZ_REF = datetime(2026, 5, 4)        # lunes 4/5 — vs techo ref. 1,6 m³/h
L11_MATRIZ_JUN = datetime(2026, 8, 10)       # lun 10/8 — consumo nocturno actual (sin control)
L11_BANOS_REF = datetime(2026, 5, 4)         # lunes 4/5 — mayor concentración 8-17h
L11_BANOS_REF2 = datetime(2026, 8, 10)       # lun 10/8 — actividad baja / inhábil
L11_NIGHT_HOURS_END = 8  # franja nocturna 0–8 h (inclusive)
L11_REF_INFORME_ANTERIOR = 1.6  # m³/h promedio nocturno informe anterior MAQ
L11_TECHO_NOCTURNO = L11_REF_INFORME_ANTERIOR

L11_SLOTS: List[Tuple[str, datetime, float, float, float, float]] = [
    ("000025-13", L11_MATRIZ_REF, 0.272, 1.130, 6.026, 2.064),
    ("000025-13", L11_MATRIZ_JUN, 0.312, 4.080, 5.946, 2.032),
    # Baños columna der. bajo panel de texto (sin tapar narrativa)
    ("000025-34", L11_BANOS_REF, 7.658, 3.475, 4.600, 1.950),
    ("000025-34", L11_BANOS_REF2, 7.658, 5.550, 4.600, 1.750),
]

L11_CAP_MATRIZ_TOP = (1.298, 3.361, 3.902, 0.266)
L11_CAP_MATRIZ_BOT = (1.574, 6.230, 3.902, 0.266)
L11_CAP_BANOS_TOP = (7.658, 5.520, 4.600, 0.266)
L11_CAP_BANOS_BOT = (7.658, 7.380, 4.600, 0.266)
L11_TEXTO_BOX = (7.174, 0.152, 6.160, 3.203)  # panel superior der. (como L07)

# L12 — MAQ perfiles horarios (Matriz + Baños; estándar L09/L11 layout)
L12_MAQ_SLOTS = L11_SLOTS
L12_MAQ_TEXTO_BOX = L11_TEXTO_BOX
L12_CAP_MATRIZ_TOP = L11_CAP_MATRIZ_TOP
L12_CAP_MATRIZ_BOT = L11_CAP_MATRIZ_BOT
L12_CAP_BANOS_TOP = L11_CAP_BANOS_TOP
L12_CAP_BANOS_BOT = L11_CAP_BANOS_BOT

# L14 — BOM perfiles horarios (layout L11: 500 izq., 300 der.)
L14_S500_REF = datetime(2026, 5, 4)          # lun 4/5
L14_S500_JUN = datetime(2026, 7, 6)          # lun 6/7
L14_S300_REF = datetime(2026, 5, 1)          # vie 1/5 — noche ~0
L14_S300_REF2 = datetime(2026, 5, 8)         # vie 8/5 — noche ~0
L14_REF_INFORME_500 = 7.0                      # m³/h informe anterior San Ignacio 500

L14_SLOTS: List[Tuple[str, datetime, float, float, float, float]] = [
    (BOM_NODE_500, L14_S500_REF, 0.272, 1.130, 6.026, 2.064),
    (BOM_NODE_500, L14_S500_JUN, 0.312, 4.080, 5.946, 2.032),
    (BOM_NODE_300, L14_S300_REF, 7.658, 3.475, 4.600, 1.950),
    (BOM_NODE_300, L14_S300_REF2, 7.658, 5.550, 4.600, 1.750),
]
L14_CAP_500_TOP = L11_CAP_MATRIZ_TOP
L14_CAP_500_BOT = L11_CAP_MATRIZ_BOT
L14_CAP_300_TOP = L11_CAP_BANOS_TOP
L14_CAP_300_BOT = L11_CAP_BANOS_BOT
L14_TEXTO_BOX = L11_TEXTO_BOX
# L15 — BOM análisis alza junio 2026 (layout horarios L12: 500 izq., 300 der.)
L15_S500_NORMAL = datetime(2026, 6, 20)   # vie 20/6 — junio previo al pico
L15_S500_PICO = datetime(2026, 6, 27)      # vie 27/6 — pico fin de mes
L15_S300_INICIO = datetime(2026, 6, 6)     # sáb 6/6 — inicio alza gradual
L15_S300_MAX = datetime(2026, 6, 27)       # sáb 27/6 — máximo del mes
L15_JUN_DESDE = "01/06/2026"
L15_JUN_HASTA = "30/06/2026"
L15_MAY_DESDE = "01/05/2026"
L15_MAY_HASTA = "31/05/2026"
L15_S500_SPIKE_START = datetime(2026, 6, 26)

L15_SLOTS: List[Tuple[str, datetime, float, float, float, float]] = [
    (BOM_NODE_500, L15_S500_NORMAL, 0.272, 1.130, 6.026, 2.064),
    (BOM_NODE_500, L15_S500_PICO, 0.312, 4.080, 5.946, 2.032),
    (BOM_NODE_300, L15_S300_INICIO, 7.658, 3.475, 4.600, 1.950),
    (BOM_NODE_300, L15_S300_MAX, 7.658, 5.550, 4.600, 1.750),
]
L15_CAP_500_TOP = L11_CAP_MATRIZ_TOP
L15_CAP_500_BOT = L11_CAP_MATRIZ_BOT
L15_CAP_300_TOP = L11_CAP_BANOS_TOP
L15_CAP_300_BOT = L11_CAP_BANOS_BOT
L15_TEXTO_BOX = L11_TEXTO_BOX

# L16 — BOM patrón nocturno 0–6 h (ranking mensual S500 + perfiles S300)
L16_NIGHT_END = 6  # horas 0–6 inclusive (00:00–06:00)
L16_ENE_MAR = L04_RANK_ENE_MAR
L16_ABR_JUN = L04_RANK_ABR_JUN
L16_RANK_DESDE = "01/01/2026"
L16_RANK_HASTA = "16/08/2026"
L16_RANK_S500_SLOT = (0.272, 1.130, 6.026, 4.150)
L16_CAP_RANK_S500 = (0.272, 5.330, 6.026, 0.280)
L16_S300_SLOTS: List[Tuple[str, str, str, float, float, float, float]] = [
    (BOM_NODE_300, *L16_ENE_MAR, 7.658, 3.475, 4.600, 1.950),
    (BOM_NODE_300, *L16_ABR_JUN, 7.658, 5.550, 4.600, 1.750),
]
L16_CAP_300_TOP = L11_CAP_BANOS_TOP
L16_CAP_300_BOT = L11_CAP_BANOS_BOT
L16_TEXTO_BOX = L11_TEXTO_BOX
L16_MESES_CORTO = ["Ene", "Feb", "Mar", "Abr", "May", "Jun", "Jul", "Ago", "Sep", "Oct", "Nov", "Dic"]
L16_TARIFA_CLP_M3 = 1400.0

# L17 — AEB análisis consumos (estándar L04: ranking dual + % + 2 diarios)
L17_AEB_PICS: List[Tuple[str, str, float, float, float, float]] = [
    ("total", "ALL", *LAYOUT_RANKING),
    ("diario", AEB_NODE_MATRIZ, *LAYOUT_LEFT_MID),
    ("diario", AEB_NODE_ANILLO, *LAYOUT_RIGHT_MID),
]
L17_AEB_TEXTO_BOX = LAYOUT_ANALISIS

# L18 — AEB comparativo nocturno (layout horarios L12: Matriz izq., Anillo der.)
L18_MATRIZ_REF = datetime(2026, 6, 8)       # lun 8/6 — consumo base nocturno
L18_MATRIZ_PEAK = datetime(2026, 8, 10)     # lun 10/8 — noche actual (sin control)
L18_ANILLO_REF = datetime(2026, 6, 8)
L18_ANILLO_PEAK = datetime(2026, 8, 10)
L18_AEB_SLOTS: List[Tuple[str, datetime, float, float, float, float]] = [
    (AEB_NODE_MATRIZ, L18_MATRIZ_REF, 0.272, 1.130, 6.026, 2.064),
    (AEB_NODE_MATRIZ, L18_MATRIZ_PEAK, 0.312, 4.080, 5.946, 2.032),
    (AEB_NODE_ANILLO, L18_ANILLO_REF, 7.658, 3.475, 4.600, 1.950),
    (AEB_NODE_ANILLO, L18_ANILLO_PEAK, 7.658, 5.550, 4.600, 1.750),
]
L18_CAP_MATRIZ_TOP = L11_CAP_MATRIZ_TOP
L18_CAP_MATRIZ_BOT = L11_CAP_MATRIZ_BOT
L18_CAP_ANILLO_TOP = L11_CAP_BANOS_TOP
L18_CAP_ANILLO_BOT = L11_CAP_BANOS_BOT
L18_AEB_TEXTO_BOX = L11_TEXTO_BOX

# L19 — CUR portada (divisoria; solo actualiza texto de instalación)
# L20 — CUR análisis consumos (período único 18/5–16/6; sin comparación Ene–Mar)
CUR_NODES = ["000025-38", "000025-37"]  # Anillo Norte, Anillo Sur
CUR_NODE_NORTE = "000025-38"
CUR_NODE_SUR = "000025-37"
CUR_LABELS = {
    "000025-37": "Anillo Sur",
    "000025-38": "Anillo Norte",
}
L20_CUR_DESDE = "01/06/2026"
L20_CUR_HASTA = "16/08/2026"
L20_CUR_PERIODO_CAP = "1/6/2026 a 16/8/2026"
L20_CUR_BOLETA_ESVAL = "2368516"  # Boleta CUR Curauma — lecturas 18/5–16/6/2026
L20_CUR_FACTURA_M3 = 621.0  # m³ boleta Esval (periodo lecturas 18/5–16/6)
L20_CUR_LECTURA_INICIO_HORA = 12  # h — lectura física inicial ~12:00 el 18/5
L20_CUR_PICS: List[Tuple[str, str, float, float, float, float]] = [
    ("total", "ALL", *LAYOUT_RANKING),
    ("diario", CUR_NODE_NORTE, *LAYOUT_LEFT_MID),
    ("diario", CUR_NODE_SUR, *LAYOUT_RIGHT_MID),
]
L20_CUR_TEXTO_BOX = LAYOUT_ANALISIS

# L22 — PAK portada (divisoria; no regenerar con script)
# L23 — PAK ranking dual (página completa, sin diarios)
# L24 — PAK cadena abastecimiento DL (22→27→35/36 + aporte 28)
# L25 — PAK consumo diario puntos 1–5 (fila de 5 gráficos)
# L26 — PAK consumo diario puntos 6–10
# L27 — PAK patrón nocturno 0–8 h (cadena DL + ranking)
PAK_NODES = [
    "000025-20",
    "000025-21",
    "000025-22",
    "000025-23",
    "000025-24",
    "000025-27",
    "000025-28",
    "000025-29",
    "000025-35",
    "000025-36",
]
LAYOUT_RANKING_FULL = (0.350, 1.050, 12.500, 6.000)
PAK_CHAIN_SOURCE_ANTIGUA = "000025-22"
PAK_CHAIN_SOURCE_NUEVA = "000025-28"
PAK_CHAIN_DISTRITO = "000025-27"
PAK_CHAIN_BAZAR = "000025-35"
PAK_CHAIN_KENNEDY = "000025-36"
PAK_CHAIN_NODES = [
    PAK_CHAIN_SOURCE_ANTIGUA,
    PAK_CHAIN_SOURCE_NUEVA,
    PAK_CHAIN_DISTRITO,
    PAK_CHAIN_BAZAR,
    PAK_CHAIN_KENNEDY,
]
LAYOUT_PAK_CADENA_BARRAS = (0.350, 1.050, 6.400, 2.750)
LAYOUT_PAK_CADENA_DIARIO = (0.350, 4.050, 6.400, 2.850)
L24_PAK_TEXTO_BOX = LAYOUT_ANALISIS
LAYOUT_RANKING_FULL_PAK_INFO = (0.350, 1.050, 12.500, 5.350)
L23_PAK_INFO_BOX = (0.350, 6.500, 12.500, 0.720)
PAK_NIGHT_END = 8
L27_PAK_REF = datetime(2026, 6, 8)
LAYOUT_PAK_RANK_NOCT = (0.350, 1.050, 6.400, 2.850)
L27_PAK_HOR_LEFT = (0.350, 4.050, 3.100, 2.900)
L27_PAK_HOR_RIGHT = (3.600, 4.050, 3.100, 2.900)
L27_PAK_TEXTO_BOX = LAYOUT_ANALISIS

# L05 — perfiles horarios MAE (posiciones del layout original)
DIAS_ES = ["Lunes", "Martes", "Miércoles", "Jueves", "Viernes", "Sábado", "Domingo"]
DIAS_CAP = ["lunes", "martes", "miercoles", "jueves", "viernes", "sabado", "domingo"]

# L05 — perfiles horarios MAE (solo Estanque Norte y Sur)
L05_MON_REF = datetime(2026, 5, 4)           # lunes 4/5 — sin control
L05_CTRL_DIA = datetime(2026, 8, 10)         # 10/8 — controles verificados (Norte 05/08, Sur corte mantención)
L05_EVENTO_SUR = date(2026, 6, 10)           # reparación tuberías lado sur
L05_JUN_INI = date(2026, 6, 1)
L05_JUN_FIN = date(2026, 6, 30)
L05_TEXTO_BOX = (7.174, 0.152, 6.160, 4.050)

# Horarios 2×2: Norte | Sur (lun 4/5 arriba, 18/6 abajo)
L05_HOR_SLOTS: List[Tuple[str, datetime, float, float, float, float]] = [
    ("000025-01", L05_MON_REF, 0.272, 1.050, 3.250, 1.420),
    ("000025-19", L05_MON_REF, 3.650, 1.050, 3.250, 1.420),
    ("000025-01", L05_CTRL_DIA, 0.272, 2.600, 3.250, 1.420),
    ("000025-19", L05_CTRL_DIA, 3.650, 2.600, 3.250, 1.420),
]
L05_DIARIO_SLOTS: List[Tuple[str, float, float, float, float]] = [
    ("000025-01", 0.272, 4.200, 3.250, 2.450),
    ("000025-19", 3.650, 4.200, 3.250, 2.450),
]

L05_CAP_NORTE_HOR = (0.272, 2.520, 3.250, 0.266)
L05_CAP_SUR_HOR = (3.650, 2.520, 3.250, 0.266)
L05_CAP_NORTE_JUN = (0.272, 6.700, 3.250, 0.266)
L05_CAP_SUR_JUN = (3.650, 6.700, 3.250, 0.266)

# L06 — Baños y Pizza Hut (más espacio; insertar como lámina 6)
L06_CTRL_DIA = L05_CTRL_DIA
L06_TEXTO_BOX = (7.174, 0.152, 6.160, 6.500)
L06_SLOTS: List[Tuple[str, datetime, float, float, float, float]] = [
    ("000025-04", L06_CTRL_DIA, 0.272, 1.130, 6.026, 2.750),
    ("000025-07", L06_CTRL_DIA, 0.272, 4.150, 6.026, 2.750),
]
L06_CAP_BANOS = (0.272, 3.930, 6.026, 0.266)
L06_CAP_PIZZA = (0.272, 6.950, 6.026, 0.266)
L06_NODE_BANOS = "000025-04"
L06_NODE_PIZZA = "000025-07"

_datos: Dict[str, dict] = {}
_horario_cache: Dict[str, Dict[date, Dict[int, float]]] = {}


def _datos_mae() -> dict:
    key = "MAE"
    if key not in _datos:
        print("[data] Descargando MAE ...", flush=True)
        _datos[key] = obtener_datos_agregados(MAE_NODES, DESDE, HASTA)
    return _datos[key]


def _datos_mam() -> dict:
    key = "MAM"
    if key not in _datos:
        print("[data] Descargando MAM ...", flush=True)
        _datos[key] = obtener_datos_agregados(MAM_NODES, DESDE, HASTA)
    return _datos[key]


def _datos_bom() -> dict:
    key = "BOM"
    if key not in _datos:
        print("[data] Descargando BOM ...", flush=True)
        _datos[key] = obtener_datos_agregados(BOM_NODES, DESDE, HASTA)
    return _datos[key]


def _datos_maq() -> dict:
    key = "MAQ"
    if key not in _datos:
        print("[data] Descargando MAQ ...", flush=True)
        _datos[key] = obtener_datos_agregados(MAQ_NODES, DESDE, HASTA)
    return _datos[key]


def _datos_aeb() -> dict:
    key = "AEB"
    if key not in _datos:
        print("[data] Descargando AEB ...", flush=True)
        _datos[key] = obtener_datos_agregados(AEB_NODES, DESDE, HASTA)
    return _datos[key]


def _datos_cur_periodo() -> dict:
    key = f"CUR:{L20_CUR_DESDE}:{L20_CUR_HASTA}"
    if key not in _datos:
        print(f"[data] Descargando CUR {L20_CUR_DESDE}–{L20_CUR_HASTA} ...", flush=True)
        _datos[key] = obtener_datos_agregados(CUR_NODES, L20_CUR_DESDE, L20_CUR_HASTA)
    return _datos[key]


def _datos_pak() -> dict:
    key = "PAK"
    if key not in _datos:
        print("[data] Descargando PAK ...", flush=True)
        _datos[key] = obtener_datos_agregados(PAK_NODES, DESDE, HASTA)
    return _datos[key]


def _datos_nodo(node_id: str) -> dict:
    if node_id not in _datos:
        print(f"[data] Descargando {node_id} ...", flush=True)
        _datos[node_id] = obtener_datos_agregados([node_id], DESDE, HASTA)
    return _datos[node_id]


def _datos_nodo_rango(node_id: str, desde: str, hasta: str) -> dict:
    key = f"N:{node_id}:{desde}:{hasta}"
    if key not in _datos:
        print(f"[data] Descargando {node_id} {desde}–{hasta} ...", flush=True)
        _datos[key] = obtener_datos_agregados([node_id], desde, hasta)
    return _datos[key]


def _figsize(w_in: float, h_in: float) -> Tuple[float, float]:
    return (max(4.0, w_in * 1.15), max(2.2, h_in * 1.15))


def _datos_mae_rango(desde: str, hasta: str) -> dict:
    return _datos_rango(MAE_NODES, desde, hasta)


def _datos_rango(nodes: List[str], desde: str, hasta: str) -> dict:
    key = f"R:{','.join(sorted(nodes))}:{desde}:{hasta}"
    if key not in _datos:
        print(f"[data] Descargando {len(nodes)} nodo(s) {desde}–{hasta} ...", flush=True)
        _datos[key] = obtener_datos_agregados(nodes, desde, hasta)
    return _datos[key]


def _labels_nodos(datos: dict, nodes: List[str]) -> Dict[str, str]:
    by_id = {n["node_id"]: n["node_name"] for n in datos.get("nodes_summary") or []}
    return {nid: by_id.get(nid, nid) for nid in nodes}


def _totales_nodos(datos: dict, nodes: List[str]) -> Dict[str, float]:
    by_id = {
        n["node_id"]: float((n["summary"] or {}).get("total") or 0)
        for n in datos.get("nodes_summary") or []
    }
    return {nid: by_id.get(nid, 0.0) for nid in nodes}


def _pct_totales(totales: Dict[str, float], nodes: List[str]) -> Dict[str, float]:
    total = sum(totales.get(n, 0.0) for n in nodes)
    return {n: (totales.get(n, 0.0) / total * 100.0 if total else 0.0) for n in nodes}


def _pct_rango(nodes: List[str], desde: str, hasta: str) -> Dict[str, float]:
    datos = _datos_rango(nodes, desde, hasta)
    return _pct_totales(_totales_nodos(datos, nodes), nodes)


def chart_ranking_dual(
    nodes: List[str],
    labels: Dict[str, str],
    out: Path,
    w: float,
    h: float,
    *,
    desde_q1: str = RANK_ENE_MAR[0],
    hasta_q1: str = RANK_ENE_MAR[1],
    desde_q2: str = RANK_ABR_JUN[0],
    hasta_q2: str = RANK_ABR_JUN[1],
    label_q1: str = "Ene–Mar 2026",
    label_q2: str = RANK_Q2_LABEL,
) -> Path:
    """Ranking dual: dos barras por punto (Ene–Mar celeste, segundo trimestre azul)."""
    d1 = _datos_rango(nodes, desde_q1, hasta_q1)
    d2 = _datos_rango(nodes, desde_q2, hasta_q2)
    t1 = _totales_nodos(d1, nodes)
    t2 = _totales_nodos(d2, nodes)
    order = sorted(nodes, key=lambda n: -t2.get(n, 0.0))
    lbls = [labels.get(n, n) for n in order]
    v1 = [t1.get(n, 0.0) for n in order]
    v2 = [t2.get(n, 0.0) for n in order]

    fig, ax = plt.subplots(figsize=(max(4.2, w * 1.05), max(1.9, h * 1.05)))
    y = list(range(len(order)))
    bh = 0.34
    ax.barh([i - bh / 2 for i in y], v1, height=bh, color="#B8D4E8", label=label_q1)
    bars2 = ax.barh([i + bh / 2 for i in y], v2, height=bh, color="#4A90E2", label=label_q2)
    ax.set_yticks(y)
    ax.set_yticklabels(lbls, fontsize=7)
    ax.invert_yaxis()
    ax.set_xlabel("m³", fontsize=7)
    ax.tick_params(axis="x", labelsize=6)
    ax.grid(True, alpha=0.3, axis="x")
    mx = max(v1 + v2) if (v1 or v2) else 1.0
    ax.set_xlim(0, mx * 1.18)
    for bar, val in zip(bars2, v2):
        if val <= 0:
            continue
        ax.text(
            bar.get_width() + mx * 0.01,
            bar.get_y() + bar.get_height() / 2,
            format_number_chilean(val, 0),
            va="center",
            ha="left",
            fontsize=6,
            color="#333333",
        )
    ax.legend(fontsize=6, loc="lower right", framealpha=0.9)
    fig.subplots_adjust(left=0.28, right=0.97, top=0.92, bottom=0.12)
    out.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(out, dpi=200, facecolor="white", pad_inches=0.03)
    plt.close(fig)
    return out


def chart_ranking_dual_mae(
    out: Path,
    w: float,
    h: float,
    *,
    desde_q1: str = L04_RANK_ENE_MAR[0],
    hasta_q1: str = L04_RANK_ENE_MAR[1],
    desde_q2: str = L04_RANK_ABR_JUN[0],
    hasta_q2: str = L04_RANK_ABR_JUN[1],
) -> Path:
    return chart_ranking_dual(
        MAE_NODES, L04_MAE_LABELS, out, w, h,
        desde_q1=desde_q1, hasta_q1=hasta_q1, desde_q2=desde_q2, hasta_q2=hasta_q2,
    )


def chart_total_ranking(out: Path, w: float, h: float) -> Path:
    """Ranking MAE L04 — comparativo Ene–Mar vs Abr–Jun."""
    return chart_ranking_dual_mae(out, w, h)


def chart_total_ranking_mam_dual(out: Path, w: float, h: float) -> Path:
    """Ranking MAM L08 — comparativo Ene–Mar vs Abr–Jun (estándar L04)."""
    d2 = _datos_rango(MAM_NODES, *RANK_ABR_JUN)
    labels = _labels_nodos(d2, MAM_NODES)
    return chart_ranking_dual(MAM_NODES, labels, out, w, h)


def chart_total_ranking_mam(out: Path, w: float, h: float) -> Path:
    """Ranking MAM: Placa Bancaria, Ripley y puntos menores (sin Falabella)."""
    datos = _datos_mam()
    crear_grafico_ranking_nodos(datos["nodes_summary"], out)
    return out


def chart_total_ranking_bom_dual(
    out: Path,
    w: float,
    h: float,
    *,
    desde_q2: str = RANK_ABR_JUN[0],
    hasta_q2: str = RANK_ABR_JUN[1],
    label_q2: str = RANK_Q2_LABEL,
) -> Path:
    """Ranking BOM L14 — comparativo Ene–Mar vs segundo período (estándar L04)."""
    d2 = _datos_rango(BOM_NODES, desde_q2, hasta_q2)
    labels = _labels_nodos(d2, BOM_NODES)
    return chart_ranking_dual(
        BOM_NODES, labels, out, w, h,
        desde_q2=desde_q2, hasta_q2=hasta_q2, label_q2=label_q2,
    )


def chart_total_ranking_bom(out: Path, w: float, h: float) -> Path:
    """Ranking BOM — barras horizontales mayor a menor."""
    datos = _datos_bom()
    crear_grafico_ranking_nodos(datos["nodes_summary"], out)
    return out


def chart_total_ranking_maq_dual(out: Path, w: float, h: float) -> Path:
    """Ranking MAQ L11 — comparativo Ene–Mar vs Abr–Jun (estándar L04)."""
    d2 = _datos_rango(MAQ_NODES, *RANK_ABR_JUN)
    labels = _labels_nodos(d2, MAQ_NODES)
    return chart_ranking_dual(MAQ_NODES, labels, out, w, h)


def chart_total_ranking_maq(out: Path, w: float, h: float) -> Path:
    """Ranking MAQ — barras horizontales mayor a menor."""
    datos = _datos_maq()
    crear_grafico_ranking_nodos(datos["nodes_summary"], out)
    return out


def chart_total_ranking_aeb_dual(out: Path, w: float, h: float) -> Path:
    """Ranking AEB L17 — comparativo Ene–Mar vs Abr–Jun (estándar L04)."""
    d2 = _datos_rango(AEB_NODES, *RANK_ABR_JUN)
    labels = _labels_aeb(d2)
    return chart_ranking_dual(AEB_NODES, labels, out, w, h)


def chart_total_ranking_cur_periodo(out: Path, w: float, h: float) -> Path:
    """Ranking CUR L20 — consumo total en el período 18/5–16/6 (sin dual Ene–Mar)."""
    datos = _datos_cur_periodo()
    ns = []
    for n in datos.get("nodes_summary") or []:
        nid = n["node_id"]
        row = dict(n)
        row["node_name"] = CUR_LABELS.get(nid, n.get("node_name", nid))
        ns.append(row)
    crear_grafico_ranking_nodos(ns, out)
    return out


def chart_total_ranking_pak_dual(out: Path, w: float, h: float) -> Path:
    """Ranking PAK L23 — comparativo Ene–Mar vs Abr–Jun (estándar L04)."""
    d2 = _datos_rango(PAK_NODES, *RANK_ABR_JUN)
    labels = _labels_pak(d2)
    return chart_ranking_dual(PAK_NODES, labels, out, w, h)


def _daily_series_pak(node_id: str) -> Dict[date, float]:
    d = _datos_pak()
    for n in d.get("nodes_summary") or []:
        if n["node_id"] != node_id:
            continue
        daily: Dict[date, float] = {}
        for m in n.get("measures") or []:
            d0 = m.date.date()
            daily[d0] = daily.get(d0, 0.0) + m.total_m3
        return daily
    return {}


def _corr_diaria_pak(nid_a: str, nid_b: str) -> float:
    import statistics

    da = _daily_series_pak(nid_a)
    db = _daily_series_pak(nid_b)
    dates = sorted(set(da) | set(db))
    if len(dates) < 2:
        return 0.0
    xs = [da.get(dt, 0.0) for dt in dates]
    ys = [db.get(dt, 0.0) for dt in dates]
    return statistics.correlation(xs, ys)


def _totales_pak_chain() -> Dict[str, float]:
    return _totales_nodos(_datos_pak(), PAK_CHAIN_NODES)


def chart_pak_cadena_totales(out: Path, w: float, h: float) -> Path:
    """Barras horizontales — totales cadena 22/28 → 27 → 35/36."""
    names = _nombres_pak()
    tot = _totales_pak_chain()
    items = [
        (PAK_CHAIN_SOURCE_ANTIGUA, "Fuente: Sandía Antigua", "#5B9BD5"),
        (PAK_CHAIN_SOURCE_NUEVA, "Fuente: Sandía Nueva", "#9DC3E6"),
        (PAK_CHAIN_DISTRITO, "Nodo: Distrito de Lujo", "#70AD47"),
        (PAK_CHAIN_BAZAR, "Aguas abajo: Bazar Gourmet", "#ED7D31"),
        (PAK_CHAIN_KENNEDY, "Aguas abajo: DL Kennedy", "#FFC000"),
    ]
    labels = [lbl for _, lbl, _ in items]
    values = [tot.get(nid, 0.0) for nid, _, _ in items]
    colors = [c for _, _, c in items]

    fig, ax = plt.subplots(figsize=(max(4.5, w * 1.05), max(2.4, h * 1.05)))
    y = list(range(len(items)))
    bars = ax.barh(y, values, color=colors, height=0.55)
    ax.set_yticks(y)
    ax.set_yticklabels(labels, fontsize=7)
    ax.invert_yaxis()
    ax.set_xlabel("m³ período", fontsize=7)
    ax.tick_params(axis="x", labelsize=6)
    ax.grid(True, alpha=0.3, axis="x")
    mx = max(values) if values else 1.0
    ax.set_xlim(0, mx * 1.22)
    for bar, val, (nid, _, _) in zip(bars, values, items):
        if val <= 0:
            continue
        ax.text(
            bar.get_width() + mx * 0.01,
            bar.get_y() + bar.get_height() / 2,
            format_number_chilean(val, 1),
            va="center",
            ha="left",
            fontsize=6,
            color="#333333",
        )
    ax.set_title(
        "Cadena declarada: 22 → 27 → 35 / 36  (+ aporte 28 al DL)",
        fontsize=8,
        pad=6,
    )
    fig.subplots_adjust(left=0.34, right=0.97, top=0.90, bottom=0.12)
    out.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(out, dpi=200, facecolor="white", pad_inches=0.03)
    plt.close(fig)
    return out


def _linreg(xs: List[float], ys: List[float]) -> Tuple[float, float]:
    n = len(xs)
    if n < 2:
        return 0.0, 0.0
    mx = sum(xs) / n
    my = sum(ys) / n
    num = sum((x - mx) * (y - my) for x, y in zip(xs, ys))
    den = sum((x - mx) ** 2 for x in xs)
    slope = num / den if den else 0.0
    return slope, my - slope * mx


def chart_pak_cadena_correlacion(out: Path, w: float, h: float) -> Path:
    """Dos scatter con tendencia: Antigua vs DL y Nueva vs DL."""
    names = _nombres_pak()
    s22 = _daily_series_pak(PAK_CHAIN_SOURCE_ANTIGUA)
    s27 = _daily_series_pak(PAK_CHAIN_DISTRITO)
    s28 = _daily_series_pak(PAK_CHAIN_SOURCE_NUEVA)

    fig, (ax1, ax2) = plt.subplots(
        1, 2,
        figsize=(max(5.8, w * 1.08), max(2.6, h * 1.08)),
    )

    dates_a = sorted(set(s22) & set(s27))
    x22 = [s22[d] for d in dates_a]
    y27a = [s27[d] for d in dates_a]
    r22 = _corr_diaria_pak(PAK_CHAIN_SOURCE_ANTIGUA, PAK_CHAIN_DISTRITO)
    ax1.scatter(x22, y27a, alpha=0.65, s=22, color="#5B9BD5", edgecolors="#2E5C8A", linewidths=0.4)
    if len(x22) >= 2:
        m, b = _linreg(x22, y27a)
        xline = [min(x22), max(x22)]
        ax1.plot(xline, [m * x + b for x in xline], color="#70AD47", linewidth=1.4, linestyle="--")
    ax1.set_xlabel("Sandía Antigua (m³/día)", fontsize=7)
    ax1.set_ylabel("Distrito de Lujo (m³/día)", fontsize=7)
    ax1.set_title(f"Antigua → DL  (r = {r22:.2f})", fontsize=8, pad=5)
    ax1.grid(True, alpha=0.25)
    ax1.tick_params(labelsize=6)

    dates_b = sorted(set(s28) & set(s27))
    x28 = [s28[d] for d in dates_b]
    y27b = [s27[d] for d in dates_b]
    r28 = _corr_diaria_pak(PAK_CHAIN_SOURCE_NUEVA, PAK_CHAIN_DISTRITO)
    ax2.scatter(x28, y27b, alpha=0.65, s=22, color="#9DC3E6", edgecolors="#5B7FA6", linewidths=0.4)
    if len(x28) >= 2:
        m, b = _linreg(x28, y27b)
        xline = [min(x28), max(x28)]
        ax2.plot(xline, [m * x + b for x in xline], color="#ED7D31", linewidth=1.4, linestyle="--")
    ax2.set_xlabel("Sandía Nueva (m³/día)", fontsize=7)
    ax2.set_ylabel("Distrito de Lujo (m³/día)", fontsize=7)
    ax2.set_title(f"Nueva → DL  (r = {r28:.2f})", fontsize=8, pad=5)
    ax2.grid(True, alpha=0.25)
    ax2.tick_params(labelsize=6)

    fig.subplots_adjust(left=0.10, right=0.98, top=0.88, bottom=0.18, wspace=0.32)
    out.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(out, dpi=200, facecolor="white", pad_inches=0.04)
    plt.close(fig)
    return out


def _nocturno_acumulado_periodo(
    node_id: str,
    desde: str = DESDE,
    hasta: str = HASTA,
    *,
    night_end: int = PAK_NIGHT_END,
) -> float:
    cache = _cargar_horario_cache(node_id, desde, hasta, night_end=night_end)
    d0 = parse_date(desde).date()
    d1 = parse_date(hasta, end_of_day=True).date()
    total = 0.0
    for d, hmap in cache.items():
        if d < d0 or d > d1:
            continue
        total += sum(hmap.get(i, 0.0) for i in range(night_end + 1))
    return total


def _dia_mayor_nocturno(
    node_id: str,
    desde: str = DESDE,
    hasta: str = HASTA,
    *,
    night_end: int = PAK_NIGHT_END,
) -> date | None:
    cache = _cargar_horario_cache(node_id, desde, hasta, night_end=night_end)
    d0 = parse_date(desde).date()
    d1 = parse_date(hasta, end_of_day=True).date()
    best_d: date | None = None
    best_v = -1.0
    for d, hmap in cache.items():
        if d < d0 or d > d1:
            continue
        v = sum(hmap.get(i, 0.0) for i in range(night_end + 1))
        if v > best_v:
            best_d, best_v = d, v
    return best_d


def _pct_nocturno_periodo(
    node_id: str,
    desde: str = DESDE,
    hasta: str = HASTA,
    *,
    night_end: int = PAK_NIGHT_END,
) -> float:
    noct = _nocturno_acumulado_periodo(node_id, desde, hasta, night_end=night_end)
    tot = _totales_nodos(_datos_pak(), [node_id]).get(node_id, 0.0)
    if not tot:
        d = _datos_nodo_rango(node_id, desde, hasta)
        tot = float(((d.get("nodes_summary") or [{}])[0].get("summary") or {}).get("total") or 0)
    return noct / tot * 100.0 if tot else 0.0


def chart_ranking_nocturno_pak(
    out: Path,
    w: float,
    h: float,
    *,
    nodes: List[str] | None = None,
) -> Path:
    """Ranking acumulado consumo nocturno 0–8 h — puntos PAK (cadena + top)."""
    focus = nodes or PAK_CHAIN_NODES
    names = _nombres_pak()
    vals = {nid: _nocturno_acumulado_periodo(nid) for nid in focus}
    order = sorted(focus, key=lambda n: -vals.get(n, 0.0))
    labels = [names.get(n, n) for n in order]
    values = [vals.get(n, 0.0) for n in order]

    fig, ax = plt.subplots(figsize=(max(4.5, w * 1.05), max(2.4, h * 1.05)))
    y = list(range(len(order)))
    bars = ax.barh(y, values, color="#DAA520", height=0.55, hatch="///", edgecolor="#B8860B")
    ax.set_yticks(y)
    ax.set_yticklabels(labels, fontsize=7)
    ax.invert_yaxis()
    ax.set_xlabel(f"m³ acumulados 0–{PAK_NIGHT_END} h", fontsize=7)
    ax.tick_params(axis="x", labelsize=6)
    ax.grid(True, alpha=0.3, axis="x")
    mx = max(values) if values else 1.0
    ax.set_xlim(0, mx * 1.22)
    for bar, val in zip(bars, values):
        if val <= 0:
            continue
        ax.text(
            bar.get_width() + mx * 0.01,
            bar.get_y() + bar.get_height() / 2,
            format_number_chilean(val, 1),
            va="center",
            ha="left",
            fontsize=6,
            color="#333333",
        )
    ax.set_title(
        f"Consumo nocturno acumulado ({DESDE}–{HASTA}) — franja 0–{PAK_NIGHT_END} h",
        fontsize=8,
        pad=6,
    )
    fig.subplots_adjust(left=0.34, right=0.97, top=0.90, bottom=0.12)
    out.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(out, dpi=200, facecolor="white", pad_inches=0.03)
    plt.close(fig)
    return out


def chart_perfil_nocturno_pak(
    node_id: str,
    out: Path,
    w_in: float,
    h_in: float,
    *,
    desde: str = DESDE,
    hasta: str = HASTA,
    night_end: int = PAK_NIGHT_END,
) -> Path:
    """Perfil promedio 24 h PAK con franja nocturna resaltada."""
    values, dias = _promedio_horario_rango(node_id, desde, hasta)
    stats = _stats_nocturno_periodo(node_id, desde, hasta, night_end=night_end)
    names = _nombres_pak()
    nm = names.get(node_id, node_id)
    plab = _label_periodo_corto(desde, hasta)
    hours = list(range(24))
    noct_vals = values[: night_end + 1]

    fig, ax = plt.subplots(figsize=(max(4.2, w_in * 1.05), max(1.9, h_in * 1.05)))
    ax.plot(hours, values, color="#4A90E2", linewidth=1.5, marker="o", markersize=2.5)
    ax.fill_between(hours, values, 0, color="#87CEEB", alpha=0.4)
    if any(v > 0 for v in noct_vals):
        ax.fill_between(
            range(night_end + 1), noct_vals, 0,
            color="#FFD700", alpha=0.35, hatch="///",
            edgecolor="#FFA500", linewidth=1.0,
        )
        ax.axvline(x=0, color="orange", linestyle="--", linewidth=1, alpha=0.7)
        ax.axvline(x=night_end, color="orange", linestyle="--", linewidth=1, alpha=0.7)

    ax.set_title(
        f"{nm} — promedio 24 h ({plab}, {int(dias)} días)",
        fontsize=8, fontweight="bold", pad=2,
    )
    ax.set_xlabel("Hora del día", fontsize=7)
    ax.set_ylabel("m³/h", fontsize=7)
    ax.set_xlim(0, 23)
    ax.set_xticks(range(0, 24, 4))
    ax.tick_params(axis="x", labelsize=6)
    ax.tick_params(axis="y", labelsize=6)
    ax.grid(True, alpha=0.3, axis="y")
    _ajustar_eje_y(ax, values)

    pct_cero = (stats["dias_cero"] / stats["dias"] * 100.0) if stats["dias"] else 0.0
    pct_noc = _pct_nocturno_periodo(node_id, desde, hasta, night_end=night_end)
    notas = [
        f"Suma 0–{night_end} h: {format_number_chilean(stats['sum_prom'], 2)} m³/noche",
        f"Prom. 0–{night_end} h: {format_number_chilean(stats['h_prom'], 2)} m³/h",
        f"% nocturno/total: {format_number_chilean(pct_noc, 1)}%",
        f"Noches ≈0: {format_number_chilean(pct_cero, 0)}%",
    ]
    ax.text(
        0.98, 0.97, "\n".join(notas),
        transform=ax.transAxes, ha="right", va="top", fontsize=6,
        bbox=dict(boxstyle="round,pad=0.3", facecolor="white", alpha=0.85, edgecolor="#4A90E2"),
    )

    fig.subplots_adjust(left=0.14, right=0.97, top=0.86, bottom=0.20)
    out.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(out, dpi=200, facecolor="white", pad_inches=0.03)
    plt.close(fig)
    return out


def _pak_daily_row_slots(n: int = 5) -> List[Tuple[float, float, float, float]]:
    """Grilla 3 arriba + 2 abajo (centrados) para 5 diarios legibles con título."""
    slide_w, slide_h = 13.333, 7.5
    margin_l, margin_r = 0.22, 0.22
    margin_t, margin_b = 1.05, 0.22
    gap_x, gap_y = 0.16, 0.18
    usable_w = slide_w - margin_l - margin_r
    usable_h = slide_h - margin_t - margin_b
    row_h = (usable_h - gap_y) / 2.0
    cw3 = (usable_w - gap_x * 2) / 3.0
    cw2 = (usable_w - gap_x) / 2.0
    # ancho uniforme (promedio) para que los 5 se vean pares
    cw = min(cw3, cw2)
    slots: List[Tuple[float, float, float, float]] = []
    # fila superior: 3
    top1 = margin_t
    total_w3 = 3 * cw + 2 * gap_x
    start1 = margin_l + (usable_w - total_w3) / 2.0
    for c in range(3):
        if len(slots) >= n:
            break
        slots.append((start1 + c * (cw + gap_x), top1, cw, row_h))
    # fila inferior: 2 centrados
    top2 = margin_t + row_h + gap_y
    total_w2 = 2 * cw + gap_x
    start2 = margin_l + (usable_w - total_w2) / 2.0
    for c in range(2):
        if len(slots) >= n:
            break
        slots.append((start2 + c * (cw + gap_x), top2, cw, row_h))
    return slots[:n]


def chart_nota_falabella(out: Path, w: float, h: float) -> Path:
    """Nota estática Falabella — en slot Pizza (L04/L07) usa caja completa."""
    compacto = h < 1.15
    fig, ax = plt.subplots(figsize=(max(3.0, w * 0.92), max(0.8, h * 0.9)))
    ax.axis("off")
    if compacto:
        texto = "Impulsión Falabella\nA la espera de la OC para el cambio de equipo."
        fs, pad = 8, 0.35
    else:
        texto = "Impulsión Falabella\n\nA la espera de la OC para\nel cambio de equipo."
        fs, pad = 11, 0.6
    ax.text(
        0.5, 0.5, texto,
        ha="center", va="center", fontsize=fs, color="#444444",
        bbox=dict(boxstyle=f"round,pad={pad}", fc="#f2f6fb", ec="#4A90E2", lw=1.0),
    )
    out.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(out, dpi=200, bbox_inches="tight", facecolor="white", pad_inches=0.02)
    plt.close(fig)
    return out


def chart_diario_linea(
    node_id: str,
    out: Path,
    w: float,
    h: float,
    *,
    desde: str = DESDE,
    hasta: str = HASTA,
    compact: bool = False,
    titulo: str | None = None,
) -> Path:
    """Consumo diario en línea con relleno celeste (mismo estilo que lámina 5 horario)."""
    if desde == DESDE and hasta == HASTA:
        d = _datos_nodo(node_id)
    else:
        d = _datos_nodo_rango(node_id, desde, hasta)
    ns = d.get("nodes_summary") or []
    measures = ns[0]["measures"] if ns else []
    daily: Dict = {}
    for m in measures:
        daily[m.date.date()] = daily.get(m.date.date(), 0.0) + m.total_m3
    dates = sorted(daily.keys())
    values = [daily[x] for x in dates]
    x = list(range(len(dates)))

    fig, ax = plt.subplots(figsize=_figsize(w, h))
    lw = 1.1 if compact else 1.8
    if values:
        ax.plot(x, values, linestyle="-", color="#4A90E2", linewidth=lw)
        ax.fill_between(x, values, 0, color="#87CEEB", alpha=0.4)
    fs_y = 6 if compact else 8
    fs_tick = 5 if compact else 7
    ax.set_ylabel("m³", fontsize=fs_y)
    ax.grid(True, alpha=0.3, axis="y")
    ax.set_ylim(bottom=0)
    # Etiquetas de fecha espaciadas
    max_ticks = 5 if compact else 8
    if len(dates) > max_ticks:
        step = max(1, len(dates) // max_ticks)
        ticks = x[::step]
        labels = [dates[i].strftime("%d/%m") for i in range(0, len(dates), step)]
    else:
        ticks = x
        labels = [d.strftime("%d/%m") for d in dates]
    ax.set_xticks(ticks)
    ax.set_xticklabels(labels, rotation=45, ha="right", fontsize=fs_tick)
    ax.tick_params(axis="y", labelsize=fs_tick)
    if titulo:
        ax.set_title(titulo, fontsize=9 if not compact else 7, fontweight="bold", pad=4, color="#1F4E79")
    else:
        nm = (ns[0].get("node_name") if ns else node_id) or node_id
        short = nm.replace("Sala de Bomba ", "S.B. ").replace("Impulsión ", "Imp. ")
        if len(short) > 42:
            short = short[:40] + "…"
        ax.set_title(short, fontsize=9 if not compact else 7, fontweight="bold", pad=4, color="#1F4E79")
    fig.tight_layout(rect=(0, 0, 1, 0.96))
    fig.savefig(out, dpi=200, bbox_inches="tight", facecolor="white", pad_inches=0.08)
    plt.close(fig)
    return out


def chart_diario_junio(
    node_id: str,
    out: Path,
    w: float,
    h: float,
    *,
    marcar: date | None = None,
) -> Path:
    """Consumo diario junio con techo ajustado y marca en día de evento (10/06)."""
    ms = _datos_nodo(node_id)["nodes_summary"][0]["measures"]
    daily: Dict[date, float] = {}
    for m in ms:
        d0 = m.date.date()
        if L05_JUN_INI <= d0 <= L05_JUN_FIN:
            daily[d0] = daily.get(d0, 0.0) + m.total_m3
    dates = sorted(daily.keys())
    values = [daily[d] for d in dates]
    x = list(range(len(dates)))

    fig, ax = plt.subplots(figsize=_figsize(w, h))
    if values:
        ax.plot(x, values, linestyle="-", color="#4A90E2", linewidth=1.8)
        ax.fill_between(x, values, 0, color="#87CEEB", alpha=0.4)
        ymax = max(values) * 1.10
        if marcar and marcar in daily:
            ymax = max(ymax, daily[marcar] * 1.35)
        ax.set_ylim(0, ymax)

    if marcar and marcar in daily:
        idx = dates.index(marcar)
        val = daily[marcar]
        ax.axvline(x=idx, color="#C0504D", linestyle="--", linewidth=1.2, alpha=0.9)
        ax.plot(idx, val, marker="o", markersize=6, color="#C0504D", zorder=5)
        ax.annotate(
            f"10/06: {format_number_chilean(val, 1)} m³",
            xy=(idx, val),
            xytext=(idx + 1.5, val + ymax * 0.06),
            fontsize=7,
            color="#C0504D",
            arrowprops=dict(arrowstyle="->", color="#C0504D", lw=0.8),
        )

    ax.set_ylabel("m³/día", fontsize=8)
    ax.grid(True, alpha=0.3, axis="y")
    step = max(1, len(dates) // 6)
    ticks = x[::step]
    labels = [dates[i].strftime("%d/%m") for i in range(0, len(dates), step)]
    ax.set_xticks(ticks)
    ax.set_xticklabels(labels, rotation=45, ha="right", fontsize=7)
    ax.tick_params(axis="y", labelsize=7)
    ax.set_title("Junio 2026", fontsize=8, fontweight="bold", pad=3)
    fig.tight_layout()
    out.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(out, dpi=200, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return out


def _match_pic(shapes, target_l: float, target_t: float, tol: float = 0.25):
    best, bd = None, tol
    for sh in shapes:
        if sh.shape_type != 13:
            continue
        if Emu(sh.width).inches < 0.6:
            continue
        l, t = Emu(sh.left).inches, Emu(sh.top).inches
        d = abs(l - target_l) + abs(t - target_t)
        if d < bd:
            bd, best = d, sh
    return best


def _replace_pic(pic, png: Path) -> None:
    """Reemplazo inseguro si varias shapes comparten rId — usar _poner_grafico_fresco."""
    rId = pic._element.blip_rId
    img_part = pic.part.related_part(rId)
    with open(png, "rb") as f:
        img_part._blob = f.read()


def _borrar_imagenes_grafico(slide) -> None:
    """Elimina todas las imágenes de gráfico (evita rId compartidos de Google Slides)."""
    eliminar = []
    for sh in slide.shapes:
        if sh.shape_type == 13 and Emu(sh.width).inches > 0.6:
            eliminar.append(sh._element)
    for el in eliminar:
        el.getparent().remove(el)


def _poner_grafico_fresco(slide, png: Path, left: float, top: float, w: float, h: float):
    """Inserta imagen nueva (cada gráfico con su propio blob)."""
    return slide.shapes.add_picture(str(png), Inches(left), Inches(top), width=Inches(w), height=Inches(h))


def _set_titulo_izq(slide, texto: str) -> None:
    for sh in slide.shapes:
        if sh.has_text_frame and sh.top < Inches(0.5) and Emu(sh.left).inches < 2.0:
            sh.text_frame.paragraphs[0].text = texto
            return
    tb = slide.shapes.add_textbox(Inches(0.17), Inches(0.27), Inches(8.11), Inches(0.86))
    tb.text_frame.paragraphs[0].text = texto


def _titulo_slide(slide) -> str:
    candidatos: List[Tuple[float, str]] = []
    for sh in slide.shapes:
        if sh.has_text_frame and sh.top < Inches(0.5):
            t = sh.text_frame.text.strip()
            if t and len(t) < 120:
                candidatos.append((Emu(sh.left).inches, t))
    if not candidatos:
        return ""
    candidatos.sort(key=lambda x: x[0])
    return candidatos[0][1]


def _eliminar_slide(prs, index: int) -> None:
    r_id = prs.slides._sldIdLst[index].rId
    prs.part.drop_rel(r_id)
    del prs.slides._sldIdLst[index]


def _es_slide_l06(slide) -> bool:
    t = _titulo_slide(slide).lower()
    if ("baño" in t or "bano" in t or "pizza" in t) and "perfiles" in t:
        return True
    tiene_pizza = False
    tiene_banos_mae = False
    tiene_estanque = False
    tiene_matriz = False
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        tx = sh.text_frame.text.lower()
        if "grafico:" in tx and "pizza hut" in tx:
            tiene_pizza = True
        if "grafico:" in tx and ("baños públicos" in tx or "banos publicos" in tx):
            tiene_banos_mae = True
        if "estanque" in tx:
            tiene_estanque = True
        if "matriz principal" in tx or "alimentación baños" in tx:
            tiene_matriz = True
    return tiene_pizza and tiene_banos_mae and not tiene_estanque and not tiene_matriz


def _lamina_6_insertada(prs) -> bool:
    return any(_es_slide_l06(s) for s in prs.slides)


def _limpiar_huerfanos_mae(prs) -> None:
    """Elimina láminas rotas entre L05 y L06 (solo índices 5–7)."""
    eliminar: List[int] = []
    for i in range(5, min(8, len(prs.slides))):
        slide = prs.slides[i]
        pics = sum(
            1 for sh in slide.shapes
            if sh.shape_type == 13 and Emu(sh.width).inches > 0.6
        )
        tiene_est = any(
            sh.has_text_frame and "estanque sur" in sh.text_frame.text.lower()
            for sh in slide.shapes
        )
        titulo = _titulo_slide(slide).lower()
        if pics == 2 and tiene_est and "estanques" not in titulo:
            eliminar.append(i)
    for idx in sorted(eliminar, reverse=True):
        _eliminar_slide(prs, idx)
        print(f"[OK] Lámina huérfana eliminada (índice {idx + 1})")


def _limpiar_duplicados_l06(prs) -> None:
    """Quita L06 duplicadas; conserva la primera válida tras L05."""
    indices = [i for i in range(5, len(prs.slides)) if _es_slide_l06(prs.slides[i])]
    if len(indices) <= 1:
        return
    conservar = indices[0]
    for idx in sorted(indices[1:], reverse=True):
        if idx != conservar:
            _eliminar_slide(prs, idx)
            print(f"[OK] L06 duplicada eliminada (índice {idx + 1})")


def _slide_tiene_keyword(prs, keyword: str) -> bool:
    kw = keyword.lower()
    for slide in prs.slides:
        for sh in slide.shapes:
            if sh.has_text_frame and kw in sh.text_frame.text.lower():
                return True
    return False


def _insertar_slide_copia(prs, src_index: int, dest_index: int) -> None:
    src = prs.slides[src_index]
    layout_idx = min(5, len(prs.slide_layouts) - 1)
    blank = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    for shape in src.shapes:
        new_el = deepcopy(shape.element)
        blank.shapes._spTree.insert_element_before(new_el, "p:extLst")
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[-1])
    xml_slides.insert(dest_index, slides[-1])


def _indice_seccion_maq(prs) -> int | None:
    for i, slide in enumerate(prs.slides):
        if _titulo_slide(slide).strip().upper() == "MAQ":
            return i
    return None


def _indice_maq_analisis(prs) -> int | None:
    for i, slide in enumerate(prs.slides):
        t = _titulo_slide(slide).lower()
        if "maq" in t and ("análisis" in t or "analisis" in t):
            return i
    return None


def _ensure_lamina_maq_analisis(prs) -> int:
    """Inserta lámina MAQ análisis (clon L04) después de la sección MAQ si no existe."""
    idx = _indice_maq_analisis(prs)
    if idx is not None:
        return idx
    idx_maq = _indice_seccion_maq(prs)
    if idx_maq is None:
        raise RuntimeError("No se encontró lámina sección MAQ en el deck")
    dest = idx_maq + 1
    _insertar_slide_copia(prs, 3, dest)
    slide = prs.slides[dest]
    _set_titulo_izq(slide, "MAQ - ANÁLISIS CONSUMOS")
    for sh in list(slide.shapes):
        if not sh.has_text_frame:
            continue
        if sh.top < Inches(0.5):
            continue
        tx = sh.text_frame.text.lower()
        if "estanque" in tx or "pizza" in tx or "baños públicos" in tx:
            sh._element.getparent().remove(sh._element)
    print(f"[OK] Lámina MAQ análisis insertada en posición {dest + 1} (clon L04)")
    return dest


def _indice_maq_horarios(prs) -> int | None:
    for i, slide in enumerate(prs.slides):
        t = _titulo_slide(slide).lower()
        if "maq" in t and "horario" in t:
            return i
        if _slide_tiene_keyword_por_slide(slide, "horario matriz principal"):
            return i
    return None


def _ensure_lamina_maq_horarios(prs) -> int:
    """Inserta lámina MAQ horarios (clon layout 4 gráficos) después de MAQ análisis."""
    idx = _indice_maq_horarios(prs)
    if idx is not None:
        return idx
    idx_an = _indice_maq_analisis(prs)
    if idx_an is None:
        raise RuntimeError("Falta lámina MAQ análisis — ejecutar --lamina 11 primero")
    dest = idx_an + 1
    _insertar_slide_copia(prs, 4, dest)
    slide = prs.slides[dest]
    _set_titulo_izq(slide, "MAQ - PERFILES HORARIOS")
    for sh in list(slide.shapes):
        if not sh.has_text_frame:
            continue
        if sh.top < Inches(0.5):
            continue
        tx = sh.text_frame.text.lower()
        if (
            "estanque" in tx
            or "comparación nocturna" in tx
            or "comparacion nocturna" in tx
            or ("control on/off" in tx and "matriz" not in tx)
        ):
            sh._element.getparent().remove(sh._element)
    _limpiar_clon_horarios_mae(slide)
    _limpiar_clon_horarios_mae(slide)
    print(f"[OK] Lámina MAQ horarios insertada en posición {dest + 1} (clon layout horarios)")
    return dest


def _indice_seccion_bom(prs) -> int | None:
    for i, slide in enumerate(prs.slides):
        if _titulo_slide(slide).strip().upper() == "BOM":
            return i
    return None


def _indice_bom_analisis(prs) -> int | None:
    for i, slide in enumerate(prs.slides):
        t = _titulo_slide(slide).lower()
        if "bom" in t and ("análisis" in t or "analisis" in t):
            return i
    return None


def _ensure_lamina_bom_analisis(prs) -> int:
    """Asegura lámina BOM análisis (estándar L04) en la posición tras sección BOM."""
    idx = _indice_bom_analisis(prs)
    if idx is not None:
        return idx
    idx_bom = _indice_seccion_bom(prs)
    if idx_bom is None:
        raise RuntimeError("No se encontró lámina sección BOM en el deck")
    dest = idx_bom + 1
    if dest < len(prs.slides):
        slide = prs.slides[dest]
        titulo = _titulo_slide(slide).strip().lower()
        if not titulo or "aeb" in titulo or len(slide.shapes) <= 4:
            _set_titulo_izq(slide, "BOM - ANÁLISIS CONSUMOS")
            print(f"[OK] Lámina {dest + 1} preparada para BOM análisis")
            return dest
    _insertar_slide_copia(prs, 3, dest)
    slide = prs.slides[dest]
    _set_titulo_izq(slide, "BOM - ANÁLISIS CONSUMOS")
    for sh in list(slide.shapes):
        if not sh.has_text_frame:
            continue
        if sh.top < Inches(0.5):
            continue
        tx = sh.text_frame.text.lower()
        if "estanque" in tx or "pizza" in tx or "baños públicos" in tx:
            sh._element.getparent().remove(sh._element)
    print(f"[OK] Lámina BOM análisis insertada en posición {dest + 1} (clon L04)")
    return dest


def _indice_bom_junio(prs) -> int | None:
    for i, slide in enumerate(prs.slides):
        t = _titulo_slide(slide).lower()
        if "bom" in t and "junio" in t:
            return i
        if _slide_tiene_keyword_por_slide(slide, "alza consumo junio"):
            return i
    return None


def _ensure_lamina_bom_junio(prs) -> int:
    """Asegura lámina BOM junio (clon layout horarios L12) tras BOM análisis."""
    idx = _indice_bom_junio(prs)
    if idx is not None:
        return idx
    idx_an = _indice_bom_analisis(prs)
    if idx_an is None:
        raise RuntimeError("Falta lámina BOM análisis — ejecutar --lamina 14 primero")
    dest = idx_an + 1
    src = _indice_maq_horarios(prs)
    if src is None:
        src = 11
    if dest < len(prs.slides):
        slide = prs.slides[dest]
        titulo = _titulo_slide(slide).strip().lower()
        if (
            not titulo
            or "aeb" in titulo
            or "pak" in titulo
            or len(slide.shapes) <= 4
            or ("bom" in titulo and "análisis" not in titulo and "analisis" not in titulo)
        ):
            _set_titulo_izq(slide, "BOM - ANÁLISIS JUNIO 2026")
            _limpiar_clon_horarios_mae(slide)
            print(f"[OK] Lámina {dest + 1} preparada para BOM junio")
            return dest
    _insertar_slide_copia(prs, src, dest)
    slide = prs.slides[dest]
    _set_titulo_izq(slide, "BOM - ANÁLISIS JUNIO 2026")
    _limpiar_clon_horarios_mae(slide)
    for sh in list(slide.shapes):
        if not sh.has_text_frame:
            continue
        if sh.top < Inches(0.5):
            continue
        tx = sh.text_frame.text.lower()
        if "matriz principal" in tx or "alimentación baños" in tx or "alimentacion banos" in tx:
            sh._element.getparent().remove(sh._element)
    print(f"[OK] Lámina BOM junio insertada en posición {dest + 1} (clon layout horarios)")
    return dest


def _indice_bom_nocturno(prs) -> int | None:
    for i, slide in enumerate(prs.slides):
        t = _titulo_slide(slide).lower()
        if "bom" in t and "nocturn" in t:
            return i
        if _slide_tiene_keyword_por_slide(slide, "ranking consumo nocturno"):
            return i
        if _slide_tiene_keyword_por_slide(slide, "patron nocturno 0-6"):
            return i
        if _slide_tiene_keyword_por_slide(slide, "control wes instalado"):
            return i
    idx_jun = _indice_bom_junio(prs)
    if idx_jun is not None and idx_jun + 1 < len(prs.slides):
        return idx_jun + 1
    return None


def _ensure_lamina_bom_nocturno(prs) -> int:
    """Asegura lámina BOM nocturno tras BOM junio (reutiliza L16 existente si hay)."""
    idx = _indice_bom_nocturno(prs)
    if idx is not None:
        return idx
    idx_jun = _indice_bom_junio(prs)
    if idx_jun is None:
        raise RuntimeError("Falta lámina BOM junio — ejecutar --lamina 15 primero")
    dest = idx_jun + 1
    src = _indice_maq_horarios(prs) or idx_jun
    _insertar_slide_copia(prs, src, dest)
    slide = prs.slides[dest]
    _set_titulo_izq(slide, "BOM - PATRÓN NOCTURNO (0–6 H)")
    _limpiar_clon_horarios_mae(slide)
    for sh in list(slide.shapes):
        if not sh.has_text_frame:
            continue
        if sh.top < Inches(0.5):
            continue
        tx = sh.text_frame.text.lower()
        if "matriz principal" in tx or "alimentación baños" in tx or "alimentacion banos" in tx:
            sh._element.getparent().remove(sh._element)
    print(f"[OK] Lámina BOM nocturno insertada en posición {dest + 1} (clon layout horarios)")
    return dest


def _indice_aeb_analisis(prs) -> int | None:
    for i, slide in enumerate(prs.slides):
        t = _titulo_slide(slide).lower()
        if "aeb" in t and ("análisis" in t or "analisis" in t):
            if "nocturn" not in t and "comparativo" not in t:
                return i
    return None


def _indice_aeb_nocturno(prs) -> int | None:
    for i, slide in enumerate(prs.slides):
        t = _titulo_slide(slide).lower()
        if "aeb" in t and ("nocturn" in t or "comparativo" in t):
            return i
    idx = _indice_aeb_analisis(prs)
    if idx is not None and idx + 1 < len(prs.slides):
        t2 = _titulo_slide(prs.slides[idx + 1]).lower()
        if "aeb" in t2:
            return idx + 1
    return None


def _ensure_lamina_aeb_analisis(prs) -> int:
    """Asegura lámina AEB análisis (estándar L04) tras BOM nocturno."""
    idx = _indice_aeb_analisis(prs)
    if idx is not None:
        return idx
    idx_bom = _indice_bom_nocturno(prs)
    if idx_bom is None:
        raise RuntimeError("Falta lámina BOM nocturno — ejecutar --lamina 16 primero")
    dest = idx_bom + 1
    if dest < len(prs.slides):
        slide = prs.slides[dest]
        titulo = _titulo_slide(slide).strip().lower()
        if not titulo or "aeb" in titulo or "pak" in titulo or len(slide.shapes) <= 4:
            _set_titulo_izq(slide, "AEB - ANÁLISIS CONSUMOS")
            print(f"[OK] Lámina {dest + 1} preparada para AEB análisis")
            return dest
    _insertar_slide_copia(prs, 3, dest)
    slide = prs.slides[dest]
    _set_titulo_izq(slide, "AEB - ANÁLISIS CONSUMOS")
    for sh in list(slide.shapes):
        if not sh.has_text_frame:
            continue
        if sh.top < Inches(0.5):
            continue
        tx = sh.text_frame.text.lower()
        if "estanque" in tx or "pizza" in tx or "san ignacio" in tx:
            sh._element.getparent().remove(sh._element)
    print(f"[OK] Lámina AEB análisis insertada en posición {dest + 1} (clon L04)")
    return dest


def _ensure_lamina_aeb_nocturno(prs) -> int:
    """Asegura lámina AEB comparativo nocturno tras AEB análisis."""
    idx = _indice_aeb_nocturno(prs)
    if idx is not None:
        return idx
    idx_an = _indice_aeb_analisis(prs)
    if idx_an is None:
        raise RuntimeError("Falta lámina AEB análisis — ejecutar --lamina 17 primero")
    dest = idx_an + 1
    src = _indice_maq_horarios(prs) or idx_an
    if dest < len(prs.slides):
        slide = prs.slides[dest]
        titulo = _titulo_slide(slide).strip().lower()
        if (
            not titulo
            or "aeb" in titulo
            or "pak" in titulo
            or len(slide.shapes) <= 4
        ):
            _set_titulo_izq(slide, "AEB - COMPARATIVO NOCTURNO")
            _limpiar_clon_horarios_mae(slide)
            print(f"[OK] Lámina {dest + 1} preparada para AEB comparativo nocturno")
            return dest
    _insertar_slide_copia(prs, src, dest)
    slide = prs.slides[dest]
    _set_titulo_izq(slide, "AEB - COMPARATIVO NOCTURNO")
    _limpiar_clon_horarios_mae(slide)
    for sh in list(slide.shapes):
        if not sh.has_text_frame:
            continue
        if sh.top < Inches(0.5):
            continue
        tx = sh.text_frame.text.lower()
        if "matriz principal" in tx or "anillo" in tx or "matriz aa" in tx:
            sh._element.getparent().remove(sh._element)
    print(f"[OK] Lámina AEB comparativo nocturno insertada en posición {dest + 1}")
    return dest


def _indice_seccion_cur(prs) -> int | None:
    for i, slide in enumerate(prs.slides):
        if _titulo_slide(slide).strip().upper() == "CUR":
            return i
    return None


def _indice_cur_analisis(prs) -> int | None:
    for i, slide in enumerate(prs.slides):
        t = _titulo_slide(slide).lower()
        if "cur" in t and ("análisis" in t or "analisis" in t or "consumo" in t):
            return i
    idx = _indice_seccion_cur(prs)
    if idx is not None and idx + 1 < len(prs.slides):
        t2 = _titulo_slide(prs.slides[idx + 1]).lower()
        if "cur" in t2 or "curauma" in t2:
            return idx + 1
    return None


def _ensure_lamina_cur_analisis(prs) -> int:
    """Inserta lámina CUR análisis (clon L04) inmediatamente después de la portada CUR."""
    idx = _indice_cur_analisis(prs)
    if idx is not None:
        t = _titulo_slide(prs.slides[idx]).lower()
        if "pak" not in t:
            return idx
    idx_cur = _indice_seccion_cur(prs)
    if idx_cur is None:
        raise RuntimeError("No se encontró lámina sección CUR en el deck")
    dest = idx_cur + 1
    if dest < len(prs.slides):
        slide = prs.slides[dest]
        titulo = _titulo_slide(slide).strip().lower()
        if titulo == "pak" or "pak" in titulo:
            src = _indice_aeb_analisis(prs) or 3
            _insertar_slide_copia(prs, src, dest)
            slide = prs.slides[dest]
            _set_titulo_izq(slide, "CUR - ANÁLISIS CONSUMOS")
            for sh in list(slide.shapes):
                if not sh.has_text_frame or sh.top < Inches(0.5):
                    continue
                tx = sh.text_frame.text.lower()
                if "estanque" in tx or "aeb" in tx or "matriz principal" in tx:
                    sh._element.getparent().remove(sh._element)
            print(f"[OK] Lámina CUR análisis insertada en posición {dest + 1} (tras portada CUR)")
            return dest
        if not titulo or "cur" in titulo or len(slide.shapes) <= 4:
            _set_titulo_izq(slide, "CUR - ANÁLISIS CONSUMOS")
            print(f"[OK] Lámina {dest + 1} preparada para CUR análisis")
            return dest
    src = _indice_aeb_analisis(prs) or idx_cur
    _insertar_slide_copia(prs, src, dest)
    slide = prs.slides[dest]
    _set_titulo_izq(slide, "CUR - ANÁLISIS CONSUMOS")
    for sh in list(slide.shapes):
        if not sh.has_text_frame or sh.top < Inches(0.5):
            continue
        tx = sh.text_frame.text.lower()
        if "estanque" in tx or "aeb" in tx:
            sh._element.getparent().remove(sh._element)
    print(f"[OK] Lámina CUR análisis insertada en posición {dest + 1}")
    return dest


def _indice_seccion_pak(prs) -> int | None:
    for i, slide in enumerate(prs.slides):
        if _titulo_slide(slide).strip().upper() == "PAK":
            return i
    return None


def _indice_pak_analisis(prs) -> int | None:
    for i, slide in enumerate(prs.slides):
        t = _titulo_slide(slide).lower()
        if "pak" in t and ("análisis" in t or "analisis" in t or "consumo" in t):
            if "cadena" not in t and "diario" not in t:
                return i
    return None


def _indice_pak_cadena(prs) -> int | None:
    for i, slide in enumerate(prs.slides):
        t = _titulo_slide(slide).lower()
        if "pak" in t and ("cadena" in t or "abastecimiento" in t):
            return i
    return None


def _indices_pak_diarios_slides(prs) -> List[int]:
    """Slides PAK con >=5 gráficos grandes (diarios), excluye cadena/nocturno."""
    found: List[int] = []
    for i, slide in enumerate(prs.slides):
        t = _titulo_slide(slide).lower()
        if "pak" not in t:
            continue
        body = " ".join(
            sh.text_frame.text.lower()
            for sh in slide.shapes
            if sh.has_text_frame
        )
        if "cadena declarada" in body or "patrón nocturno" in body or "patron nocturno" in body:
            continue
        if "cadena" in t or "nocturn" in t or "abastecimiento" in t:
            continue
        # ranking página completa suele tener 1–2 imágenes grandes
        big = [
            sh for sh in slide.shapes
            if sh.shape_type == 13 and Emu(sh.width).inches > 2.0
        ]
        if len(big) >= 5 or "diario" in t or "1/2" in t or "2/2" in t:
            # ranking L23 también es PAK análisis con pocas imgs — filtrar
            if len(big) >= 5:
                found.append(i)
            elif "diario" in t or "1/2" in t or "2/2" in t:
                found.append(i)
    return found


def _indice_pak_diarios(prs) -> int | None:
    idxs = _indices_pak_diarios_slides(prs)
    return idxs[0] if idxs else None


def _indice_pak_diarios_a(prs) -> int | None:
    idxs = _indices_pak_diarios_slides(prs)
    return idxs[0] if idxs else None


def _indice_pak_diarios_b(prs) -> int | None:
    idxs = _indices_pak_diarios_slides(prs)
    return idxs[1] if len(idxs) > 1 else None


def _indice_pak_nocturno(prs) -> int | None:
    for i, slide in enumerate(prs.slides):
        t = _titulo_slide(slide).lower()
        if "pak" in t and "nocturn" in t:
            return i
    return None


def _ensure_lamina_pak_analisis(prs) -> int:
    """Inserta lámina PAK análisis (clon L04) inmediatamente después de la portada PAK."""
    idx = _indice_pak_analisis(prs)
    if idx is not None:
        return idx
    idx_pak = _indice_seccion_pak(prs)
    if idx_pak is None:
        raise RuntimeError("No se encontró lámina sección PAK en el deck")
    dest = idx_pak + 1
    src = _indice_cur_analisis(prs) or _indice_aeb_analisis(prs) or 3
    if dest >= len(prs.slides):
        _insertar_slide_copia(prs, src, dest)
    else:
        t = _titulo_slide(prs.slides[dest]).strip().lower()
        if not ("pak" in t and ("análisis" in t or "analisis" in t)):
            _insertar_slide_copia(prs, src, dest)
    slide = prs.slides[dest]
    _set_titulo_izq(slide, "PAK - ANÁLISIS CONSUMOS")
    _limpiar_clon_pak_analisis(slide)
    print(f"[OK] Lámina PAK análisis en posición {dest + 1} (tras portada PAK)")
    return dest


def _ensure_lamina_pak_cadena(prs) -> int:
    """Inserta lámina PAK cadena abastecimiento tras el ranking."""
    idx = _indice_pak_cadena(prs)
    if idx is not None:
        return idx
    _ensure_lamina_pak_analisis(prs)
    idx_rank = _indice_pak_analisis(prs)
    dest = idx_rank + 1
    src = idx_rank
    if dest >= len(prs.slides):
        _insertar_slide_copia(prs, src, dest)
    else:
        t = _titulo_slide(prs.slides[dest]).strip().lower()
        if "cadena" not in t and "abastecimiento" not in t:
            _insertar_slide_copia(prs, src, dest)
    slide = prs.slides[dest]
    _set_titulo_izq(slide, "PAK - CADENA ABASTECIMIENTO DL")
    _limpiar_clon_pak_analisis(slide)
    print(f"[OK] Lámina PAK cadena en posición {dest + 1}")
    return dest


def _ensure_lamina_pak_diarios_a(prs) -> int:
    """Inserta lámina PAK diarios parte 1/2 (5 puntos) tras cadena abastecimiento."""
    idx = _indice_pak_diarios_a(prs)
    if idx is not None:
        return idx
    idx_old = _indice_pak_diarios(prs)
    if idx_old is not None:
        slide = prs.slides[idx_old]
        _set_titulo_izq(slide, "PAK - CONSUMO DIARIO PUNTOS (1/2)")
        return idx_old
    _ensure_lamina_pak_cadena(prs)
    idx_cadena = _indice_pak_cadena(prs)
    dest = idx_cadena + 1
    src = _indice_pak_analisis(prs) or idx_cadena
    if dest >= len(prs.slides):
        _insertar_slide_copia(prs, src, dest)
    else:
        t = _titulo_slide(prs.slides[dest]).strip().lower()
        if "diario" not in t:
            _insertar_slide_copia(prs, src, dest)
    slide = prs.slides[dest]
    _set_titulo_izq(slide, "PAK - CONSUMO DIARIO PUNTOS (1/2)")
    _limpiar_clon_pak_analisis(slide)
    print(f"[OK] Lámina PAK diarios 1/2 en posición {dest + 1}")
    return dest


def _ensure_lamina_pak_diarios_b(prs) -> int:
    """Inserta lámina PAK diarios parte 2/2 (5 puntos) tras la primera."""
    idx = _indice_pak_diarios_b(prs)
    if idx is not None:
        return idx
    _ensure_lamina_pak_diarios_a(prs)
    idx_a = _indice_pak_diarios_a(prs)
    dest = idx_a + 1
    src = idx_a
    if dest >= len(prs.slides):
        _insertar_slide_copia(prs, src, dest)
    else:
        t = _titulo_slide(prs.slides[dest]).strip().lower()
        if not ("diario" in t and "2/2" in t):
            _insertar_slide_copia(prs, src, dest)
    slide = prs.slides[dest]
    _set_titulo_izq(slide, "PAK - CONSUMO DIARIO PUNTOS (2/2)")
    _limpiar_clon_pak_analisis(slide)
    print(f"[OK] Lámina PAK diarios 2/2 en posición {dest + 1}")
    return dest


def _ensure_lamina_pak_nocturno(prs) -> int:
    """Inserta lámina PAK patrón nocturno tras diarios 2/2."""
    idx = _indice_pak_nocturno(prs)
    if idx is not None:
        return idx
    _ensure_lamina_pak_diarios_b(prs)
    idx_di = _indice_pak_diarios_b(prs)
    dest = idx_di + 1
    src = _indice_pak_analisis(prs) or idx_di
    if dest >= len(prs.slides):
        _insertar_slide_copia(prs, src, dest)
    else:
        t = _titulo_slide(prs.slides[dest]).strip().lower()
        if "nocturn" not in t:
            _insertar_slide_copia(prs, src, dest)
    slide = prs.slides[dest]
    _set_titulo_izq(slide, "PAK - PATRÓN NOCTURNO (0–8 H)")
    _limpiar_clon_pak_analisis(slide)
    print(f"[OK] Lámina PAK nocturno en posición {dest + 1}")
    return dest


def _limpiar_clon_pak_analisis(slide) -> None:
    """Quita narrativa heredada al clonar layout L04 para PAK."""
    for sh in list(slide.shapes):
        if not sh.has_text_frame or _es_titulo_slide_shape(sh):
            continue
        tx = sh.text_frame.text.strip().lower()
        if not tx or "grafico:" in tx:
            continue
        if (
            tx.startswith("cur")
            or tx.startswith("aeb")
            or "anillo" in tx
            or "estanque" in tx
            or "cuenta de agua" in tx
        ):
            sh._element.getparent().remove(sh._element)


def _limpiar_clon_cur_analisis(slide) -> None:
    """Quita títulos/narrativa heredados al clonar layout L04 para CUR."""
    for sh in list(slide.shapes):
        if not sh.has_text_frame or _es_titulo_slide_shape(sh):
            continue
        tx = sh.text_frame.text.strip().lower()
        if not tx or "grafico:" in tx:
            continue
        if (
            tx.startswith("aeb")
            or "estanque" in tx
            or "matriz principal" in tx and "anillo" not in tx
            or "san ignacio" in tx
        ):
            sh._element.getparent().remove(sh._element)


def _es_titulo_slide_shape(sh) -> bool:
    return sh.has_text_frame and sh.top < Inches(0.12) and Emu(sh.left).inches < 2.0


def _limpiar_clon_horarios_mae(slide) -> None:
    """Quita narrativa MAE/MAQ heredada al clonar layout de perfiles horarios."""
    keywords = (
        "comparación nocturna",
        "comparacion nocturna",
        "estanque sur",
        "estanque norte",
        "matriz principal",
        "alimentación baños",
        "alimentacion banos",
        "baños públicos",
        "banos publicos",
        "pizza hut",
        "anillo plaza",
        "matriz aa",
        "san ignacio",
    )
    for sh in list(slide.shapes):
        if not sh.has_text_frame or _es_titulo_slide_shape(sh):
            continue
        tx = sh.text_frame.text.lower()
        if "grafico:" in tx:
            continue
        if any(k in tx for k in keywords):
            sh._element.getparent().remove(sh._element)
            continue
        if "control on/off" in tx and "matriz" not in tx:
            sh._element.getparent().remove(sh._element)


def _limpiar_panel_narrativa_bom(slide, box: Tuple[float, float, float, float]) -> None:
    """Elimina cajas narrativas en el panel derecho antes de escribir texto BOM."""
    l, t, w, h = box
    x1, y1 = Inches(l), Inches(t)
    x2, y2 = Inches(l + w), Inches(t + h)
    for sh in list(slide.shapes):
        if not sh.has_text_frame or _es_titulo_slide_shape(sh):
            continue
        tx = sh.text_frame.text.lower()
        if "grafico:" in tx:
            continue
        sh_x2 = sh.left + sh.width
        sh_y2 = sh.top + sh.height
        if sh.left < x2 and sh_x2 > x1 and sh.top < y2 and sh_y2 > y1:
            sh._element.getparent().remove(sh._element)


def _eliminar_duplicado_tras_bom_nocturno(prs, idx_noc: int) -> None:
    """Quita lámina duplicada de BOM junio si quedó inmediatamente después de L16."""
    dup = idx_noc + 1
    if dup >= len(prs.slides):
        return
    slide = prs.slides[dup]
    if not _slide_tiene_keyword_por_slide(slide, "alza consumo junio"):
        return
    if _slide_tiene_keyword_por_slide(slide, "patrón nocturno") or _slide_tiene_keyword_por_slide(slide, "patron nocturno"):
        return
    _eliminar_slide(prs, dup)
    print(f"[OK] Lámina duplicada BOM junio eliminada (posición {dup + 1})")


def _reparar_maq_si_falta(prs) -> None:
    """Restaura láminas MAQ análisis/horarios si se perdieron."""
    if _slide_tiene_keyword(prs, "grafico: consumo diario matriz principal"):
        return
    idx_maq = next(
        (i for i, s in enumerate(prs.slides) if _titulo_slide(s).strip().upper() == "MAQ"),
        None,
    )
    if idx_maq is None:
        print("[WARN] No se encontró lámina sección MAQ")
        return
    idx_mam_an = next(
        (
            i for i, s in enumerate(prs.slides)
            if "mam" in _titulo_slide(s).lower()
            and "análisis" in _titulo_slide(s).lower()
            and _slide_tiene_keyword_por_slide(s, "ripley periodo")
        ),
        8,
    )
    idx_mam_hor = idx_mam_an + 1
    dest = idx_maq + 1
    _insertar_slide_copia(prs, idx_mam_an, dest)
    for sh in prs.slides[dest].shapes:
        if sh.has_text_frame and sh.top < Inches(0.5):
            for p in sh.text_frame.paragraphs:
                if p.text.strip():
                    p.text = "MAQ - ANÁLISIS CONSUMOS"
                    break
    _insertar_slide_copia(prs, idx_mam_hor, dest + 1)
    for sh in prs.slides[dest + 1].shapes:
        if sh.has_text_frame and sh.top < Inches(0.5):
            for p in sh.text_frame.paragraphs:
                if p.text.strip():
                    p.text = "MAQ - ANÁLISIS CONSUMOS"
                    break
    print("[OK] Láminas MAQ análisis y horarios restauradas")


def _slide_tiene_keyword_por_slide(slide, keyword: str) -> bool:
    kw = keyword.lower()
    return any(sh.has_text_frame and kw in sh.text_frame.text.lower() for sh in slide.shapes)


def _slide_index(prs, lamina: int) -> int:
    if lamina == 6:
        for i, slide in enumerate(prs.slides):
            if _es_slide_l06(slide):
                return i
        return 5
    if lamina == 8:
        for i, slide in enumerate(prs.slides):
            if _slide_tiene_keyword_por_slide(slide, "consumo diario ripley periodo"):
                return i
            t = _titulo_slide(slide).lower()
            if "mam" in t and "análisis" in t:
                if not _slide_tiene_keyword_por_slide(slide, "horario placa bancaria"):
                    return i
        return 7
    if lamina == 9:
        for i, slide in enumerate(prs.slides):
            if _slide_tiene_keyword_por_slide(slide, "horario placa bancaria"):
                return i
        return 8
    if lamina == 10:
        for i, slide in enumerate(prs.slides):
            t = _titulo_slide(slide).strip().upper()
            if t == "MAQ":
                return i
        return 9
    if lamina == 11:
        for i, slide in enumerate(prs.slides):
            t = _titulo_slide(slide).lower()
            if "maq" in t and ("análisis" in t or "analisis" in t):
                return i
    if lamina == 12:
        for i, slide in enumerate(prs.slides):
            t = _titulo_slide(slide).lower()
            if "maq" in t and "horario" in t:
                return i
            if _slide_tiene_keyword_por_slide(slide, "horario matriz principal"):
                return i
            if _slide_tiene_keyword_por_slide(slide, "horario alimentación baños"):
                return i
            if _slide_tiene_keyword_por_slide(slide, "horario alimentacion banos"):
                return i
    if lamina == 13:
        for i, slide in enumerate(prs.slides):
            if _titulo_slide(slide).strip().upper() == "BOM":
                return i
        return 12
    if lamina == 14:
        for i, slide in enumerate(prs.slides):
            t = _titulo_slide(slide).lower()
            if "bom" in t and ("análisis" in t or "analisis" in t):
                return i
    if lamina == 15:
        for i, slide in enumerate(prs.slides):
            t = _titulo_slide(slide).lower()
            if "bom" in t and "junio" in t:
                return i
            if _slide_tiene_keyword_por_slide(slide, "alza consumo junio"):
                return i
            if "bom" in t and "horario" in t:
                return i
            if _slide_tiene_keyword_por_slide(slide, "horario san ignacio"):
                return i
    if lamina == 16:
        for i, slide in enumerate(prs.slides):
            t = _titulo_slide(slide).lower()
            if "bom" in t and "nocturn" in t:
                return i
            if _slide_tiene_keyword_por_slide(slide, "patron nocturno 0-6"):
                return i
    if lamina == 17:
        idx = _indice_aeb_analisis(prs)
        if idx is not None:
            return idx
        return 16
    if lamina == 18:
        idx = _indice_aeb_nocturno(prs)
        if idx is not None:
            return idx
        return 17
    if lamina == 19:
        idx = _indice_seccion_cur(prs)
        if idx is not None:
            return idx
        return 19
    if lamina == 20:
        idx = _indice_cur_analisis(prs)
        if idx is not None:
            return idx
        return 20
    if lamina == 22:
        idx = _indice_seccion_pak(prs)
        if idx is not None:
            return idx
        return 21
    if lamina == 23:
        idx = _indice_pak_analisis(prs)
        if idx is not None:
            return idx
        return 22
    if lamina == 24:
        idx = _indice_pak_cadena(prs)
        if idx is not None:
            return idx
        return 23
    if lamina == 25:
        idx = _indice_pak_diarios_a(prs)
        if idx is not None:
            return idx
        return 24
    if lamina == 26:
        idx = _indice_pak_diarios_b(prs)
        if idx is not None:
            return idx
        return 25
    if lamina == 27:
        idx = _indice_pak_nocturno(prs)
        if idx is not None:
            return idx
        return 26
    off = 1 if _lamina_6_insertada(prs) and lamina >= 7 else 0
    base = {
        4: 3, 5: 4, 6: 5, 7: 6, 8: 7, 10: 9, 11: 10, 12: 11,
        13: 12, 14: 13, 15: 14, 16: 15, 17: 16, 18: 17, 19: 19, 20: 20,
        22: 21, 23: 22, 24: 23, 25: 24, 26: 25, 27: 26,
    }
    return base[lamina] + off


def _insertar_lamina_6(prs) -> None:
    """Inserta lámina MAE (Baños/Pizza) en posición 6 si aún no existe."""
    if _lamina_6_insertada(prs):
        print("[OK] Lámina 6 MAE ya existe")
        return
    src = prs.slides[4]
    layout_idx = min(5, len(prs.slide_layouts) - 1)
    blank = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    for shape in src.shapes:
        new_el = deepcopy(shape.element)
        blank.shapes._spTree.insert_element_before(new_el, "p:extLst")
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[-1])
    xml_slides.insert(5, slides[-1])
    slide = prs.slides[5]
    for sh in slide.shapes:
        if sh.has_text_frame and sh.top < Inches(0.5):
            for p in sh.text_frame.paragraphs:
                if p.text.strip():
                    p.text = "MAE - PERFILES HORARIOS (BAÑOS Y PIZZA)"
                    break
            break
    else:
        tb = slide.shapes.add_textbox(Inches(0.17), Inches(0.27), Inches(8.11), Inches(0.86))
        tb.text_frame.paragraphs[0].text = "MAE - PERFILES HORARIOS (BAÑOS Y PIZZA)"
    _borrar_imagenes_grafico(slide)
    for sh in list(slide.shapes):
        if not sh.has_text_frame:
            continue
        if sh.top < Inches(0.5):
            continue
        tx = sh.text_frame.text.lower()
        if "estanque" in tx or "reparación" in tx or "control on/off" in tx or "mantención" in tx:
            sh._element.getparent().remove(sh._element)
    print("[OK] Lámina 6 MAE insertada (Baños y Pizza Hut)")


def _pct_mam_may_jun() -> Dict[str, float]:
    return _pct_rango(MAM_NODES, *PCT_TRIM_REVIEW)


def _pct_mam_anual() -> Dict[str, float]:
    return _pct_rango(MAM_NODES, *PCT_ANUAL)


def _pct_mae_may_jun() -> Dict[str, float]:
    datos = _datos_mae_rango(*L04_PCT_MAY_JUN)
    return _pct_totales(_totales_nodos(datos, MAE_NODES), MAE_NODES)


def _pct_mae_anual() -> Dict[str, float]:
    datos = _datos_mae_rango(*L04_PCT_ANUAL)
    return _pct_totales(_totales_nodos(datos, MAE_NODES), MAE_NODES)


def _pct_mae() -> Dict[str, float]:
    datos = _datos_mae()
    by_id = {n["node_id"]: float((n["summary"] or {}).get("total") or 0) for n in datos["nodes_summary"]}
    total = sum(by_id.values())
    return {k: (v / total * 100 if total else 0) for k, v in by_id.items()}


def _pct_mam() -> Dict[str, float]:
    datos = _datos_mam()
    by_id = {n["node_id"]: float((n["summary"] or {}).get("total") or 0) for n in datos["nodes_summary"]}
    total = sum(by_id.values())
    return {k: (v / total * 100 if total else 0) for k, v in by_id.items()}


def _nombres_mam() -> Dict[str, str]:
    datos = _datos_mam()
    return {n["node_id"]: n["node_name"] for n in datos["nodes_summary"]}


def _pct_maq_may_jun() -> Dict[str, float]:
    return _pct_rango(MAQ_NODES, *PCT_TRIM_REVIEW)


def _pct_maq_anual() -> Dict[str, float]:
    return _pct_rango(MAQ_NODES, *PCT_ANUAL)


def _pct_maq() -> Dict[str, float]:
    datos = _datos_maq()
    by_id = {n["node_id"]: float((n["summary"] or {}).get("total") or 0) for n in datos["nodes_summary"]}
    total = sum(by_id.values())
    return {k: (v / total * 100 if total else 0) for k, v in by_id.items()}


def _nombres_maq() -> Dict[str, str]:
    datos = _datos_maq()
    return {n["node_id"]: n["node_name"] for n in datos["nodes_summary"]}


def _pct_bom_may_jun() -> Dict[str, float]:
    return _pct_rango(BOM_NODES, *PCT_TRIM_REVIEW)


def _pct_bom_anual() -> Dict[str, float]:
    return _pct_rango(BOM_NODES, *PCT_ANUAL)


def _pct_bom() -> Dict[str, float]:
    datos = _datos_bom()
    by_id = {n["node_id"]: float((n["summary"] or {}).get("total") or 0) for n in datos["nodes_summary"]}
    total = sum(by_id.values())
    return {k: (v / total * 100 if total else 0) for k, v in by_id.items()}


def _nombres_bom() -> Dict[str, str]:
    datos = _datos_bom()
    return {n["node_id"]: n["node_name"] for n in datos["nodes_summary"]}


def _labels_aeb(datos: dict) -> Dict[str, str]:
    base = _labels_nodos(datos, AEB_NODES)
    return {nid: AEB_LABELS.get(nid, base[nid]) for nid in AEB_NODES}


def _nombres_aeb() -> Dict[str, str]:
    datos = _datos_aeb()
    names = {n["node_id"]: n["node_name"] for n in datos["nodes_summary"]}
    return {nid: AEB_LABELS.get(nid, names.get(nid, nid)) for nid in AEB_NODES}


def _pct_aeb_may_jun() -> Dict[str, float]:
    return _pct_rango(AEB_NODES, *PCT_TRIM_REVIEW)


def _pct_aeb_anual() -> Dict[str, float]:
    return _pct_rango(AEB_NODES, *PCT_ANUAL)


def _ratio_anillo_sobre_matriz_aeb() -> float:
    datos = _datos_aeb()
    tot = _totales_nodos(datos, AEB_NODES)
    mat = tot.get(AEB_NODE_MATRIZ, 0.0)
    ani = tot.get(AEB_NODE_ANILLO, 0.0)
    return (ani / mat * 100.0) if mat > 0 else 0.0


def _labels_cur(datos: dict) -> Dict[str, str]:
    base = _labels_nodos(datos, CUR_NODES)
    return {nid: CUR_LABELS.get(nid, base[nid]) for nid in CUR_NODES}


def _nombres_cur() -> Dict[str, str]:
    datos = _datos_cur_periodo()
    names = {n["node_id"]: n["node_name"] for n in datos["nodes_summary"]}
    return {nid: CUR_LABELS.get(nid, names.get(nid, nid)) for nid in CUR_NODES}


def _pct_cur_periodo() -> Dict[str, float]:
    datos = _datos_cur_periodo()
    return _pct_totales(_totales_nodos(datos, CUR_NODES), CUR_NODES)


def _total_wes_cur_periodo() -> float:
    datos = _datos_cur_periodo()
    return sum(_totales_nodos(datos, CUR_NODES).values())


def _labels_pak(datos: dict) -> Dict[str, str]:
    return _labels_nodos(datos, PAK_NODES)


def _nombres_pak() -> Dict[str, str]:
    datos = _datos_pak()
    return {n["node_id"]: n["node_name"] for n in datos["nodes_summary"]}


def _pct_pak_may_jun() -> Dict[str, float]:
    return _pct_rango(PAK_NODES, *PCT_TRIM_REVIEW)


def _pct_pak_anual() -> Dict[str, float]:
    return _pct_rango(PAK_NODES, *PCT_ANUAL)


def _top_pak_nodes(n: int = 2) -> List[str]:
    datos = _datos_pak()
    tot = _totales_nodos(datos, PAK_NODES)
    return sorted(PAK_NODES, key=lambda nid: -tot.get(nid, 0.0))[:n]


def _promedio_diario(node_id: str) -> float:
    d = _datos_nodo(node_id)
    ns = d.get("nodes_summary") or []
    measures = ns[0]["measures"] if ns else []
    daily: Dict[date, float] = {}
    for m in measures:
        d0 = m.date.date()
        daily[d0] = daily.get(d0, 0.0) + m.total_m3
    vals = list(daily.values())
    return sum(vals) / len(vals) if vals else 0.0


def _texto_analisis_l13() -> str:
    fn = format_number_chilean
    p = _pct_bom()
    names = _nombres_bom()
    avg300 = _promedio_diario(BOM_NODE_300)
    avg500 = _promedio_diario(BOM_NODE_500)
    ratio = (avg500 / avg300) if avg300 > 0 else 0.0
    nm300 = names.get(BOM_NODE_300, "San Ignacio 300")
    nm500 = names.get(BOM_NODE_500, "San Ignacio 500")
    return "\n".join([
        "Respecto del total monitoreado",
        f"{nm300} = {fn(p.get(BOM_NODE_300, 0), 1)}%",
        f"{nm500} = {fn(p.get(BOM_NODE_500, 0), 1)}%",
        "",
        f"Consumo diario promedio ({DESDE}–{HASTA}):",
        f"{nm300} = {fn(avg300, 1)} m³/día",
        f"{nm500} = {fn(avg500, 1)} m³/día",
        f"{nm500} registra consumo {fn(ratio, 1)}x superior al {nm300}.",
    ])


def _buscar_o_crear_narrativa_l13(slide, texto: str) -> None:
    l, t, w, h = L13_TEXTO_BOX
    target = None
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        tx = sh.text_frame.text.lower()
        if (
            "respecto del total monitoreado" in tx
            or "san ignacio 500" in tx
            or "san ignacio 300" in tx
            or "consumo diario promedio" in tx
        ):
            target = sh
            break
    if target is None:
        target = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    else:
        target.left = Inches(l)
        target.top = Inches(t)
        target.width = Inches(w)
        target.height = Inches(h)
    _set_texto_parrafos(target, texto)
    _aplicar_negrita_solo_porcentajes(target)


def _actualizar_textos_l13(slide) -> None:
    _buscar_o_crear_narrativa_l13(slide, _texto_analisis_l13())
    print("[OK] Panel análisis L13 actualizado (BOM, párrafos limpios)")

    _remover_captions_viejos(slide)
    names = _nombres_bom()
    periodo = PERIODO_CAP
    caps = [
        (f"Grafico: Consumo total periodo {periodo}", LAYOUT_CAP_TOTAL),
        (
            f"Grafico: Consumo diario {names.get(BOM_NODE_300, 'San Ignacio 300')} periodo {periodo}",
            LAYOUT_CAP_LEFT_MID,
        ),
        (
            f"Grafico: Consumo diario {names.get(BOM_NODE_500, 'San Ignacio 500')} periodo {periodo}",
            LAYOUT_CAP_RIGHT_MID,
        ),
    ]
    for txt, pos in caps:
        _agregar_caption_fijo(slide, txt, pos)
    print("[OK] L13 — ranking horizontal + 2 diarios (estilo L10)")


def _texto_analisis_l04() -> str:
    fn = format_number_chilean
    p_mj = _pct_mae_may_jun()
    p_an = _pct_mae_anual()
    order = sorted(MAE_NODES, key=lambda n: -p_mj.get(n, 0.0))
    lines = [f"Respecto al total monitoreable entre {PCT_TRIM_LABEL}"]
    for nid in order:
        lines.append(f"{L04_MAE_LABELS[nid]} = {fn(p_mj[nid], 1)}%")
    lines.append("")
    lines.append("Respecto al total monitoreable anual")
    order_an = sorted(MAE_NODES, key=lambda n: -p_an.get(n, 0.0))
    for nid in order_an:
        lines.append(f"{L04_MAE_LABELS[nid]} = {fn(p_an[nid], 1)}%")
    lines.extend([
        "",
        "Noche con control: no se lee como fuga.",
        "Estanque Sur: corte on/off a cargo de mantención nocturna (no automático).",
        "Pizza Hut: control 00:00–06:00 desde el 01/07/2026.",
        "Estanque Norte: control on/off desde el 05/08/2026; el alza de agosto es diurna.",
        "Baños Públicos: control instalado sin uso; noches ya ~0.",
    ])
    return "\n".join(lines)


def _buscar_o_crear_caja_analisis_l04(slide, texto: str) -> None:
    l, t, w, h = L04_TEXTO_BOX
    target = None
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        tx = sh.text_frame.text.lower()
        if "respecto" in tx and ("monitore" in tx or "estanque" in tx):
            target = sh
            break
    if target is None:
        target = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    else:
        target.left = Inches(l)
        target.top = Inches(t)
        target.width = Inches(w)
        target.height = Inches(h)
    _set_texto_parrafos(target, texto)
    _aplicar_negrita_solo_porcentajes(target)


def _placa_promedios_jun18() -> Tuple[float, float, float]:
    """Promedio diario Placa Bancaria antes/desde 18-jun-2026 y % de aumento."""
    corte = date(2026, 6, 18)
    d = _datos_nodo("000025-08")
    ns = d.get("nodes_summary") or []
    measures = ns[0]["measures"] if ns else []
    daily: Dict[date, float] = {}
    for m in measures:
        d0 = m.date.date()
        daily[d0] = daily.get(d0, 0.0) + m.total_m3
    pre = [v for k, v in daily.items() if k < corte]
    post = [v for k, v in daily.items() if k >= corte]
    avg_pre = sum(pre) / len(pre) if pre else 0.0
    avg_post = sum(post) / len(post) if post else 0.0
    pct = ((avg_post / avg_pre) - 1.0) * 100.0 if avg_pre > 0 else 0.0
    return avg_pre, avg_post, pct


def _texto_analisis_l07() -> str:
    p = _pct_mam()
    names = _nombres_mam()
    avg_pre, avg_post, pct_inc = _placa_promedios_jun18()

    lines = ["Respecto del total monitoreado \x0b"]
    for nid in sorted(MAM_NODES, key=lambda x: -p.get(x, 0)):
        nm = names.get(nid, nid)
        lines.append(f"> {nm} = {format_number_chilean(p[nid], 1)}%\x0b")

    lines.append(
        "\nRespecto a Impulsión Falabella:\n"
        "Activa desde el 11/08/2026 (OC / cambio de equipo). "
        "Junio–julio = 0 m³. No comparar con la serie previa. "
        "Sin control nocturno; baseline 2–3 semanas."
    )
    lines.append(
        "\nRespecto a la auditoría (Placa Bancaria):\n"
        f"El 18/6/2026 el consumo diario promedio cambió "
        f"{format_number_chilean(pct_inc, 1)}% "
        f"(de {format_number_chilean(avg_pre, 1)} a {format_number_chilean(avg_post, 1)} m³/día). "
        "En agosto el volumen vuelve cerca del nivel pre-auditoría. "
        "La noche de Placa no se lee como fuga."
    )
    return "".join(lines)


def _texto_analisis_l08() -> str:
    fn = format_number_chilean
    p_mj = _pct_mam_may_jun()
    p_an = _pct_mam_anual()
    names = _nombres_mam()
    avg_pre, avg_post, pct_inc = _placa_promedios_jun18()
    order = sorted(MAM_NODES, key=lambda n: -p_mj.get(n, 0.0))
    lines = [f"Respecto al total monitoreable entre {PCT_TRIM_LABEL}"]
    for nid in order:
        lines.append(f"{names.get(nid, nid)} = {fn(p_mj[nid], 1)}%")
    lines.append("")
    lines.append("Respecto al total monitoreable anual")
    order_an = sorted(MAM_NODES, key=lambda n: -p_an.get(n, 0.0))
    for nid in order_an:
        lines.append(f"{names.get(nid, nid)} = {fn(p_an[nid], 1)}%")
    lines.extend([
        "",
        "Respecto a Impulsión Falabella:",
        "Activa desde el 11/08/2026. Junio–julio = 0 m³ (equipo fuera). "
        "Sin control nocturno; seguir 2–3 semanas para fijar baseline.",
        "",
        "Respecto a la auditoría (Placa Bancaria):",
        f"Desde el 18/6/2026 el promedio diario cambió {fn(pct_inc, 1)}% "
        f"(de {fn(avg_pre, 1)} a {fn(avg_post, 1)} m³/día). "
        "Agosto revierte cerca del nivel pre-auditoría. Noche no se lee como fuga.",
    ])
    return "\n".join(lines)


def _buscar_o_crear_caja_analisis_l08(slide, texto: str) -> None:
    l, t, w, h = L08_TEXTO_BOX
    target = None
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        tx = sh.text_frame.text.lower()
        if "respecto" in tx and ("monitore" in tx or "placa" in tx or "falabella" in tx):
            target = sh
            break
    if target is None:
        target = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    else:
        target.left = Inches(l)
        target.top = Inches(t)
        target.width = Inches(w)
        target.height = Inches(h)
    _set_texto_parrafos(target, texto)
    _aplicar_viñetas_panel_l04(target)
    _aplicar_negrita_solo_porcentajes(target)


def _actualizar_textos_l08(slide) -> None:
    _buscar_o_crear_caja_analisis_l08(slide, _texto_analisis_l08())
    print("[OK] Texto L08 — % mayo–jun, % anual, Falabella y auditoría Placa")

    _remover_notas_falabella(slide)
    _remover_captions_viejos(slide)
    periodo = PERIODO_CAP
    caps = [
        (RANK_CAP, LAYOUT_CAP_TOTAL),
        (f"Grafico: Consumo diario Ripley periodo {periodo}", L07_CAP_RIPLEY),
        (f"Grafico: Consumo diario Placa Bancaria periodo {periodo}", L07_CAP_PLACA),
        (f"Grafico: Consumo diario Impulsión Falabella periodo {periodo} (activa 11/08)", L07_CAP_FALABELLA),
    ]
    for txt, pos in caps:
        _agregar_caption_fijo(slide, txt, pos)
    print("[OK] L08 — ranking dual + diarios (estándar L04)")


def _texto_analisis_l11_maq() -> str:
    fn = format_number_chilean
    p_mj = _pct_maq_may_jun()
    p_an = _pct_maq_anual()
    names = _nombres_maq()
    order = sorted(MAQ_NODES, key=lambda n: -p_mj.get(n, 0.0))
    lines = [f"Respecto al total monitoreable entre {PCT_TRIM_LABEL}"]
    for nid in order:
        lines.append(f"{names.get(nid, nid)} = {fn(p_mj[nid], 1)}%")
    lines.append("")
    lines.append("Respecto al total monitoreable anual")
    order_an = sorted(MAQ_NODES, key=lambda n: -p_an.get(n, 0.0))
    for nid in order_an:
        lines.append(f"{names.get(nid, nid)} = {fn(p_an[nid], 1)}%")
    lines.extend([
        "",
        "No hay control nocturno en Matriz Principal: la madrugada SÍ entra al análisis.",
        "Pedido: on/off 00:00–08:00 (mismo criterio estanques MAE). ~24 m³/noche el 10/08.",
    ])
    return "\n".join(lines)


def _buscar_o_crear_caja_analisis_l11_maq(slide, texto: str) -> None:
    l, t, w, h = L11_MAQ_TEXTO_BOX
    target = None
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        tx = sh.text_frame.text.lower()
        if "respecto" in tx and ("monitore" in tx or "matriz" in tx or "baño" in tx or "bano" in tx):
            target = sh
            break
    if target is None:
        target = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    else:
        target.left = Inches(l)
        target.top = Inches(t)
        target.width = Inches(w)
        target.height = Inches(h)
    _set_texto_parrafos(target, texto)
    _aplicar_viñetas_panel_l04(target)
    _aplicar_negrita_solo_porcentajes(target)


def _actualizar_textos_l11_maq(slide) -> None:
    _buscar_o_crear_caja_analisis_l11_maq(slide, _texto_analisis_l11_maq())
    print("[OK] Texto L11 MAQ — % mayo–jun y % anual (estándar L04)")

    _remover_captions_viejos(slide)
    names = _nombres_maq()
    periodo = PERIODO_CAP
    caps = [
        (RANK_CAP, LAYOUT_CAP_TOTAL),
        (
            f"Grafico: Consumo diario {names.get('000025-13', 'Matriz Principal')} periodo {periodo}",
            LAYOUT_CAP_LEFT_MID,
        ),
        (
            f"Grafico: Consumo diario {names.get('000025-34', 'Alimentación Baños')} periodo {periodo}",
            LAYOUT_CAP_RIGHT_MID,
        ),
    ]
    for txt, pos in caps:
        _agregar_caption_fijo(slide, txt, pos)
    print("[OK] L11 MAQ — ranking dual + 2 diarios (estándar L04)")


def _texto_analisis_l14_bom() -> str:
    fn = format_number_chilean
    p_mj = _pct_bom_may_jun()
    p_an = _pct_bom_anual()
    names = _nombres_bom()
    order = sorted(BOM_NODES, key=lambda n: -p_mj.get(n, 0.0))
    lines = [f"Respecto al total monitoreable entre {PCT_TRIM_LABEL}"]
    for nid in order:
        lines.append(f"{names.get(nid, nid)} = {fn(p_mj[nid], 1)}%")
    lines.append("")
    lines.append("Respecto al total monitoreable anual")
    order_an = sorted(BOM_NODES, key=lambda n: -p_an.get(n, 0.0))
    for nid in order_an:
        lines.append(f"{names.get(nid, nid)} = {fn(p_an[nid], 1)}%")
    lines.extend([
        "",
        "San Ignacio 500: control nocturno desde el 16/07 FUNCIONA (noche ~42 → ~2 m³). Esa madrugada no se lee como fuga.",
        "Queda alza diurna desde el 26/06. San Ignacio 300: sin control; la noche SÍ entra al análisis (~51% el 10/08).",
    ])
    return "\n".join(lines)


def _buscar_o_crear_caja_analisis_l14_bom(slide, texto: str) -> None:
    l, t, w, h = L14_BOM_TEXTO_BOX
    target = None
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        tx = sh.text_frame.text.lower()
        if "respecto" in tx and ("monitore" in tx or "san ignacio" in tx):
            target = sh
            break
    if target is None:
        target = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    else:
        target.left = Inches(l)
        target.top = Inches(t)
        target.width = Inches(w)
        target.height = Inches(h)
    _set_texto_parrafos(target, texto)
    _aplicar_viñetas_panel_l04(target)
    _aplicar_negrita_solo_porcentajes(target)


def _actualizar_textos_l14_bom(slide) -> None:
    _buscar_o_crear_caja_analisis_l14_bom(slide, _texto_analisis_l14_bom())
    print("[OK] Texto L14 BOM — % mayo–jun y % anual (estándar L04)")

    _remover_captions_viejos(slide)
    names = _nombres_bom()
    periodo = L14_BOM_PERIODO_CAP
    caps = [
        (RANK_CAP, LAYOUT_CAP_TOTAL),
        (
            f"Grafico: Consumo diario {names.get(BOM_NODE_300, 'San Ignacio 300')} periodo {periodo}",
            LAYOUT_CAP_LEFT_MID,
        ),
        (
            f"Grafico: Consumo diario {names.get(BOM_NODE_500, 'San Ignacio 500')} periodo {periodo}",
            LAYOUT_CAP_RIGHT_MID,
        ),
    ]
    for txt, pos in caps:
        _agregar_caption_fijo(slide, txt, pos)
    print("[OK] L14 BOM — ranking dual + 2 diarios (estándar L04)")


def _texto_analisis_l17_aeb() -> str:
    fn = format_number_chilean
    p_mj = _pct_aeb_may_jun()
    p_an = _pct_aeb_anual()
    names = _nombres_aeb()
    order = sorted(AEB_NODES, key=lambda n: -p_mj.get(n, 0.0))
    lines = [f"Respecto al total monitoreable entre {PCT_TRIM_LABEL}"]
    for nid in order:
        lines.append(f"{names.get(nid, nid)} = {fn(p_mj[nid], 1)}%")
    lines.append("")
    lines.append("Respecto al total monitoreable anual")
    order_an = sorted(AEB_NODES, key=lambda n: -p_an.get(n, 0.0))
    for nid in order_an:
        lines.append(f"{names.get(nid, nid)} = {fn(p_an[nid], 1)}%")
    lines.extend([
        "",
        "Puntos activos: Matriz 1° piso (000025-11) y Anillo Plaza (000025-12). Matriz A.A. = 0 m³.",
        "No hay control nocturno: la madrugada SÍ entra al análisis. Pedir on/off 00:00–08:00 en Matriz.",
    ])
    return "\n".join(lines)


def _buscar_o_crear_caja_analisis_l17_aeb(slide, texto: str) -> None:
    l, t, w, h = L17_AEB_TEXTO_BOX
    target = None
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        tx = sh.text_frame.text.lower()
        if (
            "respecto" in tx and ("monitore" in tx or "matriz" in tx or "anillo" in tx)
            or "consumo diario promedio" in tx
        ):
            target = sh
            break
    if target is None:
        target = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    else:
        target.left = Inches(l)
        target.top = Inches(t)
        target.width = Inches(w)
        target.height = Inches(h)
    _set_texto_parrafos(target, texto)
    _aplicar_viñetas_panel_l04(target)
    _aplicar_negrita_solo_porcentajes(target)


def _actualizar_textos_l17_aeb(slide) -> None:
    _buscar_o_crear_caja_analisis_l17_aeb(slide, _texto_analisis_l17_aeb())
    print("[OK] Texto L17 AEB — % mayo–jun y % anual (estándar L04)")

    _remover_captions_viejos(slide)
    names = _nombres_aeb()
    periodo = PERIODO_CAP
    caps = [
        (RANK_CAP, LAYOUT_CAP_TOTAL),
        (
            f"Grafico: Consumo diario {names.get(AEB_NODE_MATRIZ, 'Matriz Principal')} periodo {periodo}",
            LAYOUT_CAP_LEFT_MID,
        ),
        (
            f"Grafico: Consumo diario {names.get(AEB_NODE_ANILLO, 'Anillo (Primer Piso)')} periodo {periodo}",
            LAYOUT_CAP_RIGHT_MID,
        ),
    ]
    for txt, pos in caps:
        _agregar_caption_fijo(slide, txt, pos)
    print("[OK] L17 AEB — ranking dual + 2 diarios (estándar L04)")


def _texto_matriz_l18_aeb() -> str:
    fn = format_number_chilean
    prom_noche = _promedio_nocturno_diario(AEB_NODE_MATRIZ)
    prom_h = _promedio_nocturno_h(AEB_NODE_MATRIZ)
    s_ref = _stats_horario_dia(AEB_NODE_MATRIZ, L18_MATRIZ_REF) or {}
    s_peak = _stats_horario_dia(AEB_NODE_MATRIZ, L18_MATRIZ_PEAK) or {}
    return "\n".join([
        "Matriz Principal — consumo nocturno (0–8 h)",
        f"Promedio del período: {fn(prom_noche, 1)} m³/noche ({fn(prom_h, 2)} m³/h); "
        "consumo base sostenido en horario inhábil.",
        f"Lun 8/6: nocturno = {fn(s_ref.get('night_sum', 0), 1)} m³ "
        f"(pico {fn(s_ref.get('night_max', 0), 2)} m³/h).",
        f"Sáb 27/6: nocturno = {fn(s_peak.get('night_sum', 0), 1)} m³ "
        f"(pico {fn(s_peak.get('night_max', 0), 2)} m³/h).",
        "",
        "Oportunidad:",
        "Implementar control on/off en horario inhábil (0–8 h), "
        "replicando la gestión de estanques MAE.",
    ])


def _texto_anillo_l18_aeb() -> str:
    fn = format_number_chilean
    prom_noche = _promedio_nocturno_diario(AEB_NODE_ANILLO)
    prom_h = _promedio_nocturno_h(AEB_NODE_ANILLO)
    ratio = _ratio_anillo_sobre_matriz_aeb()
    s_ref = _stats_horario_dia(AEB_NODE_ANILLO, L18_ANILLO_REF) or {}
    s_peak = _stats_horario_dia(AEB_NODE_ANILLO, L18_ANILLO_PEAK) or {}
    pct_noche_peak = (
        float(s_peak.get("night_sum", 0)) / float(s_peak.get("day_sum", 1)) * 100.0
        if s_peak.get("day_sum")
        else 0.0
    )
    return "\n".join([
        "Anillo (Primer Piso) — consumo nocturno (0–8 h)",
        f"Promedio del período: {fn(prom_noche, 1)} m³/noche ({fn(prom_h, 2)} m³/h).",
        f"Representa {fn(ratio, 0)}% del consumo de Matriz Principal en el período.",
        f"Lun 8/6: nocturno = {fn(s_ref.get('night_sum', 0), 1)} m³; "
        f"Sáb 27/6: nocturno = {fn(s_peak.get('night_sum', 0), 1)} m³ "
        f"({fn(pct_noche_peak, 0)}% del día).",
        "",
        "Oportunidad:",
        "Revisar equipos y llaves del anillo primer piso para reducir "
        "el consumo base nocturno.",
    ])


def _texto_analisis_l18_aeb() -> str:
    return _texto_matriz_l18_aeb() + "\n\n" + _texto_anillo_l18_aeb()


def _encontrar_caja_narrativa_aeb(slide):
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        tx = sh.text_frame.text.lower()
        if (
            "matriz principal" in tx
            or "anillo" in tx
            or "matriz aa" in tx
            or "consumo nocturno" in tx
        ):
            return sh
    return None


def _buscar_o_crear_narrativa_l18_aeb(slide, texto: str) -> None:
    target = _encontrar_caja_narrativa_aeb(slide)
    l, t, w, h = L18_AEB_TEXTO_BOX
    if target is None:
        target = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    else:
        target.left = Inches(l)
        target.top = Inches(t)
        target.width = Inches(w)
        target.height = Inches(h)
    _set_texto_parrafos(target, texto)
    _aplicar_viñetas_panel_l04(target)
    _aplicar_negrita_solo_porcentajes(target)


def _actualizar_textos_l18_aeb(slide) -> None:
    _buscar_o_crear_narrativa_l18_aeb(slide, _texto_analisis_l18_aeb())
    print("[OK] Narrativa L18 AEB — comparativo nocturno Matriz + Anillo")

    _remover_captions_viejos(slide)
    names = _nombres_aeb()
    nm_matriz = names.get(AEB_NODE_MATRIZ, "Matriz Principal")
    nm_anillo = names.get(AEB_NODE_ANILLO, "Anillo (Primer Piso)")
    caps = [
        (
            f"Grafico: Consumos horario {nm_matriz} - {DIAS_CAP[L18_MATRIZ_REF.weekday()]} "
            f"{L18_MATRIZ_REF.day}/{L18_MATRIZ_REF.month}",
            L18_CAP_MATRIZ_TOP,
        ),
        (
            f"Grafico: Consumos horario {nm_matriz} - {DIAS_CAP[L18_MATRIZ_PEAK.weekday()]} "
            f"{L18_MATRIZ_PEAK.day}/{L18_MATRIZ_PEAK.month}",
            L18_CAP_MATRIZ_BOT,
        ),
        (
            f"Grafico: Consumos horario {nm_anillo} - {DIAS_CAP[L18_ANILLO_REF.weekday()]} "
            f"{L18_ANILLO_REF.day}/{L18_ANILLO_REF.month}",
            L18_CAP_ANILLO_TOP,
        ),
        (
            f"Grafico: Consumos horario {nm_anillo} - {DIAS_CAP[L18_ANILLO_PEAK.weekday()]} "
            f"{L18_ANILLO_PEAK.day}/{L18_ANILLO_PEAK.month}",
            L18_CAP_ANILLO_BOT,
        ),
    ]
    for txt, pos in caps:
        _agregar_caption_fijo(slide, txt, pos)
    print("[OK] L18 AEB — 4 perfiles horarios comparativos")


def _texto_portada_cur() -> str:
    return "\n".join([
        "Puntos intervenidos: 2",
        "Anillo Norte",
        "Anillo Sur",
        "",
        "Recepción trabajos: 27/11/2025",
        "Fecha capacitación: 12/12/2025",
        "Reconfiguración de red: monitoreo actual en Anillo Norte y Anillo Sur",
        "(puntos anteriores Matriz Principal / Baños no comparables con el período actual).",
        "",
        "Usuarios habilitados:",
        "Joceline Lazo: Jefe de Operaciones.",
        "Equipo: medioambiente.dcl@parauco.com",
        "Tamara Martínez: Supervisora Linkes.",
        "Constanza Vilches: Analista Ambiental.",
    ])


def _actualizar_textos_l19_cur(slide) -> None:
    target = None
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        tx = sh.text_frame.text.lower()
        if "puntos intervenidos" in tx or "matriz principal" in tx or "anillo norte" in tx:
            target = sh
            break
    if target is None:
        print("[WARN] No se encontró caja de instalación CUR")
        return
    _set_texto_parrafos(target, _texto_portada_cur())
    print("[OK] Texto L19 CUR portada — Anillo Norte / Anillo Sur")


def _texto_analisis_l20_cur(*, factura_m3: float | None = None) -> str:
    fn = format_number_chilean
    p = _pct_cur_periodo()
    names = _nombres_cur()
    total_wes = _total_wes_cur_periodo()
    lines = [
        f"Período analizado: {L20_CUR_PERIODO_CAP}",
        "Monitoreo actual: Anillo Norte y Anillo Sur.",
        "Sin comparación con enero–marzo (reconfiguración de la red).",
        "",
        "Respecto al total monitoreado en el período",
    ]
    for nid in sorted(CUR_NODES, key=lambda n: -p.get(n, 0.0)):
        lines.append(f"{names.get(nid, nid)} = {fn(p[nid], 1)}%")
    lines.append("")
    lines.append(f"Consumo WES agregado: {fn(total_wes, 1)} m³")
    lines.append("")
    lines.append(
        f"Cuadre histórico vs boleta Esval N° {L20_CUR_BOLETA_ESVAL} "
        f"(lecturas 18/5–16/6, {fn(L20_CUR_FACTURA_M3, 0)} m³): la diferencia se explicó "
        "por desfase de lecturas (WES 24 h vs corte físico ~12:00), no por pérdida."
    )
    lines.append(
        "Pedir la(s) boleta(s) Esval de junio–agosto para repetir el cuadre WES vs cuenta "
        "en este período."
    )
    return "\n".join(lines)


def _buscar_o_crear_caja_analisis_l20_cur(slide, texto: str) -> None:
    l, t, w, h = L20_CUR_TEXTO_BOX
    target = None
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        tx = sh.text_frame.text.lower()
        if (
            "respecto" in tx
            or "período analizado" in tx
            or "periodo analizado" in tx
            or "cuenta de agua" in tx
            or "anillo" in tx
        ):
            target = sh
            break
    if target is None:
        target = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    else:
        target.left = Inches(l)
        target.top = Inches(t)
        target.width = Inches(w)
        target.height = Inches(h)
    _set_texto_parrafos(target, texto)
    _aplicar_viñetas_panel_l04(target)
    _aplicar_negrita_solo_porcentajes(target)


def _actualizar_textos_l20_cur(slide, *, factura_m3: float | None = None) -> None:
    _buscar_o_crear_caja_analisis_l20_cur(slide, _texto_analisis_l20_cur(factura_m3=factura_m3))
    print("[OK] Texto L20 CUR — período único + comparación WES vs cuenta")

    _remover_captions_viejos(slide)
    names = _nombres_cur()
    caps = [
        (f"Grafico: Consumo total periodo {L20_CUR_PERIODO_CAP}", LAYOUT_CAP_TOTAL),
        (
            f"Grafico: Consumo diario {names.get(CUR_NODE_NORTE, 'Anillo Norte')} "
            f"periodo {L20_CUR_PERIODO_CAP}",
            LAYOUT_CAP_LEFT_MID,
        ),
        (
            f"Grafico: Consumo diario {names.get(CUR_NODE_SUR, 'Anillo Sur')} "
            f"periodo {L20_CUR_PERIODO_CAP}",
            LAYOUT_CAP_RIGHT_MID,
        ),
    ]
    for txt, pos in caps:
        _agregar_caption_fijo(slide, txt, pos)
    print("[OK] L20 CUR — ranking período + 2 diarios (sin dual Ene–Mar)")


def _texto_resumen_l23_pak() -> str:
    fn = format_number_chilean
    datos = _datos_pak()
    totales = _totales_nodos(datos, PAK_NODES)
    total = sum(totales.values())
    p_mj = _pct_pak_may_jun()
    names = _nombres_pak()
    top3 = sorted(PAK_NODES, key=lambda n: -totales.get(n, 0.0))[:3]
    parts = [
        f"Total período {DESDE}–{HASTA}: {fn(total, 1)} m³ en 10 puntos.",
        "Mayores consumos (may–jun): "
        + ", ".join(f"{names.get(n, n)} {fn(p_mj[n], 1)}%" for n in top3),
        "Cadena DL, correlación y patrón nocturno en láminas siguientes.",
    ]
    return "\n".join(parts)


def _texto_analisis_l23_pak() -> str:
    return _texto_resumen_l23_pak()


def _actualizar_textos_l23_pak(slide) -> None:
    """L23 — ranking a página completa + resumen informativo inferior."""
    for sh in list(slide.shapes):
        if not sh.has_text_frame or _es_titulo_slide_shape(sh):
            continue
        tx = sh.text_frame.text.strip().lower()
        if tx and ("respecto" in tx or "monitore" in tx or "cadena" in tx or "total período" in tx):
            sh._element.getparent().remove(sh._element)
    l, t, w, h = L23_PAK_INFO_BOX
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    _set_texto_parrafos(box, _texto_resumen_l23_pak())
    for par in box.text_frame.paragraphs:
        for run in par.runs:
            run.font.size = Pt(8)
    _remover_captions_viejos(slide)
    _agregar_caption_fijo(
        slide,
        RANK_CAP,
        (0.272, 6.280, 12.500, 0.280),
    )
    print("[OK] L23 PAK — ranking dual + resumen período")


def _texto_analisis_l24_pak() -> str:
    fn = format_number_chilean
    tot = _totales_pak_chain()
    t22 = tot.get(PAK_CHAIN_SOURCE_ANTIGUA, 0.0)
    t28 = tot.get(PAK_CHAIN_SOURCE_NUEVA, 0.0)
    t27 = tot.get(PAK_CHAIN_DISTRITO, 0.0)
    t35 = tot.get(PAK_CHAIN_BAZAR, 0.0)
    t36 = tot.get(PAK_CHAIN_KENNEDY, 0.0)
    diff_dl_ant = t27 - t22
    pct_dl_ant = (diff_dl_ant / t22 * 100.0) if t22 else 0.0
    down = t35 + t36
    diff_int = t27 - down
    diff_fuentes = (t22 + t28) - t27
    r22 = _corr_diaria_pak(PAK_CHAIN_SOURCE_ANTIGUA, PAK_CHAIN_DISTRITO)
    r28 = _corr_diaria_pak(PAK_CHAIN_SOURCE_NUEVA, PAK_CHAIN_DISTRITO)
    p_mj = _pct_pak_may_jun()
    p_an = _pct_pak_anual()
    datos = _datos_pak()
    tot_pak = sum(_totales_nodos(datos, PAK_NODES).values())
    names = _nombres_pak()
    lines = [
        "Cadena declarada:",
        f"• {names.get(PAK_CHAIN_SOURCE_ANTIGUA, '22')} alimenta "
        f"{names.get(PAK_CHAIN_DISTRITO, '27')}",
        f"• {names.get(PAK_CHAIN_SOURCE_NUEVA, '28')} también surte al DL",
        f"• DL se reparte en {names.get(PAK_CHAIN_BAZAR, '35')} y "
        f"{names.get(PAK_CHAIN_KENNEDY, '36')}",
        "",
        f"Período {DESDE} a {HASTA} (total PAK {fn(tot_pak, 1)} m³):",
        f"• DL ({fn(t27, 1)} m³) supera a Sandía Antigua ({fn(t22, 1)} m³) "
        f"en +{fn(diff_dl_ant, 1)} m³ (+{fn(pct_dl_ant, 1)}%). "
        "El distrito recibe aporte adicional desde Sandía Nueva "
        f"({fn(t28, 1)} m³).",
        f"• Correlación diaria Antigua↔DL: r={r22:.2f}. Nueva↔DL: r={r28:.2f}.",
        f"• Aguas abajo 35+36: {fn(down, 1)} m³ "
        f"({fn(down / t27 * 100 if t27 else 0, 1)}% del DL). "
        f"Consumo interno DL: {fn(diff_int, 1)} m³.",
        f"• Suma fuentes 22+28 vs DL: excedente {fn(diff_fuentes, 1)} m³.",
        "",
        f"Distribución % {PCT_TRIM_LABEL} (top 5):",
    ]
    for nid in sorted(PAK_NODES, key=lambda n: -p_mj.get(n, 0.0))[:5]:
        lines.append(f"• {names.get(nid, nid)} = {fn(p_mj[nid], 1)}%")
    lines.append("")
    lines.append("Distribución % anual (top 5):")
    for nid in sorted(PAK_NODES, key=lambda n: -p_an.get(n, 0.0))[:5]:
        lines.append(f"• {names.get(nid, nid)} = {fn(p_an[nid], 1)}%")
    lines.append("")
    lines.append(f"Consumo nocturno 0–{PAK_NIGHT_END} h (cadena):")
    for nid in PAK_CHAIN_NODES:
        s = _stats_nocturno_periodo(nid, DESDE, HASTA, night_end=PAK_NIGHT_END)
        pct_n = _pct_nocturno_periodo(nid)
        lines.append(
            f"• {names.get(nid, nid)}: {fn(s['sum_prom'], 1)} m³/noche prom. "
            f"({fn(pct_n, 1)}% del total)"
        )
    lines.append("Detalle en lámina PATRÓN NOCTURNO.")
    return "\n".join(lines)


def _buscar_o_crear_caja_analisis_l24_pak(slide, texto: str) -> None:
    l, t, w, h = L24_PAK_TEXTO_BOX
    target = None
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        tx = sh.text_frame.text.lower()
        if "cadena declarada" in tx or ("respecto" in tx and "monitore" in tx):
            target = sh
            break
    if target is None:
        target = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    else:
        target.left = Inches(l)
        target.top = Inches(t)
        target.width = Inches(w)
        target.height = Inches(h)
    _set_texto_parrafos(target, texto)
    for par in target.text_frame.paragraphs:
        for run in par.runs:
            if run.font.size is None or run.font.size > Pt(9):
                run.font.size = Pt(7)


def _actualizar_textos_l24_pak(slide) -> None:
    _buscar_o_crear_caja_analisis_l24_pak(slide, _texto_analisis_l24_pak())
    print("[OK] Texto L24 PAK — cadena abastecimiento DL")

    _remover_captions_viejos(slide)
    periodo = PERIODO_CAP
    caps = [
        (f"Grafico: Totales cadena abastecimiento período {periodo}", LAYOUT_CAP_TOTAL),
        (
            f"Grafico: Dispersión consumo diario vs DL — correlación {periodo}",
            (0.350, 3.750, 6.400, 0.280),
        ),
    ]
    for txt, pos in caps:
        _agregar_caption_fijo(slide, txt, pos)
    print("[OK] L24 PAK — barras cadena + scatter correlación + narrativa")


def _actualizar_textos_l25_pak(slide, nodes: List[str]) -> None:
    """No modifica título de lámina; limpia captions viejos. Títulos van en el gráfico."""
    _remover_captions_viejos(slide)
    print(f"[OK] L25 PAK — {len(nodes)} diarios (título lámina conservado)")


def _actualizar_textos_l26_pak(slide, nodes: List[str]) -> None:
    """No modifica título de lámina; limpia captions viejos. Títulos van en el gráfico."""
    _remover_captions_viejos(slide)
    print(f"[OK] L26 PAK — {len(nodes)} diarios (título lámina conservado)")


def _texto_analisis_l27_pak(peak_day: date) -> str:
    fn = format_number_chilean
    names = _nombres_pak()
    lines = [
        f"PAK Kennedy — patrón nocturno (0–{PAK_NIGHT_END} h)",
        f"Período {DESDE} a {HASTA}.",
        "",
        "Resumen cadena DL:",
    ]
    for nid in PAK_CHAIN_NODES:
        s = _stats_nocturno_periodo(nid, DESDE, HASTA, night_end=PAK_NIGHT_END)
        acum = _nocturno_acumulado_periodo(nid)
        pct_n = _pct_nocturno_periodo(nid)
        lines.append(
            f"• {names.get(nid, nid)}: {fn(acum, 1)} m³ acum. nocturno "
            f"({fn(pct_n, 1)}% del total); prom. {fn(s['sum_prom'], 1)} m³/noche, "
            f"pico {fn(s['pico_max'], 2)} m³/h."
        )

    s_dl_ref = _stats_horario_dia(PAK_CHAIN_DISTRITO, L27_PAK_REF) or {}
    s_dl_peak = _stats_horario_dia(
        PAK_CHAIN_DISTRITO, datetime.combine(peak_day, datetime.min.time()),
    ) or {}
    prom_dl = _promedio_nocturno_diario(PAK_CHAIN_DISTRITO, night_end=PAK_NIGHT_END)
    prom_h_dl = _promedio_nocturno_h(PAK_CHAIN_DISTRITO, night_end=PAK_NIGHT_END)

    lines.extend([
        "",
        f"{names.get(PAK_CHAIN_DISTRITO, 'DL')} — análisis:",
        f"Promedio período: {fn(prom_dl, 1)} m³/noche ({fn(prom_h_dl, 2)} m³/h).",
        f"Lun {L27_PAK_REF.strftime('%d/%m')}: nocturno = "
        f"{fn(s_dl_ref.get('night_sum', 0), 1)} m³ "
        f"(pico {fn(s_dl_ref.get('night_max', 0), 2)} m³/h).",
        f"Día pico nocturno ({DIAS_ES[peak_day.weekday()]} {peak_day.strftime('%d/%m')}): "
        f"{fn(s_dl_peak.get('night_sum', 0), 1)} m³ "
        f"(pico {fn(s_dl_peak.get('night_max', 0), 2)} m³/h).",
        "",
        "Hallazgos:",
        "• El DL concentra el mayor consumo nocturno de la cadena.",
        "• Sandía Antigua (bomba) muestra operación nocturna asociada al abastecimiento.",
        "• Bazar Gourmet y Sandía Nueva aportan menor fracción nocturna relativa.",
        "",
        "Oportunidad:",
        "Evaluar control on/off en horario inhábil (0–8 h) en DL y salas de bomba, "
        "alineado a la gestión MAE/AEB.",
    ])
    return "\n".join(lines)


def _buscar_o_crear_caja_analisis_l27_pak(slide, texto: str) -> None:
    l, t, w, h = L27_PAK_TEXTO_BOX
    target = None
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        tx = sh.text_frame.text.lower()
        if "patrón nocturno" in tx or "patron nocturno" in tx or "pak kennedy" in tx:
            target = sh
            break
    if target is None:
        target = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    else:
        target.left = Inches(l)
        target.top = Inches(t)
        target.width = Inches(w)
        target.height = Inches(h)
    _set_texto_parrafos(target, texto)
    for par in target.text_frame.paragraphs:
        for run in par.runs:
            if run.font.size is None or run.font.size > Pt(9):
                run.font.size = Pt(7)


def _actualizar_textos_l27_pak(slide, peak_day: date) -> None:
    _buscar_o_crear_caja_analisis_l27_pak(slide, _texto_analisis_l27_pak(peak_day))
    print("[OK] Texto L27 PAK — patrón nocturno cadena DL")

    _remover_captions_viejos(slide)
    names = _nombres_pak()
    periodo = PERIODO_CAP
    caps = [
        (
            f"Grafico: Ranking consumo nocturno 0–{PAK_NIGHT_END} h — cadena DL — {periodo}",
            LAYOUT_CAP_TOTAL,
        ),
        (
            f"Grafico: Perfil horario {names.get(PAK_CHAIN_DISTRITO, 'DL')} "
            f"lun {L27_PAK_REF.strftime('%d/%m/%Y')}",
            (0.350, 3.850, 3.100, 0.280),
        ),
        (
            f"Grafico: Perfil horario {names.get(PAK_CHAIN_DISTRITO, 'DL')} "
            f"día pico nocturno {peak_day.strftime('%d/%m/%Y')}",
            (3.600, 3.850, 3.100, 0.280),
        ),
    ]
    for txt, pos in caps:
        _agregar_caption_fijo(slide, txt, pos)
    print("[OK] L27 PAK — ranking nocturno + 2 perfiles horarios DL + narrativa")


def _poblar_diarios_pak(slide, nodes: List[str], tag: str) -> None:
    names = _nombres_pak()
    slots = _pak_daily_row_slots(len(nodes))
    for node_id, (sl, st, sw, sh) in zip(nodes, slots):
        png = CHARTS / f"{tag}_pak_{node_id}_diario_linea.png"
        nm = names.get(node_id, node_id)
        chart_diario_linea(node_id, png, sw, sh, compact=False, titulo=nm)
        _poner_grafico_fresco(slide, png, sl, st, sw, sh)
        print(f"[OK] DIARIO {node_id} ({nm}) @ L={sl:.2f} T={st:.2f} W={sw:.2f} H={sh:.2f}")


def _texto_analisis_l10() -> str:
    p = _pct_maq()
    names = _nombres_maq()
    lines = ["Respecto del total monitoreado \x0b"]
    for nid in sorted(MAQ_NODES, key=lambda x: -p.get(x, 0)):
        nm = names.get(nid, nid)
        lines.append(f"> {nm} = {format_number_chilean(p[nid], 1)}%\x0b")
    return "".join(lines)


def _buscar_o_crear_caja_analisis_maq(slide, texto: str, box=L10_TEXTO_BOX) -> None:
    """Panel % MAQ — reemplaza caja narrativa antigua si existe."""
    target = None
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        t = sh.text_frame.text.lower()
        if (
            "respecto del total monitoreado" in t
            or "matriz principal" in t
            or "consumo de ba" in t
            or "consumo diario promedio" in t
        ):
            target = sh
            break
    l, t, w, h = box
    if target is None:
        target = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        print("[OK] Caja de análisis MAQ creada (columna derecha)")
    else:
        target.left = Inches(l)
        target.top = Inches(t)
        target.width = Inches(w)
        target.height = Inches(h)
    _set_texto_caja(target, texto)


def _actualizar_textos_l10(slide) -> None:
    """Conserva panel de análisis; negrita solo en %; viñetas estándar."""
    periodo = PERIODO_CAP
    _conservar_caja_analisis(slide, L10_TEXTO_BOX, "Respecto del total monitoreado")
    print("[OK] Panel análisis L10 conservado (negrita en % estandarizada)")

    _remover_captions_viejos(slide)
    names = _nombres_maq()
    caps = [
        (f"Grafico: Consumo total periodo {periodo}", LAYOUT_CAP_TOTAL),
        (
            f"Grafico: Consumo diario {names.get('000025-13', 'Matriz Principal')} periodo {periodo}",
            LAYOUT_CAP_LEFT_MID,
        ),
        (
            f"Grafico: Consumo diario {names.get('000025-34', 'Baños')} periodo {periodo}",
            LAYOUT_CAP_RIGHT_MID,
        ),
    ]
    for txt, pos in caps:
        _agregar_caption_fijo(slide, txt, pos)
    print("[OK] L10 — ranking horizontal + viñetas (formato L04/L07)")


def _buscar_o_crear_caja_analisis(slide, texto: str, box=L07_TEXTO_BOX) -> None:
    """Solo actualiza el texto del panel de %; no cambia tipografía ni posición."""
    target = None
    for sh in slide.shapes:
        if sh.has_text_frame and "Respecto del total monitoreado" in sh.text_frame.text:
            target = sh
            break
    if target is None:
        l, t, w, h = box
        target = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        print("[OK] Caja de análisis creada (columna derecha)")
    else:
        l, t, w, h = box
        target.left = Inches(l)
        target.top = Inches(t)
        target.width = Inches(w)
        target.height = Inches(h)
    _set_texto_caja(target, texto)


def _set_texto_parrafos(shape, texto: str) -> None:
    """Un párrafo por línea; evita \\x0b (_x000B_) y prefijos > visibles."""
    texto = texto.replace("\x0b", "\n").replace("_x000B_", "\n")
    lineas = [ln.lstrip("> ").strip() if ln.strip() else "" for ln in texto.split("\n")]
    tf = shape.text_frame
    tf.word_wrap = True
    while len(tf.paragraphs) < len(lineas):
        tf.add_paragraph()
    for i, ln in enumerate(lineas):
        tf.paragraphs[i].text = ln
    for par in tf.paragraphs[len(lineas) :]:
        par.text = ""


def _set_texto_caja(shape, texto: str) -> None:
    tf = shape.text_frame
    tf.word_wrap = True
    if tf.paragraphs:
        tf.paragraphs[0].text = texto
        for par in tf.paragraphs[1:]:
            par.text = ""


def _estilo_caption(shape, *, font_name: str = "", font_size: int = 10) -> None:
    """Títulos bajo gráficos — negrita cursiva centrada (10 pt estándar)."""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    for para in tf.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        para.space_before = Pt(0)
        para.space_after = Pt(0)
        for run in para.runs:
            run.font.size = Pt(font_size)
            run.font.bold = True
            run.font.italic = True
            if font_name:
                run.font.name = font_name


def _aplicar_negrita_solo_porcentajes(shape) -> None:
    """Negrita únicamente en valores con % (estándar L07/L10)."""
    for para in shape.text_frame.paragraphs:
        text = para.text
        if "%" not in text:
            continue
        partes = re.split(r"([\d.,]+%)", text)
        para.text = ""
        for parte in partes:
            if not parte:
                continue
            run = para.add_run()
            run.text = parte
            if parte.endswith("%"):
                run.font.bold = True


def _es_encabezado_viñeta_l04(texto: str) -> bool:
    """Encabezados con viñeta circular (estándar panel L04 / narrativas horarios)."""
    t = texto.strip()
    if not t:
        return False
    if t.startswith("Respecto") or t.startswith("Recomendación"):
        return True
    if t.startswith("En ") and "=" not in t:
        return True
    if "— consumo nocturno" in t:
        return True
    return False


def _poner_viñeta_parrafo(para, char: str = "\uf0b7") -> None:
    pPr = para._p.get_or_add_pPr()
    for tag in ("a:buChar", "a:buAutoNum", "a:buNone"):
        el = pPr.find(qn(tag))
        if el is not None:
            pPr.remove(el)
    bu = OxmlElement("a:buChar")
    bu.set("char", char)
    pPr.insert(0, bu)


def _quitar_viñeta_parrafo(para) -> None:
    pPr = para._p.get_or_add_pPr()
    for tag in ("a:buChar", "a:buAutoNum"):
        el = pPr.find(qn(tag))
        if el is not None:
            pPr.remove(el)
    if pPr.find(qn("a:buNone")) is None:
        pPr.insert(0, OxmlElement("a:buNone"))


def _aplicar_viñetas_panel_l04(shape) -> None:
    """Viñetas en encabezados de sección; sub-ítems (líneas con =) sin viñeta."""
    for para in shape.text_frame.paragraphs:
        if _es_encabezado_viñeta_l04(para.text):
            _poner_viñeta_parrafo(para)
        else:
            _quitar_viñeta_parrafo(para)


def _conservar_caja_analisis(slide, box: Tuple[float, float, float, float], marcador: str) -> None:
    """Reposiciona caja de análisis sin cambiar texto; aplica negrita en %."""
    l, t, w, h = box
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        if marcador not in sh.text_frame.text:
            continue
        sh.left = Inches(l)
        sh.top = Inches(t)
        sh.width = Inches(w)
        sh.height = Inches(h)
        _aplicar_negrita_solo_porcentajes(sh)
        return


def _caption_bajo_grafico(chart_l: float, chart_t: float, chart_w: float, chart_h: float) -> Tuple[float, float, float]:
    """Posición del título centrado justo debajo del gráfico."""
    return chart_l, chart_t + chart_h + CAPTION_GAP, chart_w


def _configurar_caption(
    shape,
    texto: str,
    left: float,
    top: float,
    width: float,
    height: float = CAPTION_H,
    *,
    font_name: str = "",
    font_size: int = 10,
) -> None:
    shape.left = Inches(left)
    shape.top = Inches(top)
    shape.width = Inches(width)
    shape.height = Inches(height)
    shape.text_frame.paragraphs[0].text = texto
    _estilo_caption(shape, font_name=font_name, font_size=font_size)


def _agregar_caption_fijo(
    slide,
    texto: str,
    layout: Tuple[float, float, float, float],
    *,
    font_name: str = "",
    font_size: int = 10,
) -> None:
    l, t, w, h = layout
    sh = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    _configurar_caption(sh, texto, l, t, w, h, font_name=font_name, font_size=font_size)


def _remover_captions_viejos(slide) -> None:
    """Quita captions 'Grafico:' previos para evitar duplicados al reformatear."""
    eliminar = []
    for sh in slide.shapes:
        if sh.has_text_frame and "grafico:" in sh.text_frame.text.lower():
            eliminar.append(sh._element)
    for el in eliminar:
        el.getparent().remove(el)
    if eliminar:
        print(f"[OK] {len(eliminar)} caption(s) antiguo(s) eliminado(s)")


def _actualizar_textos_l04(slide) -> None:
    narrativo = _texto_analisis_l04()
    _buscar_o_crear_caja_analisis_l04(slide, narrativo)
    print("[OK] Texto L04 — % mayo–jun y % anual")

    _remover_captions_viejos(slide)
    periodo = PERIODO_CAP
    caps = [
        (RANK_CAP, LAYOUT_CAP_TOTAL),
        (f"Grafico: Consumo diario Estanque Sur periodo {periodo}", LAYOUT_CAP_LEFT_MID),
        (f"Grafico: Consumo diario Baños Públicos periodo {periodo}", LAYOUT_CAP_RIGHT_MID),
        (f"Grafico: Consumo diario Estanque Norte Locales Mall periodo {periodo}", LAYOUT_CAP_LEFT_BOT),
        (f"Grafico: Consumo diario Pizza Hut periodo {periodo}", LAYOUT_CAP_RIGHT_BOT),
    ]
    for txt, pos in caps:
        _agregar_caption_fijo(slide, txt, pos)
    print("[OK] Títulos L04 — posiciones fijas (igual deck editado)")


def _remover_notas_falabella(slide) -> None:
    """Quita cajas de texto Falabella sueltas (evita duplicados al regenerar L07)."""
    eliminar = []
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        t = sh.text_frame.text.lower()
        if "falabella" in t and "respecto" not in t:
            eliminar.append(sh._element)
    for el in eliminar:
        el.getparent().remove(el)


def _actualizar_textos_l07(slide) -> None:
    """Panel % período + Falabella activa + auditoría Placa."""
    periodo = PERIODO_CAP
    _buscar_o_crear_caja_analisis_l08(slide, _texto_analisis_l07())
    print("[OK] Panel análisis L07 — % período, Falabella activa, Placa")

    _remover_notas_falabella(slide)
    _remover_captions_viejos(slide)
    caps = [
        (f"Grafico: Consumo total periodo {periodo}", LAYOUT_CAP_TOTAL),
        (f"Grafico: Consumo diario Ripley periodo {periodo}", L07_CAP_RIPLEY),
        (f"Grafico: Consumo diario Placa Bancaria periodo {periodo}", L07_CAP_PLACA),
        (f"Grafico: Consumo diario Impulsión Falabella periodo {periodo} (activa 11/08)", L07_CAP_FALABELLA),
    ]
    for txt, pos in caps:
        _agregar_caption_fijo(slide, txt, pos)
    print("[OK] L07 — layout + Falabella activa 11/08")


def _ajustar_eje_y(ax, values: List[float]) -> None:
    """Escala Y que refleje consumos bajos nocturnos (2–3 decimales como la app)."""
    mx = max(values) if values else 0.1
    noct = values[:7] if len(values) >= 7 else values
    mx_n = max(noct) if noct else 0.0

    if mx <= 1.2:
        ymax = max(0.5, round(mx * 1.15 + 0.05, 2))
        step = 0.1 if ymax <= 1.0 else 0.2
        ax.set_ylim(0, ymax)
        ax.yaxis.set_major_locator(MultipleLocator(step))
        ax.yaxis.set_major_formatter(FormatStrFormatter("%.2f"))
    elif mx <= 6.0:
        ymax = max(mx * 1.12, mx_n * 3 if mx_n > 0 else mx * 1.12)
        ax.set_ylim(0, ymax)
        ax.yaxis.set_major_locator(MultipleLocator(0.5 if ymax <= 3 else 1.0))
        ax.yaxis.set_major_formatter(FormatStrFormatter("%.2f"))
    else:
        ax.set_ylim(0, mx * 1.08)
        ax.yaxis.set_major_formatter(FormatStrFormatter("%.1f"))


def chart_horario_compact(
    node_id: str,
    day: datetime,
    out: Path,
    w_in: float,
    h_in: float,
    *,
    anotar: bool = False,
    techo_ref: float | None = None,
) -> Path:
    """Perfil horario compacto: título español, área nocturna achurada, escala precisa."""
    serie = get_hourly_measures_for_day(node_id, day)
    vals_map = {h: v for h, v in serie}
    hours = list(range(24))
    values = [vals_map.get(h, 0.0) for h in hours]

    fig, ax = plt.subplots(figsize=(max(4.2, w_in * 1.05), max(1.9, h_in * 1.05)))
    ax.plot(hours, values, color="#4A90E2", linewidth=1.5, marker="o", markersize=2.5)
    ax.fill_between(hours, values, 0, color="#87CEEB", alpha=0.4)

    noct_vals = [values[h] for h in range(7)]
    marcar_nocturno = any(v > 0 for v in noct_vals) or node_id == "000025-07"
    if marcar_nocturno:
        ax.fill_between(
            range(7), noct_vals, 0,
            color="#FFD700", alpha=0.3, hatch="///",
            edgecolor="#FFA500", linewidth=1.0,
        )
        ax.axvline(x=0, color="orange", linestyle="--", linewidth=1, alpha=0.7)
        ax.axvline(x=6, color="orange", linestyle="--", linewidth=1, alpha=0.7)

    titulo = f"{DIAS_ES[day.weekday()]} {day.strftime('%d/%m/%Y')}"
    ax.set_title(titulo, fontsize=8, fontweight="bold", pad=2)
    ax.set_xlabel("Hora del día", fontsize=7)
    ax.set_ylabel("m³/h", fontsize=7)
    ax.set_xlim(0, 23)
    ax.set_xticks(range(0, 24, 4))
    ax.tick_params(axis="x", labelsize=6)
    ax.tick_params(axis="y", labelsize=6)
    ax.grid(True, alpha=0.3, axis="y")
    _ajustar_eje_y(ax, values)

    if techo_ref is not None:
        ax.axhline(
            y=techo_ref, color="#C0504D", linestyle="--", linewidth=1.1, alpha=0.85,
        )
        ax.text(
            0.5, techo_ref, f"  Techo {format_number_chilean(techo_ref, 1)} m³/h",
            fontsize=6, color="#C0504D", va="bottom", ha="left",
        )
        ymax = ax.get_ylim()[1]
        if techo_ref > ymax * 0.85:
            ax.set_ylim(0, max(ymax, techo_ref * 1.12))

    if anotar and values:
        mx_n = max(noct_vals) if noct_vals else 0.0
        sm_n = sum(noct_vals)
        mx_d = max(values)
        notas: List[str] = []

        if techo_ref is not None:
            notas.append(f"Techo informe: {format_number_chilean(techo_ref, 1)} m³/h")
            notas.append(f"Pico medido 0-6h: {format_number_chilean(mx_n, 2)} m³/h")
            if mx_n > techo_ref:
                exc = mx_n - techo_ref
                pct = (exc / techo_ref * 100.0) if techo_ref else 0.0
                notas.append(
                    f"Sobre el techo: +{format_number_chilean(exc, 2)} m³/h "
                    f"({format_number_chilean(pct, 0)}%)"
                )
            else:
                notas.append(f"Bajo el techo: −{format_number_chilean(techo_ref - mx_n, 2)} m³/h")
            notas.append(f"Suma nocturna 0-6h: {format_number_chilean(sm_n, 1)} m³")
            # Marca hora y valor del pico nocturno en el gráfico
            if mx_n > 0:
                hora_pico = noct_vals.index(mx_n)
                ax.plot(hora_pico, mx_n, marker="o", markersize=5, color="#C0504D", zorder=5)
                ax.annotate(
                    f"{format_number_chilean(mx_n, 2)} m³/h",
                    xy=(hora_pico, mx_n),
                    xytext=(hora_pico + 1.2, mx_n + max(mx_n * 0.08, 0.15)),
                    fontsize=6,
                    color="#C0504D",
                    arrowprops=dict(arrowstyle="->", color="#C0504D", lw=0.8),
                )
        else:
            notas = [
                f"Pico 0-6h: {format_number_chilean(mx_n, 2)} m³/h",
                f"Suma 0-6h: {format_number_chilean(sm_n, 1)} m³",
                f"Pico día: {format_number_chilean(mx_d, 2)} m³/h",
            ]

        ax.text(
            0.98, 0.97, "\n".join(notas),
            transform=ax.transAxes, ha="right", va="top", fontsize=6,
            bbox=dict(boxstyle="round,pad=0.3", facecolor="white", alpha=0.85, edgecolor="#4A90E2"),
        )

    fig.subplots_adjust(left=0.14, right=0.97, top=0.86, bottom=0.20)
    out.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(out, dpi=200, facecolor="white", pad_inches=0.03)
    plt.close(fig)
    return out


def _promedio_horario_24h(node_id: str) -> Tuple[List[float], int]:
    """Promedio m³/h por hora (0–23) en el período del informe."""
    desde = parse_date(DESDE)
    hasta = parse_date(HASTA, end_of_day=True)
    sums = [0.0] * 24
    counts = [0] * 24
    dias = 0
    d = desde
    while d <= hasta:
        serie = get_hourly_measures_for_day(node_id, d)
        if not serie:
            d += timedelta(days=1)
            continue
        vals_map = {h: v for h, v in serie}
        if len(vals_map) < 20:
            d += timedelta(days=1)
            continue
        dias += 1
        for h in range(24):
            if h in vals_map:
                sums[h] += float(vals_map[h])
                counts[h] += 1
        d += timedelta(days=1)
    avgs = [sums[h] / counts[h] if counts[h] else 0.0 for h in range(24)]
    return avgs, dias


def _promedio_horario_rango(node_id: str, desde: str, hasta: str) -> Tuple[List[float], int]:
    """Promedio m³/h por hora (0–23) en un rango de fechas."""
    cache = _cargar_horario_cache(node_id, desde, hasta)
    d0 = parse_date(desde).date()
    d1 = parse_date(hasta, end_of_day=True).date()
    sums = [0.0] * 24
    counts = [0] * 24
    dias = 0
    for d, hmap in cache.items():
        if d < d0 or d > d1:
            continue
        dias += 1
        for h in range(24):
            if h in hmap:
                sums[h] += hmap[h]
                counts[h] += 1
    avgs = [sums[h] / counts[h] if counts[h] else 0.0 for h in range(24)]
    return avgs, dias


def chart_ranking_nocturno_mensual_s500(
    out: Path,
    w_in: float,
    h_in: float,
    *,
    desde: str = L16_RANK_DESDE,
    hasta: str = L16_RANK_HASTA,
    night_end: int = L16_NIGHT_END,
) -> Path:
    """Ranking mensual consumo nocturno 0–6 h — San Ignacio 500 (barras amarillas + ahorro control)."""
    names = _nombres_bom()
    nm = names.get(BOM_NODE_500, "San Ignacio 500")
    monthly = _suma_nocturna_por_mes(BOM_NODE_500, desde, hasta, night_end=night_end)
    d0 = parse_date(desde)
    d1 = parse_date(hasta, end_of_day=True)
    order: List[Tuple[int, int]] = []
    y, m = d0.year, d0.month
    while (y, m) <= (d1.year, d1.month):
        order.append((y, m))
        m += 1
        if m > 12:
            m = 1
            y += 1
    vals = [monthly.get(k, 0.0) for k in order]
    ahorros = vals  # control WES → consumo nocturno ~0 (100% recuperable)
    total_noc = sum(vals)
    total_ahorro = sum(ahorros)
    rank_order = sorted(range(len(order)), key=lambda i: -vals[i])
    labels = [f"{L16_MESES_CORTO[order[i][1] - 1]} {order[i][0]}" for i in rank_order]
    v_rank = [vals[i] for i in rank_order]
    a_rank = [ahorros[i] for i in rank_order]

    fig, ax = plt.subplots(figsize=(max(4.5, w_in * 1.08), max(2.4, h_in * 1.08)))
    y_pos = list(range(len(rank_order)))
    bars = ax.barh(
        y_pos, v_rank, height=0.52,
        color="#FFD700", edgecolor="#FFA500", linewidth=0.8,
        hatch="///", label="Consumo nocturno 0–6 h",
    )
    ax.set_yticks(y_pos)
    ax.set_yticklabels(labels, fontsize=7)
    ax.invert_yaxis()
    ax.set_xlabel("m³ (suma mensual 0–6 h)", fontsize=7)
    ax.tick_params(axis="x", labelsize=6)
    ax.grid(True, alpha=0.3, axis="x")
    mx = max(v_rank) if v_rank else 1.0
    ax.set_xlim(0, mx * 1.55)
    fn = format_number_chilean
    for bar, v, a in zip(bars, v_rank, a_rank):
        if v <= 0:
            continue
        ax.text(
            bar.get_width() + mx * 0.02,
            bar.get_y() + bar.get_height() / 2,
            f"{fn(v, 0)} m³  |  Ahorro: {fn(a, 0)} m³",
            va="center", ha="left", fontsize=6, color="#333333",
        )
    ax.set_title(
        f"{nm} — ranking nocturno por mes (0–{night_end} h)",
        fontsize=8, fontweight="bold", pad=3,
    )
    ax.text(
        0.98, 0.02,
        f"Total Ene–Jun: {fn(total_noc, 0)} m³  |  Ahorro si activa control WES: {fn(total_ahorro, 0)} m³",
        transform=ax.transAxes, ha="right", va="bottom", fontsize=6,
        bbox=dict(boxstyle="round,pad=0.3", facecolor="#FFF8DC", alpha=0.95, edgecolor="#FFA500"),
    )
    fig.subplots_adjust(left=0.22, right=0.97, top=0.90, bottom=0.14)
    out.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(out, dpi=200, facecolor="white", pad_inches=0.03)
    plt.close(fig)
    return out


def _cargar_horario_cache(
    node_id: str,
    desde: str,
    hasta: str,
    *,
    night_end: int = L16_NIGHT_END,
) -> Dict[date, Dict[int, float]]:
    """Cache diaria de medidas horarias (evita llamadas API repetidas en L16)."""
    key = f"H:{node_id}:{desde}:{hasta}"
    if key in _horario_cache:
        return _horario_cache[key]
    d0 = parse_date(desde)
    d1 = parse_date(hasta, end_of_day=True)
    out: Dict[date, Dict[int, float]] = {}
    d = d0
    while d <= d1:
        serie = get_hourly_measures_for_day(node_id, d)
        if serie:
            hmap = {int(h): float(v) for h, v in serie}
            if len(hmap) >= 20:
                out[d.date()] = hmap
        d += timedelta(days=1)
    _horario_cache[key] = out
    print(f"[data] Cache horario {node_id} {desde}–{hasta}: {len(out)} días", flush=True)
    return out


def _suma_nocturna_dia(hmap: Dict[int, float], night_end: int = L16_NIGHT_END) -> float:
    return sum(hmap.get(i, 0.0) for i in range(night_end + 1))


def _suma_nocturna_por_mes(
    node_id: str,
    desde: str,
    hasta: str,
    *,
    night_end: int = L16_NIGHT_END,
) -> Dict[Tuple[int, int], float]:
    cache = _cargar_horario_cache(node_id, desde, hasta, night_end=night_end)
    monthly: Dict[Tuple[int, int], float] = {}
    for d, hmap in cache.items():
        k = (d.year, d.month)
        monthly[k] = monthly.get(k, 0.0) + _suma_nocturna_dia(hmap, night_end)
    return monthly


def _stats_nocturno_periodo(
    node_id: str,
    desde: str,
    hasta: str,
    *,
    night_end: int = L16_NIGHT_END,
) -> Dict[str, float]:
    """Estadísticas nocturnas 0–night_end h en un período."""
    cache = _cargar_horario_cache(node_id, desde, hasta, night_end=night_end)
    d0 = parse_date(desde).date()
    d1 = parse_date(hasta, end_of_day=True).date()
    night_sums: List[float] = []
    night_peaks: List[float] = []
    for d, hmap in sorted(cache.items()):
        if d < d0 or d > d1:
            continue
        noct = [hmap.get(i, 0.0) for i in range(night_end + 1)]
        night_sums.append(sum(noct))
        night_peaks.append(max(noct) if noct else 0.0)
    n = len(night_sums)
    if not n:
        return {
            "dias": 0,
            "sum_prom": 0.0,
            "h_prom": 0.0,
            "pico_prom": 0.0,
            "pico_max": 0.0,
            "dias_cero": 0,
        }
    dias_cero = sum(1 for s in night_sums if s < 0.05)
    return {
        "dias": float(n),
        "sum_prom": sum(night_sums) / n,
        "h_prom": (sum(night_sums) / n) / (night_end + 1),
        "pico_prom": sum(night_peaks) / n,
        "pico_max": max(night_peaks),
        "dias_cero": float(dias_cero),
    }


def _label_periodo_corto(desde: str, hasta: str) -> str:
    d0 = parse_date(desde)
    d1 = parse_date(hasta)
    meses = ["Ene", "Feb", "Mar", "Abr", "May", "Jun", "Jul", "Ago", "Sep", "Oct", "Nov", "Dic"]
    if d0.year == d1.year and d0.month == d1.month:
        return f"{meses[d0.month - 1]} {d0.year}"
    if d0.year == d1.year:
        return f"{meses[d0.month - 1]}–{meses[d1.month - 1]} {d0.year}"
    return f"{desde}–{hasta}"


def chart_perfil_nocturno_periodo(
    node_id: str,
    desde: str,
    hasta: str,
    out: Path,
    w_in: float,
    h_in: float,
    *,
    night_end: int = L16_NIGHT_END,
) -> Path:
    """Perfil promedio 24 h del período con franja nocturna 0–6 h resaltada."""
    values, dias = _promedio_horario_rango(node_id, desde, hasta)
    stats = _stats_nocturno_periodo(node_id, desde, hasta, night_end=night_end)
    names = _nombres_bom()
    nm = names.get(node_id, node_id)
    plab = _label_periodo_corto(desde, hasta)
    hours = list(range(24))
    noct_vals = values[: night_end + 1]

    fig, ax = plt.subplots(figsize=(max(4.2, w_in * 1.05), max(1.9, h_in * 1.05)))
    ax.plot(hours, values, color="#4A90E2", linewidth=1.5, marker="o", markersize=2.5)
    ax.fill_between(hours, values, 0, color="#87CEEB", alpha=0.4)
    if any(v > 0 for v in noct_vals):
        ax.fill_between(
            range(night_end + 1), noct_vals, 0,
            color="#FFD700", alpha=0.35, hatch="///",
            edgecolor="#FFA500", linewidth=1.0,
        )
        ax.axvline(x=0, color="orange", linestyle="--", linewidth=1, alpha=0.7)
        ax.axvline(x=night_end, color="orange", linestyle="--", linewidth=1, alpha=0.7)

    ax.set_title(
        f"{nm} — promedio 24 h ({plab}, {int(dias)} días)",
        fontsize=8, fontweight="bold", pad=2,
    )
    ax.set_xlabel("Hora del día", fontsize=7)
    ax.set_ylabel("m³/h", fontsize=7)
    ax.set_xlim(0, 23)
    ax.set_xticks(range(0, 24, 4))
    ax.tick_params(axis="x", labelsize=6)
    ax.tick_params(axis="y", labelsize=6)
    ax.grid(True, alpha=0.3, axis="y")
    _ajustar_eje_y(ax, values)

    pct_cero = (stats["dias_cero"] / stats["dias"] * 100.0) if stats["dias"] else 0.0
    notas = [
        f"Suma 0–{night_end} h: {format_number_chilean(stats['sum_prom'], 2)} m³/noche",
        f"Prom. 0–{night_end} h: {format_number_chilean(stats['h_prom'], 2)} m³/h",
        f"Pico 0–{night_end} h: {format_number_chilean(stats['pico_max'], 2)} m³/h",
        f"Noches ≈0: {format_number_chilean(pct_cero, 0)}%",
    ]
    ax.text(
        0.98, 0.97, "\n".join(notas),
        transform=ax.transAxes, ha="right", va="top", fontsize=6,
        bbox=dict(boxstyle="round,pad=0.3", facecolor="white", alpha=0.85, edgecolor="#4A90E2"),
    )

    fig.subplots_adjust(left=0.14, right=0.97, top=0.86, bottom=0.20)
    out.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(out, dpi=200, facecolor="white", pad_inches=0.03)
    plt.close(fig)
    return out


def chart_horario_promedio(
    node_id: str,
    out: Path,
    w_in: float,
    h_in: float,
    *,
    night_end: int = L11_NIGHT_HOURS_END,
) -> Path:
    """Perfil horario promedio de 24 h del período (estilo chart_horario_compact)."""
    values, dias = _promedio_horario_24h(node_id)
    hours = list(range(24))

    fig, ax = plt.subplots(figsize=(max(4.2, w_in * 1.05), max(1.9, h_in * 1.05)))
    ax.plot(hours, values, color="#4A90E2", linewidth=1.5, marker="o", markersize=2.5)
    ax.fill_between(hours, values, 0, color="#87CEEB", alpha=0.4)

    noct_vals = values[: night_end + 1]
    ax.fill_between(
        range(night_end + 1), noct_vals, 0,
        color="#FFD700", alpha=0.3, hatch="///",
        edgecolor="#FFA500", linewidth=1.0,
    )
    ax.axvline(x=0, color="orange", linestyle="--", linewidth=1, alpha=0.7)
    ax.axvline(x=night_end, color="orange", linestyle="--", linewidth=1, alpha=0.7)

    names = _nombres_mam() if node_id in MAM_NODES else {}
    nm = names.get(node_id, node_id)
    ax.set_title(
        f"{nm} — promedio 24 h ({dias} días, {DESDE}–{HASTA})",
        fontsize=8, fontweight="bold", pad=2,
    )
    ax.set_xlabel("Hora del día", fontsize=7)
    ax.set_ylabel("m³/h", fontsize=7)
    ax.set_xlim(0, 23)
    ax.set_xticks(range(0, 24, 4))
    ax.tick_params(axis="x", labelsize=6)
    ax.tick_params(axis="y", labelsize=6)
    ax.grid(True, alpha=0.3, axis="y")
    _ajustar_eje_y(ax, values)

    prom_noc = sum(noct_vals) / len(noct_vals) if noct_vals else 0.0
    ax.text(
        0.98, 0.97,
        f"Prom. 0–{night_end} h: {format_number_chilean(prom_noc, 2)} m³/h",
        transform=ax.transAxes, ha="right", va="top", fontsize=6,
        bbox=dict(boxstyle="round,pad=0.3", facecolor="white", alpha=0.85, edgecolor="#4A90E2"),
    )

    fig.subplots_adjust(left=0.14, right=0.97, top=0.86, bottom=0.20)
    out.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(out, dpi=200, facecolor="white", pad_inches=0.03)
    plt.close(fig)
    return out


def _daily_junio_node(node_id: str) -> Dict[date, float]:
    ms = _datos_nodo(node_id)["nodes_summary"][0]["measures"]
    daily: Dict[date, float] = {}
    for m in ms:
        d0 = m.date.date()
        if L05_JUN_INI <= d0 <= L05_JUN_FIN:
            daily[d0] = daily.get(d0, 0.0) + m.total_m3
    return daily


def _evento_jun10_estanques() -> Dict[str, Dict[str, float]]:
    """Métricas diarias Norte y Sur alrededor del 10/06."""
    out: Dict[str, Dict[str, float]] = {}
    for nid in ("000025-01", "000025-19"):
        daily = _daily_junio_node(nid)
        pre_vals = [v for k, v in daily.items() if k < L05_EVENTO_SUR]
        avg_pre = sum(pre_vals) / len(pre_vals) if pre_vals else 0.0
        dia10 = daily.get(L05_EVENTO_SUR, 0.0)
        pct_baja = ((dia10 - avg_pre) / avg_pre * 100.0) if avg_pre else 0.0
        out[nid] = {"avg_pre": avg_pre, "dia10": dia10, "pct_baja": pct_baja}
    return out


def _texto_analisis_l05() -> str:
    fn = format_number_chilean
    ev = _evento_jun10_estanques()
    ev_n = ev["000025-01"]
    ev_s = ev["000025-19"]
    s_sur_ref = _stats_horario_dia("000025-19", L05_MON_REF) or {}
    s_sur_ctrl = _stats_horario_dia("000025-19", L05_CTRL_DIA) or {}
    s_norte_ref = _stats_horario_dia("000025-01", L05_MON_REF) or {}
    s_norte_ctrl = _stats_horario_dia("000025-01", L05_CTRL_DIA) or {}
    pct_sur_noc = (
        (float(s_sur_ctrl.get("night_sum", 0)) - float(s_sur_ref.get("night_sum", 0)))
        / float(s_sur_ref.get("night_sum", 1))
        * 100.0
    ) if s_sur_ref.get("night_sum") else 0.0
    pct_norte_noc = (
        (float(s_norte_ctrl.get("night_sum", 0)) - float(s_norte_ref.get("night_sum", 0)))
        / float(s_norte_ref.get("night_sum", 1))
        * 100.0
    ) if s_norte_ref.get("night_sum") else 0.0
    return "\n".join([
        "Reparación tuberías — 10/06/2026 (validado con mantención MAE)",
        "Al revelar resultados WES, ese día inicia la baja de consumo en ambos estanques.",
        f"Estanque Sur: promedio 1–9/jun {fn(ev_s['avg_pre'], 1)} m³/día → "
        f"10/jun {fn(ev_s['dia10'], 1)} m³/día ({fn(ev_s['pct_baja'], 0)}%).",
        f"Estanque Norte: promedio 1–9/jun {fn(ev_n['avg_pre'], 1)} m³/día → "
        f"10/jun {fn(ev_n['dia10'], 1)} m³/día ({fn(ev_n['pct_baja'], 0)}%).",
        "Comportamiento proporcional: el Sur es abastecido por la salida del Norte; "
        "al corregirse la red lado sur, baja el llenado y el consumo en ambos puntos.",
        "",
        "Comparación nocturna — lun 4/5 (sin control) vs 10/08 (noche con control)",
        f"Estanque Sur 0–8 h: {fn(s_sur_ref.get('night_sum', 0), 1)} m³ → "
        f"{fn(s_sur_ctrl.get('night_sum', 0), 1)} m³ ({fn(pct_sur_noc, 0)}%). "
        "Corte on/off de mantención nocturna: esa madrugada no se lee como fuga.",
        f"Estanque Norte 0–8 h: {fn(s_norte_ref.get('night_sum', 0), 1)} m³ → "
        f"{fn(s_norte_ctrl.get('night_sum', 0), 1)} m³ ({fn(pct_norte_noc, 0)}%). "
        "Control desde el 05/08/2026: esa madrugada no se lee como fuga "
        "(el alza de agosto es diurna).",
    ])


def _texto_analisis_l06() -> str:
    fn = format_number_chilean
    s_banos = _stats_horario_dia(L06_NODE_BANOS, L06_CTRL_DIA) or {}
    s_pizza = _stats_horario_dia(L06_NODE_PIZZA, L06_CTRL_DIA) or {}
    s_banos_ref = _stats_horario_dia(L06_NODE_BANOS, L05_MON_REF) or {}
    s_pizza_ref = _stats_horario_dia(L06_NODE_PIZZA, L05_MON_REF) or {}
    return "\n".join([
        "Baños Públicos — 10/08/2026",
        "Sistema validado: consumo concentrado en horario hábil.",
        f"Nocturno 0–8 h = {fn(s_banos.get('night_sum', 0), 1)} m³ "
        f"(lun 4/5: {fn(s_banos_ref.get('night_sum', 0), 1)} m³).",
        "Control instalado sin uso; noches ya ~0. No se interpreta como fuga.",
        "",
        "Pizza Hut — 10/08/2026 (control activo desde 01/07)",
        f"Control on/off 00:00–06:00; nocturno 0–8 h = {fn(s_pizza.get('night_sum', 0), 1)} m³ "
        f"(lun 4/5 sin control: {fn(s_pizza_ref.get('night_sum', 0), 1)} m³).",
        f"Pico nocturno 0–8 h: {fn(s_pizza.get('night_max', 0), 2)} m³/h.",
        "Esa noche no se lee como fuga. El alza de julio fue diurna; agosto baja.",
    ])


def _buscar_o_crear_narrativa_l05(slide, texto: str) -> None:
    l, t, w, h = L05_TEXTO_BOX
    target = None
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        tx = sh.text_frame.text.lower()
        if "estanque sur" in tx or "control on/off" in tx or "reparación" in tx or "mantención" in tx:
            target = sh
            break
    if target is None:
        target = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    else:
        target.left = Inches(l)
        target.top = Inches(t)
        target.width = Inches(w)
        target.height = Inches(h)
    _set_texto_parrafos(target, texto)
    _aplicar_negrita_solo_porcentajes(target)


def _buscar_o_crear_narrativa_l06(slide, texto: str) -> None:
    l, t, w, h = L06_TEXTO_BOX
    target = None
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        tx = sh.text_frame.text.lower()
        if "baño" in tx or "bano" in tx or "pizza" in tx or "validado" in tx:
            target = sh
            break
    if target is None:
        target = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    else:
        target.left = Inches(l)
        target.top = Inches(t)
        target.width = Inches(w)
        target.height = Inches(h)
    _set_texto_parrafos(target, texto)
    _aplicar_negrita_solo_porcentajes(target)


def chart_horario(
    node_id: str,
    day: datetime,
    out: Path,
    w: float,
    h: float,
    *,
    anotar: bool = False,
    techo_ref: float | None = None,
) -> Path:
    return chart_horario_compact(node_id, day, out, w, h, anotar=anotar, techo_ref=techo_ref)


def _actualizar_captions_l05(slide) -> None:
    _remover_captions_viejos(slide)
    caps = [
        (
            f"Grafico: Consumos horario Estanque Norte - {DIAS_CAP[L05_MON_REF.weekday()]} "
            f"{L05_MON_REF.day}/{L05_MON_REF.month} y {DIAS_CAP[L05_CTRL_DIA.weekday()]} "
            f"{L05_CTRL_DIA.day}/{L05_CTRL_DIA.month}",
            L05_CAP_NORTE_HOR,
        ),
        (
            f"Grafico: Consumos horario Estanque Sur - {DIAS_CAP[L05_MON_REF.weekday()]} "
            f"{L05_MON_REF.day}/{L05_MON_REF.month} y {DIAS_CAP[L05_CTRL_DIA.weekday()]} "
            f"{L05_CTRL_DIA.day}/{L05_CTRL_DIA.month}",
            L05_CAP_SUR_HOR,
        ),
        (
            "Grafico: Consumo diario Estanque Norte — junio 2026 (evento 10/06)",
            L05_CAP_NORTE_JUN,
        ),
        (
            "Grafico: Consumo diario Estanque Sur — junio 2026 (evento 10/06)",
            L05_CAP_SUR_JUN,
        ),
    ]
    for txt, pos in caps:
        _agregar_caption_fijo(slide, txt, pos)
    print("[OK] L05 — estanques: horarios lun 4/5 + 18/6 y diarios junio")


def _actualizar_captions_l06(slide) -> None:
    _remover_captions_viejos(slide)
    caps = [
        (
            f"Grafico: Consumos horario Baños Públicos - {DIAS_CAP[L06_CTRL_DIA.weekday()]} "
            f"{L06_CTRL_DIA.day}/{L06_CTRL_DIA.month} (sistema validado)",
            L06_CAP_BANOS,
        ),
        (
            f"Grafico: Consumos horario Pizza Hut - {DIAS_CAP[L06_CTRL_DIA.weekday()]} "
            f"{L06_CTRL_DIA.day}/{L06_CTRL_DIA.month} (control activo)",
            L06_CAP_PIZZA,
        ),
    ]
    for txt, pos in caps:
        _agregar_caption_fijo(slide, txt, pos)
    print("[OK] L06 — Baños y Pizza Hut 18/6 (gráficos amplios)")


def _texto_placa_l09() -> str:
    fn = format_number_chilean
    prom_noche = _promedio_nocturno_diario(L09_NODE_PLACA)
    prom_h = _promedio_nocturno_h(L09_NODE_PLACA)
    _, dias = _promedio_horario_24h(L09_NODE_PLACA)
    return "\n".join([
        "Placa Bancaria — consumo nocturno (0–8 h)",
        f"Promedio nocturno del período: {fn(prom_noche, 1)} m³/día "
        f"({fn(prom_h, 2)} m³/h de promedio horario).",
        "En el período analizado, el consumo nocturno se mantiene sostenido "
        "en todas las noches (sin tendencia a cero).",
        f"Gráfico: perfil promedio de 24 h ({dias} días, {DESDE}–{HASTA}).",
        "",
        "Recomendación:",
        "Implementar control on/off en horario inhábil (0–8 h), "
        "gestionado por el servicio de cámaras del mall.",
    ])


def _encontrar_panel_l09(slide):
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        tx = sh.text_frame.text.lower()
        if "ripley" in tx or "placa bancaria" in tx:
            return sh
    return None


def _append_placa_texto_l09(slide) -> None:
    """Conserva texto Ripley manual; agrega o actualiza bloque Placa Bancaria."""
    l, t, w, h = L09_TEXTO_BOX
    target = _encontrar_panel_l09(slide)
    if target is None:
        target = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        _set_texto_parrafos(target, _texto_placa_l09())
    else:
        target.left = Inches(l)
        target.top = Inches(t)
        target.width = Inches(w)
        target.height = Inches(h)
        existing = target.text_frame.text
        marker = "placa bancaria"
        placa_part = _texto_placa_l09()
        idx = existing.lower().find(marker)
        if idx >= 0:
            ripley_part = existing[:idx].rstrip()
            _set_texto_parrafos(target, ripley_part + "\n\n" + placa_part)
        else:
            _set_texto_parrafos(target, existing.rstrip() + "\n\n" + placa_part)
    _aplicar_viñetas_panel_l04(target)
    _aplicar_negrita_solo_porcentajes(target)


def _remover_placa_l09(slide) -> None:
    """Quita solo gráfico y caption de Placa (no toca Ripley)."""
    eliminar = []
    for sh in slide.shapes:
        if sh.has_text_frame:
            tx = sh.text_frame.text.lower()
            if "grafico:" in tx and "placa" in tx:
                eliminar.append(sh._element)
        elif sh.shape_type == 13 and Emu(sh.width).inches > 0.6:
            left = Emu(sh.left).inches
            top = Emu(sh.top).inches
            if left > 6.5 and top > 3.0:
                eliminar.append(sh._element)
    for el in eliminar:
        el.getparent().remove(el)
    if eliminar:
        print(f"[OK] {len(eliminar)} elemento(s) Placa anterior(es) eliminado(s)")


def _agregar_caption_placa_l09(slide) -> None:
    names = _nombres_mam()
    nm = names.get(L09_NODE_PLACA, "Placa Bancaria")
    txt = f"Grafico: Consumos horario {nm} — promedio 24 h periodo {DESDE} a {HASTA}"
    _agregar_caption_fijo(slide, txt, L09_CAP_PLACA)


def _actualizar_textos_l09(slide) -> None:
    """Solo agrega/actualiza bloque Placa en panel; no modifica Ripley ni sus captions."""
    _append_placa_texto_l09(slide)
    print("[OK] Narrativa L09 — bloque Placa Bancaria + control on/off")


def _actualizar_captions_l09(slide) -> None:
    """Alias — narrativa Placa en _actualizar_textos_l09."""
    _actualizar_textos_l09(slide)


def _stats_horario_dia(
    node_id: str,
    day: datetime,
    *,
    night_end: int = L11_NIGHT_HOURS_END,
) -> Dict[str, float] | None:
    serie = get_hourly_measures_for_day(node_id, day)
    vals = [v for _, v in serie]
    if len(vals) < night_end + 2:
        return None
    noct = vals[: night_end + 1]
    work = vals[8:18]
    return {
        "night_max": max(noct) if noct else 0.0,
        "night_avg": (sum(noct) / len(noct)) if noct else 0.0,
        "night_sum": sum(noct),
        "work_sum": sum(work),
        "day_sum": sum(vals),
    }


def _lunes_banos_may_jun() -> List[Tuple[datetime, Dict[str, float]]]:
    filas: List[Tuple[datetime, Dict[str, float]]] = []
    d = datetime(2026, 5, 4)
    fin = datetime(2026, 6, 30)
    while d <= fin:
        if d.weekday() == 0:
            s = _stats_horario_dia("000025-34", d)
            if s:
                filas.append((d, s))
        d += timedelta(days=1)
    return filas


def _banos_dias_nocturno_cero(node_id: str = "000025-34") -> Tuple[int, int]:
    """Días con data en el período y cuántos tienen suma nocturna 0–8 h ≈ 0."""
    desde = parse_date(DESDE)
    hasta = parse_date(HASTA, end_of_day=True)
    d = desde
    total = 0
    cero = 0
    while d <= hasta:
        s = _stats_horario_dia(node_id, d)
        if s is not None:
            total += 1
            if float(s.get("night_sum") or 0) < 0.05:
                cero += 1
        d += timedelta(days=1)
    return total, cero


def _texto_banos_l11() -> str:
    return "\n".join([
        "Alimentación Baños — consumo nocturno (0–8 h)",
        "En el período analizado, el consumo nocturno tiende a cero en todos los días.",
        "No se registra agua corriendo en la red de baños en horario inhábil.",
        "El consumo se concentra en horario de apertura del mall (8–17 h).",
    ])


def _texto_matriz_l12_maq() -> str:
    fn = format_number_chilean
    ref = L11_REF_INFORME_ANTERIOR
    s_may = _stats_horario_dia("000025-13", L11_MATRIZ_REF) or {}
    s_jun = _stats_horario_dia("000025-13", L11_MATRIZ_JUN) or {}
    prom_noche = _promedio_nocturno_diario("000025-13")
    prom_h = _promedio_nocturno_h("000025-13")
    pico_jun = float(s_jun.get("night_max") or 0)
    pct_ref = ((pico_jun - ref) / ref * 100.0) if ref else 0.0
    return "\n".join([
        "Matriz Principal — consumo nocturno (0–8 h)",
        f"Promedio del período: {fn(prom_noche, 1)} m³/noche ({fn(prom_h, 2)} m³/h); "
        "sin noches en cero.",
        f"Lun 4/5: nocturno = {fn(s_may.get('night_sum', 0), 1)} m³ "
        f"(pico {fn(s_may.get('night_max', 0), 2)} m³/h).",
        f"Lun 8/6: nocturno = {fn(s_jun.get('night_sum', 0), 1)} m³ "
        f"(pico {fn(pico_jun, 2)} m³/h).",
        f"Referencia informe anterior: {fn(ref, 1)} m³/h "
        f"(pico lun 8/6 +{fn(pct_ref, 0)}% sobre techo).",
        "",
        "Recomendación:",
        "Implementar control on/off en horario inhábil (0–8 h), "
        "replicando la gestión de estanques MAE.",
    ])


def _texto_banos_l12_maq() -> str:
    fn = format_number_chilean
    total, cero = _dias_nocturno_cero("000025-34")
    pct_cero = (cero / total * 100.0) if total else 0.0
    s_may = _stats_horario_dia("000025-34", L11_BANOS_REF) or {}
    s_jun = _stats_horario_dia("000025-34", L11_BANOS_REF2) or {}
    return "\n".join([
        "Alimentación Baños — consumo nocturno (0–8 h)",
        f"En el período analizado, el consumo nocturno tiende a cero "
        f"({fn(pct_cero, 0)}% de las noches sin flujo significativo).",
        f"Lun 4/5: nocturno = {fn(s_may.get('night_sum', 0), 2)} m³; "
        f"hábil 8–17 h = {fn(s_may.get('work_sum', 0), 1)} m³.",
        f"Lun 8/6: nocturno = {fn(s_jun.get('night_sum', 0), 2)} m³; "
        "patrón esperado en día de baja actividad del mall.",
        "La red de baños no requiere control on/off; el consumo se concentra en apertura.",
    ])


def _texto_analisis_l12_maq() -> str:
    return _texto_matriz_l12_maq() + "\n\n" + _texto_banos_l12_maq()


def _texto_matriz_l11() -> str:
    fn = format_number_chilean
    ref_ant = L11_REF_INFORME_ANTERIOR
    s_jun = _stats_horario_dia("000025-13", L11_MATRIZ_JUN) or {}
    actual = float(s_jun.get("night_max") or 0)
    pct_aum = ((actual - ref_ant) / ref_ant * 100.0) if ref_ant else 0.0
    return "\n".join([
        "Matriz Principal — consumo nocturno (0–8 h)",
        f"Informe anterior: promedio nocturno de {fn(ref_ant, 1)} m³/h.",
        "Período actual: se mantiene consumo nocturno sostenido en horario inhábil.",
        f"Incremento respecto al informe anterior: +{fn(pct_aum, 0)}%.",
    ])


def _texto_analisis_l11() -> str:
    return _texto_matriz_l11() + "\n\n" + _texto_banos_l11()


def _encontrar_caja_narrativa_l11(slide):
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        tx = sh.text_frame.text.lower()
        if "matriz principal" in tx or "alimentación baños" in tx or "alimentacion banos" in tx:
            return sh
    return None


def _buscar_o_crear_narrativa_l11(slide, texto: str, *, conservar_matriz: bool = False) -> None:
    target = _encontrar_caja_narrativa_l11(slide)
    if target is None:
        l, t, w, h = L11_TEXTO_BOX
        target = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        _set_texto_parrafos(target, texto)
    elif conservar_matriz:
        existing = target.text_frame.text
        marker = "alimentación baños"
        idx = existing.lower().find(marker)
        if idx < 0:
            marker = "alimentacion banos"
            idx = existing.lower().find(marker)
        if idx >= 0:
            matriz_part = existing[:idx].rstrip()
            banos_part = _texto_banos_l11()
            _set_texto_parrafos(target, matriz_part + "\n\n" + banos_part)
        else:
            _set_texto_parrafos(target, existing.rstrip() + "\n\n" + _texto_banos_l11())
    else:
        l, t, w, h = L11_TEXTO_BOX
        target.left = Inches(l)
        target.top = Inches(t)
        target.width = Inches(w)
        target.height = Inches(h)
        _set_texto_parrafos(target, texto)
    _aplicar_negrita_solo_porcentajes(target)


def _resumen_lunes_banos_may_jun() -> str:
    """Revisión de todos los lunes may–jun para Baños (patrón 8–17h, noche ~0)."""
    filas = [(d, s["night_sum"], s["work_sum"], s["day_sum"]) for d, s in _lunes_banos_may_jun()]
    filas.sort(key=lambda x: -x[2])
    mejor = filas[0] if filas else None
    bajo = min(filas, key=lambda x: x[3]) if filas else None
    msg = "[INFO] Baños — lunes revisados may/jun: "
    if mejor:
        msg += (
            f"mayor uso horario hábil {mejor[0].strftime('%d/%m')} "
            f"(8–17h: {format_number_chilean(mejor[2], 1)} m³, noche: {format_number_chilean(mejor[1], 2)} m³); "
        )
    if bajo:
        msg += (
            f"menor actividad {bajo[0].strftime('%d/%m')} "
            f"(total {format_number_chilean(bajo[3], 1)} m³)."
        )
    return msg


def _buscar_o_crear_narrativa_l12_maq(slide, texto: str) -> None:
    target = _encontrar_caja_narrativa_l11(slide)
    l, t, w, h = L12_MAQ_TEXTO_BOX
    if target is None:
        target = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    else:
        target.left = Inches(l)
        target.top = Inches(t)
        target.width = Inches(w)
        target.height = Inches(h)
    _set_texto_parrafos(target, texto)
    _aplicar_viñetas_panel_l04(target)
    _aplicar_negrita_solo_porcentajes(target)


def _actualizar_textos_l12_maq(slide) -> None:
    _buscar_o_crear_narrativa_l12_maq(slide, _texto_analisis_l12_maq())
    print("[OK] Narrativa L12 MAQ — Matriz nocturna + Baños (viñetas estándar)")

    _remover_captions_viejos(slide)
    names = _nombres_maq()
    nm_matriz = names.get("000025-13", "Matriz Principal")
    nm_banos = names.get("000025-34", "Alimentación Baños")
    caps = [
        (
            f"Grafico: Consumos horario {nm_matriz} - {DIAS_CAP[L11_MATRIZ_REF.weekday()]} "
            f"{L11_MATRIZ_REF.day}/{L11_MATRIZ_REF.month}",
            L12_CAP_MATRIZ_TOP,
        ),
        (
            f"Grafico: Consumos horario {nm_matriz} - {DIAS_CAP[L11_MATRIZ_JUN.weekday()]} "
            f"{L11_MATRIZ_JUN.day}/{L11_MATRIZ_JUN.month}",
            L12_CAP_MATRIZ_BOT,
        ),
        (
            f"Grafico: Consumos horario {nm_banos} - {DIAS_CAP[L11_BANOS_REF.weekday()]} "
            f"{L11_BANOS_REF.day}/{L11_BANOS_REF.month}",
            L12_CAP_BANOS_TOP,
        ),
        (
            f"Grafico: Consumos horario {nm_banos} - {DIAS_CAP[L11_BANOS_REF2.weekday()]} "
            f"{L11_BANOS_REF2.day}/{L11_BANOS_REF2.month}",
            L12_CAP_BANOS_BOT,
        ),
    ]
    for txt, pos in caps:
        _agregar_caption_fijo(slide, txt, pos)
    print(_resumen_lunes_banos_may_jun())
    print("[OK] L12 MAQ — 4 horarios + panel (formato estándar)")


def _actualizar_textos_l11(slide) -> None:
    """Reescribe panel derecho: Matriz sin valor puntual; Baños con patrón nocturno cero."""
    _buscar_o_crear_narrativa_l11(slide, _texto_analisis_l11(), conservar_matriz=False)
    print("[OK] Narrativa L11 — texto Matriz y Baños reordenado")

    _remover_captions_viejos(slide)
    names = _nombres_maq()
    nm_matriz = names.get("000025-13", "Matriz Principal")
    nm_banos = names.get("000025-34", "Alimentación Baños")
    caps = [
        (
            f"Grafico: Consumos horario {nm_matriz} - {DIAS_CAP[L11_MATRIZ_REF.weekday()]} "
            f"{L11_MATRIZ_REF.day}/{L11_MATRIZ_REF.month}",
            L11_CAP_MATRIZ_TOP,
        ),
        (
            f"Grafico: Consumos horario {nm_matriz} - {DIAS_CAP[L11_MATRIZ_JUN.weekday()]} "
            f"{L11_MATRIZ_JUN.day}/{L11_MATRIZ_JUN.month}",
            L11_CAP_MATRIZ_BOT,
        ),
        (
            f"Grafico: Consumos horario {nm_banos} - {DIAS_CAP[L11_BANOS_REF.weekday()]} "
            f"{L11_BANOS_REF.day}/{L11_BANOS_REF.month}",
            L11_CAP_BANOS_TOP,
        ),
        (
            f"Grafico: Consumos horario {nm_banos} - {DIAS_CAP[L11_BANOS_REF2.weekday()]} "
            f"{L11_BANOS_REF2.day}/{L11_BANOS_REF2.month}",
            L11_CAP_BANOS_BOT,
        ),
    ]
    for txt, pos in caps:
        _agregar_caption_fijo(slide, txt, pos)
    print("[OK] L11 — viñetas 10 pt (estándar L08); gráficos sin tapar panel")


def _promedio_nocturno_diario(node_id: str, *, night_end: int = L11_NIGHT_HOURS_END) -> float:
    d = parse_date(DESDE)
    fin = parse_date(HASTA, end_of_day=True)
    sums: List[float] = []
    while d <= fin:
        s = _stats_horario_dia(node_id, d, night_end=night_end)
        if s is not None:
            sums.append(float(s.get("night_sum") or 0))
        d += timedelta(days=1)
    return sum(sums) / len(sums) if sums else 0.0


def _promedio_nocturno_h(node_id: str, *, night_end: int = L11_NIGHT_HOURS_END) -> float:
    d = parse_date(DESDE)
    fin = parse_date(HASTA, end_of_day=True)
    avgs: List[float] = []
    while d <= fin:
        s = _stats_horario_dia(node_id, d, night_end=night_end)
        if s is not None:
            avgs.append(float(s.get("night_avg") or 0))
        d += timedelta(days=1)
    return sum(avgs) / len(avgs) if avgs else 0.0


def _dias_nocturno_cero(node_id: str, *, night_end: int = L11_NIGHT_HOURS_END) -> Tuple[int, int]:
    d = parse_date(DESDE)
    fin = parse_date(HASTA, end_of_day=True)
    total = 0
    cero = 0
    while d <= fin:
        s = _stats_horario_dia(node_id, d, night_end=night_end)
        if s is not None:
            total += 1
            if float(s.get("night_sum") or 0) < 0.05:
                cero += 1
        d += timedelta(days=1)
    return total, cero


def _total_diario_fecha(node_id: str, day: datetime) -> float:
    fecha = day.strftime("%d/%m/%Y")
    d = _datos_nodo_rango(node_id, fecha, fecha)
    ms = (d.get("nodes_summary") or [{}])[0].get("measures") or []
    return float(ms[0].total_m3) if ms else 0.0


def _promedio_diario_rango(node_id: str, desde: str, hasta: str) -> float:
    d = _datos_nodo_rango(node_id, desde, hasta)
    vals = [float(m.total_m3) for m in (d.get("nodes_summary") or [{}])[0].get("measures") or []]
    return sum(vals) / len(vals) if vals else 0.0


def _franjas_horario_dia(node_id: str, day: datetime) -> Dict[str, float]:
    serie = get_hourly_measures_for_day(node_id, day)
    h = {int(hr): float(v) for hr, v in serie}
    diurno = sum(h.get(i, 0.0) for i in range(8, 18))
    vespertino = sum(h.get(i, 0.0) for i in range(18, 23))
    nocturno = sum(h.get(i, 0.0) for i in list(range(0, 8)) + [23])
    total = sum(h.values()) or 1.0
    peak_h, peak_v = max(h.items(), key=lambda x: x[1]) if h else (0, 0.0)
    return {
        "diurno": diurno,
        "vespertino": vespertino,
        "nocturno": nocturno,
        "total": total,
        "pct_diurno": diurno / total * 100.0,
        "pct_vespertino": vespertino / total * 100.0,
        "pct_nocturno": nocturno / total * 100.0,
        "peak_h": float(peak_h),
        "peak_v": peak_v,
    }


def _texto_s500_l15_bom() -> str:
    fn = format_number_chilean
    names = _nombres_bom()
    nm = names.get(BOM_NODE_500, "San Ignacio 500")
    may_avg = _promedio_diario_rango(BOM_NODE_500, L15_MAY_DESDE, L15_MAY_HASTA)
    jun_avg = _promedio_diario_rango(BOM_NODE_500, L15_JUN_DESDE, L15_JUN_HASTA)
    pre_spike = _promedio_diario_rango(BOM_NODE_500, "01/06/2026", "25/06/2026")
    post_spike = _promedio_diario_rango(BOM_NODE_500, "26/06/2026", L15_JUN_HASTA)
    f_norm = _franjas_horario_dia(BOM_NODE_500, L15_S500_NORMAL)
    f_pico = _franjas_horario_dia(BOM_NODE_500, L15_S500_PICO)
    tot_norm = _total_diario_fecha(BOM_NODE_500, L15_S500_NORMAL)
    tot_pico = _total_diario_fecha(BOM_NODE_500, L15_S500_PICO)
    return "\n".join([
        f"{nm} — alza consumo junio 2026",
        f"Mayo: {fn(may_avg, 1)} m³/d promedio; junio: {fn(jun_avg, 1)} m³/d "
        f"(+{fn((jun_avg - may_avg) / may_avg * 100 if may_avg else 0, 0)}% vs mayo).",
        f"1–25/6: {fn(pre_spike, 1)} m³/d; 26–30/6: {fn(post_spike, 1)} m³/d "
        f"(inicio alza sostenida el {L15_S500_SPIKE_START.strftime('%d/%m')}).",
        f"{DIAS_ES[L15_S500_NORMAL.weekday()]} {L15_S500_NORMAL.day}/{L15_S500_NORMAL.month} "
        f"(referencia): {fn(tot_norm, 1)} m³ — diurno {fn(f_norm['pct_diurno'], 0)}%, "
        f"vespertino {fn(f_norm['pct_vespertino'], 0)}%.",
        f"{DIAS_ES[L15_S500_PICO.weekday()]} {L15_S500_PICO.day}/{L15_S500_PICO.month} "
        f"(pico): {fn(tot_pico, 1)} m³ — concentración 15–20 h "
        f"(pico {int(f_pico['peak_h']):02d} h = {fn(f_pico['peak_v'], 1)} m³/h).",
        "Patrón: diurno–vespertino; validar causa operacional en terreno (última semana).",
    ])


def _texto_s300_l15_bom() -> str:
    fn = format_number_chilean
    names = _nombres_bom()
    nm = names.get(BOM_NODE_300, "San Ignacio 300")
    may_avg = _promedio_diario_rango(BOM_NODE_300, L15_MAY_DESDE, L15_MAY_HASTA)
    jun_avg = _promedio_diario_rango(BOM_NODE_300, L15_JUN_DESDE, L15_JUN_HASTA)
    pre = _promedio_diario_rango(BOM_NODE_300, "01/06/2026", "05/06/2026")
    post = _promedio_diario_rango(BOM_NODE_300, "06/06/2026", L15_JUN_HASTA)
    f_ini = _franjas_horario_dia(BOM_NODE_300, L15_S300_INICIO)
    f_max = _franjas_horario_dia(BOM_NODE_300, L15_S300_MAX)
    tot_ini = _total_diario_fecha(BOM_NODE_300, L15_S300_INICIO)
    tot_max = _total_diario_fecha(BOM_NODE_300, L15_S300_MAX)
    return "\n".join([
        f"{nm} — alza consumo junio 2026",
        f"Mayo: {fn(may_avg, 1)} m³/d; junio: {fn(jun_avg, 1)} m³/d "
        f"(+{fn((jun_avg - may_avg) / may_avg * 100 if may_avg else 0, 0)}% vs mayo).",
        f"Inicio alza sostenida: {L15_S300_INICIO.strftime('%d/%m')} "
        f"(1–5/6: {fn(pre, 1)} m³/d → desde 6/6: {fn(post, 1)} m³/d).",
        f"{DIAS_ES[L15_S300_INICIO.weekday()]} {L15_S300_INICIO.day}/{L15_S300_INICIO.month} "
        f"(primer pico): {fn(tot_ini, 1)} m³; "
        f"{DIAS_ES[L15_S300_MAX.weekday()]} {L15_S300_MAX.day}/{L15_S300_MAX.month} "
        f"(máximo): {fn(tot_max, 1)} m³.",
        f"Franja dominante: diurno 8–17 h ({fn(f_max['pct_diurno'], 0)}% en día pico); "
        f"refuerzo vespertino 17–20 h.",
        f"Horas pico recurrentes: {int(f_max['peak_h']):02d} h "
        f"({fn(f_max['peak_v'], 1)} m³/h el 27/6).",
    ])


def _texto_analisis_l15_bom() -> str:
    return _texto_s500_l15_bom() + "\n\n" + _texto_s300_l15_bom()


def _encontrar_caja_narrativa_l15(slide):
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        tx = sh.text_frame.text.lower()
        if (
            "san ignacio 500" in tx
            or "san ignacio 300" in tx
            or "alza consumo junio" in tx
            or "matriz principal" in tx
        ):
            return sh
    return None


def _buscar_o_crear_narrativa_l15(slide, texto: str) -> None:
    target = _encontrar_caja_narrativa_l15(slide)
    l, t, w, h = L15_TEXTO_BOX
    if target is None:
        target = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    else:
        target.left = Inches(l)
        target.top = Inches(t)
        target.width = Inches(w)
        target.height = Inches(h)
    _set_texto_parrafos(target, texto)
    _aplicar_viñetas_panel_l04(target)
    _aplicar_negrita_solo_porcentajes(target)


def _actualizar_textos_l15_bom(slide) -> None:
    _buscar_o_crear_narrativa_l15(slide, _texto_analisis_l15_bom())
    print("[OK] Narrativa L15 BOM — alza junio San Ignacio 500 y 300 (viñetas estándar)")

    _remover_captions_viejos(slide)
    names = _nombres_bom()
    nm500 = names.get(BOM_NODE_500, "San Ignacio 500")
    nm300 = names.get(BOM_NODE_300, "San Ignacio 300")
    caps = [
        (
            f"Grafico: Consumos horario {nm500} - {DIAS_CAP[L15_S500_NORMAL.weekday()]} "
            f"{L15_S500_NORMAL.day}/{L15_S500_NORMAL.month} (junio previo al pico)",
            L15_CAP_500_TOP,
        ),
        (
            f"Grafico: Consumos horario {nm500} - {DIAS_CAP[L15_S500_PICO.weekday()]} "
            f"{L15_S500_PICO.day}/{L15_S500_PICO.month} (pico fin de mes)",
            L15_CAP_500_BOT,
        ),
        (
            f"Grafico: Consumos horario {nm300} - {DIAS_CAP[L15_S300_INICIO.weekday()]} "
            f"{L15_S300_INICIO.day}/{L15_S300_INICIO.month} (inicio alza)",
            L15_CAP_300_TOP,
        ),
        (
            f"Grafico: Consumos horario {nm300} - {DIAS_CAP[L15_S300_MAX.weekday()]} "
            f"{L15_S300_MAX.day}/{L15_S300_MAX.month} (máximo junio)",
            L15_CAP_300_BOT,
        ),
    ]
    for txt, pos in caps:
        _agregar_caption_fijo(slide, txt, pos)
    print("[OK] L15 BOM — 4 horarios + panel (formato estándar L12)")


def _clp_desde_m3(m3: float, tarifa: float = L16_TARIFA_CLP_M3) -> str:
    return f"CLP ${format_number_chilean(m3 * tarifa, 0)}"


def _texto_nodo_l16_noc(node_id: str) -> str:
    fn = format_number_chilean
    names = _nombres_bom()
    nm = names.get(node_id, node_id)
    s1 = _stats_nocturno_periodo(node_id, *L16_ENE_MAR)
    s2 = _stats_nocturno_periodo(node_id, *L16_ABR_JUN)
    pct_sum = ((s2["sum_prom"] - s1["sum_prom"]) / s1["sum_prom"] * 100.0) if s1["sum_prom"] else 0.0
    pct_h = ((s2["h_prom"] - s1["h_prom"]) / s1["h_prom"] * 100.0) if s1["h_prom"] else 0.0
    pct_cero1 = (s1["dias_cero"] / s1["dias"] * 100.0) if s1["dias"] else 0.0
    pct_cero2 = (s2["dias_cero"] / s2["dias"] * 100.0) if s2["dias"] else 0.0
    tendencia = "aumenta" if s2["sum_prom"] > s1["sum_prom"] * 1.05 else (
        "disminuye" if s2["sum_prom"] < s1["sum_prom"] * 0.95 else "se mantiene"
    )
    return "\n".join([
        f"{nm} — patrón nocturno (0–6 h)",
        f"Ene–Mar: {fn(s1['sum_prom'], 2)} m³/noche ({fn(s1['h_prom'], 2)} m³/h); "
        f"pico {fn(s1['pico_max'], 2)} m³/h; noches ≈0: {fn(pct_cero1, 0)}%.",
        f"Abr–Jun: {fn(s2['sum_prom'], 2)} m³/noche ({fn(s2['h_prom'], 2)} m³/h); "
        f"pico {fn(s2['pico_max'], 2)} m³/h; noches ≈0: {fn(pct_cero2, 0)}%.",
        f"Variación Abr–Jun vs Ene–Mar: {fn(pct_sum, 0)}% en suma nocturna; "
        f"el consumo nocturno {tendencia} ({fn(pct_h, 0)}% m³/h promedio).",
    ])


def _texto_s500_l16_ranking() -> str:
    fn = format_number_chilean
    names = _nombres_bom()
    nm = names.get(BOM_NODE_500, "San Ignacio 500")
    monthly = _suma_nocturna_por_mes(BOM_NODE_500, L16_RANK_DESDE, L16_RANK_HASTA)
    total = sum(monthly.values())
    n_meses = max(len(monthly), 1)
    prom_mes = total / n_meses
    peak_m = max(monthly, key=monthly.get) if monthly else (2026, 1)
    peak_v = monthly.get(peak_m, 0.0)
    mes_lab = L16_MESES_CORTO[peak_m[1] - 1]
    s1 = _stats_nocturno_periodo(BOM_NODE_500, *L16_ENE_MAR)
    s2 = _stats_nocturno_periodo(BOM_NODE_500, *L16_ABR_JUN)
    tarifa = L16_TARIFA_CLP_M3
    return "\n".join([
        f"{nm} — consumo nocturno (0–6 h)",
        "Control WES activo desde el 16/07/2026: la noche bajó de ~42 m³ a ~2 m³. "
        "Esa madrugada no se lee como fuga. Queda alza operacional de día.",
        f"Ranking ene–ago 2026: {fn(total, 0)} m³ acumulados en horario 0–6 h.",
        f"Valor consumo nocturno (tarifa ${fn(tarifa, 0)}/m³): {_clp_desde_m3(total)}.",
        f"Mes de mayor consumo nocturno: {mes_lab} {peak_m[0]} "
        f"({fn(peak_v, 0)} m³ = {_clp_desde_m3(peak_v)}).",
        "El ahorro nocturno del 500 ya está capturado; el volumen diurno no volvió a la base de junio.",
        f"Referencia mensual: ≈ {fn(prom_mes, 0)} m³/mes = {_clp_desde_m3(prom_mes)}/mes.",
        f"Promedio nocturno Ene–Mar: {fn(s1['sum_prom'], 2)} m³/noche; "
        f"Abr–Jun: {fn(s2['sum_prom'], 2)} m³/noche.",
    ])


def _texto_analisis_l16_bom() -> str:
    return _texto_s500_l16_ranking()


def _encontrar_caja_narrativa_l16(slide):
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        tx = sh.text_frame.text.lower()
        if "patrón nocturno" in tx or "patron nocturno" in tx or "san ignacio" in tx:
            return sh
    return None


def _buscar_o_crear_narrativa_l16(slide, texto: str) -> None:
    _limpiar_panel_narrativa_bom(slide, L16_TEXTO_BOX)
    l, t, w, h = L16_TEXTO_BOX
    target = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    _set_texto_parrafos(target, texto)
    _aplicar_viñetas_panel_l04(target)
    _aplicar_negrita_solo_porcentajes(target)


def _actualizar_textos_l16_bom(slide, *, solo_panel: bool = False) -> None:
    if not solo_panel:
        _limpiar_clon_horarios_mae(slide)
    _limpiar_panel_narrativa_bom(slide, L16_TEXTO_BOX)
    _buscar_o_crear_narrativa_l16(slide, _texto_analisis_l16_bom())
    print("[OK] Narrativa L16 BOM — San Ignacio 500 + valor CLP (tarifa 1.400/m³)")

    if solo_panel:
        return

    _remover_captions_viejos(slide)
    names = _nombres_bom()
    nm500 = names.get(BOM_NODE_500, "San Ignacio 500")
    caps = [
        (
            f"Grafico: Ranking consumo nocturno 0–6 h {nm500} por mes (ene–ago 2026)",
            L16_CAP_RANK_S500,
        ),
    ]
    for txt, pos in caps:
        _agregar_caption_fijo(slide, txt, pos)
    print("[OK] L16 BOM — caption ranking S500")


def _texto_s500_l14() -> str:
    fn = format_number_chilean
    ref_ant = L14_REF_INFORME_500
    actual = _promedio_nocturno_h(BOM_NODE_500)
    pct_red = ((ref_ant - actual) / ref_ant * 100.0) if ref_ant else 0.0
    names = _nombres_bom()
    nm = names.get(BOM_NODE_500, "San Ignacio 500")
    return "\n".join([
        f"{nm} — consumo nocturno (0–8 h)",
        f"Informe anterior: consumos nocturnos de hasta {fn(ref_ant, 1)} m³/h.",
        f"Período actual: consumo nocturno promedio de {fn(actual, 1)} m³/h.",
        f"Reducción respecto al informe anterior: −{fn(pct_red, 0)}%.",
    ])


def _texto_s300_l14() -> str:
    names = _nombres_bom()
    nm = names.get(BOM_NODE_300, "San Ignacio 300")
    return "\n".join([
        f"{nm} — consumo nocturno (0–8 h)",
        "En el período analizado, el consumo nocturno tiende a cero.",
        "No se registra agua corriendo en la red en horario inhábil.",
        "El consumo se concentra en horario de apertura del mall (8–17 h).",
    ])


def _texto_analisis_l14() -> str:
    return _texto_s500_l14() + "\n\n" + _texto_s300_l14()


def _encontrar_caja_narrativa_l14(slide):
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        tx = sh.text_frame.text.lower()
        if "san ignacio 500" in tx or "san ignacio 300" in tx:
            return sh
    return None


def _buscar_o_crear_narrativa_l14(slide, texto: str) -> None:
    target = _encontrar_caja_narrativa_l14(slide)
    l, t, w, h = L14_TEXTO_BOX
    if target is None:
        target = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    else:
        target.left = Inches(l)
        target.top = Inches(t)
        target.width = Inches(w)
        target.height = Inches(h)
    _set_texto_parrafos(target, texto)
    _aplicar_negrita_solo_porcentajes(target)


def _actualizar_textos_l14(slide) -> None:
    _buscar_o_crear_narrativa_l14(slide, _texto_analisis_l14())
    print("[OK] Narrativa L14 — San Ignacio 500 y 300 (párrafos limpios)")

    _remover_captions_viejos(slide)
    names = _nombres_bom()
    nm500 = names.get(BOM_NODE_500, "San Ignacio 500")
    nm300 = names.get(BOM_NODE_300, "San Ignacio 300")
    caps = [
        (
            f"Grafico: Consumos horario {nm500} - {DIAS_CAP[L14_S500_REF.weekday()]} "
            f"{L14_S500_REF.day}/{L14_S500_REF.month}",
            L14_CAP_500_TOP,
        ),
        (
            f"Grafico: Consumos horario {nm500} - {DIAS_CAP[L14_S500_JUN.weekday()]} "
            f"{L14_S500_JUN.day}/{L14_S500_JUN.month}",
            L14_CAP_500_BOT,
        ),
        (
            f"Grafico: Consumos horario {nm300} - {DIAS_CAP[L14_S300_REF.weekday()]} "
            f"{L14_S300_REF.day}/{L14_S300_REF.month}",
            L14_CAP_300_TOP,
        ),
        (
            f"Grafico: Consumos horario {nm300} - {DIAS_CAP[L14_S300_REF2.weekday()]} "
            f"{L14_S300_REF2.day}/{L14_S300_REF2.month}",
            L14_CAP_300_BOT,
        ),
    ]
    for txt, pos in caps:
        _agregar_caption_fijo(slide, txt, pos)
    print("[OK] L14 — viñetas 10 pt; layout L11 (500 izq., 300 der.)")


def editar_lamina_5() -> None:
    if not PPT.is_file():
        raise FileNotFoundError(PPT)
    prs = Presentation(str(PPT))
    slide = prs.slides[_slide_index(prs, 5)]

    _borrar_imagenes_grafico(slide)
    print("[OK] Imágenes anteriores eliminadas — L05 solo estanques")

    _set_titulo_izq(slide, "MAE - PERFILES HORARIOS (ESTANQUES)")

    for node_id, day, l, t, w, h in L05_HOR_SLOTS:
        png = CHARTS / f"l05_{node_id}_{day.strftime('%Y%m%d')}.png"
        chart_horario(node_id, day, png, w, h)
        _poner_grafico_fresco(slide, png, l, t, w, h)
        print(f"[OK] horario {node_id} {DIAS_ES[day.weekday()]} {day.strftime('%d/%m/%Y')}")

    for node_id, l, t, w, h in L05_DIARIO_SLOTS:
        png = CHARTS / f"l05_{node_id}_jun2026.png"
        chart_diario_junio(node_id, png, w, h, marcar=L05_EVENTO_SUR)
        _poner_grafico_fresco(slide, png, l, t, w, h)
        print(f"[OK] diario junio {node_id}")

    _buscar_o_crear_narrativa_l05(slide, _texto_analisis_l05())
    print("[OK] Panel L05 — evento 10/6 y control 18/6")
    _actualizar_captions_l05(slide)
    prs.save(str(PPT))
    print(f"\n[OK] Guardado in-place: {PPT}")


def _techo_nocturno_ref(node_id: str) -> float | None:
    s = _stats_horario_dia(node_id, L05_MON_REF)
    if not s:
        return None
    mx = float(s.get("night_max") or 0)
    return mx if mx > 0 else None


def editar_lamina_6() -> None:
    if not PPT.is_file():
        raise FileNotFoundError(PPT)
    prs = Presentation(str(PPT))
    _limpiar_huerfanos_mae(prs)
    _limpiar_duplicados_l06(prs)
    if not _lamina_6_insertada(prs):
        _insertar_lamina_6(prs)
    slide = prs.slides[_slide_index(prs, 6)]

    _borrar_imagenes_grafico(slide)
    print("[OK] Imágenes anteriores eliminadas — L06 Baños y Pizza")

    _set_titulo_izq(slide, "MAE - PERFILES HORARIOS (BAÑOS Y PIZZA)")

    for sh in list(slide.shapes):
        if not sh.has_text_frame or sh.top < Inches(0.5):
            continue
        tx = sh.text_frame.text.lower()
        if "estanque" in tx or "reparación tuberías" in tx:
            sh._element.getparent().remove(sh._element)

    for node_id, day, l, t, w, h in L06_SLOTS:
        techo = _techo_nocturno_ref(node_id)
        png = CHARTS / f"l06_{node_id}_{day.strftime('%Y%m%d')}.png"
        chart_horario(node_id, day, png, w, h, techo_ref=techo)
        _poner_grafico_fresco(slide, png, l, t, w, h)
        print(f"[OK] {node_id} {DIAS_ES[day.weekday()]} {day.strftime('%d/%m/%Y')} techo={techo}")

    _buscar_o_crear_narrativa_l06(slide, _texto_analisis_l06())
    _actualizar_captions_l06(slide)
    prs.save(str(PPT))
    print(f"\n[OK] Guardado in-place: {PPT}")


def editar_lamina_4() -> None:
    if not PPT.is_file():
        raise FileNotFoundError(PPT)
    prs = Presentation(str(PPT))
    slide = prs.slides[3]
    shapes = list(slide.shapes)

    for role, node_id, l, t, w, h in L04_PICS:
        pic = _match_pic(shapes, l, t)
        if not pic:
            print(f"[WARN] No encontré imagen en L={l} T={t}")
            continue
        pic.left = Inches(l)
        pic.top = Inches(t)
        pic.width = Inches(w)
        pic.height = Inches(h)
        if role == "total":
            png = CHARTS / "l04_total_ranking.png"
            chart_total_ranking(png, w, h)
            label = "TOTAL ranking MAE"
        else:
            png = CHARTS / f"l04_{node_id}_diario_linea.png"
            chart_diario_linea(node_id, png, w, h)
            label = f"DIARIO línea {node_id}"
        _replace_pic(pic, png)
        print(f"[OK] {label} -> L={l:.2f} T={t:.2f}")

    _actualizar_textos_l04(slide)

    prs.save(str(PPT))
    print(f"\n[OK] Guardado in-place: {PPT}")


def editar_lamina_7() -> None:
    if not PPT.is_file():
        raise FileNotFoundError(PPT)
    prs = Presentation(str(PPT))
    slide = prs.slides[_slide_index(prs, 7)]

    _borrar_imagenes_grafico(slide)
    print("[OK] Imágenes anteriores eliminadas — layout L07 limpio")

    for role, node_id, l, t, w, h in L07_PICS:
        if role == "total":
            png = CHARTS / "l07_total_ranking.png"
            chart_total_ranking_mam(png, w, h)
            label = "TOTAL ranking MAM"
        elif role == "nota":
            png = CHARTS / "l07_falabella_nota.png"
            chart_nota_falabella(png, w, h)
            label = "NOTA Falabella (OC)"
        else:
            png = CHARTS / f"l07_{node_id}_diario_linea.png"
            chart_diario_linea(node_id, png, w, h)
            label = f"DIARIO línea {node_id}"
        _poner_grafico_fresco(slide, png, l, t, w, h)
        print(f"[OK] {label} -> L={l:.2f} T={t:.2f} W={w:.2f} H={h:.2f}")

    _actualizar_textos_l07(slide)
    prs.save(str(PPT))
    print(f"\n[OK] Guardado in-place: {PPT}")


def editar_lamina_8() -> None:
    """MAM análisis consumos — estándar L04 (ranking dual + % trimestre + anual). FIJO."""
    if not PPT.is_file():
        raise FileNotFoundError(PPT)
    prs = Presentation(str(PPT))
    slide = prs.slides[_slide_index(prs, 8)]

    _borrar_imagenes_grafico(slide)
    print("[OK] Imágenes anteriores eliminadas — layout L08 MAM análisis")

    for role, node_id, l, t, w, h in L08_PICS:
        if role == "total":
            png = CHARTS / "l08_total_ranking.png"
            chart_total_ranking_mam_dual(png, w, h)
            label = "TOTAL ranking MAM dual"
        elif role == "nota":
            png = CHARTS / "l08_falabella_nota.png"
            chart_nota_falabella(png, w, h)
            label = "NOTA Falabella (OC)"
        else:
            png = CHARTS / f"l08_{node_id}_diario_linea.png"
            chart_diario_linea(node_id, png, w, h)
            label = f"DIARIO línea {node_id}"
        _poner_grafico_fresco(slide, png, l, t, w, h)
        print(f"[OK] {label} -> L={l:.2f} T={t:.2f} W={w:.2f} H={h:.2f}")

    _actualizar_textos_l08(slide)
    prs.save(str(PPT))
    print(f"\n[OK] Guardado in-place: {PPT}")


def editar_lamina_9() -> None:
    """L09 — conserva Ripley manual; agrega Placa Bancaria (promedio 24 h) + texto on/off."""
    if not PPT.is_file():
        raise FileNotFoundError(PPT)
    prs = Presentation(str(PPT))
    slide = prs.slides[_slide_index(prs, 9)]

    _remover_placa_l09(slide)

    l, t, w, h = L09_PLACA_PIC
    png = CHARTS / f"l09_{L09_NODE_PLACA}_promedio_24h.png"
    chart_horario_promedio(L09_NODE_PLACA, png, w, h)
    _poner_grafico_fresco(slide, png, l, t, w, h)
    print(f"[OK] Placa Bancaria promedio 24 h @ L={l:.2f} T={t:.2f}")

    _actualizar_textos_l09(slide)
    _agregar_caption_placa_l09(slide)
    print("[OK] L09 — Ripley conservado; Placa promedio + control on/off")

    prs.save(str(PPT))
    print(f"\n[OK] Guardado in-place: {PPT}")


def editar_lamina_10() -> None:
    """Lámina 10 = sección MAQ (divisoria). No se modifica."""
    print("[OK] L10 — sección MAQ conservada (usar --lamina 11 para análisis consumos)")


def editar_lamina_11() -> None:
    """MAQ análisis consumos — ranking dual + diarios del período."""
    _editar_lamina_11_generar()


def _editar_lamina_11_generar() -> None:
    """Regenera lámina 11 MAQ análisis — solo si se solicita explícitamente."""
    if not PPT.is_file():
        raise FileNotFoundError(PPT)
    prs = Presentation(str(PPT))
    idx = _ensure_lamina_maq_analisis(prs)
    slide = prs.slides[idx]

    _borrar_imagenes_grafico(slide)
    print("[OK] Imágenes anteriores eliminadas — layout L11 MAQ análisis")

    for role, node_id, l, t, w, h in L11_MAQ_PICS:
        if role == "total":
            png = CHARTS / "l11_maq_total_ranking.png"
            chart_total_ranking_maq_dual(png, w, h)
            label = "TOTAL ranking MAQ dual"
        else:
            png = CHARTS / f"l11_maq_{node_id}_diario_linea.png"
            chart_diario_linea(node_id, png, w, h)
            label = f"DIARIO línea {node_id}"
        _poner_grafico_fresco(slide, png, l, t, w, h)
        print(f"[OK] {label} -> L={l:.2f} T={t:.2f} W={w:.2f} H={h:.2f}")

    _actualizar_textos_l11_maq(slide)
    prs.save(str(PPT))
    print(f"\n[OK] Guardado in-place: {PPT}")


def editar_lamina_12() -> None:
    """MAQ perfiles horarios — Matriz (on/off) + Baños (noche ~0)."""
    if not PPT.is_file():
        raise FileNotFoundError(PPT)
    prs = Presentation(str(PPT))
    idx = _ensure_lamina_maq_horarios(prs)
    slide = prs.slides[idx]

    _limpiar_clon_horarios_mae(slide)
    _set_titulo_izq(slide, "MAQ - PERFILES HORARIOS")
    _borrar_imagenes_grafico(slide)
    print("[OK] Imágenes anteriores eliminadas — layout L12 MAQ horarios")

    for node_id, day, l, t, w, h in L12_MAQ_SLOTS:
        techo = L11_TECHO_NOCTURNO if node_id == "000025-13" else None
        anotar = node_id == "000025-13" and day == L11_MATRIZ_JUN
        png = CHARTS / f"l12_{node_id}_{day.strftime('%Y%m%d')}.png"
        chart_horario(node_id, day, png, w, h, anotar=anotar, techo_ref=techo)
        _poner_grafico_fresco(slide, png, l, t, w, h)
        extra = f" techo={techo}" if techo else ""
        print(
            f"[OK] {node_id} {DIAS_ES[day.weekday()]} {day.strftime('%d/%m/%Y')} "
            f"@ L={l:.2f} T={t:.2f}{extra}"
        )

    _actualizar_textos_l12_maq(slide)
    prs.save(str(PPT))
    print(f"\n[OK] Guardado in-place: {PPT}")


def editar_lamina_14() -> None:
    """BOM análisis consumos — estándar L04 (ranking dual + % trimestre + anual)."""
    if not PPT.is_file():
        raise FileNotFoundError(PPT)
    prs = Presentation(str(PPT))
    idx = _ensure_lamina_bom_analisis(prs)
    slide = prs.slides[idx]

    _set_titulo_izq(slide, "BOM - ANÁLISIS CONSUMOS")
    _borrar_imagenes_grafico(slide)
    print("[OK] Imágenes anteriores eliminadas — layout L14 BOM análisis")

    for role, node_id, l, t, w, h in L14_BOM_PICS:
        if role == "total":
            png = CHARTS / "l14_bom_total_ranking.png"
            chart_total_ranking_bom_dual(
                png, w, h,
                desde_q2=L14_BOM_CHART_DESDE,
                hasta_q2=L14_BOM_CHART_HASTA,
                label_q2=RANK_Q2_LABEL,
            )
            label = "TOTAL ranking BOM dual (abr–ago)"
        else:
            png = CHARTS / f"l14_bom_{node_id}_diario_linea.png"
            chart_diario_linea(
                node_id, png, w, h,
                desde=L14_BOM_CHART_DESDE,
                hasta=L14_BOM_CHART_HASTA,
            )
            label = f"DIARIO línea {node_id} (may–jun)"
        _poner_grafico_fresco(slide, png, l, t, w, h)
        print(f"[OK] {label} -> L={l:.2f} T={t:.2f} W={w:.2f} H={h:.2f}")

    _actualizar_textos_l14_bom(slide)
    prs.save(str(PPT))
    print(f"\n[OK] Guardado in-place: {PPT}")


def editar_lamina_15() -> None:
    """BOM análisis alza junio — 4 perfiles horarios + narrativa (clon layout L12)."""
    if not PPT.is_file():
        raise FileNotFoundError(PPT)
    prs = Presentation(str(PPT))
    idx = _ensure_lamina_bom_junio(prs)
    slide = prs.slides[idx]

    _set_titulo_izq(slide, "BOM - ANÁLISIS JUNIO 2026")
    _borrar_imagenes_grafico(slide)
    print("[OK] Imágenes anteriores eliminadas — layout L15 BOM junio")

    for node_id, day, l, t, w, h in L15_SLOTS:
        png = CHARTS / f"l15_{node_id}_{day.strftime('%Y%m%d')}.png"
        chart_horario(node_id, day, png, w, h)
        _poner_grafico_fresco(slide, png, l, t, w, h)
        print(
            f"[OK] {node_id} {DIAS_ES[day.weekday()]} {day.strftime('%d/%m/%Y')} "
            f"@ L={l:.2f} T={t:.2f}"
        )

    _actualizar_textos_l15_bom(slide)
    prs.save(str(PPT))
    print(f"\n[OK] Guardado in-place: {PPT}")


def editar_lamina_16(*, solo_texto: bool = False) -> None:
    """BOM patrón nocturno — ranking mensual S500. Use solo_texto=True para no tocar gráficos."""
    if not PPT.is_file():
        raise FileNotFoundError(PPT)
    prs = Presentation(str(PPT))
    idx = _ensure_lamina_bom_nocturno(prs)
    slide = prs.slides[idx]

    _horario_cache.clear()
    _cargar_horario_cache(BOM_NODE_500, L16_RANK_DESDE, L16_RANK_HASTA)

    if solo_texto:
        _set_titulo_izq(slide, "BOM - PATRÓN NOCTURNO (0–6 H)")
        _actualizar_textos_l16_bom(slide, solo_panel=True)
        prs.save(str(PPT))
        print(f"\n[OK] Guardado in-place (solo texto): {PPT}")
        return

    _limpiar_clon_horarios_mae(slide)
    _limpiar_panel_narrativa_bom(slide, L16_TEXTO_BOX)
    _set_titulo_izq(slide, "BOM - PATRÓN NOCTURNO (0–6 H)")
    _borrar_imagenes_grafico(slide)
    print("[OK] Layout L16 limpio — generando ranking mensual S500")

    _cargar_horario_cache(BOM_NODE_300, L16_RANK_DESDE, L16_RANK_HASTA)

    l, t, w, h = L16_RANK_S500_SLOT
    png_rank = CHARTS / "l16_s500_ranking_nocturno_mensual.png"
    chart_ranking_nocturno_mensual_s500(png_rank, w, h)
    _poner_grafico_fresco(slide, png_rank, l, t, w, h)
    print(f"[OK] Ranking nocturno mensual S500 @ L={l:.2f} T={t:.2f}")

    for node_id, desde, hasta, sl, st, sw, sh in L16_S300_SLOTS:
        plab = _label_periodo_corto(desde, hasta).replace("–", "_").replace(" ", "")
        png = CHARTS / f"l16_{node_id}_{plab}.png"
        chart_perfil_nocturno_periodo(node_id, desde, hasta, png, sw, sh)
        _poner_grafico_fresco(slide, png, sl, st, sw, sh)
        print(f"[OK] {node_id} {desde}–{hasta} @ L={sl:.2f} T={st:.2f}")

    _actualizar_textos_l16_bom(slide)
    _eliminar_duplicado_tras_bom_nocturno(prs, idx)
    prs.save(str(PPT))
    print(f"\n[OK] Guardado in-place: {PPT}")


def editar_lamina_13() -> None:
    """Lámina 13 = sección BOM (divisoria). No se modifica."""
    print("[OK] L13 — sección BOM conservada (usar --lamina 14 para análisis consumos)")


def _editar_lamina_13_generar() -> None:
    if not PPT.is_file():
        raise FileNotFoundError(PPT)
    prs = Presentation(str(PPT))
    slide = prs.slides[_slide_index(prs, 13)]

    _borrar_imagenes_grafico(slide)
    print("[OK] Imágenes anteriores eliminadas — layout L13 limpio")

    for role, node_id, l, t, w, h in L13_PICS:
        if role == "total":
            png = CHARTS / "l13_total_ranking.png"
            chart_total_ranking_bom(png, w, h)
            label = "TOTAL ranking BOM"
        else:
            png = CHARTS / f"l13_{node_id}_diario_linea.png"
            chart_diario_linea(node_id, png, w, h)
            label = f"DIARIO línea {node_id}"
        _poner_grafico_fresco(slide, png, l, t, w, h)
        print(f"[OK] {label} -> L={l:.2f} T={t:.2f} W={w:.2f} H={h:.2f}")

    _actualizar_textos_l13(slide)
    prs.save(str(PPT))
    print(f"\n[OK] Guardado in-place: {PPT}")


def editar_lamina_17() -> None:
    """AEB análisis consumos — estándar L04 (ranking dual + % trimestre + anual + 2 diarios)."""
    if not PPT.is_file():
        raise FileNotFoundError(PPT)
    prs = Presentation(str(PPT))
    idx = _ensure_lamina_aeb_analisis(prs)
    slide = prs.slides[idx]

    _set_titulo_izq(slide, "AEB - ANÁLISIS CONSUMOS")
    _borrar_imagenes_grafico(slide)
    print("[OK] Imágenes anteriores eliminadas — layout L17 AEB análisis")

    for role, node_id, l, t, w, h in L17_AEB_PICS:
        if role == "total":
            png = CHARTS / "l17_aeb_total_ranking.png"
            chart_total_ranking_aeb_dual(png, w, h)
            label = "TOTAL ranking AEB dual"
        else:
            png = CHARTS / f"l17_aeb_{node_id}_diario_linea.png"
            chart_diario_linea(node_id, png, w, h)
            label = f"DIARIO línea {node_id}"
        _poner_grafico_fresco(slide, png, l, t, w, h)
        print(f"[OK] {label} -> L={l:.2f} T={t:.2f} W={w:.2f} H={h:.2f}")

    _actualizar_textos_l17_aeb(slide)
    prs.save(str(PPT))
    print(f"\n[OK] Guardado in-place: {PPT}")


def editar_lamina_18() -> None:
    """AEB comparativo nocturno — Matriz vs Anillo (consumo base y oportunidades)."""
    if not PPT.is_file():
        raise FileNotFoundError(PPT)
    prs = Presentation(str(PPT))
    idx = _ensure_lamina_aeb_nocturno(prs)
    slide = prs.slides[idx]

    _limpiar_clon_horarios_mae(slide)
    _set_titulo_izq(slide, "AEB - COMPARATIVO NOCTURNO")
    _borrar_imagenes_grafico(slide)
    print("[OK] Imágenes anteriores eliminadas — layout L18 AEB comparativo")

    for node_id, day, l, t, w, h in L18_AEB_SLOTS:
        png = CHARTS / f"l18_{node_id}_{day.strftime('%Y%m%d')}.png"
        chart_horario(node_id, day, png, w, h)
        _poner_grafico_fresco(slide, png, l, t, w, h)
        print(
            f"[OK] {node_id} {DIAS_ES[day.weekday()]} {day.strftime('%d/%m/%Y')} "
            f"@ L={l:.2f} T={t:.2f}"
        )

    _actualizar_textos_l18_aeb(slide)
    prs.save(str(PPT))
    print(f"\n[OK] Guardado in-place: {PPT}")


def editar_lamina_19() -> None:
    """CUR portada — actualiza texto de instalación (Anillo Norte / Anillo Sur)."""
    if not PPT.is_file():
        raise FileNotFoundError(PPT)
    prs = Presentation(str(PPT))
    idx = _indice_seccion_cur(prs)
    if idx is None:
        raise RuntimeError("No se encontró lámina sección CUR en el deck")
    slide = prs.slides[idx]
    _set_titulo_izq(slide, "CUR")
    _actualizar_textos_l19_cur(slide)
    prs.save(str(PPT))
    print(f"\n[OK] Guardado in-place: {PPT}")


def editar_lamina_20(*, factura_m3: float | None = None, solo_texto: bool = False) -> None:
    """CUR análisis consumos — período 18/5–16/6; comparación WES vs cuenta de agua."""
    if not PPT.is_file():
        raise FileNotFoundError(PPT)
    prs = Presentation(str(PPT))
    idx = _ensure_lamina_cur_analisis(prs)
    slide = prs.slides[idx]

    _set_titulo_izq(slide, "CUR - ANÁLISIS CONSUMOS")

    if solo_texto:
        _actualizar_textos_l20_cur(slide, factura_m3=factura_m3)
        prs.save(str(PPT))
        print(f"\n[OK] Guardado in-place (solo texto): {PPT}")
        return

    _limpiar_clon_cur_analisis(slide)
    _borrar_imagenes_grafico(slide)
    print("[OK] Imágenes anteriores eliminadas — layout L20 CUR análisis")

    for role, node_id, l, t, w, h in L20_CUR_PICS:
        if role == "total":
            png = CHARTS / "l20_cur_total_ranking.png"
            chart_total_ranking_cur_periodo(png, w, h)
            label = "TOTAL ranking CUR período"
        else:
            png = CHARTS / f"l20_cur_{node_id}_diario_linea.png"
            chart_diario_linea(
                node_id, png, w, h,
                desde=L20_CUR_DESDE,
                hasta=L20_CUR_HASTA,
            )
            label = f"DIARIO línea {node_id} ({L20_CUR_PERIODO_CAP})"
        _poner_grafico_fresco(slide, png, l, t, w, h)
        print(f"[OK] {label} -> L={l:.2f} T={t:.2f} W={w:.2f} H={h:.2f}")

    _actualizar_textos_l20_cur(slide, factura_m3=factura_m3)
    prs.save(str(PPT))
    print(f"\n[OK] Guardado in-place: {PPT}")


def editar_lamina_22() -> None:
    """Lámina 22 = sección PAK (divisoria). No se modifica."""
    print("[OK] L22 — sección PAK conservada (usar --lamina 23–27 para análisis PAK)")


def editar_lamina_23() -> None:
    """PAK ranking dual — lámina completa (sin diarios ni panel %)."""
    if not PPT.is_file():
        raise FileNotFoundError(PPT)
    prs = Presentation(str(PPT))
    idx = _ensure_lamina_pak_analisis(prs)
    slide = prs.slides[idx]

    _set_titulo_izq(slide, "PAK - ANÁLISIS CONSUMOS")
    _limpiar_clon_pak_analisis(slide)
    _borrar_imagenes_grafico(slide)
    print("[OK] Imágenes anteriores eliminadas — layout L23 PAK ranking")

    l, t, w, h = LAYOUT_RANKING_FULL_PAK_INFO
    png_rank = CHARTS / "l23_pak_total_ranking.png"
    chart_total_ranking_pak_dual(png_rank, w, h)
    _poner_grafico_fresco(slide, png_rank, l, t, w, h)
    print(f"[OK] TOTAL ranking PAK dual (página completa) @ L={l:.2f} T={t:.2f}")

    _actualizar_textos_l23_pak(slide)
    prs.save(str(PPT))
    print(f"\n[OK] Guardado in-place: {PPT}")


def editar_lamina_24() -> None:
    """PAK cadena abastecimiento DL — correlación 22/28→27→35/36."""
    if not PPT.is_file():
        raise FileNotFoundError(PPT)
    prs = Presentation(str(PPT))
    idx = _ensure_lamina_pak_cadena(prs)
    slide = prs.slides[idx]

    _set_titulo_izq(slide, "PAK - CADENA ABASTECIMIENTO DL")
    _limpiar_clon_pak_analisis(slide)
    _borrar_imagenes_grafico(slide)
    print("[OK] Imágenes anteriores eliminadas — layout L24 PAK cadena")

    l, t, w, h = LAYOUT_PAK_CADENA_BARRAS
    png_barras = CHARTS / "l24_pak_cadena_totales.png"
    chart_pak_cadena_totales(png_barras, w, h)
    _poner_grafico_fresco(slide, png_barras, l, t, w, h)
    print(f"[OK] Cadena totales @ L={l:.2f} T={t:.2f}")

    l2, t2, w2, h2 = LAYOUT_PAK_CADENA_DIARIO
    png_corr = CHARTS / "l24_pak_cadena_correlacion.png"
    chart_pak_cadena_correlacion(png_corr, w2, h2)
    _poner_grafico_fresco(slide, png_corr, l2, t2, w2, h2)
    print(f"[OK] Scatter correlación cadena @ L={l2:.2f} T={t2:.2f}")

    _actualizar_textos_l24_pak(slide)
    prs.save(str(PPT))
    print(f"\n[OK] Guardado in-place: {PPT}")


def editar_lamina_25() -> None:
    """PAK consumo diario — 5 puntos con mayor consumo (1/2). No toca título de lámina."""
    if not PPT.is_file():
        raise FileNotFoundError(PPT)
    prs = Presentation(str(PPT))
    idx = _indice_pak_diarios_a(prs)
    if idx is None:
        idx = _ensure_lamina_pak_diarios_a(prs)
    slide = prs.slides[idx]
    titulo_prev = _titulo_slide(slide)

    _limpiar_clon_pak_analisis(slide)
    _borrar_imagenes_grafico(slide)
    print(f"[OK] Imágenes anteriores eliminadas — L25 diarios (título conservado: {titulo_prev!r})")

    orden = _top_pak_nodes(len(PAK_NODES))
    nodes = orden[:5]
    _poblar_diarios_pak(slide, nodes, "l25")
    _actualizar_textos_l25_pak(slide, nodes)
    # Restaura título por si algún helper lo tocó
    if titulo_prev:
        _set_titulo_izq(slide, titulo_prev)
    prs.save(str(PPT))
    print(f"\n[OK] Guardado in-place: {PPT}")


def editar_lamina_26() -> None:
    """PAK consumo diario — 5 puntos restantes (2/2). No toca título de lámina."""
    if not PPT.is_file():
        raise FileNotFoundError(PPT)
    prs = Presentation(str(PPT))
    idx = _indice_pak_diarios_b(prs)
    if idx is None:
        idx = _ensure_lamina_pak_diarios_b(prs)
    slide = prs.slides[idx]
    titulo_prev = _titulo_slide(slide)

    _limpiar_clon_pak_analisis(slide)
    _borrar_imagenes_grafico(slide)
    print(f"[OK] Imágenes anteriores eliminadas — L26 diarios (título conservado: {titulo_prev!r})")

    orden = _top_pak_nodes(len(PAK_NODES))
    nodes = orden[5:]
    _poblar_diarios_pak(slide, nodes, "l26")
    _actualizar_textos_l26_pak(slide, nodes)
    if titulo_prev:
        _set_titulo_izq(slide, titulo_prev)
    prs.save(str(PPT))
    print(f"\n[OK] Guardado in-place: {PPT}")


def editar_lamina_27() -> None:
    """PAK patrón nocturno — ranking 0–8 h + perfiles horarios DL."""
    if not PPT.is_file():
        raise FileNotFoundError(PPT)
    prs = Presentation(str(PPT))
    idx = _ensure_lamina_pak_nocturno(prs)
    slide = prs.slides[idx]

    _set_titulo_izq(slide, "PAK - PATRÓN NOCTURNO (0–8 H)")
    _limpiar_clon_pak_analisis(slide)
    _limpiar_clon_horarios_mae(slide)
    _borrar_imagenes_grafico(slide)
    print("[OK] Imágenes anteriores eliminadas — layout L27 PAK nocturno")

    peak_d = _dia_mayor_nocturno(PAK_CHAIN_DISTRITO) or date(2026, 6, 14)
    peak_dt = datetime.combine(peak_d, datetime.min.time())

    l, t, w, h = LAYOUT_PAK_RANK_NOCT
    png_rank = CHARTS / "l27_pak_ranking_nocturno.png"
    chart_ranking_nocturno_pak(png_rank, w, h)
    _poner_grafico_fresco(slide, png_rank, l, t, w, h)
    print(f"[OK] Ranking nocturno PAK @ L={l:.2f} T={t:.2f}")

    sl, st, sw, sh = L27_PAK_HOR_LEFT
    png_ref = CHARTS / f"l27_pak_{PAK_CHAIN_DISTRITO}_hor_{L27_PAK_REF.strftime('%Y%m%d')}.png"
    chart_horario(PAK_CHAIN_DISTRITO, L27_PAK_REF, png_ref, sw, sh, anotar=True)
    _poner_grafico_fresco(slide, png_ref, sl, st, sw, sh)
    print(f"[OK] Horario DL ref {L27_PAK_REF.strftime('%d/%m/%Y')}")

    sr, st2, sw2, sh2 = L27_PAK_HOR_RIGHT
    png_peak = CHARTS / f"l27_pak_{PAK_CHAIN_DISTRITO}_hor_{peak_d.strftime('%Y%m%d')}.png"
    chart_horario(PAK_CHAIN_DISTRITO, peak_dt, png_peak, sw2, sh2, anotar=True)
    _poner_grafico_fresco(slide, png_peak, sr, st2, sw2, sh2)
    print(f"[OK] Horario DL pico nocturno {peak_d.strftime('%d/%m/%Y')}")

    _actualizar_textos_l27_pak(slide, peak_d)
    prs.save(str(PPT))
    print(f"\n[OK] Guardado in-place: {PPT}")


def editar_lamina_1() -> None:
    """Portada y cierre: fecha de emisión."""
    if not PPT.is_file():
        raise FileNotFoundError(PPT)
    prs = Presentation(str(PPT))
    for idx in (0, len(prs.slides) - 1):
        slide = prs.slides[idx]
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            raw = (sh.text_frame.text or "").strip()
            if "2026" not in raw or len(raw) > 40:
                continue
            low = raw.lower()
            if not any(m in low for m in ("enero", "febrero", "marzo", "abril", "mayo", "junio", "julio", "agosto")):
                continue
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    if "2026" in (r.text or ""):
                        r.text = FECHA_PORTADA
            print(f"[OK] Fecha slide {idx + 1} → {FECHA_PORTADA}")
    prs.save(str(PPT))
    print(f"\n[OK] Guardado in-place: {PPT}")


def editar_lamina_3() -> None:
    """MAE divisoria: estado de controles actualizado."""
    if not PPT.is_file():
        raise FileNotFoundError(PPT)
    prs = Presentation(str(PPT))
    slide = prs.slides[2]
    reemplazos = [
        (
            "Estanque Norte Locales Mall Control Activo via On/Off",
            "Estanque Norte: control on/off desde 05/08/2026 (noche no se lee como fuga)",
        ),
        (
            "Baños Públicos Cotro, instalado sin funcionamiento",
            "Baños Públicos: control instalado sin uso; noches ~0",
        ),
        (
            "Pizza Hut Control Activo 00:00 a 06:00",
            "Pizza Hut: control 00:00–06:00 desde 01/07/2026 (noche no se lee como fuga)",
        ),
        (
            "Sala de Bomba Estanque Sur Control Activo via On/Off",
            "Estanque Sur: corte on/off de mantención nocturna (noche no se lee como fuga)",
        ),
    ]
    n = 0
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                t = r.text or ""
                nt = t
                for a, b in reemplazos:
                    if a in nt:
                        nt = nt.replace(a, b)
                if nt != t:
                    r.text = nt
                    n += 1
    print(f"[OK] L3 MAE — {n} reemplazo(s) de control")
    prs.save(str(PPT))
    print(f"\n[OK] Guardado in-place: {PPT}")


def _editar_todas(*, factura_m3: float | None = None) -> int:
    orden = [1, 3, 4, 5, 6, 7, 8, 9, 11, 12, 14, 15, 16, 17, 18, 19, 20, 23, 24, 25, 26, 27]
    dispatch = {
        1: editar_lamina_1,
        3: editar_lamina_3,
        4: editar_lamina_4,
        5: editar_lamina_5,
        6: editar_lamina_6,
        7: editar_lamina_7,
        8: editar_lamina_8,
        9: editar_lamina_9,
        11: editar_lamina_11,
        12: editar_lamina_12,
        14: editar_lamina_14,
        15: editar_lamina_15,
        16: lambda: editar_lamina_16(solo_texto=False),
        17: editar_lamina_17,
        18: editar_lamina_18,
        19: editar_lamina_19,
        20: lambda: editar_lamina_20(factura_m3=factura_m3, solo_texto=False),
        23: editar_lamina_23,
        24: editar_lamina_24,
        25: editar_lamina_25,
        26: editar_lamina_26,
        27: editar_lamina_27,
    }
    errores = 0
    for n in orden:
        print(f"\n===== LÁMINA {n} =====", flush=True)
        try:
            dispatch[n]()
        except Exception as exc:
            errores += 1
            print(f"[ERROR] Lámina {n}: {exc}", flush=True)
    print(f"\n[OK] Regeneración completa. Errores: {errores}", flush=True)
    return 1 if errores else 0


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--lamina", type=int, default=None)
    ap.add_argument(
        "--todas",
        action="store_true",
        help="Regenera portada y láminas de análisis del informe 7 Malls",
    )
    ap.add_argument(
        "--solo-texto",
        action="store_true",
        help="Solo actualiza narrativa/captions (no regenera gráficos)",
    )
    ap.add_argument(
        "--factura-m3",
        type=float,
        default=None,
        help="Consumo m³ de la cuenta de agua (CUR L20) para calcular diferencia y variación",
    )
    args = ap.parse_args()
    if args.todas:
        return _editar_todas(factura_m3=args.factura_m3)
    if args.lamina is None:
        ap.error("indique --lamina N o --todas")
    if args.lamina == 1:
        editar_lamina_1()
        return 0
    if args.lamina == 3:
        editar_lamina_3()
        return 0
    if args.lamina == 4:
        editar_lamina_4()
        return 0
    if args.lamina == 5:
        editar_lamina_5()
        return 0
    if args.lamina == 6:
        editar_lamina_6()
        return 0
    if args.lamina == 7:
        editar_lamina_7()
        return 0
    if args.lamina == 8:
        editar_lamina_8()
        return 0
    if args.lamina == 9:
        editar_lamina_9()
        return 0
    if args.lamina == 10:
        editar_lamina_10()
        return 0
    if args.lamina == 11:
        editar_lamina_11()
        return 0
    if args.lamina == 12:
        editar_lamina_12()
        return 0
    if args.lamina == 13:
        editar_lamina_13()
        return 0
    if args.lamina == 14:
        editar_lamina_14()
        return 0
    if args.lamina == 15:
        editar_lamina_15()
        return 0
    if args.lamina == 16:
        editar_lamina_16(solo_texto=args.solo_texto)
        return 0
    if args.lamina == 17:
        editar_lamina_17()
        return 0
    if args.lamina == 18:
        editar_lamina_18()
        return 0
    if args.lamina == 19:
        editar_lamina_19()
        return 0
    if args.lamina == 20:
        editar_lamina_20(factura_m3=args.factura_m3, solo_texto=args.solo_texto)
        return 0
    if args.lamina == 22:
        editar_lamina_22()
        return 0
    if args.lamina == 23:
        editar_lamina_23()
        return 0
    if args.lamina == 24:
        editar_lamina_24()
        return 0
    if args.lamina == 25:
        editar_lamina_25()
        return 0
    if args.lamina == 26:
        editar_lamina_26()
        return 0
    if args.lamina == 27:
        editar_lamina_27()
        return 0
    print(f"Lámina {args.lamina}: aún no implementada.")
    return 1


if __name__ == "__main__":
    raise SystemExit(main())
