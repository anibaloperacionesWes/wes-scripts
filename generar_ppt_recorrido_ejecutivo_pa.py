# -*- coding: utf-8 -*-
"""
Recorrido ejecutivo Parque Arauco — PPT aparte de las fichas.

Misma estética (navy / gold / cards / fondo PA). Mall por mall.

  1) Presentación — equipos en tiras + gráfico mayo → fecha + peso de julio
  2) Hallazgos — se arma después, recinto por recinto (MAE ya tiene la suya)

Uso:
  python3 generar_ppt_recorrido_ejecutivo_pa.py
  python3 generar_ppt_recorrido_ejecutivo_pa.py --skip-refresh
"""

from __future__ import annotations

import json
import sys
from concurrent.futures import ThreadPoolExecutor, as_completed
from datetime import date, datetime, timedelta
from pathlib import Path
from typing import Any, Dict, List, Tuple

import matplotlib

matplotlib.use("Agg")
import matplotlib.dates as mdates
import matplotlib.pyplot as plt
import numpy as np
from matplotlib.colors import ListedColormap
from matplotlib.patches import Patch
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGB
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches as PptInches, Pt as PptPt

from generar_reporte_word import format_number_chilean

ROOT = Path(__file__).resolve().parent
OUT_DIR = ROOT / "reports" / "Parque_Arauco" / "TMP_7MALLS" / "entrega_diego_anibal"
CHARTS = OUT_DIR / "charts_recorrido_mae"
JSON_ALL = OUT_DIR / "datos_recorrido_may_ago.json"
JSON_DATOS = OUT_DIR / "datos_mae_may_ago.json"
JSON_MAM = OUT_DIR / "datos_mam_may_ago.json"
JSON_NOCHES = OUT_DIR / "noches_control_mae.json"
JSON_PAK_CADENA = OUT_DIR / "noches_cadena_pak.json"
JSON_MAM_PLACA = OUT_DIR / "noches_mam_placa.json"
JSON_MAQ_MATRIZ = OUT_DIR / "noches_maq_matriz.json"
JSON_BOM_SI500 = OUT_DIR / "noches_bom_si500.json"
JSON_AEB = OUT_DIR / "noches_aeb.json"
JSON_CUR = OUT_DIR / "noches_cur.json"
JSON_PERFILES = OUT_DIR / "perfiles_horarios_control_mae.json"
LOGO = ROOT / "logo wes.bmp"
FONDO = ROOT / "Parque arauco fondo.jpg"

# MAE: 4 puntos del deck. 000025-02 no entra en este recorrido.
MAE_NODOS = ["000025-01", "000025-04", "000025-07", "000025-19"]
MAM_NODOS = ["000025-08", "000025-09", "000025-10", "000025-32", "000025-33"]
PAK_CADENA = ["000025-27", "000025-35", "000025-36"]
DIA_PAK_NOCHE = date(2026, 8, 10)
JUL_NOCHE_D0 = date(2026, 7, 1)
JUL_NOCHE_D1 = date(2026, 7, 31)
FALABELLA = "000025-09"
PLACA = "000025-08"
FALABELLA_DESDE = date(2026, 8, 11)
PASILLO = "000025-32"
ARROW = "000025-33"
UMBRAL_PLACA_DIA = 290.0
UMBRAL_FALABELLA_DIA = 140.0
MATRIZ_MAQ = "000025-13"
MAQ_ALZA = date(2026, 6, 22)
MAQ_NOCHE_D0 = date(2026, 6, 1)  # la noche también sube; el gráfico parte acá
UMBRAL_MAQ_DIA = 240.0
BAZAR = "000025-35"
DL_KENNEDY = "000025-36"
SI500 = "000025-18"
SI300 = "000025-17"
CTRL_SI500 = date(2026, 7, 17)
UMBRAL_SI500_DIA = 145.0
UMBRAL_SI300_DIA = 45.0
UMBRAL_SI500_NOCHE = 8.0  # solo interno: residual con control ~2 m³
MATRIZ_AEB = "000025-11"
ANILLO_AEB = "000025-12"
MATRIZ_AA_HASTA = date(2026, 5, 15)  # último día con caudal; ese día entra Matriz 1° piso
UMBRAL_ANILLO_DIA = 22.0
UMBRAL_MATRIZ_AEB_DIA = 75.0  # agosto ~58 × 1,25; solo se ofrece umbral del día
ANILLO_SUR = "000025-37"
ANILLO_NORTE = "000025-38"
CUR_NOCHE_D0 = date(2026, 6, 1)  # mayo no se grafica (cambio de anillos)
UMBRAL_SUR_DIA = 14.0  # ~11 × 1,25
UMBRAL_NORTE_DIA = 13.0  # ~10 × 1,25
UMBRAL_MAE_NORTE_DIA = 40.0  # julio ~29
UMBRAL_MAE_SUR_DIA = 35.0  # post presostatos ~28
UMBRAL_MAE_PIZZA_DIA = 50.0  # julio ~38
UMBRAL_MAE_BANOS_DIA = 10.0  # julio ~7

MALLS: List[Dict[str, Any]] = [
    {
        "code": "MAE",
        "titulo": "Mall Arauco Estación",
        "nodes": MAE_NODOS,
        "chip_order": ["000025-01", "000025-07", "000025-19", "000025-04"],
        "recepcion": "20/10/2025",
        "capacitacion": "18/02/2025",
        "usuarios": "medioambiente.dcl@parauco.com  ·  Sala de monitores MAE  ·  Sergio Fuenzalida",
        "caption": "Mayo cierra más alto por Estanque Sur, antes de la reparación del 10/06.",
    },
    {
        "code": "MAM",
        "titulo": "Mall Arauco Maipú",
        "nodes": MAM_NODOS,
        "chip_order": ["000025-08", "000025-10", "000025-09", "000025-32", "000025-33"],
        "recepcion": "06/11/2025",
        "capacitacion": "14/11/2025",
        "usuarios": "Miguel Rupayan  ·  Constanza Vilches  ·  Mantención: C. Bustamante, O. Cuevas y Supervisor Eléctrico",
        "caption": "Junio sube por Placa (18/06). Falabella activo 11/08. Lámina 2: noche, umbral y reubicación.",
    },
    {
        "code": "MAQ",
        "titulo": "Mall Arauco Quilicura",
        "nodes": ["000025-13", "000025-34"],
        "recepcion": "06/11/2025 / relocalizado 17/02/2026",
        "capacitacion": "14/11/2025",
        "usuarios": "Mario Freitez  ·  Tomás Saba  ·  Sebastián Araneda  ·  Mantención: I. Dustan, K. Varas, L. Méndez, C. Leyto",
        "caption": "Matriz concentra el recinto. Lámina 2: alza desde junio, umbral y control de noche.",
    },
    {
        "code": "BOM",
        "titulo": "Buenaventura (San Ignacio)",
        "nodes": ["000025-17", "000025-18"],
        "chip_order": ["000025-18", "000025-17"],
        "recepcion": "18/11/2025",
        "capacitacion": "11/12/2025",
        "usuarios": "Aliro Cortés  ·  Tomás Saba  ·  Sebastián Araneda",
        "caption": "San Ignacio 500 es el mayor volumen. Lámina 2: control nocturno desde 17/07, ahorro y umbral.",
    },
    {
        "code": "AEB",
        "titulo": "Arauco El Bosque",
        "nodes": ["000025-11", "000025-12"],
        "recepcion": "29/10/2025 / relocalizado 16/01/2026",
        "capacitacion": "20/11/2025",
        "usuarios": "Tamara Martínez  ·  Tomás Saba  ·  Sebastián Araneda",
        "caption": "Matriz 1° piso y Anillo Plaza. Matriz A.A. desactivada 15/05. Lámina 2: alza Anillo y umbrales.",
    },
    {
        "code": "CUR",
        "titulo": "Arauco Curauma",
        "nodes": ["000025-37", "000025-38"],
        "chip_order": ["000025-38", "000025-37"],
        "recepcion": "29/04/2026 (Anillo Norte / Anillo Sur)",
        "capacitacion": "12/12/2025",
        "usuarios": "Joceline Lazo  ·  Constanza Vilches",
        "sin_mayo": True,
        "caption": "Cambio de monitoreo: Anillo Norte y Sur desde el 19/05 (mayo no se grafica). Lámina 2: diferencia, noche y umbral.",
    },
    {
        "code": "PAK",
        "titulo": "Parque Arauco Kennedy",
        "nodes": [
            "000025-20", "000025-21", "000025-22", "000025-23", "000025-24",
            "000025-27", "000025-28", "000025-29", "000025-35", "000025-36",
        ],
        "cabecera": [
            "000025-20", "000025-21", "000025-22", "000025-23",
            "000025-24", "000025-28", "000025-29",
        ],
        "chip_order": [
            "000025-20", "000025-21", "000025-22", "000025-28", "000025-29",
            "000025-23", "000025-24", "000025-27", "000025-35", "000025-36",
        ],
        "recepcion": "12/12/2025",
        "capacitacion": "17/12/2025",
        "usuarios": "Francisco Jeldres  ·  Paula Azolas  ·  Mantención: C. Naranjo, M. Jara, R. Moreno, R. Díaz, J. Gutiérrez, H. Fierro",
        "caption": "Julio = cabecera (7 puntos). Distrito de Lujo, Bazar y DL Kennedy no se suman (lámina 2).",
    },
]

NOMBRE_CORTO = {
    "000025-01": "Estanque Norte",
    "000025-04": "Baños Públicos",
    "000025-07": "Pizza Hut",
    "000025-19": "Estanque Sur",
    "000025-08": "Placa Bancaria",
    "000025-09": "Falabella",
    "000025-10": "Ripley",
    "000025-32": "Pasillo Técnico",
    "000025-33": "Salida ARROW",
    "000025-13": "Matriz Principal",
    "000025-34": "Alim. Baños",
    "000025-17": "San Ignacio 300",
    "000025-18": "San Ignacio 500",
    "000025-11": "Matriz 1° piso",
    "000025-12": "Anillo Plaza",
    "000025-37": "Anillo Sur",
    "000025-38": "Anillo Norte",
    "000025-20": "Matriz Andén",
    "000025-21": "Locales Gast.",
    "000025-22": "Sandía Antigua",
    "000025-23": "Pileta",
    "000025-24": "Cascada",
    "000025-27": "Distrito Lujo",
    "000025-28": "Sandía Nueva",
    "000025-29": "Restaurante",
    "000025-35": "Bazar Gourmet",
    "000025-36": "DL Kennedy",
}
CHIP_NOTA = {
    "000025-01": "locales mall",
    "000025-19": "sala de bomba",
    "000025-09": "activo 11/08",
    "000025-32": "residual",
    "000025-33": "residual",
    "000025-34": "uso hábil",
    "000025-17": "monitoreo",
    "000025-11": "activo 15/05",
    "000025-22": "alimenta DL",
    "000025-28": "alimenta DL",
    "000025-27": "cadena · no suma",
    "000025-35": "sale de DL",
    "000025-36": "sale de DL",
}
PALETA = [
    (13, 59, 102),
    (31, 119, 180),
    (196, 92, 38),
    (201, 162, 39),
    (123, 163, 201),
    (90, 140, 110),
    (140, 90, 90),
    (70, 100, 140),
    (160, 120, 60),
    (80, 130, 130),
]
COLOR_NODO: Dict[str, Tuple[int, int, int]] = {
    "000025-07": (196, 92, 38),
    "000025-01": (13, 59, 102),
    "000025-19": (31, 119, 180),
    "000025-04": (123, 163, 201),
    "000025-08": (13, 59, 102),
    "000025-10": (31, 119, 180),
    "000025-09": (201, 162, 39),
    "000025-32": (123, 163, 201),
    "000025-33": (90, 140, 110),
    "000025-27": (196, 92, 38),
    "000025-35": (201, 162, 39),
    "000025-36": (90, 140, 110),
    "000025-13": (13, 59, 102),
    "000025-34": (123, 163, 201),
    "000025-18": (13, 59, 102),
    "000025-17": (196, 92, 38),
    "000025-11": (13, 59, 102),
    "000025-12": (196, 92, 38),
    "000025-37": (13, 59, 102),
    "000025-38": (201, 162, 39),
}
_ci = 0
for _mall in MALLS:
    for _nid in _mall["nodes"]:
        if _nid not in COLOR_NODO:
            COLOR_NODO[_nid] = PALETA[_ci % len(PALETA)]
            _ci += 1

DESDE = date(2026, 5, 1)
HASTA = date(2026, 8, 17)
AGO_MES = 31
PERIODO = f"{DESDE.strftime('%d/%m/%Y')} – {HASTA.strftime('%d/%m/%Y')}"
FECHA_EMISION = "18 agosto 2026"
AGO_ETQ = f"1–{HASTA.day}"

# Hallazgos MAE
SUR_REPARACION = date(2026, 6, 10)  # presostatos Estanque Sur
CTRL_PIZZA = date(2026, 7, 1)
CTRL_NORTE = date(2026, 8, 5)
HORAS_CTRL_NORTE = 5  # 00:00–05:00 (hasta las 05:00)
# 26–29/07 Pizza Hut: la noche volvió; no entra a la mediana post-control.
PIZZA_NOCHES_ATIPICAS = {date(2026, 7, 26), date(2026, 7, 27), date(2026, 7, 28), date(2026, 7, 29)}

NAVY = (13, 59, 102)
GOLD = (201, 162, 39)
TEAL = (31, 119, 180)
GRAY = (90, 90, 90)
LIGHT = (245, 247, 250)
WHITE = (255, 255, 255)
GOLD_SOFT = (232, 213, 163)
TARIFA_CLP_M3 = 1400.0  # misma tarifa del PPT 7 Malls


def fn(v: float, dec: int = 1) -> str:
    return format_number_chilean(float(v), dec)


def _etq_pct_julio(nid: str, vol: float, jul: float, un_decimal: bool = False) -> str:
    """% de julio en las lengüetas. Falabella = 0 (no operativo). Residuales no se pintan 0%."""
    if nid == FALABELLA or jul <= 0:
        return "0 %"
    pct = float(vol) / jul * 100.0
    if pct >= 1:
        return f"{fn(pct, 1 if un_decimal else 0)} %"
    if vol > 0 and pct < 0.1:
        return "<0,1 %"
    if vol > 0:
        return f"{fn(pct, 1)} %"
    return "0 %"


def _clp(m3: float) -> str:
    return f"${fn(float(m3) * TARIFA_CLP_M3, 0)}"


def _rango_horas(horas: List[int]) -> str:
    if not horas:
        return "—"
    h = sorted(horas)
    return f"{h[0]:02d}:00–{h[-1]:02d}:00"


def _rgb(t: Tuple[int, int, int]) -> PptRGB:
    return PptRGB(*t)


def _hex(t: Tuple[int, int, int]) -> str:
    return f"#{t[0]:02X}{t[1]:02X}{t[2]:02X}"


def refrescar_datos(nodos: List[str], json_path: Path, etiqueta: str) -> None:
    from generar_reportes_y_ppt_mall_maipu import guardar_datos_json, obtener_datos_agregados

    print(f"[INFO] Descargando {etiqueta} {PERIODO}…", flush=True)
    datos = obtener_datos_agregados(
        nodos,
        DESDE.strftime("%d/%m/%Y"),
        HASTA.strftime("%d/%m/%Y"),
    )
    datos["all_measures"] = []
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    guardar_datos_json(datos, json_path)


def cargar_mall(
    json_path: Path, nodos: List[str]
) -> Tuple[Dict[str, str], Dict[str, Dict[str, Any]], Dict[str, float]]:
    raw = json.loads(json_path.read_text(encoding="utf-8"))
    names: Dict[str, str] = {}
    by: Dict[str, Dict[str, Any]] = {}
    for ns in raw["nodes_summary"]:
        nid = ns["node_id"]
        names[nid] = ns["node_name"]
        row: Dict[str, Any] = {
            "may": 0.0,
            "jun": 0.0,
            "jul": 0.0,
            "ago": 0.0,
            "daily": {},
        }
        for m in ns["measures"]:
            d = str(m["date"])[:10]
            v = float(m["total_m3"])
            row["daily"][d] = row["daily"].get(d, 0.0) + v
            month = d[5:7]
            if month == "05":
                row["may"] += v
            elif month == "06":
                row["jun"] += v
            elif month == "07":
                row["jul"] += v
            elif month == "08":
                row["ago"] += v
        by[nid] = row
    tot = totales_de(by, nodos)
    return names, by, tot


def totales_de(by: Dict[str, Dict[str, Any]], nodos: List[str]) -> Dict[str, float]:
    tot = {"may": 0.0, "jun": 0.0, "jul": 0.0, "ago": 0.0}
    for nid in nodos:
        for k in tot:
            tot[k] += float((by.get(nid) or {}).get(k) or 0)
    tot["ago_d"] = tot["ago"] / HASTA.day if HASTA.day else 0.0
    tot["ago_proy"] = tot["ago_d"] * AGO_MES
    tot["jul_d"] = tot["jul"] / 31.0
    tot["may_d"] = tot["may"] / 31.0
    tot["jun_d"] = tot["jun"] / 30.0
    return tot


def nodos_totales(mall: Dict[str, Any]) -> List[str]:
    return list(mall.get("cabecera") or mall["nodes"])


def nodos_todos() -> List[str]:
    out: List[str] = []
    seen = set()
    for mall in MALLS:
        for nid in mall["nodes"]:
            if nid not in seen:
                seen.add(nid)
                out.append(nid)
    return out


def cargar_mae() -> Tuple[Dict[str, str], Dict[str, Dict[str, Any]], Dict[str, float]]:
    path = JSON_ALL if JSON_ALL.is_file() else JSON_DATOS
    return cargar_mall(path, MAE_NODOS)


def _n06_de_serie(serie) -> float:
    return sum(float(v) for h, v in (serie or []) if int(h) < 6)


def _horas_noche_de_serie(serie) -> Dict[str, float]:
    rec = {str(h): 0.0 for h in range(6)}
    for h, v in serie or []:
        hi = int(h)
        if 0 <= hi < 6:
            rec[str(hi)] = round(float(v), 3)
    return rec


def _rango_dias(d0: date, d1: date) -> List[date]:
    out = []
    d = d0
    while d <= d1:
        out.append(d)
        d += timedelta(days=1)
    return out


def refrescar_noches() -> Dict[str, Dict[str, float]]:
    """m³ 00:00–06:00 por día para Pizza Hut y Estanque Norte (controles MAE)."""
    from generar_reporte_word import get_hourly_measures_for_day

    hourly: Dict[str, Dict[str, float]] = {"000025-07": {}, "000025-01": {}}
    if JSON_NOCHES.is_file():
        try:
            hourly = json.loads(JSON_NOCHES.read_text(encoding="utf-8")).get("hourly") or hourly
        except Exception:
            pass
    pendientes: List[Tuple[str, date]] = []
    for nid, d0, d1 in (
        ("000025-07", date(2026, 6, 20), HASTA),
        ("000025-01", date(2026, 7, 25), HASTA),
    ):
        hourly.setdefault(nid, {})
        for d in _rango_dias(d0, d1):
            if d.isoformat() not in hourly[nid]:
                pendientes.append((nid, d))
    print(f"[INFO] Noches MAE a descargar: {len(pendientes)}", flush=True)

    def _uno(nid: str, d: date) -> Tuple[str, str, float]:
        serie = get_hourly_measures_for_day(nid, datetime(d.year, d.month, d.day)) or []
        return nid, d.isoformat(), round(_n06_de_serie(serie), 2)

    if pendientes:
        with ThreadPoolExecutor(max_workers=6) as pool:
            futs = [pool.submit(_uno, nid, d) for nid, d in pendientes]
            for i, fut in enumerate(as_completed(futs), 1):
                nid, iso, n06 = fut.result()
                hourly[nid][iso] = n06
                if i % 15 == 0 or i == len(pendientes):
                    print(f"  {i}/{len(pendientes)} noches", flush=True)
    JSON_NOCHES.write_text(
        json.dumps({"hourly": hourly}, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    return hourly


def cargar_noches() -> Dict[str, Dict[str, float]]:
    if not JSON_NOCHES.is_file():
        return {}
    return json.loads(JSON_NOCHES.read_text(encoding="utf-8")).get("hourly") or {}


def cargar_cadena_pak() -> Dict[str, Any]:
    if not JSON_PAK_CADENA.is_file():
        return {"perfil": {}, "n06": {}}
    raw = json.loads(JSON_PAK_CADENA.read_text(encoding="utf-8"))
    return {"perfil": raw.get("perfil") or {}, "n06": raw.get("n06") or {}}


def _mediana(vals: List[float]) -> float:
    xs = sorted(float(v) for v in vals)
    if not xs:
        return 0.0
    n = len(xs)
    mid = n // 2
    if n % 2:
        return xs[mid]
    return (xs[mid - 1] + xs[mid]) / 2.0


def refrescar_cadena_pak() -> Dict[str, Any]:
    """Perfil 10/08 y m³ 00:00–06:00 de julio (mes completo) para 27, 35 y 36."""
    from generar_reporte_word import get_hourly_measures_for_day

    data = cargar_cadena_pak()
    perfil: Dict[str, Dict[str, float]] = data.get("perfil") or {}
    n06: Dict[str, Dict[str, float]] = data.get("n06") or {}
    iso_perfil = DIA_PAK_NOCHE.isoformat()
    pendientes_p: List[str] = []
    for nid in PAK_CADENA:
        perfil.setdefault(nid, {})
        if iso_perfil not in perfil[nid]:
            pendientes_p.append(nid)
    pendientes_n: List[Tuple[str, date]] = []
    for nid in PAK_CADENA:
        n06.setdefault(nid, {})
        for d in _rango_dias(JUL_NOCHE_D0, HASTA):
            if d.isoformat() not in n06[nid]:
                pendientes_n.append((nid, d))
    print(
        f"[INFO] Cadena PAK: {len(pendientes_p)} perfiles + {len(pendientes_n)} noches",
        flush=True,
    )

    def _perfil(nid: str) -> Tuple[str, Dict[str, float]]:
        serie = get_hourly_measures_for_day(
            nid, datetime(DIA_PAK_NOCHE.year, DIA_PAK_NOCHE.month, DIA_PAK_NOCHE.day)
        ) or []
        rec = {str(int(h)): round(float(v), 3) for h, v in serie}
        return nid, rec

    def _noche(nid: str, d: date) -> Tuple[str, str, float]:
        serie = get_hourly_measures_for_day(nid, datetime(d.year, d.month, d.day)) or []
        return nid, d.isoformat(), round(_n06_de_serie(serie), 2)

    if pendientes_p:
        with ThreadPoolExecutor(max_workers=3) as pool:
            for fut in as_completed([pool.submit(_perfil, nid) for nid in pendientes_p]):
                nid, rec = fut.result()
                perfil[nid][iso_perfil] = rec
    if pendientes_n:
        with ThreadPoolExecutor(max_workers=6) as pool:
            futs = [pool.submit(_noche, nid, d) for nid, d in pendientes_n]
            for i, fut in enumerate(as_completed(futs), 1):
                nid, iso, v = fut.result()
                n06[nid][iso] = v
                if i % 15 == 0 or i == len(pendientes_n):
                    print(f"  {i}/{len(pendientes_n)} noches cadena PAK", flush=True)
    payload = {"perfil": perfil, "n06": n06}
    JSON_PAK_CADENA.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")
    return payload


def cargar_mam_placa() -> Dict[str, Any]:
    if not JSON_MAM_PLACA.is_file():
        return {"perfil": {}, "n06": {}}
    raw = json.loads(JSON_MAM_PLACA.read_text(encoding="utf-8"))
    return {"perfil": raw.get("perfil") or {}, "n06": raw.get("n06") or {}}


def refrescar_mam_placa() -> Dict[str, Any]:
    """Perfil 10/08 de Placa y m³ 00:00–06:00 (julio + agosto a la fecha) de Placa y Falabella."""
    from generar_reporte_word import get_hourly_measures_for_day

    data = cargar_mam_placa()
    perfil: Dict[str, Dict[str, Any]] = data.get("perfil") or {}
    n06: Dict[str, Dict[str, float]] = data.get("n06") or {}
    iso_perfil = date(2026, 8, 10).isoformat()
    perfil.setdefault(PLACA, {})
    pendientes_p: List[str] = []
    if iso_perfil not in perfil[PLACA]:
        pendientes_p.append(PLACA)
    pendientes_n: List[Tuple[str, date]] = []
    n06.setdefault(PLACA, {})
    n06.setdefault(FALABELLA, {})
    for d in _rango_dias(JUL_NOCHE_D0, HASTA):
        if d.isoformat() not in n06[PLACA]:
            pendientes_n.append((PLACA, d))
    for d in _rango_dias(FALABELLA_DESDE, HASTA):
        if d.isoformat() not in n06[FALABELLA]:
            pendientes_n.append((FALABELLA, d))
    print(
        f"[INFO] MAM Placa: {len(pendientes_p)} perfiles + {len(pendientes_n)} noches",
        flush=True,
    )

    def _perfil(nid: str) -> Tuple[str, Dict[str, float]]:
        serie = get_hourly_measures_for_day(nid, datetime(2026, 8, 10)) or []
        rec = {str(int(h)): round(float(v), 3) for h, v in serie}
        return nid, rec

    def _noche(nid: str, d: date) -> Tuple[str, str, float]:
        serie = get_hourly_measures_for_day(nid, datetime(d.year, d.month, d.day)) or []
        return nid, d.isoformat(), round(_n06_de_serie(serie), 2)

    if pendientes_p:
        with ThreadPoolExecutor(max_workers=2) as pool:
            for fut in as_completed([pool.submit(_perfil, nid) for nid in pendientes_p]):
                nid, rec = fut.result()
                perfil[nid][iso_perfil] = rec
    if pendientes_n:
        with ThreadPoolExecutor(max_workers=6) as pool:
            futs = [pool.submit(_noche, nid, d) for nid, d in pendientes_n]
            for i, fut in enumerate(as_completed(futs), 1):
                nid, iso, v = fut.result()
                n06[nid][iso] = v
                if i % 15 == 0 or i == len(pendientes_n):
                    print(f"  {i}/{len(pendientes_n)} noches MAM Placa", flush=True)
    payload = {"perfil": perfil, "n06": n06}
    JSON_MAM_PLACA.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")
    return payload


def cargar_maq_matriz() -> Dict[str, Any]:
    if not JSON_MAQ_MATRIZ.is_file():
        return {"n06": {}, "horas": {}}
    raw = json.loads(JSON_MAQ_MATRIZ.read_text(encoding="utf-8"))
    return {"n06": raw.get("n06") or {}, "horas": raw.get("horas") or {}}


def refrescar_maq_matriz() -> Dict[str, Any]:
    """m³ 00:00–06:00 de Matriz Principal desde 01/06, más m³/h de madrugada."""
    from generar_reporte_word import get_hourly_measures_for_day

    data = cargar_maq_matriz()
    n06: Dict[str, Dict[str, float]] = data.get("n06") or {}
    horas: Dict[str, Dict[str, float]] = data.get("horas") or {}
    n06.setdefault(MATRIZ_MAQ, {})
    pendientes: List[date] = [
        d
        for d in _rango_dias(MAQ_NOCHE_D0, HASTA)
        if d.isoformat() not in n06[MATRIZ_MAQ] or d.isoformat() not in horas
    ]
    print(f"[INFO] MAQ Matriz: {len(pendientes)} noches", flush=True)

    def _noche(d: date) -> Tuple[str, float, Dict[str, float]]:
        serie = get_hourly_measures_for_day(MATRIZ_MAQ, datetime(d.year, d.month, d.day)) or []
        return d.isoformat(), round(_n06_de_serie(serie), 2), _horas_noche_de_serie(serie)

    if pendientes:
        with ThreadPoolExecutor(max_workers=6) as pool:
            futs = [pool.submit(_noche, d) for d in pendientes]
            for i, fut in enumerate(as_completed(futs), 1):
                iso, v, rec = fut.result()
                n06[MATRIZ_MAQ][iso] = v
                horas[iso] = rec
                if i % 15 == 0 or i == len(pendientes):
                    print(f"  {i}/{len(pendientes)} noches MAQ Matriz", flush=True)
    payload = {"n06": n06, "horas": horas}
    JSON_MAQ_MATRIZ.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")
    return payload


def cargar_bom_si500() -> Dict[str, Any]:
    if not JSON_BOM_SI500.is_file():
        return {"n06": {}}
    raw = json.loads(JSON_BOM_SI500.read_text(encoding="utf-8"))
    return {"n06": raw.get("n06") or {}}


def refrescar_bom_si500() -> Dict[str, Any]:
    """m³ 00:00–06:00 de San Ignacio 500 (julio–agosto a la fecha)."""
    from generar_reporte_word import get_hourly_measures_for_day

    data = cargar_bom_si500()
    n06: Dict[str, Dict[str, float]] = data.get("n06") or {}
    pendientes: List[Tuple[str, date]] = []
    n06.setdefault(SI500, {})
    for d in _rango_dias(JUL_NOCHE_D0, HASTA):
        if d.isoformat() not in n06[SI500]:
            pendientes.append((SI500, d))
    print(f"[INFO] BOM SI500: {len(pendientes)} noches", flush=True)

    def _noche(nid: str, d: date) -> Tuple[str, str, float]:
        serie = get_hourly_measures_for_day(nid, datetime(d.year, d.month, d.day)) or []
        return nid, d.isoformat(), round(_n06_de_serie(serie), 2)

    if pendientes:
        with ThreadPoolExecutor(max_workers=6) as pool:
            futs = [pool.submit(_noche, nid, d) for nid, d in pendientes]
            for i, fut in enumerate(as_completed(futs), 1):
                nid, iso, v = fut.result()
                n06[nid][iso] = v
                if i % 15 == 0 or i == len(pendientes):
                    print(f"  {i}/{len(pendientes)} noches BOM", flush=True)
    payload = {"n06": n06}
    JSON_BOM_SI500.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")
    return payload


def cargar_aeb() -> Dict[str, Any]:
    if not JSON_AEB.is_file():
        return {"n06": {}}
    raw = json.loads(JSON_AEB.read_text(encoding="utf-8"))
    return {"n06": raw.get("n06") or {}}


def refrescar_aeb() -> Dict[str, Any]:
    """m³ 00:00–06:00 de Matriz 1° piso y Anillo Plaza (mayo–agosto a la fecha)."""
    from generar_reporte_word import get_hourly_measures_for_day

    data = cargar_aeb()
    n06: Dict[str, Dict[str, float]] = data.get("n06") or {}
    pendientes: List[Tuple[str, date]] = []
    for nid in (MATRIZ_AEB, ANILLO_AEB):
        n06.setdefault(nid, {})
        for d in _rango_dias(DESDE, HASTA):
            if d.isoformat() not in n06[nid]:
                pendientes.append((nid, d))
    print(f"[INFO] AEB: {len(pendientes)} noches", flush=True)

    def _noche(nid: str, d: date) -> Tuple[str, str, float]:
        serie = get_hourly_measures_for_day(nid, datetime(d.year, d.month, d.day)) or []
        return nid, d.isoformat(), round(_n06_de_serie(serie), 2)

    if pendientes:
        with ThreadPoolExecutor(max_workers=6) as pool:
            futs = [pool.submit(_noche, nid, d) for nid, d in pendientes]
            for i, fut in enumerate(as_completed(futs), 1):
                nid, iso, v = fut.result()
                n06[nid][iso] = v
                if i % 20 == 0 or i == len(pendientes):
                    print(f"  {i}/{len(pendientes)} noches AEB", flush=True)
                    JSON_AEB.write_text(
                        json.dumps({"n06": n06}, ensure_ascii=False, indent=2),
                        encoding="utf-8",
                    )
    payload = {"n06": n06}
    JSON_AEB.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")
    return payload


def cargar_cur() -> Dict[str, Any]:
    if not JSON_CUR.is_file():
        return {"n06": {}}
    raw = json.loads(JSON_CUR.read_text(encoding="utf-8"))
    return {"n06": raw.get("n06") or {}}


def refrescar_cur() -> Dict[str, Any]:
    """m³ 00:00–06:00 de Anillo Sur y Norte (junio–agosto; mayo no entra)."""
    from generar_reporte_word import get_hourly_measures_for_day

    data = cargar_cur()
    n06: Dict[str, Dict[str, float]] = data.get("n06") or {}
    pendientes: List[Tuple[str, date]] = []
    for nid in (ANILLO_SUR, ANILLO_NORTE):
        n06.setdefault(nid, {})
        for d in _rango_dias(CUR_NOCHE_D0, HASTA):
            if d.isoformat() not in n06[nid]:
                pendientes.append((nid, d))
    print(f"[INFO] CUR: {len(pendientes)} noches", flush=True)

    def _noche(nid: str, d: date) -> Tuple[str, str, float]:
        serie = get_hourly_measures_for_day(nid, datetime(d.year, d.month, d.day)) or []
        return nid, d.isoformat(), round(_n06_de_serie(serie), 2)

    if pendientes:
        with ThreadPoolExecutor(max_workers=6) as pool:
            futs = [pool.submit(_noche, nid, d) for nid, d in pendientes]
            for i, fut in enumerate(as_completed(futs), 1):
                nid, iso, v = fut.result()
                n06[nid][iso] = v
                if i % 20 == 0 or i == len(pendientes):
                    print(f"  {i}/{len(pendientes)} noches CUR", flush=True)
                    JSON_CUR.write_text(
                        json.dumps({"n06": n06}, ensure_ascii=False, indent=2),
                        encoding="utf-8",
                    )
    payload = {"n06": n06}
    JSON_CUR.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")
    return payload


def _prom_mes(serie: Dict[str, float], year: int, month: int, activos: Dict[str, float] | None = None) -> float:
    """Promedio diario o nocturno del mes, solo días con caudal de día."""
    vals: List[float] = []
    for iso, v in serie.items():
        d = date.fromisoformat(iso)
        if d.year != year or d.month != month or d > HASTA:
            continue
        if activos is not None and float(activos.get(iso, 0.0) or 0.0) <= 0.1:
            continue
        vals.append(float(v))
    return (sum(vals) / len(vals)) if vals else 0.0


def _stats_si500(n06: Dict[str, float]) -> Dict[str, Any]:
    """Noche típica 7 días antes del 17/07 vs con control, y residuales."""
    pre0 = CTRL_SI500 - timedelta(days=7)
    pre, post_vals = [], []
    post_altas: List[Tuple[date, float]] = []
    vispera = 0.0
    for iso, v in n06.items():
        d = date.fromisoformat(iso)
        fv = float(v)
        if d == CTRL_SI500 - timedelta(days=2):
            vispera = fv  # 15/07, última noche alta de referencia
        if pre0 <= d < CTRL_SI500:
            pre.append(fv)
        elif d >= CTRL_SI500:
            post_vals.append(fv)
            if fv >= UMBRAL_SI500_NOCHE:
                post_altas.append((d, fv))
    post_altas.sort(key=lambda x: -x[1])
    pre_m = _mediana(pre)
    post_m = _mediana(post_vals)
    ahorro = max(0.0, pre_m - post_m)
    return {
        "pre": pre_m,
        "post": post_m,
        "ahorro_noche": ahorro,
        "n_pre": float(len(pre)),
        "n_post": float(len(post_vals)),
        "ahorro_acum": ahorro * float(len(post_vals)),
        "max_post": max(post_vals) if post_vals else 0.0,
        "n_sobre_umbral": float(len(post_altas)),
        "noches_altas": post_altas[:3],
        "vispera": vispera,
    }


def _stats_maq_hora(horas: Dict[str, Dict[str, float]]) -> Dict[str, float]:
    """Mínimo y típico m³/h en 00–06 desde el alza del 22/06."""
    mins: List[float] = []
    todas: List[float] = []
    for iso, rec in (horas or {}).items():
        if date.fromisoformat(iso) < MAQ_ALZA:
            continue
        vals = [float(rec.get(str(h), 0.0) or 0.0) for h in range(6)]
        if not rec:
            continue
        mins.append(min(vals))
        todas.extend(vals)
    return {
        "min_hora": _mediana(mins),
        "med_hora": _mediana(todas),
    }


def cargar_perfiles() -> Dict[str, Dict[str, Dict[str, float]]]:
    if not JSON_PERFILES.is_file():
        return {}
    return json.loads(JSON_PERFILES.read_text(encoding="utf-8")).get("by_h") or {}


def refrescar_perfiles() -> Dict[str, Dict[str, Dict[str, float]]]:
    """Perfil horario 0–7 h en días con control (para marcar horas en cero)."""
    from generar_reporte_word import get_hourly_measures_for_day

    by_h: Dict[str, Dict[str, Dict[str, float]]] = cargar_perfiles() or {
        "000025-01": {},
        "000025-07": {},
        "000025-19": {},
    }
    pendientes: List[Tuple[str, date]] = []
    for nid, d0, d1 in (
        ("000025-01", CTRL_NORTE, HASTA),
        ("000025-07", CTRL_PIZZA, HASTA),
    ):
        by_h.setdefault(nid, {})
        for d in _rango_dias(d0, d1):
            if d.isoformat() not in by_h[nid]:
                pendientes.append((nid, d))
    print(f"[INFO] Perfiles horarios a descargar: {len(pendientes)}", flush=True)

    def _uno(nid: str, d: date) -> Tuple[str, str, Dict[str, float]]:
        serie = get_hourly_measures_for_day(nid, datetime(d.year, d.month, d.day)) or []
        rec = {str(int(h)): round(float(v), 3) for h, v in serie}
        return nid, d.isoformat(), rec

    if pendientes:
        with ThreadPoolExecutor(max_workers=6) as pool:
            futs = [pool.submit(_uno, nid, d) for nid, d in pendientes]
            for i, fut in enumerate(as_completed(futs), 1):
                nid, iso, rec = fut.result()
                by_h[nid][iso] = rec
                if i % 20 == 0 or i == len(pendientes):
                    print(f"  {i}/{len(pendientes)} perfiles", flush=True)
    JSON_PERFILES.write_text(
        json.dumps({"by_h": by_h}, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    return by_h


def _mediana(vals: List[float]) -> float:
    if not vals:
        return 0.0
    s = sorted(vals)
    n = len(s)
    if n % 2:
        return float(s[n // 2])
    return 0.5 * (s[n // 2 - 1] + s[n // 2])


def _stats_noche(hourly: Dict[str, Dict[str, float]], nid: str, ctrl: date) -> Dict[str, float]:
    serie = hourly.get(nid) or {}
    pre, post = [], []
    for iso, v in serie.items():
        d = date.fromisoformat(iso)
        if nid == "000025-07" and d in PIZZA_NOCHES_ATIPICAS:
            continue
        if d < ctrl:
            pre.append(float(v))
        else:
            post.append(float(v))
    pre_m = _mediana(pre)
    post_m = _mediana(post)
    ahorro = max(0.0, pre_m - post_m)
    return {
        "pre": pre_m,
        "post": post_m,
        "ahorro_noche": ahorro,
        "n_pre": float(len(pre)),
        "n_post": float(len(post)),
        "ahorro_acum": ahorro * float(len(post)),
    }


def _ahorro_norte_mes(st: Dict[str, float]) -> float:
    """Ahorro mensual Norte: m³ de la noche × 5 h de control × 30 días."""
    return round(float(st.get("ahorro_noche") or 0.0), 1) * HORAS_CTRL_NORTE * 30.0


def _avg_daily(daily: Dict[str, float], d0: date, d1: date) -> float:
    vals = [
        float(v)
        for iso, v in daily.items()
        if d0 <= date.fromisoformat(iso) <= d1
    ]
    return (sum(vals) / len(vals)) if vals else 0.0


def chart_sur_diario(path: Path, daily: Dict[str, float]) -> Dict[str, float]:
    """Estanque Sur m³/día. Marca la reparación de presostatos (10/06)."""
    dias = sorted(date.fromisoformat(k) for k in daily)
    vals = [float(daily[d.isoformat()]) for d in dias]
    pre = _avg_daily(daily, DESDE, SUR_REPARACION - timedelta(days=1))
    post = _avg_daily(daily, SUR_REPARACION + timedelta(days=1), HASTA)
    dia10 = float(daily.get(SUR_REPARACION.isoformat()) or 0.0)
    n_post = sum(
        1
        for iso in daily
        if SUR_REPARACION + timedelta(days=1) <= date.fromisoformat(iso) <= HASTA
    )
    baja = max(0.0, pre - post)
    m3_acum = baja * float(n_post)

    fig, ax = plt.subplots(figsize=(7.15, 3.55), dpi=160)
    ax.axvspan(datetime(2026, 5, 1), datetime(2026, 6, 10), color="#F4E6C8", alpha=0.55, zorder=0)
    ax.axvspan(datetime(2026, 6, 10), datetime(2026, 8, 18), color="#E7F1F8", alpha=0.55, zorder=0)
    ax.plot(dias, vals, color=_hex(TEAL), linewidth=1.7, zorder=3)
    ax.fill_between(dias, vals, 0, color=_hex(TEAL), alpha=0.18, zorder=2)
    ax.axvline(SUR_REPARACION, color="#C04545", linestyle="--", linewidth=1.3, zorder=4)
    ax.axhline(pre, color="#C9A227", linestyle=":", linewidth=1.2, zorder=4)
    ax.axhline(post, color=_hex(NAVY), linestyle=":", linewidth=1.2, zorder=4)
    ax.plot([SUR_REPARACION], [dia10], marker="o", markersize=7, color="#C04545", zorder=5)
    ax.annotate(
        f"10/06  {fn(dia10, 0)} m³",
        xy=(SUR_REPARACION, dia10),
        xytext=(15, 18),
        textcoords="offset points",
        fontsize=9,
        color="#C04545",
        fontweight="bold",
        arrowprops=dict(arrowstyle="->", color="#C04545", lw=0.9),
    )
    ax.text(
        datetime(2026, 5, 12),
        pre + 4,
        f"antes  {fn(pre, 0)} m³/día",
        fontsize=8.5,
        color="#8A6A12",
        fontweight="bold",
    )
    ax.text(
        datetime(2026, 7, 5),
        post + 6,
        f"después  {fn(post, 0)} m³/día",
        fontsize=8.5,
        color=_hex(NAVY),
        fontweight="bold",
    )
    ax.text(
        datetime(2026, 7, 8),
        max(vals) * 0.78 if vals else 90,
        f"−{fn(baja, 0)} m³/día\n{_clp(baja)} / día\n{_clp(baja * 30)} / mes",
        fontsize=9,
        color=_hex(NAVY),
        fontweight="bold",
        va="top",
        ha="left",
        bbox=dict(boxstyle="round,pad=0.35", facecolor="white", edgecolor=_hex(GOLD), linewidth=1.2),
    )
    ax.set_ylabel("m³/día", fontsize=10, color=_hex(NAVY))
    ax.set_xlim(datetime(2026, 5, 1), datetime(2026, 8, 18))
    ax.set_ylim(0, max(vals) * 1.18 if vals else 1)
    ax.xaxis.set_major_formatter(mdates.DateFormatter("%d/%m"))
    ax.xaxis.set_major_locator(mdates.WeekdayLocator(interval=2))
    ax.tick_params(labelsize=8, colors=_hex(NAVY))
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#C5CDD6")
    ax.spines["bottom"].set_color("#C5CDD6")
    ax.yaxis.grid(True, linestyle=":", alpha=0.5, zorder=1)
    ax.set_axisbelow(True)
    fig.autofmt_xdate(rotation=30, ha="right")
    fig.tight_layout(pad=0.25)
    path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(path, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return {
        "pre": pre,
        "post": post,
        "dia10": dia10,
        "baja": baja,
        "n_post": n_post,
        "m3_acum": m3_acum,
    }


def chart_controles_nocturnos(path: Path, hourly: Dict[str, Dict[str, float]]) -> Dict[str, Dict[str, float]]:
    """Noche típica (00:00–06:00) antes vs con control. Pizza Hut y Estanque Norte."""
    st_p = _stats_noche(hourly, "000025-07", CTRL_PIZZA)
    st_n = _stats_noche(hourly, "000025-01", CTRL_NORTE)
    labels = ["Pizza Hut\ncontrol 01/07", "Estanque Norte\ncontrol 05/08"]
    x = np.arange(len(labels))
    w = 0.34
    fig, ax = plt.subplots(figsize=(5.55, 3.55), dpi=160)
    col_piz = _hex(COLOR_NODO["000025-07"])
    ax.bar(x[0] - w / 2, st_p["pre"], w, color="#8FA4B8", zorder=3)
    ax.bar(x[0] + w / 2, st_p["post"], w, color=col_piz, zorder=3)
    ax.bar(x[1] - w / 2, st_n["pre"], w, color="#8FA4B8", zorder=3)
    ax.bar(x[1] + w / 2, st_n["post"], w, color=_hex(GOLD), zorder=3)
    ax.set_xticks(x)
    ax.set_xticklabels(labels, fontsize=9, color=_hex(NAVY))
    ax.set_ylabel("m³ / noche (00:00–06:00)", fontsize=9, color=_hex(NAVY))
    ax.tick_params(axis="y", labelsize=8, colors=_hex(NAVY))
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#C5CDD6")
    ax.spines["bottom"].set_color("#C5CDD6")
    ax.yaxis.grid(True, linestyle=":", alpha=0.5, zorder=0)
    ax.set_axisbelow(True)
    ymax = max(st_p["pre"], st_p["post"], st_n["pre"], st_n["post"], 1.0) * 1.45
    ax.set_ylim(0, ymax)
    for xi, st in ((0, st_p), (1, st_n)):
        for dx, val in ((-w / 2, st["pre"]), (w / 2, st["post"])):
            ax.text(
                xi + dx,
                val + ymax * 0.03,
                fn(val, 1),
                ha="center",
                va="bottom",
                fontsize=10,
                color=_hex(NAVY),
                fontweight="bold",
            )
    # Ahorro Norte, sobre las barras del punto
    if st_n["ahorro_noche"] >= 0.2:
        mes_n = _ahorro_norte_mes(st_n)
        ax.annotate(
            f"ahorro {fn(st_n['ahorro_noche'], 1)} m³\n"
            f"× {HORAS_CTRL_NORTE} h × 30 d\n"
            f"{fn(mes_n, 0)} m³/mes\n{_clp(mes_n)}/mes",
            xy=(1 + w / 2, st_n["post"]),
            xytext=(1.18, ymax * 0.58),
            fontsize=8,
            color=_hex(NAVY),
            fontweight="bold",
            ha="left",
            arrowprops=dict(arrowstyle="->", color=_hex(GOLD), lw=0.9),
        )
    ax.legend(
        handles=[
            Patch(facecolor="#8FA4B8", edgecolor="none", label="Noche típica antes"),
            Patch(facecolor=col_piz, edgecolor="none", label="Pizza Hut con control"),
            Patch(facecolor=_hex(GOLD), edgecolor="none", label="Norte con control"),
        ],
        frameon=False,
        fontsize=8,
        loc="upper left",
        labelcolor=_hex(NAVY),
    )
    fig.tight_layout(pad=0.25)
    path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(path, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return {"000025-07": st_p, "000025-01": st_n}


def chart_horas_en_cero(
    path: Path,
    perfiles: Dict[str, Dict[str, Dict[str, float]]],
) -> Dict[str, Any]:
    """Mapa día × hora: dorado = en cero. Estanque Norte desde el 05/08."""
    cero = 0.05
    horas = list(range(7))  # 00–06
    norte = perfiles.get("000025-01") or {}
    dias = sorted(
        date.fromisoformat(k) for k in norte if date.fromisoformat(k) >= CTRL_NORTE
    )
    if not dias:
        dias = []
    mat = np.zeros((len(dias), len(horas)))
    for i, d in enumerate(dias):
        rec = norte.get(d.isoformat()) or {}
        for j, h in enumerate(horas):
            v = float(rec.get(str(h), 0.0) or 0.0)
            mat[i, j] = 1.0 if v < cero else 0.0  # 1 = en cero

    fig, ax = plt.subplots(figsize=(5.55, 3.55), dpi=160)
    cmap = ListedColormap(["#0D3B66", "#C9A227"])
    ax.imshow(mat, aspect="auto", cmap=cmap, vmin=0, vmax=1, interpolation="nearest")
    ax.set_xticks(range(len(horas)))
    ax.set_xticklabels([f"{h:02d}" for h in horas], fontsize=9, color=_hex(NAVY))
    ax.set_yticks(range(len(dias)))
    ax.set_yticklabels([d.strftime("%d/%m") for d in dias], fontsize=8, color=_hex(NAVY))
    ax.set_xlabel("Hora", fontsize=9, color=_hex(NAVY))
    ax.tick_params(colors=_hex(NAVY))
    for spine in ax.spines.values():
        spine.set_color("#C5CDD6")
    ax.set_xticks(np.arange(-0.5, len(horas), 1), minor=True)
    ax.set_yticks(np.arange(-0.5, len(dias), 1), minor=True)
    ax.grid(which="minor", color="white", linestyle="-", linewidth=1.1)
    ax.tick_params(which="minor", bottom=False, left=False)
    ax.legend(
        handles=[
            Patch(facecolor="#C9A227", edgecolor="none", label="En cero"),
            Patch(facecolor="#0D3B66", edgecolor="none", label="Con caudal"),
        ],
        loc="lower right",
        frameon=False,
        fontsize=8,
        bbox_to_anchor=(1.02, -0.22),
        ncol=2,
        labelcolor=_hex(NAVY),
    )
    fig.tight_layout(pad=0.35)
    path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(path, bbox_inches="tight", facecolor="white")
    plt.close(fig)

    # Horas que quedan en cero casi todas las noches (≥ 80 %).
    n = len(dias) or 1
    horas_cero = []
    for j, h in enumerate(horas):
        pct = 100.0 * float(mat[:, j].sum()) / n
        if pct >= 80:
            horas_cero.append(h)
    pizza = perfiles.get("000025-07") or {}
    pizza_dias = [
        date.fromisoformat(k)
        for k in pizza
        if date.fromisoformat(k) >= CTRL_PIZZA and date.fromisoformat(k) not in PIZZA_NOCHES_ATIPICAS
    ]
    pizza_cero = []
    if pizza_dias:
        for h in horas:
            c = 0
            for d in pizza_dias:
                v = float((pizza.get(d.isoformat()) or {}).get(str(h), 0) or 0)
                if v < cero:
                    c += 1
            if 100.0 * c / len(pizza_dias) >= 80:
                pizza_cero.append(h)
    return {"norte_horas_cero": horas_cero, "pizza_horas_cero": pizza_cero, "n_norte": len(dias)}


def chart_mensual(path: Path, tot: Dict[str, float], sin_mayo: bool = False) -> None:
    """Meses cerrados + agosto apilado (a la fecha + proyección al 31)."""
    cerrados: List[float] = []
    labels: List[str] = []
    if not sin_mayo:
        cerrados.append(float(tot["may"]))
        labels.append("Mayo")
    cerrados.extend([float(tot["jun"]), float(tot["jul"])])
    labels.extend(["Junio", "Julio"])
    ago, proy = float(tot["ago"]), float(tot["ago_proy"])
    resto = max(proy - ago, 0.0)
    labels.append(f"Agosto\n({AGO_ETQ} + proy.)")
    n = len(cerrados) + 1
    last = n - 1
    x = np.arange(n)
    w = 0.58

    fig, ax = plt.subplots(figsize=(9.4, 3.85), dpi=160)
    ax.bar(
        x[:last],
        cerrados,
        width=w,
        color=_hex(NAVY),
        zorder=3,
        label="Mes cerrado",
    )
    ax.bar(
        [last],
        [ago],
        width=w,
        color=_hex(GOLD),
        zorder=3,
        label=f"Agosto {AGO_ETQ} (a la fecha)",
    )
    ax.bar(
        [last],
        [resto],
        width=w,
        bottom=[ago],
        color=_hex(GOLD_SOFT),
        edgecolor=_hex(GOLD),
        linewidth=0.8,
        zorder=3,
        label="Proyección resto de agosto",
    )

    ax.set_xticks(x)
    ax.set_xticklabels(labels, fontsize=11, color=_hex(NAVY), fontweight="bold")
    ax.set_ylabel("m³", fontsize=11, color=_hex(NAVY))
    ax.tick_params(axis="y", labelsize=10, colors=_hex(NAVY))
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#C5CDD6")
    ax.spines["bottom"].set_color("#C5CDD6")
    ax.yaxis.grid(True, linestyle=":", alpha=0.55, zorder=0)
    ax.set_axisbelow(True)

    pico = max(cerrados + [proy, ago, 1.0])
    ymax = pico * 1.16
    ax.set_ylim(0, ymax)
    ax.set_xlim(-0.62, last + 1.15)

    for i, v in enumerate(cerrados):
        ax.text(
            i,
            v + ymax * 0.02,
            fn(v, 0),
            ha="center",
            va="bottom",
            fontsize=11,
            color=_hex(NAVY),
            fontweight="bold",
        )
    if ago >= ymax * 0.14:
        ax.text(
            last,
            ago * 0.50,
            fn(ago, 0),
            ha="center",
            va="center",
            fontsize=11,
            color="white",
            fontweight="bold",
        )
    else:
        ax.text(
            last - 0.38,
            max(ago, ymax * 0.04),
            fn(ago, 0),
            ha="right",
            va="center",
            fontsize=10,
            color=_hex(GOLD),
            fontweight="bold",
        )
    ax.annotate(
        f"proy. {fn(proy, 0)}",
        xy=(last + w / 2, proy),
        xytext=(last + 0.48, proy),
        fontsize=11,
        color=_hex(NAVY),
        fontweight="bold",
        va="center",
        ha="left",
        arrowprops=dict(arrowstyle="-", color=_hex(GOLD), lw=0.8),
    )

    leg = ax.legend(
        loc="upper center",
        bbox_to_anchor=(0.48, -0.18),
        ncol=3,
        frameon=False,
        fontsize=8.5,
        labelcolor=_hex(NAVY),
    )
    for t in leg.get_texts():
        t.set_color(_hex(NAVY))

    fig.tight_layout(pad=0.25)
    fig.subplots_adjust(bottom=0.20)
    path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(path, bbox_inches="tight", facecolor="white")
    plt.close(fig)


def chart_pak_perfil_10ago(path: Path, perfil: Dict[str, Dict[str, float]]) -> Dict[str, float]:
    """Perfil 00:00–23:00 del 10/08 para 27, 35 y 36."""
    horas = list(range(24))
    fig, ax = plt.subplots(figsize=(6.55, 2.85), dpi=160)
    n06: Dict[str, float] = {}
    for nid in PAK_CADENA:
        rec = perfil.get(nid) or {}
        ys = [float(rec.get(str(h), rec.get(h, 0.0)) or 0.0) for h in horas]
        n06[nid] = sum(ys[h] for h in horas if h < 6)
        ax.plot(
            horas,
            ys,
            color=_hex(COLOR_NODO[nid]),
            linewidth=2.0,
            marker="o",
            markersize=3.5,
            label=NOMBRE_CORTO[nid],
            zorder=3,
        )
    ax.axvspan(-0.4, 5.5, color="#F4E6C8", alpha=0.55, zorder=0)
    ax.set_xlim(-0.4, 23.4)
    ax.set_xticks([0, 3, 6, 9, 12, 15, 18, 21, 23])
    ax.set_xticklabels(
        [f"{h:02d}" for h in (0, 3, 6, 9, 12, 15, 18, 21, 23)],
        fontsize=8,
        color=_hex(NAVY),
    )
    ax.set_ylabel("m³/h", fontsize=10, color=_hex(NAVY))
    ax.set_xlabel("Hora del 10/08", fontsize=9, color=_hex(NAVY))
    ax.tick_params(axis="y", labelsize=8, colors=_hex(NAVY))
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#C5CDD6")
    ax.spines["bottom"].set_color("#C5CDD6")
    ax.yaxis.grid(True, linestyle=":", alpha=0.5, zorder=1)
    ax.set_axisbelow(True)
    ax.legend(frameon=False, fontsize=8, loc="upper right", labelcolor=_hex(NAVY))
    ymax = ax.get_ylim()[1]
    ax.text(2.4, ymax * 0.92 if ymax else 1, "madrugada 00–06", fontsize=8, color="#8A6A12", fontweight="bold", ha="center")
    fig.tight_layout(pad=0.25)
    path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(path, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return n06


def chart_pak_noches_julio(path: Path, n06: Dict[str, Dict[str, float]]) -> Dict[str, float]:
    """m³ 00:00–06:00: Bazar Gourmet y DL Kennedy (julio–agosto)."""
    ramales = [BAZAR, DL_KENNEDY]
    dias = _rango_dias(JUL_NOCHE_D0, HASTA)
    fig, ax = plt.subplots(figsize=(6.55, 2.85), dpi=160)
    med: Dict[str, float] = {}
    for nid in ramales:
        ys = [float((n06.get(nid) or {}).get(d.isoformat(), 0.0) or 0.0) for d in dias]
        med[nid] = _mediana(ys)
        ax.plot(
            dias,
            ys,
            color=_hex(COLOR_NODO[nid]),
            linewidth=1.7,
            marker="o",
            markersize=3.2,
            label=NOMBRE_CORTO[nid],
            zorder=3,
        )
    ax.set_ylabel("m³ / noche (00:00–06:00)", fontsize=9, color=_hex(NAVY))
    ax.set_xlabel("Julio – agosto 2026", fontsize=9, color=_hex(NAVY))
    ax.xaxis.set_major_locator(mdates.WeekdayLocator(byweekday=mdates.MO))
    ax.xaxis.set_major_formatter(mdates.DateFormatter("%d/%m"))
    ax.tick_params(labelsize=8, colors=_hex(NAVY))
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#C5CDD6")
    ax.spines["bottom"].set_color("#C5CDD6")
    ax.yaxis.grid(True, linestyle=":", alpha=0.5, zorder=1)
    ax.set_axisbelow(True)
    ax.legend(frameon=False, fontsize=8, loc="upper left", labelcolor=_hex(NAVY))
    fig.tight_layout(pad=0.25)
    path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(path, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return med


def chart_mam_placa_perfil(path: Path, rec: Dict[str, float]) -> float:
    """Perfil 10/08 de Placa Bancaria; sombra 00:00–06:00."""
    horas = list(range(24))
    ys = [float(rec.get(str(h), rec.get(h, 0.0)) or 0.0) for h in horas]
    n06 = sum(ys[h] for h in horas if h < 6)
    fig, ax = plt.subplots(figsize=(6.55, 2.85), dpi=160)
    ax.plot(
        horas,
        ys,
        color=_hex(COLOR_NODO[PLACA]),
        linewidth=2.0,
        marker="o",
        markersize=3.5,
        zorder=3,
        label="Placa Bancaria",
    )
    ax.axvspan(-0.4, 5.5, color="#F4E6C8", alpha=0.55, zorder=0)
    ax.set_xlim(-0.4, 23.4)
    ax.set_xticks([0, 3, 6, 9, 12, 15, 18, 21, 23])
    ax.set_xticklabels(
        [f"{h:02d}" for h in (0, 3, 6, 9, 12, 15, 18, 21, 23)],
        fontsize=8,
        color=_hex(NAVY),
    )
    ax.set_ylabel("m³/h", fontsize=10, color=_hex(NAVY))
    ax.set_xlabel("Hora del 10/08", fontsize=9, color=_hex(NAVY))
    ax.tick_params(axis="y", labelsize=8, colors=_hex(NAVY))
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#C5CDD6")
    ax.spines["bottom"].set_color("#C5CDD6")
    ax.yaxis.grid(True, linestyle=":", alpha=0.5, zorder=1)
    ax.set_axisbelow(True)
    ax.legend(frameon=False, fontsize=8, loc="upper left", labelcolor=_hex(NAVY))
    ymax = ax.get_ylim()[1]
    ax.text(
        2.4,
        ymax * 0.92 if ymax else 1,
        f"madrugada 00–06: {fn(n06, 1)} m³",
        fontsize=8,
        color="#8A6A12",
        fontweight="bold",
        ha="center",
    )
    fig.tight_layout(pad=0.25)
    path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(path, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return n06


def chart_mam_placa_vs_falabella(
    path: Path, daily_p: Dict[str, float], daily_f: Dict[str, float]
) -> None:
    """m³/día Placa vs Falabella desde el 11/08."""
    dias = _rango_dias(FALABELLA_DESDE, HASTA)
    x = np.arange(len(dias))
    w = 0.36
    yp = [float(daily_p.get(d.isoformat(), 0.0) or 0.0) for d in dias]
    yf = [float(daily_f.get(d.isoformat(), 0.0) or 0.0) for d in dias]
    fig, ax = plt.subplots(figsize=(6.55, 2.85), dpi=160)
    ax.bar(x - w / 2, yp, w, color=_hex(COLOR_NODO[PLACA]), zorder=3, label="Placa Bancaria")
    ax.bar(x + w / 2, yf, w, color=_hex(COLOR_NODO[FALABELLA]), zorder=3, label="Falabella")
    ax.set_xticks(x)
    ax.set_xticklabels([d.strftime("%d/%m") for d in dias], fontsize=8, color=_hex(NAVY))
    ax.set_ylabel("m³/día", fontsize=10, color=_hex(NAVY))
    ax.tick_params(axis="y", labelsize=8, colors=_hex(NAVY))
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#C5CDD6")
    ax.spines["bottom"].set_color("#C5CDD6")
    ax.yaxis.grid(True, linestyle=":", alpha=0.5, zorder=1)
    ax.set_axisbelow(True)
    ax.legend(frameon=False, fontsize=8, loc="upper left", labelcolor=_hex(NAVY))
    fig.tight_layout(pad=0.25)
    path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(path, bbox_inches="tight", facecolor="white")
    plt.close(fig)


def chart_maq_matriz_diario(path: Path, daily: Dict[str, float]) -> None:
    """m³/día de Matriz Principal desde junio: el alza del 22/06 se sostiene."""
    dias = _rango_dias(date(2026, 6, 1), HASTA)
    ys = [float(daily.get(d.isoformat(), 0.0) or 0.0) for d in dias]
    fig, ax = plt.subplots(figsize=(6.55, 2.85), dpi=160)
    ax.axvspan(MAQ_ALZA, dias[-1] + timedelta(days=1), color="#F4E6C8", alpha=0.45, zorder=0)
    ax.plot(
        dias,
        ys,
        color=_hex(COLOR_NODO[MATRIZ_MAQ]),
        linewidth=1.8,
        zorder=3,
    )
    ax.axvline(MAQ_ALZA, color=_hex(GOLD), linestyle="--", linewidth=1.1, zorder=4)
    ax.set_ylabel("m³/día", fontsize=10, color=_hex(NAVY))
    ax.set_xlabel("Junio – agosto 2026", fontsize=9, color=_hex(NAVY))
    ax.xaxis.set_major_locator(mdates.WeekdayLocator(byweekday=mdates.MO))
    ax.xaxis.set_major_formatter(mdates.DateFormatter("%d/%m"))
    ax.tick_params(labelsize=8, colors=_hex(NAVY))
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#C5CDD6")
    ax.spines["bottom"].set_color("#C5CDD6")
    ax.yaxis.grid(True, linestyle=":", alpha=0.5, zorder=1)
    ax.set_axisbelow(True)
    ymax = ax.get_ylim()[1]
    ax.text(
        MAQ_ALZA + timedelta(days=2),
        ymax * 0.12 if ymax else 1,
        "22/06 alza",
        fontsize=8,
        color="#8A6A12",
        fontweight="bold",
    )
    fig.tight_layout(pad=0.25)
    path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(path, bbox_inches="tight", facecolor="white")
    plt.close(fig)


def chart_maq_matriz_noches(path: Path, n06: Dict[str, float]) -> Dict[str, float]:
    """m³ 00:00–06:00 de Matriz Principal desde 01/06: se ve el alza nocturna."""
    dias = _rango_dias(MAQ_NOCHE_D0, HASTA)
    ys = [float(n06.get(d.isoformat(), 0.0) or 0.0) for d in dias]
    pre = [v for d, v in zip(dias, ys) if d < MAQ_ALZA]
    post = [v for d, v in zip(dias, ys) if d >= MAQ_ALZA]
    med_pre = _mediana(pre)
    med_post = _mediana(post)
    fig, ax = plt.subplots(figsize=(6.55, 2.85), dpi=160)
    ax.axvspan(MAQ_ALZA, dias[-1] + timedelta(days=1), color="#F4E6C8", alpha=0.45, zorder=0)
    ax.plot(
        dias,
        ys,
        color=_hex(COLOR_NODO[MATRIZ_MAQ]),
        linewidth=1.7,
        marker="o",
        markersize=2.6,
        zorder=3,
        label="Matriz Principal",
    )
    ax.axvline(MAQ_ALZA, color=_hex(GOLD), linestyle="--", linewidth=1.1, zorder=4)
    ax.set_ylabel("m³ / noche (00:00–06:00)", fontsize=9, color=_hex(NAVY))
    ax.set_xlabel("Junio – agosto 2026", fontsize=9, color=_hex(NAVY))
    ax.xaxis.set_major_locator(mdates.WeekdayLocator(byweekday=mdates.MO))
    ax.xaxis.set_major_formatter(mdates.DateFormatter("%d/%m"))
    ax.tick_params(labelsize=8, colors=_hex(NAVY))
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#C5CDD6")
    ax.spines["bottom"].set_color("#C5CDD6")
    ax.yaxis.grid(True, linestyle=":", alpha=0.5, zorder=1)
    ax.set_axisbelow(True)
    ymax = ax.get_ylim()[1]
    ax.text(
        MAQ_ALZA + timedelta(days=2),
        ymax * 0.12 if ymax else 1,
        "22/06 alza",
        fontsize=8,
        color="#8A6A12",
        fontweight="bold",
    )
    ax.legend(frameon=False, fontsize=8, loc="upper right", labelcolor=_hex(NAVY))
    fig.tight_layout(pad=0.25)
    path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(path, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return {"med": med_post, "med_pre": med_pre, "med_post": med_post}


def chart_bom_si500_noches(path: Path, n06_500: Dict[str, float]) -> None:
    """m³ 00:00–06:00 de San Ignacio 500: el corte del 17/07 se sostiene."""
    dias = _rango_dias(JUL_NOCHE_D0, HASTA)
    ys500 = [float(n06_500.get(d.isoformat(), 0.0) or 0.0) for d in dias]
    fig, ax = plt.subplots(figsize=(6.55, 2.85), dpi=160)
    ax.axvspan(CTRL_SI500, dias[-1] + timedelta(days=1), color="#E7F1F8", alpha=0.65, zorder=0)
    ax.plot(
        dias,
        ys500,
        color=_hex(COLOR_NODO[SI500]),
        linewidth=1.8,
        marker="o",
        markersize=3.0,
        zorder=3,
        label="San Ignacio 500",
    )
    ax.axvline(CTRL_SI500, color=_hex(GOLD), linestyle="--", linewidth=1.2, zorder=4)
    ax.set_ylabel("m³ / noche (00:00–06:00)", fontsize=9, color=_hex(NAVY))
    ax.set_xlabel("Julio – agosto 2026", fontsize=9, color=_hex(NAVY))
    ax.xaxis.set_major_locator(mdates.WeekdayLocator(byweekday=mdates.MO))
    ax.xaxis.set_major_formatter(mdates.DateFormatter("%d/%m"))
    ax.tick_params(labelsize=8, colors=_hex(NAVY))
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#C5CDD6")
    ax.spines["bottom"].set_color("#C5CDD6")
    ax.yaxis.grid(True, linestyle=":", alpha=0.5, zorder=1)
    ax.set_axisbelow(True)
    ymax = ax.get_ylim()[1]
    ax.text(
        CTRL_SI500 + timedelta(days=1),
        ymax * 0.88 if ymax else 1,
        "17/07 control",
        fontsize=8,
        color="#8A6A12",
        fontweight="bold",
    )
    ax.legend(frameon=False, fontsize=7.5, loc="upper right", labelcolor=_hex(NAVY))
    fig.tight_layout(pad=0.25)
    path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(path, bbox_inches="tight", facecolor="white")
    plt.close(fig)


def chart_bom_si500_ahorro(path: Path, st: Dict[str, float]) -> None:
    """Noche típica SI500 antes vs con control, con ahorro en $."""
    fig, ax = plt.subplots(figsize=(6.55, 2.85), dpi=160)
    labels = ["Antes\n(10–16/07)", "Con control\n(17/07–17/08)"]
    vals = [st["pre"], st["post"]]
    cols = ["#8FA4B8", _hex(COLOR_NODO[SI500])]
    bars = ax.bar(labels, vals, color=cols, width=0.55, zorder=3)
    ax.set_ylabel("m³ / noche (00:00–06:00)", fontsize=9, color=_hex(NAVY))
    ax.tick_params(labelsize=9, colors=_hex(NAVY))
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#C5CDD6")
    ax.spines["bottom"].set_color("#C5CDD6")
    ax.yaxis.grid(True, linestyle=":", alpha=0.5, zorder=0)
    ax.set_axisbelow(True)
    ymax = max(vals + [1.0]) * 1.55
    ax.set_ylim(0, ymax)
    for bar, val in zip(bars, vals):
        ax.text(
            bar.get_x() + bar.get_width() / 2,
            val + ymax * 0.03,
            fn(val, 1),
            ha="center",
            va="bottom",
            fontsize=12,
            color=_hex(NAVY),
            fontweight="bold",
        )
    if st["ahorro_noche"] >= 0.2:
        ax.annotate(
            f"ahorro {fn(st['ahorro_noche'], 1)} m³/noche\n"
            f"{_clp(st['ahorro_noche'])}/noche\n"
            f"{_clp(st['ahorro_noche'] * 30)}/mes",
            xy=(1, st["post"]),
            xytext=(1.18, ymax * 0.55),
            fontsize=9,
            color=_hex(NAVY),
            fontweight="bold",
            ha="left",
            arrowprops=dict(arrowstyle="->", color=_hex(GOLD), lw=0.9),
        )
    fig.tight_layout(pad=0.25)
    path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(path, bbox_inches="tight", facecolor="white")
    plt.close(fig)


def chart_aeb_dia_noche(
    path: Path,
    daily: Dict[str, float],
    n06: Dict[str, float],
    *,
    color_dia: Tuple[int, int, int],
    umbral: float,
    umbral_noche: bool,
    umbral_etq: str,
) -> Dict[str, List[float]]:
    """Barras por mes: promedio día + promedio noche 00–06, con umbral."""
    meses = [
        (2026, 5, "Mayo"),
        (2026, 6, "Junio"),
        (2026, 7, "Julio"),
        (2026, 8, f"Agosto\n1–{HASTA.day}"),
    ]
    dias_avg = [_prom_mes(daily, y, m) for y, m, _ in meses]
    noc_avg = [_prom_mes(n06, y, m, activos=daily) for y, m, _ in meses]
    labels = [lab for _, _, lab in meses]
    x = np.arange(len(meses))
    w = 0.36
    fig, ax = plt.subplots(figsize=(6.55, 2.85), dpi=160)
    ax.bar(x - w / 2, dias_avg, w, color=_hex(color_dia), zorder=3, label="Promedio día")
    ax.bar(x + w / 2, noc_avg, w, color=_hex(GOLD), zorder=3, label="Promedio noche 00–06")
    ax.axhline(
        umbral,
        color="#C04545" if umbral_noche else _hex(GOLD),
        linestyle=":",
        linewidth=1.2,
        zorder=4,
        label=umbral_etq,
    )
    ax.set_xticks(x)
    ax.set_xticklabels(labels, fontsize=9, color=_hex(NAVY))
    ax.set_ylabel("m³ / día  ·  m³ / noche", fontsize=9, color=_hex(NAVY))
    ax.tick_params(axis="y", labelsize=8, colors=_hex(NAVY))
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#C5CDD6")
    ax.spines["bottom"].set_color("#C5CDD6")
    ax.yaxis.grid(True, linestyle=":", alpha=0.5, zorder=0)
    ax.set_axisbelow(True)
    ymax = max(dias_avg + noc_avg + [umbral, 1.0]) * 1.28
    ax.set_ylim(0, ymax)
    for xi, d, n in zip(x, dias_avg, noc_avg):
        ax.text(xi - w / 2, d + ymax * 0.02, fn(d, 0), ha="center", va="bottom", fontsize=8, color=_hex(NAVY), fontweight="bold")
        ax.text(xi + w / 2, n + ymax * 0.02, fn(n, 1), ha="center", va="bottom", fontsize=8, color=_hex(NAVY), fontweight="bold")
    ax.legend(frameon=False, fontsize=7.5, loc="upper left", labelcolor=_hex(NAVY))
    fig.tight_layout(pad=0.25)
    path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(path, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return {"dia": dias_avg, "noche": noc_avg}


def chart_cur_noches(path: Path, n06_s: Dict[str, float], n06_n: Dict[str, float]) -> Dict[str, float]:
    """m³ 00:00–06:00: Anillo Sur y Norte, junio–agosto."""
    dias = _rango_dias(CUR_NOCHE_D0, HASTA)
    ys = [float(n06_s.get(d.isoformat(), 0.0) or 0.0) for d in dias]
    yn = [float(n06_n.get(d.isoformat(), 0.0) or 0.0) for d in dias]
    fig, ax = plt.subplots(figsize=(6.55, 2.85), dpi=160)
    ax.plot(
        dias,
        ys,
        color=_hex(COLOR_NODO[ANILLO_SUR]),
        linewidth=1.7,
        marker="o",
        markersize=2.6,
        zorder=3,
        label="Anillo Sur",
    )
    ax.plot(
        dias,
        yn,
        color=_hex(COLOR_NODO[ANILLO_NORTE]),
        linewidth=1.7,
        marker="o",
        markersize=2.6,
        zorder=3,
        label="Anillo Norte",
    )
    ax.set_ylabel("m³ / noche (00:00–06:00)", fontsize=9, color=_hex(NAVY))
    ax.set_xlabel("Junio – agosto 2026", fontsize=9, color=_hex(NAVY))
    ax.xaxis.set_major_locator(mdates.WeekdayLocator(byweekday=mdates.MO))
    ax.xaxis.set_major_formatter(mdates.DateFormatter("%d/%m"))
    ax.tick_params(labelsize=8, colors=_hex(NAVY))
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#C5CDD6")
    ax.spines["bottom"].set_color("#C5CDD6")
    ax.yaxis.grid(True, linestyle=":", alpha=0.5, zorder=1)
    ax.set_axisbelow(True)
    ax.legend(frameon=False, fontsize=8, loc="upper right", labelcolor=_hex(NAVY))
    fig.tight_layout(pad=0.25)
    path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(path, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return {"sur": _mediana(ys), "norte": _mediana(yn)}


def chart_cur_dia(path: Path, daily_s: Dict[str, float], daily_n: Dict[str, float]) -> Dict[str, List[float]]:
    """Promedio día por mes: Sur un poco sobre Norte. Umbrales del total del día."""
    meses = [
        (2026, 6, "Junio"),
        (2026, 7, "Julio"),
        (2026, 8, f"Agosto\n1–{HASTA.day}"),
    ]
    sur = [_prom_mes(daily_s, y, m) for y, m, _ in meses]
    nor = [_prom_mes(daily_n, y, m) for y, m, _ in meses]
    labels = [lab for _, _, lab in meses]
    x = np.arange(len(meses))
    w = 0.36
    fig, ax = plt.subplots(figsize=(6.55, 2.85), dpi=160)
    ax.bar(x - w / 2, sur, w, color=_hex(COLOR_NODO[ANILLO_SUR]), zorder=3, label="Anillo Sur")
    ax.bar(x + w / 2, nor, w, color=_hex(COLOR_NODO[ANILLO_NORTE]), zorder=3, label="Anillo Norte")
    ax.axhline(UMBRAL_SUR_DIA, color=_hex(COLOR_NODO[ANILLO_SUR]), linestyle=":", linewidth=1.1, zorder=4, label=f"Umbral Sur {fn(UMBRAL_SUR_DIA, 0)}")
    ax.axhline(UMBRAL_NORTE_DIA, color=_hex(COLOR_NODO[ANILLO_NORTE]), linestyle=":", linewidth=1.1, zorder=4, label=f"Umbral Norte {fn(UMBRAL_NORTE_DIA, 0)}")
    ax.set_xticks(x)
    ax.set_xticklabels(labels, fontsize=9, color=_hex(NAVY))
    ax.set_ylabel("m³ / día", fontsize=9, color=_hex(NAVY))
    ax.tick_params(axis="y", labelsize=8, colors=_hex(NAVY))
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#C5CDD6")
    ax.spines["bottom"].set_color("#C5CDD6")
    ax.yaxis.grid(True, linestyle=":", alpha=0.5, zorder=0)
    ax.set_axisbelow(True)
    ymax = max(sur + nor + [UMBRAL_SUR_DIA, 1.0]) * 1.35
    ax.set_ylim(0, ymax)
    for xi, s, n in zip(x, sur, nor):
        ax.text(xi - w / 2, s + ymax * 0.02, fn(s, 1), ha="center", va="bottom", fontsize=8, color=_hex(NAVY), fontweight="bold")
        ax.text(xi + w / 2, n + ymax * 0.02, fn(n, 1), ha="center", va="bottom", fontsize=8, color=_hex(NAVY), fontweight="bold")
    ax.legend(frameon=False, fontsize=7.5, loc="upper right", labelcolor=_hex(NAVY), ncol=2)
    fig.tight_layout(pad=0.25)
    path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(path, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return {"sur": sur, "norte": nor}


def chart_mam_placa_noches(
    path: Path, n06_p: Dict[str, float], n06_f: Dict[str, float]
) -> float:
    """m³ 00:00–06:00: Placa (jul–ago) y Falabella (desde el 11/08)."""
    dias = _rango_dias(JUL_NOCHE_D0, HASTA)
    ys_p = [float(n06_p.get(d.isoformat(), 0.0) or 0.0) for d in dias]
    ys_f = [
        float(n06_f.get(d.isoformat(), 0.0) or 0.0) if d >= FALABELLA_DESDE else float("nan")
        for d in dias
    ]
    med = _mediana(ys_p)
    fig, ax = plt.subplots(figsize=(6.55, 2.85), dpi=160)
    ax.plot(
        dias,
        ys_p,
        color=_hex(COLOR_NODO[PLACA]),
        linewidth=1.7,
        marker="o",
        markersize=3.0,
        zorder=3,
        label="Placa Bancaria",
    )
    ax.plot(
        dias,
        ys_f,
        color=_hex(COLOR_NODO[FALABELLA]),
        linewidth=1.8,
        marker="o",
        markersize=3.5,
        zorder=4,
        label="Falabella",
    )
    ax.set_ylabel("m³ / noche (00:00–06:00)", fontsize=9, color=_hex(NAVY))
    ax.set_xlabel("Julio – agosto 2026", fontsize=9, color=_hex(NAVY))
    ax.xaxis.set_major_locator(mdates.WeekdayLocator(byweekday=mdates.MO))
    ax.xaxis.set_major_formatter(mdates.DateFormatter("%d/%m"))
    ax.tick_params(labelsize=8, colors=_hex(NAVY))
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#C5CDD6")
    ax.spines["bottom"].set_color("#C5CDD6")
    ax.yaxis.grid(True, linestyle=":", alpha=0.5, zorder=1)
    ax.set_axisbelow(True)
    ax.legend(frameon=False, fontsize=8, loc="upper right", labelcolor=_hex(NAVY))
    fig.tight_layout(pad=0.25)
    path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(path, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return med


def _set_run(run, text: str, size: int, bold: bool = False, color=NAVY) -> None:
    run.text = text
    run.font.size = PptPt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)
    run.font.name = "Calibri"


def _caja(slide, l, t, w, h, fill=LIGHT, line=None):
    sh = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        PptInches(l),
        PptInches(t),
        PptInches(w),
        PptInches(h),
    )
    sh.fill.solid()
    sh.fill.fore_color.rgb = _rgb(fill)
    sh.line.fill.background()
    if line:
        sh.line.fill.solid()
        sh.line.color.rgb = _rgb(line)
        sh.line.width = Emu(6350)
    try:
        sh.adjustments[0] = 0.08
    except Exception:
        pass
    return sh


def _tb(slide, l, t, w, h, lines: List[Tuple[str, int, bool, Tuple[int, int, int]]], align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(PptInches(l), PptInches(t), PptInches(w), PptInches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, (text, size, bold, color) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = PptPt(3)
        run = p.add_run()
        _set_run(run, text, size, bold, color)
    return box


def _vineta_hallazgos(slide, lineas: List[Tuple[str, int, bool, Tuple[int, int, int]]], h: float = 1.18) -> None:
    """Caja superior de hallazgos: título naranja HALLAZGOS + texto."""
    _caja(slide, 0.22, 1.08, 12.88, h, fill=(255, 249, 235), line=GOLD)
    _tb(slide, 0.40, 1.12, 12.55, 0.20, [("HALLAZGOS", 11, True, GOLD)])
    _tb(slide, 0.40, 1.32, 12.55, max(h - 0.28, 0.36), lineas)


def _umbral_dia_de(vol_mes: float, dias: int = 31, factor: float = 1.25) -> float:
    """Umbral del total de 24 h: ~1,25 × promedio día, redondeado."""
    if vol_mes <= 0 or dias <= 0:
        return 0.0
    crudo = (float(vol_mes) / float(dias)) * factor
    if crudo >= 40:
        return float(int(round(crudo / 5.0) * 5.0))
    return float(int(round(crudo)))


def _header_bar(slide, prs, titulo: str, sub: str) -> None:
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, PptInches(0.92)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb(NAVY)
    bar.line.fill.background()
    gold = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, PptInches(0.92), prs.slide_width, PptInches(0.06)
    )
    gold.fill.solid()
    gold.fill.fore_color.rgb = _rgb(GOLD)
    gold.line.fill.background()
    if LOGO.is_file():
        slide.shapes.add_picture(str(LOGO), PptInches(11.85), PptInches(0.16), width=PptInches(1.25))
    _tb(slide, 0.28, 0.12, 11.6, 0.42, [(titulo, 20, True, WHITE)])
    _tb(slide, 0.28, 0.50, 11.6, 0.36, [(sub, 11, False, (220, 230, 240))])


def _portada(prs) -> None:
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    if FONDO.is_file():
        pic = sl.shapes.add_picture(str(FONDO), 0, 0, width=prs.slide_width, height=prs.slide_height)
        spTree = sl.shapes._spTree
        spTree.remove(pic.element)
        spTree.insert(2, pic.element)
    veil = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, PptInches(3.55), prs.slide_width, PptInches(3.95))
    veil.fill.solid()
    veil.fill.fore_color.rgb = _rgb(NAVY)
    veil.line.fill.background()
    _tb(sl, 0.6, 3.75, 12, 0.5, [("WES  ·  Parque Arauco", 16, True, GOLD)])
    _tb(sl, 0.6, 4.18, 12, 0.7, [("Recorrido ejecutivo por recinto", 32, True, WHITE)])
    _tb(
        sl,
        0.6,
        4.95,
        12,
        1.2,
        [
            ("Una lámina de presentación por recinto", 16, False, WHITE),
            (f"Período {PERIODO}   |   Emisión {FECHA_EMISION}", 15, False, (220, 230, 240)),
            ("MAE  ·  MAM  ·  MAQ  ·  BOM  ·  AEB  ·  CUR  ·  PAK", 16, False, GOLD),
        ],
    )
    if LOGO.is_file():
        sl.shapes.add_picture(str(LOGO), PptInches(11.70), PptInches(6.85), width=PptInches(1.35))


def _slide_hallazgos(
    prs,
    by: Dict[str, Dict[str, Any]],
    hourly: Dict[str, Dict[str, float]],
) -> None:
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _header_bar(
        sl,
        prs,
        "MAE  ·  Hallazgos",
        "Estanque Sur: reparación de presostatos (10/06)  ·  Estanque Norte: control nocturno",
    )

    daily_sur = (by.get("000025-19") or {}).get("daily") or {}
    ch_sur = CHARTS / "mae_sur_diario_presostatos.png"
    ch_noc = CHARTS / "mae_controles_nocturnos.png"
    st_sur = chart_sur_diario(ch_sur, daily_sur)
    st_noc = chart_controles_nocturnos(ch_noc, hourly)
    baja = st_sur["baja"]
    pct = (baja / st_sur["pre"] * 100) if st_sur["pre"] else 0.0
    st_n = st_noc["000025-01"]
    st_p = st_noc["000025-07"]

    _vineta_hallazgos(
        sl,
        [
            (
                f"Estanque Norte (control 05/08, 00:00–05:00): noche típica "
                f"{fn(st_n['pre'], 1)} → {fn(st_n['post'], 1)} m³, ahorro "
                f"{fn(st_n['ahorro_noche'], 1)} m³. × {HORAS_CTRL_NORTE} h de control × 30 días = "
                f"{fn(_ahorro_norte_mes(st_n), 0)} m³/mes = {_clp(_ahorro_norte_mes(st_n))}/mes "
                f"(tarifa ${fn(TARIFA_CLP_M3, 0)}/m³). Es el que más ahorra de los controles MAE.",
                12,
                True,
                NAVY,
            ),
            (
                f"Estanque Sur: el 10/06 se repararon los presostatos. Venía en "
                f"{fn(st_sur['pre'], 0)} m³/día (mayo–9/jun) y ese día ya marca "
                f"{fn(st_sur['dia10'], 0)} m³. Desde el 11/06 se sostiene en "
                f"{fn(st_sur['post'], 0)} m³/día (−{fn(baja, 0)} m³/día, {fn(pct, 0)}%). "
                f"A ${fn(TARIFA_CLP_M3, 0)}/m³: {_clp(baja)}/día  ·  {_clp(baja * 30)}/mes  ·  "
                f"acumulado 11/06–{HASTA.strftime('%d/%m')}: {_clp(st_sur['m3_acum'])} "
                f"({fn(st_sur['m3_acum'], 0)} m³ en {int(st_sur['n_post'])} días).",
                12,
                False,
                NAVY,
            ),
            (
                f"Umbrales a activar (total 24 h): Estanque Norte {fn(UMBRAL_MAE_NORTE_DIA, 0)} m³/día  ·  "
                f"Estanque Sur {fn(UMBRAL_MAE_SUR_DIA, 0)} m³/día  ·  "
                f"Pizza Hut {fn(UMBRAL_MAE_PIZZA_DIA, 0)} m³/día  ·  "
                f"Baños Públicos {fn(UMBRAL_MAE_BANOS_DIA, 0)} m³/día.",
                12,
                True,
                NAVY,
            ),
        ],
        h=1.62,
    )

    _caja(sl, 0.22, 2.80, 5.78, 3.32)
    _tb(sl, 0.36, 2.84, 5.50, 0.22, [("CONTROLES NOCTURNOS — noche típica antes vs con control", 11, True, TEAL)])
    sl.shapes.add_picture(
        str(ch_noc), PptInches(0.36), PptInches(3.08), width=PptInches(5.48), height=PptInches(2.92)
    )

    _caja(sl, 6.14, 2.80, 6.96, 3.32)
    _tb(sl, 6.28, 2.84, 6.68, 0.22, [("ESTANQUE SUR — m³/día y costo evitado", 11, True, TEAL)])
    sl.shapes.add_picture(
        str(ch_sur), PptInches(6.28), PptInches(3.08), width=PptInches(6.68), height=PptInches(2.92)
    )

    _tb(
        sl,
        0.28,
        6.22,
        12.8,
        1.10,
        [
            (
                f"Pizza Hut (control 01/07): noche típica {fn(st_p['pre'], 1)} → "
                f"{fn(st_p['post'], 1)} m³.",
                12,
                False,
                NAVY,
            ),
            (
                "Baños Guardias: control instalado; no está solicitado el uso.",
                12,
                False,
                NAVY,
            ),
        ],
    )


def _chip_layout(n: int) -> Tuple[float, List[Tuple[float, float, float, float]]]:
    """Posiciones de tiras (x, y, w, h) y Y de inicio del gráfico."""
    y0, left, total_w, gap = 1.12, 0.22, 12.90, 0.08
    if n <= 5:
        h = 0.72
        w = (total_w - gap * (n - 1)) / max(n, 1)
        pos = [(left + i * (w + gap), y0, w, h) for i in range(n)]
        return 1.96, pos
    n1 = (n + 1) // 2
    n2 = n - n1
    h = 0.68 if n >= 8 else 0.58
    gap = 0.10 if n >= 8 else gap
    w1 = (total_w - gap * (n1 - 1)) / n1
    pos = [(left + i * (w1 + gap), y0, w1, h) for i in range(n1)]
    y1 = y0 + h + 0.08
    w2 = (total_w - gap * max(n2 - 1, 0)) / max(n2, 1)
    row2_w = n2 * w2 + max(n2 - 1, 0) * gap
    x2 = left + (total_w - row2_w) / 2
    pos += [(x2 + i * (w2 + gap), y1, w2, h) for i in range(n2)]
    return y1 + h + 0.10, pos


def _fit_picture(slide, path: Path, l: float, t: float, max_w: float, max_h: float) -> None:
    """Inserta la imagen dentro del recuadro, sin recortar ni pisar el pie."""
    from PIL import Image

    with Image.open(path) as im:
        iw, ih = im.size
    aspect = ih / float(iw) if iw else 1.0
    w = max_w
    h = w * aspect
    if h > max_h:
        h = max_h
        w = h / aspect if aspect else max_w
    left = l + (max_w - w) / 2
    top = t + (max_h - h) / 2
    slide.shapes.add_picture(str(path), PptInches(left), PptInches(top), width=PptInches(w), height=PptInches(h))


def _slide_presentacion(
    prs,
    mall: Dict[str, Any],
    by: Dict[str, Dict[str, Any]],
    tot: Dict[str, float],
) -> None:
    """Equipos en tiras + gráfico mensual + peso de julio."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    chips_ids = list(mall.get("chip_order") or mall["nodes"])
    rank_ids = nodos_totales(mall)
    n_pts = len(chips_ids)
    sub = f"{mall['titulo']}   |   {n_pts} puntos WES   |   {PERIODO}"
    if mall["code"] == "PAK":
        sub = f"{mall['titulo']}   |   10 puntos WES   |   julio = cabecera (sin cadena DL)   |   {PERIODO}"
    _header_bar(
        sl,
        prs,
        f"{mall['code']}  ·  Equipos y consumo",
        sub,
    )

    chart_top, positions = _chip_layout(n_pts)
    compact = n_pts >= 8
    for nid, (x, y, w, h) in zip(chips_ids, positions):
        nota = CHIP_NOTA.get(nid, "")
        borde = COLOR_NODO.get(nid, TEAL)
        fill = (255, 249, 235) if nid in PAK_CADENA else LIGHT
        _caja(sl, x, y, w, h, fill=fill, line=borde)
        name_sz_c, nota_sz = (12, 8) if compact else (13, 10)
        _tb(
            sl,
            x + 0.08,
            y + 0.10,
            w - 0.14,
            0.32,
            [(NOMBRE_CORTO.get(nid, nid), name_sz_c, True, NAVY)],
        )
        if nota:
            _tb(sl, x + 0.08, y + h - 0.22, w - 0.14, 0.18, [(nota, nota_sz, False, GRAY)])

    ch_mes = CHARTS / f"{mall['code'].lower()}_mensual_may_ago.png"
    chart_mensual(ch_mes, tot, sin_mayo=bool(mall.get("sin_mayo")))
    chart_h = 6.14 - chart_top
    _caja(sl, 0.22, chart_top, 8.72, chart_h)
    _fit_picture(sl, ch_mes, 0.30, chart_top + 0.06, 8.56, chart_h - 0.12)

    _caja(sl, 9.08, chart_top, 4.02, chart_h, fill=WHITE, line=TEAL)
    jul_tit = "JULIO · cabecera" if mall["code"] == "PAK" else "JULIO · último mes cerrado"
    _tb(sl, 9.22, chart_top + 0.05, 3.74, 0.18, [(jul_tit, 11, True, TEAL)])
    _tb(
        sl,
        9.22,
        chart_top + 0.23,
        3.74,
        0.22,
        [(f"El recinto sumó {fn(tot['jul'], 0)} m³", 12, False, GRAY)],
    )
    ranked = sorted(rank_ids, key=lambda n: -float((by.get(n) or {}).get("jul") or 0))
    n_rank = max(len(ranked), 1)
    y0 = chart_top + 0.46
    avail = (chart_top + chart_h - 0.08) - y0
    step = avail / n_rank
    card_h = min(0.84, max(0.36, step - 0.04))
    name_sz = 10 if n_rank >= 6 else 13
    val_sz = 12 if n_rank >= 6 else 16
    y = y0
    for nid in ranked:
        v = float((by.get(nid) or {}).get("jul") or 0)
        extra = " · 11/08" if nid == FALABELLA else ""
        _caja(sl, 9.22, y, 3.74, card_h, fill=LIGHT, line=COLOR_NODO.get(nid, TEAL))
        _tb(
            sl,
            9.34,
            y + 0.03,
            3.50,
            min(0.18, card_h * 0.42),
            [(NOMBRE_CORTO.get(nid, nid), name_sz, True, COLOR_NODO.get(nid, NAVY))],
        )
        _tb(
            sl,
            9.34,
            y + card_h * 0.42,
            3.50,
            min(0.26, card_h * 0.52),
            [(f"{fn(v, 0)} m³    {_etq_pct_julio(nid, v, tot['jul'], un_decimal=n_rank >= 6)}{extra}", val_sz, True, NAVY)],
        )
        y += step

    _caja(sl, 0.22, 6.24, 12.90, 0.48, fill=(255, 249, 235), line=GOLD)
    _tb(
        sl,
        0.36,
        6.32,
        12.62,
        0.34,
        [
            (
                f"Recepción {mall['recepcion']}  ·  Capacitación {mall['capacitacion']}  ·  "
                f"{mall['usuarios']}",
                11,
                False,
                NAVY,
            )
        ],
    )
    _tb(
        sl,
        0.28,
        6.76,
        12.8,
        0.62,
        [
            (
                ("" if mall.get("sin_mayo") else f"Mayo {fn(tot['may'], 0)}   ·   ")
                + f"Junio {fn(tot['jun'], 0)}   ·   "
                f"Julio {fn(tot['jul'], 0)}   ·   Agosto {AGO_ETQ}: {fn(tot['ago'], 0)}   ·   "
                f"proy. ago {fn(tot['ago_proy'], 0)} m³.",
                12,
                False,
                NAVY,
            ),
            (mall["caption"], 11, False, GRAY),
        ],
    )


def _slide_pak_cadena(prs, by: Dict[str, Dict[str, Any]], cadena: Dict[str, Any]) -> None:
    """Sandía Antigua/Nueva → Distrito de Lujo → Bazar y DL Kennedy. Control en el tronco."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _header_bar(
        sl,
        prs,
        "PAK  ·  2. Cadena Sandía / Distrito de Lujo",
        "Sandía Antigua y Sandía Nueva alimentan Distrito de Lujo  ·  se divide en Bazar Gourmet y DL Kennedy",
    )
    perfil = {
        nid: ((cadena.get("perfil") or {}).get(nid) or {}).get(DIA_PAK_NOCHE.isoformat()) or {}
        for nid in PAK_CADENA
    }
    n06 = cadena.get("n06") or {}
    ch_p = CHARTS / "pak_cadena_perfil_20260810.png"
    ch_n = CHARTS / "pak_cadena_noches_jul.png"
    chart_pak_perfil_10ago(ch_p, perfil)
    med_n = chart_pak_noches_julio(ch_n, n06)
    med_bazar = float(med_n.get(BAZAR) or 0.0)
    med_ken = float(med_n.get(DL_KENNEDY) or 0.0)
    if med_bazar >= med_ken:
        onoff_nid, onoff_med = BAZAR, med_bazar
        otro_nom, otro_med = NOMBRE_CORTO[DL_KENNEDY], med_ken
    else:
        onoff_nid, onoff_med = DL_KENNEDY, med_ken
        otro_nom, otro_med = NOMBRE_CORTO[BAZAR], med_bazar
    onoff_nom = NOMBRE_CORTO[onoff_nid]
    ahorro_mes = onoff_med * 30.0
    jul_bazar = float((by.get(BAZAR) or {}).get("jul") or 0)
    jul_ken = float((by.get(DL_KENNEDY) or {}).get("jul") or 0)
    jul_dl = float((by.get("000025-27") or {}).get("jul") or 0)
    umb_bazar = _umbral_dia_de(jul_bazar)
    umb_ken = _umbral_dia_de(jul_ken)
    umb_dl = _umbral_dia_de(jul_dl)

    _vineta_hallazgos(
        sl,
        [
            (
                "Sandía Antigua y Sandía Nueva alimentan Distrito de Lujo. "
                "Desde ahí el caudal se divide en Bazar Gourmet y DL Kennedy. "
                "Por eso esos tres puntos no se suman a la cabecera: sería doble conteo.",
                12,
                False,
                NAVY,
            ),
            (
                f"Propuesta: on/off 00–06 en {onoff_nom} (el que más gasta de noche: "
                f"{fn(onoff_med, 1)} m³ vs {fn(otro_med, 1)} m³). "
                f"{fn(onoff_med, 1)} m³/noche = {_clp(onoff_med)}/noche  ·  "
                f"{fn(ahorro_mes, 0)} m³/mes = {_clp(ahorro_mes)}/mes "
                f"(tarifa ${fn(TARIFA_CLP_M3, 0)}/m³).",
                12,
                True,
                NAVY,
            ),
            (
                f"Umbrales a activar (total 24 h): Distrito de Lujo {fn(umb_dl, 0)} m³/día  ·  "
                f"Bazar Gourmet {fn(umb_bazar, 0)} m³/día  ·  DL Kennedy {fn(umb_ken, 0)} m³/día.",
                12,
                True,
                NAVY,
            ),
        ],
        h=1.22,
    )

    flujo = [
        (0.22, 2.12, "Sandía Antigua", "000025-22", LIGHT),
        (2.70, 2.12, "Sandía Nueva", "000025-28", LIGHT),
        (5.32, 2.40, "Distrito de Lujo", "000025-27", (255, 249, 235)),
        (8.22, 2.12, "Bazar Gourmet", "000025-35", LIGHT),
        (10.44, 2.12, "DL Kennedy", "000025-36", LIGHT),
    ]
    yf = 2.36
    hf = 0.42
    for x, w, nom, nid, fill in flujo:
        _caja(sl, x, yf, w, hf, fill=fill, line=COLOR_NODO.get(nid, TEAL))
        _tb(sl, x + 0.08, yf + 0.12, w - 0.16, 0.28, [(nom, 12, True, NAVY)], align=PP_ALIGN.CENTER)
    _tb(sl, 2.34, yf + 0.08, 0.36, 0.34, [("+", 16, True, GOLD)], align=PP_ALIGN.CENTER)
    _tb(sl, 4.82, yf + 0.08, 0.50, 0.34, [("→", 18, True, GOLD)], align=PP_ALIGN.CENTER)
    _tb(sl, 7.72, yf + 0.08, 0.50, 0.34, [("→", 18, True, GOLD)], align=PP_ALIGN.CENTER)

    _caja(sl, 0.22, 2.84, 6.38, 3.22)
    _tb(sl, 0.36, 2.88, 6.10, 0.22, [("10/08 · perfil 24 h  (sombra = 00:00–06:00)", 11, True, TEAL)])
    _fit_picture(sl, ch_p, 0.32, 3.12, 6.18, 2.84)

    _caja(sl, 6.74, 2.84, 6.36, 3.22)
    _tb(sl, 6.88, 2.88, 6.08, 0.22, [("NOCHE 00–06  ·  Bazar Gourmet y DL Kennedy", 11, True, TEAL)])
    _fit_picture(sl, ch_n, 6.84, 3.12, 6.16, 2.84)

    _caja(sl, 0.22, 6.16, 12.88, 1.16, fill=(255, 249, 235), line=GOLD)
    _tb(
        sl,
        0.40,
        6.20,
        12.55,
        1.06,
        [
            (
                f"Noche típica 00–06: Bazar Gourmet {fn(med_bazar, 1)} m³  ·  "
                f"DL Kennedy {fn(med_ken, 1)} m³. {otro_nom} queda en {fn(otro_med, 1)} m³/noche.",
                12,
                False,
                NAVY,
            ),
        ],
    )


def _slide_mam_placa(prs, by: Dict[str, Dict[str, Any]], mam: Dict[str, Any]) -> None:
    """Placa Bancaria: noche, umbral, reubicación de residuales y cruce vs Falabella."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _header_bar(
        sl,
        prs,
        "MAM  ·  2. Placa Bancaria",
        "Cuando se inyecta Falabella, se corta Placa y el mall sale por Falabella",
    )
    daily_p = (by.get(PLACA) or {}).get("daily") or {}
    daily_f = (by.get(FALABELLA) or {}).get("daily") or {}
    n06_p = (mam.get("n06") or {}).get(PLACA) or {}
    n06_f = (mam.get("n06") or {}).get(FALABELLA) or {}
    ch_n = CHARTS / "mam_placa_noches_jul_ago.png"
    ch_f = CHARTS / "mam_placa_vs_falabella.png"
    chart_mam_placa_noches(ch_n, n06_p, n06_f)
    chart_mam_placa_vs_falabella(ch_f, daily_p, daily_f)

    jul_p = float((by.get(PLACA) or {}).get("jul") or 0)
    jul_tot = sum(float((by.get(n) or {}).get("jul") or 0) for n in MAM_NODOS)
    pct_jul = (jul_p / jul_tot * 100.0) if jul_tot else 0.0
    jul_pas = float((by.get(PASILLO) or {}).get("jul") or 0)
    jul_arr = float((by.get(ARROW) or {}).get("jul") or 0)
    jul_dia = jul_p / 31.0
    fala_ahora = [
        float(daily_f.get(d.isoformat(), 0.0) or 0.0)
        for d in _rango_dias(date(2026, 8, 15), HASTA)
    ]
    fala_dia = _mediana(fala_ahora) if fala_ahora else 0.0

    _vineta_hallazgos(
        sl,
        [
            (
                f"Placa Bancaria se lleva {fn(pct_jul, 0)} % de julio ({fn(jul_p, 0)} m³). "
                f"Pasillo Técnico ({fn(jul_pas, 0)} m³) y ARROW ({fn(jul_arr, 1)} m³) casi no miden: "
                "son los puntos para reubicar y partir la línea.",
                13,
                False,
                NAVY,
            ),
            (
                f"Umbrales a activar (total 24 h): Placa {fn(UMBRAL_PLACA_DIA, 0)} m³/día "
                f"(julio andaba en {fn(jul_dia, 0)})  ·  Falabella {fn(UMBRAL_FALABELLA_DIA, 0)} m³/día "
                f"(ahora ~{fn(fala_dia, 0)}).",
                13,
                True,
                NAVY,
            ),
        ],
        h=1.10,
    )

    _caja(sl, 0.22, 2.28, 6.38, 3.78)
    _tb(sl, 0.36, 2.32, 6.10, 0.22, [("DESDE EL 15/08 EL MALL SALE POR FALABELLA", 11, True, TEAL)])
    _fit_picture(sl, ch_f, 0.32, 2.56, 6.18, 3.40)

    _caja(sl, 6.74, 2.28, 6.36, 3.78)
    _tb(sl, 6.88, 2.32, 6.08, 0.22, [("NOCHE 00–06  ·  mall cerrado", 11, True, TEAL)])
    _fit_picture(sl, ch_n, 6.84, 2.56, 6.16, 3.40)

    _caja(sl, 0.22, 6.16, 12.88, 1.16, fill=(255, 249, 235), line=GOLD)
    _tb(
        sl,
        0.40,
        6.22,
        12.55,
        1.04,
        [
            (
                "Desde el 15/08 Placa se fue a 0. Ya pasó: se inyecta Falabella, "
                "se corta Placa y el mall sale por Falabella. "
                "Siguiente paso: subdividir la línea de Falabella para ver hacia dónde va el consumo de noche.",
                14,
                False,
                NAVY,
            ),
        ],
    )


def _slide_maq_matriz(prs, by: Dict[str, Dict[str, Any]], maq: Dict[str, Any]) -> None:
    """Alza desde junio, umbrales y control nocturno en Matriz Principal."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _header_bar(
        sl,
        prs,
        "MAQ  ·  2. Matriz Principal",
        "Alza desde junio  ·  umbrales  ·  la noche pide un control",
    )
    daily = (by.get(MATRIZ_MAQ) or {}).get("daily") or {}
    n06 = (maq.get("n06") or {}).get(MATRIZ_MAQ) or {}
    horas = maq.get("horas") or {}
    ch_d = CHARTS / "maq_matriz_diario_jun_ago.png"
    ch_n = CHARTS / "maq_matriz_noches_jun_ago.png"
    chart_maq_matriz_diario(ch_d, daily)
    st_n = chart_maq_matriz_noches(ch_n, n06)
    st_h = _stats_maq_hora(horas)

    jul = float((by.get(MATRIZ_MAQ) or {}).get("jul") or 0)
    ago = float((by.get(MATRIZ_MAQ) or {}).get("ago") or 0)
    jun = float((by.get(MATRIZ_MAQ) or {}).get("jun") or 0)
    jul_tot = sum(float((by.get(n) or {}).get("jul") or 0) for n in ["000025-13", "000025-34"])
    pct = (jul / jul_tot * 100.0) if jul_tot else 0.0
    med_n = float(st_n.get("med_post") or st_n.get("med") or 0.0)
    min_h = float(st_h.get("min_hora") or 0.0)
    if min_h <= 0 and med_n:
        min_h = med_n / 6.0
    ahorro_noche = med_n
    ahorro_mes = ahorro_noche * 30.0

    _vineta_hallazgos(
        sl,
        [
            (
                f"Matriz Principal se lleva {fn(pct, 0)} % de julio. "
                f"Desde el 22/06 el día se duplicó (junio {fn(jun / 30.0, 0)} → julio {fn(jul / 31.0, 0)} "
                f"→ agosto {fn(ago / 17.0, 0)} m³/día) y se quedó arriba. "
                "Baños es uso hábil: no es el problema.",
                13,
                False,
                NAVY,
            ),
            (
                f"Umbral a activar (total 24 h): Matriz Principal {fn(UMBRAL_MAQ_DIA, 0)} m³/día "
                f"(julio ~{fn(jul / 31.0, 0)}).",
                13,
                True,
                NAVY,
            ),
        ],
        h=1.10,
    )

    _caja(sl, 0.22, 2.28, 6.38, 3.78)
    _tb(sl, 0.36, 2.32, 6.10, 0.22, [("DÍA · el alza desde junio se sostiene", 11, True, TEAL)])
    _fit_picture(sl, ch_d, 0.32, 2.56, 6.18, 3.40)

    _caja(sl, 6.74, 2.28, 6.36, 3.78)
    _tb(sl, 6.88, 2.32, 6.08, 0.22, [("NOCHE 00–06  ·  desde el 01/06", 11, True, TEAL)])
    _fit_picture(sl, ch_n, 6.84, 2.56, 6.16, 3.40)

    _caja(sl, 0.22, 6.16, 12.88, 1.16, fill=(255, 249, 235), line=GOLD)
    _tb(
        sl,
        0.40,
        6.22,
        12.55,
        1.04,
        [
            (
                f"Desde el 22/06 la madrugada anda en {fn(med_n, 1)} m³ "
                f"(antes {fn(st_n.get('med_pre') or 0.0, 1)}). "
                f"Mínimo típico por hora 00–06: {fn(min_h, 1)} m³/h.",
                13,
                False,
                NAVY,
            ),
            (
                f"On/off 00–06 en Matriz: {fn(ahorro_noche, 1)} m³/noche = "
                f"{_clp(ahorro_noche)}/noche  ·  {fn(ahorro_mes, 0)} m³/mes = "
                f"{_clp(ahorro_mes)}/mes (tarifa ${fn(TARIFA_CLP_M3, 0)}/m³).",
                13,
                True,
                NAVY,
            ),
        ],
    )


def _slide_bom_control(prs, by: Dict[str, Dict[str, Any]], bom: Dict[str, Any]) -> None:
    """Control operativo en San Ignacio 500 desde 17/07: noche, ahorro y umbral."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _header_bar(
        sl,
        prs,
        "BOM  ·  2. San Ignacio 500",
        "Control nocturno operativo  ·  ahorro  ·  umbral",
    )
    n06_500 = (bom.get("n06") or {}).get(SI500) or {}
    st = _stats_si500(n06_500)
    ch_n = CHARTS / "bom_si500_noches_jul_ago.png"
    ch_a = CHARTS / "bom_si500_ahorro.png"
    chart_bom_si500_noches(ch_n, n06_500)
    chart_bom_si500_ahorro(ch_a, st)

    jul = float((by.get(SI500) or {}).get("jul") or 0)
    ago = float((by.get(SI500) or {}).get("ago") or 0)
    jul_tot = sum(float((by.get(n) or {}).get("jul") or 0) for n in [SI500, SI300])
    pct = (jul / jul_tot * 100.0) if jul_tot else 0.0
    n_post = int(st["n_post"])
    if st["n_sobre_umbral"] < 1:
        residual = (
            f"Desde el 17/07 no hay noche alta: residual típico {fn(st['post'], 1)} m³ "
            f"(máximo {fn(st['max_post'], 1)}). No se lee como fuga."
        )
    else:
        altas = ", ".join(
            f"{d.strftime('%d/%m')} {fn(v, 1)} m³" for d, v in st["noches_altas"]
        )
        residual = (
            f"Quedan {fn(st['n_sobre_umbral'], 0)} noches ≥ {fn(UMBRAL_SI500_NOCHE, 0)} m³ "
            f"({altas}). No se leen como fuga: el control puede haber aflojado esas madrugadas."
        )

    _vineta_hallazgos(
        sl,
        [
            (
                f"San Ignacio 500 se lleva {fn(pct, 0)} % de julio. "
                f"Control on/off operativo desde el 17/07: la madrugada pasó de "
                f"{fn(st['pre'], 0)} a {fn(st['post'], 1)} m³ "
                f"(el 15/07 eran {fn(st['vispera'], 0)} m³) y se sostiene hasta hoy.",
                13,
                False,
                NAVY,
            ),
            (
                f"Ahorro del corte 00–06: {fn(st['ahorro_noche'], 1)} m³/noche = "
                f"{_clp(st['ahorro_noche'])}/noche  ·  {_clp(st['ahorro_noche'] * 30)}/mes. "
                f"Acumulado 17/07–{HASTA.strftime('%d/%m')} ({n_post} noches): "
                f"{fn(st['ahorro_acum'], 0)} m³ = {_clp(st['ahorro_acum'])} "
                f"(tarifa ${fn(TARIFA_CLP_M3, 0)}/m³).",
                13,
                True,
                NAVY,
            ),
            (
                f"Umbrales a activar (total 24 h): San Ignacio 500 {fn(UMBRAL_SI500_DIA, 0)} m³/día "
                f"(agosto ~{fn(ago / 17.0, 0)})  ·  San Ignacio 300 {fn(UMBRAL_SI300_DIA, 0)} m³/día.",
                13,
                True,
                NAVY,
            ),
        ],
        h=1.32,
    )

    _caja(sl, 0.22, 2.50, 6.38, 3.56)
    _tb(sl, 0.36, 2.54, 6.10, 0.22, [("AHORRO  ·  noche típica antes vs con control", 11, True, TEAL)])
    _fit_picture(sl, ch_a, 0.32, 2.78, 6.18, 3.18)

    _caja(sl, 6.74, 2.50, 6.36, 3.56)
    _tb(sl, 6.88, 2.54, 6.08, 0.22, [("NOCHE 00–06  ·  el corte se ve y se sostiene", 11, True, TEAL)])
    _fit_picture(sl, ch_n, 6.84, 2.78, 6.16, 3.18)

    _caja(sl, 0.22, 6.16, 12.88, 1.16, fill=(255, 249, 235), line=GOLD)
    _tb(
        sl,
        0.40,
        6.22,
        12.55,
        1.04,
        [
            (
                f"{residual}",
                13,
                False,
                NAVY,
            ),
        ],
    )


def _slide_aeb_anillo(prs, by: Dict[str, Dict[str, Any]], aeb: Dict[str, Any]) -> None:
    """Matriz A.A. desactivada 15/05, alza Anillo Plaza y umbrales."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _header_bar(
        sl,
        prs,
        "AEB  ·  2. Anillo Plaza y Matriz",
        "Matriz A.A. desactivada 15/05  ·  alza Anillo  ·  umbrales",
    )
    daily_a = (by.get(ANILLO_AEB) or {}).get("daily") or {}
    daily_m = (by.get(MATRIZ_AEB) or {}).get("daily") or {}
    n06_a = (aeb.get("n06") or {}).get(ANILLO_AEB) or {}
    n06_m = (aeb.get("n06") or {}).get(MATRIZ_AEB) or {}
    ch_a = CHARTS / "aeb_anillo_dia_noche.png"
    ch_m = CHARTS / "aeb_matriz_dia_noche.png"
    st_a = chart_aeb_dia_noche(
        ch_a,
        daily_a,
        n06_a,
        color_dia=COLOR_NODO[ANILLO_AEB],
        umbral=UMBRAL_ANILLO_DIA,
        umbral_noche=False,
        umbral_etq=f"Umbral día {fn(UMBRAL_ANILLO_DIA, 0)} m³",
    )
    st_m = chart_aeb_dia_noche(
        ch_m,
        daily_m,
        n06_m,
        color_dia=COLOR_NODO[MATRIZ_AEB],
        umbral=UMBRAL_MATRIZ_AEB_DIA,
        umbral_noche=False,
        umbral_etq=f"Umbral día {fn(UMBRAL_MATRIZ_AEB_DIA, 0)} m³",
    )
    may_a, jun_a, jul_a, ago_a = st_a["dia"]
    noc_a_jul, noc_a_ago = st_a["noche"][2], st_a["noche"][3]
    noc_m_jul, noc_m_ago = st_m["noche"][2], st_m["noche"][3]
    noc_a = (noc_a_jul + noc_a_ago) / 2.0
    noc_m = (noc_m_jul + noc_m_ago) / 2.0
    ahorro_a_mes = noc_a * 30.0
    ahorro_m_mes = noc_m * 30.0
    ahorro_tot_mes = ahorro_a_mes + ahorro_m_mes
    pico = max((float(v) for iso, v in daily_a.items() if iso.startswith("2026-07")), default=0.0)
    pico_iso = max(
        ((iso, float(v)) for iso, v in daily_a.items() if iso.startswith("2026-07")),
        key=lambda x: x[1],
        default=("", 0.0),
    )
    pico_d = date.fromisoformat(pico_iso[0]).strftime("%d/%m") if pico_iso[0] else "09/07"

    _vineta_hallazgos(
        sl,
        [
            (
                f"On/off 00–06: Anillo Plaza {fn(noc_a, 1)} m³/noche = "
                f"{_clp(noc_a)}/noche · {_clp(ahorro_a_mes)}/mes. "
                f"Matriz 1° piso {fn(noc_m, 1)} m³/noche = "
                f"{_clp(noc_m)}/noche · {_clp(ahorro_m_mes)}/mes. "
                f"Los dos: {fn(noc_a + noc_m, 1)} m³/noche = {fn(ahorro_tot_mes, 0)} m³/mes = "
                f"{_clp(ahorro_tot_mes)}/mes (tarifa ${fn(TARIFA_CLP_M3, 0)}/m³).",
                12,
                True,
                NAVY,
            ),
            (
                f"Umbrales a activar (total 24 h): Anillo Plaza {fn(UMBRAL_ANILLO_DIA, 0)} m³/día "
                f"(julio {fn(jul_a, 0)}; el 09/07 llegó a {fn(pico, 0)}). "
                f"Matriz 1° piso {fn(UMBRAL_MATRIZ_AEB_DIA, 0)} m³/día "
                f"(agosto ~{fn(st_m['dia'][-1], 0)} × 1,25).",
                12,
                False,
                NAVY,
            ),
        ],
        h=1.10,
    )

    _caja(sl, 0.22, 2.28, 6.38, 3.78)
    _tb(sl, 0.36, 2.32, 6.10, 0.22, [("ANILLO PLAZA  ·  promedio día vs noche", 11, True, TEAL)])
    _fit_picture(sl, ch_a, 0.32, 2.56, 6.18, 3.40)

    _caja(sl, 6.74, 2.28, 6.36, 3.78)
    _tb(sl, 6.88, 2.32, 6.08, 0.22, [("MATRIZ 1° PISO  ·  promedio día vs noche", 11, True, TEAL)])
    _fit_picture(sl, ch_m, 6.84, 2.56, 6.16, 3.40)

    _caja(sl, 0.22, 6.16, 12.88, 1.16, fill=(255, 249, 235), line=GOLD)
    _tb(
        sl,
        0.40,
        6.22,
        12.55,
        1.04,
        [
            (
                f"Matriz A.A. se desactivó el {MATRIZ_AA_HASTA.strftime('%d/%m')}: "
                "desde ese día el recinto se lee en Matriz 1° piso y Anillo Plaza. "
                f"Anillo Plaza subió (mayo {fn(may_a, 0)} → junio {fn(jun_a, 0)} → julio {fn(jul_a, 0)} m³/día; "
                f"pico {fn(pico, 0)} el {pico_d}). Agosto {fn(ago_a, 0)}: ¿hay un trabajo en curso?",
                13,
                False,
                NAVY,
            ),
        ],
    )


def _slide_cur_anillos(prs, by: Dict[str, Dict[str, Any]], cur: Dict[str, Any]) -> None:
    """Diferencia chica Sur vs Norte, noche 00–06 y umbrales del día."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _header_bar(
        sl,
        prs,
        "CUR  ·  2. Anillo Sur y Norte",
        "Diferencia chica  ·  noche 00–06  ·  umbral del día",
    )
    daily_s = (by.get(ANILLO_SUR) or {}).get("daily") or {}
    daily_n = (by.get(ANILLO_NORTE) or {}).get("daily") or {}
    n06_s = (cur.get("n06") or {}).get(ANILLO_SUR) or {}
    n06_n = (cur.get("n06") or {}).get(ANILLO_NORTE) or {}
    ch_n = CHARTS / "cur_noches_jun_ago.png"
    ch_d = CHARTS / "cur_dia_sur_norte.png"
    st_n = chart_cur_noches(ch_n, n06_s, n06_n)
    st_d = chart_cur_dia(ch_d, daily_s, daily_n)
    jul_s = float((by.get(ANILLO_SUR) or {}).get("jul") or 0)
    jul_n = float((by.get(ANILLO_NORTE) or {}).get("jul") or 0)
    jul_tot = jul_s + jul_n
    pct_s = (jul_s / jul_tot * 100.0) if jul_tot else 0.0
    pct_n = (jul_n / jul_tot * 100.0) if jul_tot else 0.0
    dif = st_d["sur"][1] - st_d["norte"][1]  # julio m³/día

    _vineta_hallazgos(
        sl,
        [
            (
                f"Anillo Sur anda un poco sobre Norte: julio {fn(pct_s, 0)} % vs {fn(pct_n, 0)} % "
                f"(unos {fn(dif, 1)} m³/día). Los dos cubren el mall; la diferencia es chica y se sostiene.",
                13,
                False,
                NAVY,
            ),
            (
                f"Umbrales a activar (total 24 h): Anillo Sur {fn(UMBRAL_SUR_DIA, 0)} m³/día "
                f"(junio–agosto ~{fn(sum(st_d['sur']) / 3.0, 0)} × 1,25)  ·  "
                f"Anillo Norte {fn(UMBRAL_NORTE_DIA, 0)} m³/día "
                f"(junio–agosto ~{fn(sum(st_d['norte']) / 3.0, 0)} × 1,25).",
                13,
                True,
                NAVY,
            ),
        ],
        h=1.10,
    )

    _caja(sl, 0.22, 2.28, 6.38, 3.78)
    _tb(sl, 0.36, 2.32, 6.10, 0.22, [("NOCHE 00–06  ·  los dos anillos", 11, True, TEAL)])
    _fit_picture(sl, ch_n, 0.32, 2.56, 6.18, 3.40)

    _caja(sl, 6.74, 2.28, 6.36, 3.78)
    _tb(sl, 6.88, 2.32, 6.08, 0.22, [("DÍA  ·  promedio mes  ·  Sur un poco sobre Norte", 11, True, TEAL)])
    _fit_picture(sl, ch_d, 6.84, 2.56, 6.16, 3.40)

    _caja(sl, 0.22, 6.16, 12.88, 1.16, fill=(255, 249, 235), line=GOLD)
    _tb(
        sl,
        0.40,
        6.22,
        12.55,
        1.04,
        [
            (
                f"Noche típica 00–06: Sur {fn(st_n['sur'], 1)} m³  ·  Norte {fn(st_n['norte'], 1)} m³. "
                "La madrugada sigue el mismo patrón chico; no se lee como fuga.",
                13,
                False,
                NAVY,
            ),
        ],
    )


def _filas_propuesta(
    by: Dict[str, Dict[str, Any]],
    hourly: Dict[str, Dict[str, float]],
    cadena_pak: Dict[str, Any],
    maq_matriz: Dict[str, Any],
    bom_si500: Dict[str, Any],
    aeb: Dict[str, Any],
    cur: Dict[str, Any],
) -> List[Dict[str, str]]:
    """Qué proponemos en cada recinto: on/off (ahorro) + umbral 24 h."""
    st_n = _stats_noche(hourly, "000025-01", CTRL_NORTE)
    mes_n = _ahorro_norte_mes(st_n)

    n06_mq = (maq_matriz.get("n06") or {}).get(MATRIZ_MAQ) or {}
    noc_maq = _mediana(
        [
            float(v)
            for iso, v in n06_mq.items()
            if date.fromisoformat(iso) >= MAQ_ALZA
        ]
    )
    mes_maq = noc_maq * 30.0

    st_b = _stats_si500((bom_si500.get("n06") or {}).get(SI500) or {})
    mes_bom = float(st_b.get("ahorro_noche") or 0.0) * 30.0

    daily_a = (by.get(ANILLO_AEB) or {}).get("daily") or {}
    daily_m = (by.get(MATRIZ_AEB) or {}).get("daily") or {}
    n06_a = (aeb.get("n06") or {}).get(ANILLO_AEB) or {}
    n06_m = (aeb.get("n06") or {}).get(MATRIZ_AEB) or {}
    noc_a = (
        _prom_mes(n06_a, 2026, 7, activos=daily_a) + _prom_mes(n06_a, 2026, 8, activos=daily_a)
    ) / 2.0
    noc_m = (
        _prom_mes(n06_m, 2026, 7, activos=daily_m) + _prom_mes(n06_m, 2026, 8, activos=daily_m)
    ) / 2.0
    mes_aeb = (noc_a + noc_m) * 30.0

    n06_cs = (cur.get("n06") or {}).get(ANILLO_SUR) or {}
    n06_cn = (cur.get("n06") or {}).get(ANILLO_NORTE) or {}
    med_sur = _mediana([float(v) for v in n06_cs.values()])
    med_nor = _mediana([float(v) for v in n06_cn.values()])

    n06_pak = cadena_pak.get("n06") or {}
    med_bazar = _mediana([float(v) for v in (n06_pak.get(BAZAR) or {}).values()])
    med_ken = _mediana([float(v) for v in (n06_pak.get(DL_KENNEDY) or {}).values()])
    if med_bazar >= med_ken:
        pak_nom, pak_med, pak_otro, pak_otro_m = "Bazar Gourmet", med_bazar, "DL Kennedy", med_ken
    else:
        pak_nom, pak_med, pak_otro, pak_otro_m = "DL Kennedy", med_ken, "Bazar Gourmet", med_bazar
    mes_pak = pak_med * 30.0
    umb_dl = _umbral_dia_de(float((by.get("000025-27") or {}).get("jul") or 0))
    umb_bazar = _umbral_dia_de(float((by.get(BAZAR) or {}).get("jul") or 0))
    umb_ken = _umbral_dia_de(float((by.get(DL_KENNEDY) or {}).get("jul") or 0))

    return [
        {
            "code": "MAE",
            "nombre": "Mall Arauco Estación",
            "onoff": (
                f"On/off: Estanque Norte ya operativo (05/08, 00:00–05:00) "
                f"{fn(st_n['pre'], 1)} → {fn(st_n['post'], 1)} m³ = "
                f"{fn(mes_n, 0)} m³/mes = {_clp(mes_n)}/mes. Pizza Hut control 01/07."
            ),
            "umbral": (
                f"Umbrales 24 h: Norte {fn(UMBRAL_MAE_NORTE_DIA, 0)}  ·  "
                f"Sur {fn(UMBRAL_MAE_SUR_DIA, 0)}  ·  Pizza Hut {fn(UMBRAL_MAE_PIZZA_DIA, 0)}  ·  "
                f"Baños {fn(UMBRAL_MAE_BANOS_DIA, 0)} m³/día."
            ),
        },
        {
            "code": "MAM",
            "nombre": "Mall Arauco Maipú",
            "onoff": (
                "On/off: no se propone corte ahora. Desde el 15/08 el mall sale por Falabella. "
                "Siguiente paso: partir esa línea para ver a dónde va la noche."
            ),
            "umbral": (
                f"Umbrales 24 h: Placa Bancaria {fn(UMBRAL_PLACA_DIA, 0)}  ·  "
                f"Falabella {fn(UMBRAL_FALABELLA_DIA, 0)} m³/día."
            ),
        },
        {
            "code": "MAQ",
            "nombre": "Mall Arauco Quilicura",
            "onoff": (
                f"On/off a proponer 00–06 en Matriz Principal: "
                f"{fn(noc_maq, 1)} m³/noche = {fn(mes_maq, 0)} m³/mes = {_clp(mes_maq)}/mes."
            ),
            "umbral": f"Umbral 24 h: Matriz Principal {fn(UMBRAL_MAQ_DIA, 0)} m³/día.",
        },
        {
            "code": "BOM",
            "nombre": "Buenaventura (San Ignacio)",
            "onoff": (
                f"On/off: San Ignacio 500 ya operativo desde 17/07. "
                f"Ahorro {fn(st_b['ahorro_noche'], 1)} m³/noche = "
                f"{fn(mes_bom, 0)} m³/mes = {_clp(mes_bom)}/mes."
            ),
            "umbral": (
                f"Umbrales 24 h: San Ignacio 500 {fn(UMBRAL_SI500_DIA, 0)}  ·  "
                f"San Ignacio 300 {fn(UMBRAL_SI300_DIA, 0)} m³/día."
            ),
        },
        {
            "code": "AEB",
            "nombre": "Arauco El Bosque",
            "onoff": (
                f"On/off a proponer 00–06 en Anillo Plaza ({fn(noc_a, 1)} m³) y Matriz 1° piso "
                f"({fn(noc_m, 1)} m³): {fn(noc_a + noc_m, 1)} m³/noche = "
                f"{fn(mes_aeb, 0)} m³/mes = {_clp(mes_aeb)}/mes."
            ),
            "umbral": (
                f"Umbrales 24 h: Anillo Plaza {fn(UMBRAL_ANILLO_DIA, 0)}  ·  "
                f"Matriz 1° piso {fn(UMBRAL_MATRIZ_AEB_DIA, 0)} m³/día."
            ),
        },
        {
            "code": "CUR",
            "nombre": "Arauco Curauma",
            "onoff": (
                f"On/off: no se propone. Noche típica Sur {fn(med_sur, 1)} · Norte {fn(med_nor, 1)} m³. "
                "Es chica; no se lee como fuga."
            ),
            "umbral": (
                f"Umbrales 24 h: Anillo Sur {fn(UMBRAL_SUR_DIA, 0)}  ·  "
                f"Anillo Norte {fn(UMBRAL_NORTE_DIA, 0)} m³/día."
            ),
        },
        {
            "code": "PAK",
            "nombre": "Parque Arauco Kennedy",
            "onoff": (
                f"On/off a proponer 00–06 en {pak_nom} "
                f"({fn(pak_med, 1)} vs {fn(pak_otro_m, 1)} m³ de {pak_otro}): "
                f"{fn(pak_med, 1)} m³/noche = {fn(mes_pak, 0)} m³/mes = {_clp(mes_pak)}/mes."
            ),
            "umbral": (
                f"Umbrales 24 h: Distrito de Lujo {fn(umb_dl, 0)}  ·  "
                f"Bazar Gourmet {fn(umb_bazar, 0)}  ·  DL Kennedy {fn(umb_ken, 0)} m³/día."
            ),
        },
    ]


def _portada_resumen(prs) -> None:
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    if FONDO.is_file():
        pic = sl.shapes.add_picture(str(FONDO), 0, 0, width=prs.slide_width, height=prs.slide_height)
        spTree = sl.shapes._spTree
        spTree.remove(pic.element)
        spTree.insert(2, pic.element)
    veil = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, PptInches(3.55), prs.slide_width, PptInches(3.95))
    veil.fill.solid()
    veil.fill.fore_color.rgb = _rgb(NAVY)
    veil.line.fill.background()
    _tb(sl, 0.6, 3.75, 12, 0.5, [("WES  ·  Parque Arauco", 16, True, GOLD)])
    _tb(sl, 0.6, 4.18, 12, 0.7, [("Propuesta: on/off y umbrales", 32, True, WHITE)])
    _tb(
        sl,
        0.6,
        4.95,
        12,
        1.2,
        [
            ("Qué pedimos en cada recinto  ·  ahorro de madrugada y umbral del día (24 h)", 16, False, WHITE),
            (f"Período {PERIODO}   |   Emisión {FECHA_EMISION}", 15, False, (220, 230, 240)),
            ("MAE  ·  MAM  ·  MAQ  ·  BOM  ·  AEB  ·  CUR  ·  PAK", 16, False, GOLD),
        ],
    )
    if LOGO.is_file():
        sl.shapes.add_picture(str(LOGO), PptInches(11.70), PptInches(6.85), width=PptInches(1.35))


def _slide_resumen_propuestas(prs, filas: List[Dict[str, str]]) -> None:
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _header_bar(
        sl,
        prs,
        "Propuesta por recinto  ·  on/off y umbrales",
        "Umbral = total de las 24 h (promedio operativo × 1,25)  ·  On/off = corte 00:00–06:00",
    )
    y0, gap, bottom = 1.08, 0.05, 7.38
    n = max(len(filas), 1)
    h = (bottom - y0 - gap * (n - 1)) / n
    y = y0
    for fila in filas:
        _caja(sl, 0.22, y, 12.88, h, fill=(255, 249, 235), line=GOLD)
        _tb(sl, 0.36, y + 0.04, 12.60, 0.20, [(f"{fila['code']}  ·  {fila['nombre']}", 12, True, GOLD)])
        _tb(
            sl,
            0.36,
            y + 0.24,
            12.60,
            h - 0.30,
            [
                (fila["onoff"], 11, False, NAVY),
                (fila["umbral"], 11, True, NAVY),
            ],
        )
        y += h + gap


def build_ppt_resumen(
    by: Dict[str, Dict[str, Any]],
    hourly: Dict[str, Dict[str, float]],
    cadena_pak: Dict[str, Any],
    maq_matriz: Dict[str, Any],
    bom_si500: Dict[str, Any],
    aeb: Dict[str, Any],
    cur: Dict[str, Any],
) -> Path:
    filas = _filas_propuesta(by, hourly, cadena_pak, maq_matriz, bom_si500, aeb, cur)
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)
    _portada_resumen(prs)
    _slide_resumen_propuestas(prs, filas)
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    path = OUT_DIR / f"Propuesta_onoff_umbrales_PA_7malls_{HASTA.strftime('%Y%m%d')}.pptx"
    prs.save(str(path))
    print(f"[OK] PPT resumen {path}")
    return path


def build_ppt(
    by: Dict[str, Dict[str, Any]],
    hourly: Dict[str, Dict[str, float]],
    cadena_pak: Dict[str, Any],
    mam_placa: Dict[str, Any],
    maq_matriz: Dict[str, Any],
    bom_si500: Dict[str, Any],
    aeb: Dict[str, Any],
    cur: Dict[str, Any],
) -> Path:
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)
    _portada(prs)
    for mall in MALLS:
        ids = nodos_totales(mall)
        tot = totales_de(by, ids)
        _slide_presentacion(prs, mall, by, tot)
        if mall["code"] == "MAE":
            _slide_hallazgos(prs, by, hourly)
        if mall["code"] == "MAM":
            _slide_mam_placa(prs, by, mam_placa)
        if mall["code"] == "MAQ":
            _slide_maq_matriz(prs, by, maq_matriz)
        if mall["code"] == "BOM":
            _slide_bom_control(prs, by, bom_si500)
        if mall["code"] == "AEB":
            _slide_aeb_anillo(prs, by, aeb)
        if mall["code"] == "CUR":
            _slide_cur_anillos(prs, by, cur)
        if mall["code"] == "PAK":
            _slide_pak_cadena(prs, by, cadena_pak)
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    path = OUT_DIR / f"Recorrido_ejecutivo_PA_MAE_{HASTA.strftime('%Y%m%d')}.pptx"
    prs.save(str(path))
    print(f"[OK] PPT {path}")
    return path


def main() -> int:
    skip = "--skip-refresh" in sys.argv
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    CHARTS.mkdir(parents=True, exist_ok=True)
    todos = nodos_todos()
    if not skip or not JSON_ALL.is_file():
        refrescar_datos(todos, JSON_ALL, "7 malls")
    if not skip or not JSON_NOCHES.is_file():
        refrescar_noches()
    if not skip or not JSON_PAK_CADENA.is_file():
        refrescar_cadena_pak()
    if not skip or not JSON_MAM_PLACA.is_file():
        refrescar_mam_placa()
    if not skip or not JSON_MAQ_MATRIZ.is_file():
        refrescar_maq_matriz()
    if not skip or not JSON_BOM_SI500.is_file():
        refrescar_bom_si500()
    if not skip or not JSON_AEB.is_file():
        refrescar_aeb()
    if not skip or not JSON_CUR.is_file():
        refrescar_cur()
    _names, by, _tot = cargar_mall(JSON_ALL, todos)
    hourly = cargar_noches()
    if not (hourly.get("000025-07") and hourly.get("000025-01")):
        hourly = refrescar_noches()
    cadena_pak = cargar_cadena_pak()
    n06_27 = ((cadena_pak.get("n06") or {}).get("000025-27") or {})
    n06_35 = ((cadena_pak.get("n06") or {}).get(BAZAR) or {})
    if (
        not ((cadena_pak.get("perfil") or {}).get("000025-27") or {}).get(DIA_PAK_NOCHE.isoformat())
        or JUL_NOCHE_D0.isoformat() not in n06_27
        or JUL_NOCHE_D0.isoformat() not in n06_35
        or HASTA.isoformat() not in n06_35
    ):
        cadena_pak = refrescar_cadena_pak()
    mam_placa = cargar_mam_placa()
    n06_pl = ((mam_placa.get("n06") or {}).get(PLACA) or {})
    if JUL_NOCHE_D0.isoformat() not in n06_pl or HASTA.isoformat() not in n06_pl:
        mam_placa = refrescar_mam_placa()
    maq_matriz = cargar_maq_matriz()
    n06_mq = ((maq_matriz.get("n06") or {}).get(MATRIZ_MAQ) or {})
    horas_mq = maq_matriz.get("horas") or {}
    if (
        MAQ_NOCHE_D0.isoformat() not in n06_mq
        or HASTA.isoformat() not in n06_mq
        or MAQ_NOCHE_D0.isoformat() not in horas_mq
        or HASTA.isoformat() not in horas_mq
    ):
        maq_matriz = refrescar_maq_matriz()
    bom_si500 = cargar_bom_si500()
    n06_b500 = ((bom_si500.get("n06") or {}).get(SI500) or {})
    if (
        JUL_NOCHE_D0.isoformat() not in n06_b500
        or HASTA.isoformat() not in n06_b500
        or CTRL_SI500.isoformat() not in n06_b500
    ):
        bom_si500 = refrescar_bom_si500()
    aeb = cargar_aeb()
    n06_ae11 = ((aeb.get("n06") or {}).get(MATRIZ_AEB) or {})
    n06_ae12 = ((aeb.get("n06") or {}).get(ANILLO_AEB) or {})
    if (
        MATRIZ_AA_HASTA.isoformat() not in n06_ae11
        or HASTA.isoformat() not in n06_ae11
        or HASTA.isoformat() not in n06_ae12
        or DESDE.isoformat() not in n06_ae12
    ):
        aeb = refrescar_aeb()
    cur = cargar_cur()
    n06_cs = ((cur.get("n06") or {}).get(ANILLO_SUR) or {})
    n06_cn = ((cur.get("n06") or {}).get(ANILLO_NORTE) or {})
    if (
        CUR_NOCHE_D0.isoformat() not in n06_cs
        or HASTA.isoformat() not in n06_cs
        or CUR_NOCHE_D0.isoformat() not in n06_cn
        or HASTA.isoformat() not in n06_cn
    ):
        cur = refrescar_cur()
    ppt = build_ppt(by, hourly, cadena_pak, mam_placa, maq_matriz, bom_si500, aeb, cur)
    ppt_res = build_ppt_resumen(by, hourly, cadena_pak, maq_matriz, bom_si500, aeb, cur)
    print("\n=== SALIDA ===")
    print(ppt)
    print(ppt_res)
    for mall in MALLS:
        ids = nodos_totales(mall)
        tot = totales_de(by, ids)
        ranked = sorted(ids, key=lambda n: -float((by.get(n) or {}).get("jul") or 0))
        print(
            mall["code"],
            "julio",
            round(tot["jul"], 1),
            {NOMBRE_CORTO.get(n, n): _etq_pct_julio(n, float((by.get(n) or {}).get("jul") or 0), tot["jul"]) for n in ranked},
        )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
