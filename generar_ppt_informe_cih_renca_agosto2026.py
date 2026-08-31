"""
Informa de consumos Renca en formato CIH (PPT 16:9 → PDF).

Toma el «Informe de consumos CIH Nº 2» (portada, intros satelitales, paleta) y lo
rellena con la auditoría homologada de agosto 2026 (4 puntos 24–26 vs 17–19;
ICCP al final: 10–16 con WES vs 17–23 sin WES).

Uso:
  python generar_ppt_informe_cih_renca_agosto2026.py
  python generar_ppt_informe_cih_renca_agosto2026.py --run-dir "reports/reporte de auditoria/auditoria_renca_semana_pasada_vs_esta_20260826_1806"
"""
from __future__ import annotations

import argparse
import json
import shutil
import subprocess
from dataclasses import dataclass
from datetime import date, datetime
from pathlib import Path
from typing import List, Optional, Sequence, Tuple

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

import generar_graficos_comparativos_desde_excel_consolidado as gxlsx
from generar_graficos_comparativos_desde_excel_consolidado import (
    _dibujar_comparativo_area_lineas,
    leer_matriz_consolidado,
)
from generar_reporte_word import format_number_chilean

ROOT = Path(__file__).resolve().parent
TZ_NAME = "America/Santiago"
TARIFA = 1300.0
CIH_PDF_DEFAULT = Path("/tmp/drive_pdfs/Informe_consumos_CIH_N2.pdf")
CIH_PAGES = Path("/tmp/drive_pdfs/cih_pages")
CIH_ASSETS = Path("/tmp/drive_pdfs/cih_assets")

AZUL = RGBColor(0x1D, 0x53, 0x72)
AZUL_OSCURO = RGBColor(0x14, 0x3D, 0x56)
BLANCO = RGBColor(255, 255, 255)
GRIS = RGBColor(0x1F, 0x1F, 0x1F)
GRIS_TITULO = RGBColor(0x66, 0x66, 0x66)
ROJO = RGBColor(0xC0, 0x50, 0x4D)
VERDE = RGBColor(0x54, 0x82, 0x35)
COLOR_CON = "#2a6fad"
COLOR_SIN = "#a0503d"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Páginas 1-based del PDF CIH Nº 2 (intros satelitales).
INTRO_PDF = {
    "000017-05": 4,   # Gimnasio
    "000017-06": 7,   # Piscina
    "000017-04": 10,  # Escuela Lo Velásquez
    "000017-08": 13,  # ICCO (oriente)
    "000017-07": 16,  # ICCP (poniente)
}

FONT_SANS = "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf"
FONT_SANS_B = "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf"


@dataclass
class Recinto:
    nid: str
    titulo: str
    corto: str
    direccion: str
    folder_glob: str
    iccp: bool
    comentario: str


RECINTOS: Tuple[Recinto, ...] = (
    Recinto(
        "000017-08",
        "INSTITUTO CUMBRE DE CÓNDORES ORIENTE",
        "ICCO",
        "Av. Dorsal 1365, Renca",
        "Auditoria_Colegio_ICCO_Renca_000017-08",
        False,
        "Control nocturno 00:01–06:00. El % del día mezcla ocupación escolar. "
        "No se reconstruye con la fuga del domingo 16.",
    ),
    Recinto(
        "000017-04",
        "ESCUELA LO VELÁSQUEZ",
        "Escuela Lo Velásquez",
        "Av. José Miguel Infante 7401, Renca",
        "Auditoria_Esc._Lo_Velásquez_000017-04",
        False,
        "Bypass abierto desde jue 27 ~17:00: de noche ~0,55 m³/h parejo (antes 0). "
        "El control quedó fuera de circuito; el % negativo no es ocupación.",
    ),
    Recinto(
        "000017-05",
        "GIMNASIO MUNICIPAL PONIENTE",
        "Gimnasio Municipal",
        "Av. Vicuña Mackenna 7836, Renca",
        "Auditoria_Gimnasio_municipal_000017-05",
        False,
        "Tras la regulación lun 24 el tope deja de ser 0,54 m³/h parejo y cubre picos reales.",
    ),
    Recinto(
        "000017-06",
        "PISCINA MUNICIPAL LO VELÁSQUEZ",
        "Piscina Municipal",
        "Av. José Miguel Infante 6502, Renca",
        "Auditoria_Piscina_municipal_000017-06",
        False,
        "Mayor aporte al ahorro de los 4 puntos con control. El control corta basal fuera de uso.",
    ),
    Recinto(
        "000017-07",
        "INSTITUTO CUMBRE DE CÓNDORES PONIENTE",
        "ICCP",
        "Av. Brasil 7965, Renca",
        "Auditoria_ICCP_000017-07",
        True,
        "Punto sí o sí. Semana completa lun 10–dom 16 con WES vs lun 17–dom 23 sin WES. "
        "Esta semana el control está off (OT 2282); no se usa 24–26. No entra al % de los 4.",
    ),
)


def _emu(inches: float) -> int:
    return int(Inches(inches))


def _font(path: str, size: int) -> ImageFont.FreeTypeFont:
    return ImageFont.truetype(path, size)


def _set_run(p, text: str, size: int, bold: bool = False, color: RGBColor = GRIS, align=None) -> None:
    p.clear()
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    if align is not None:
        p.alignment = align


def _textbox(slide, l, t, w, h, text: str, size: int, bold: bool = False, color: RGBColor = GRIS, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    _set_run(tf.paragraphs[0], text, size, bold, color, align)
    return box


def _shape_fill(slide, l, t, w, h, fill: RGBColor, text: str = "", size: int = 18, tcolor: RGBColor = BLANCO, bold: bool = True, align=PP_ALIGN.LEFT):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    tf = sh.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf._txBody.bodyPr.set("anchor", "ctr")
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    _set_run(p, text, size, bold, tcolor, align)
    # padding
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    return sh


def _add_logo_renca(slide) -> None:
    logo = CIH_ASSETS / "renca_logo.png"
    if logo.is_file():
        slide.shapes.add_picture(str(logo), Inches(11.15), Inches(0.12), width=Inches(1.95))


def _banner(slide, titulo: str) -> None:
    _shape_fill(slide, 0.0, 0.16, 8.15, 0.58, AZUL, titulo, size=20, align=PP_ALIGN.LEFT)
    _add_logo_renca(slide)


def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _cell(table, r, c, text: str, size: int = 10, bold: bool = False, color: RGBColor = GRIS, fill: Optional[RGBColor] = None, align=PP_ALIGN.CENTER):
    cell = table.cell(r, c)
    cell.text = ""
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    _set_run(p, text, size, bold, color, align)
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    return cell


def _n_horas(n_dias: int, hora_corte: int) -> int:
    if n_dias < 1:
        return 1
    if hora_corte >= 23:
        return n_dias * 24
    return (n_dias - 1) * 24 + (hora_corte + 1)


def _ahorro_mes(m3_con: float, m3_sin: float, n_horas: int) -> Tuple[float, float, float, float]:
    """(m³/día con, m³/día sin, ahorro mes m³, ahorro mes CLP)."""
    h = max(1, n_horas)
    d_con = m3_con * 24.0 / h
    d_sin = m3_sin * 24.0 / h
    mes = (d_sin - d_con) * 30.0
    return d_con, d_sin, mes, mes * TARIFA


def _pct(m3_con: float, m3_sin: float) -> float:
    if abs(m3_sin) < 1e-9:
        return 0.0
    return (m3_sin - m3_con) / m3_sin * 100.0


def _color_pct(pct: float) -> RGBColor:
    if pct >= 0:
        return AZUL
    return ROJO


def _ensure_cih_pages(pdf: Path) -> None:
    CIH_PAGES.mkdir(parents=True, exist_ok=True)
    CIH_ASSETS.mkdir(parents=True, exist_ok=True)
    need = any(not (CIH_PAGES / f"p{i:02d}.png").is_file() for i in (1, 4, 7, 10, 13, 16, 20))
    if need:
        import fitz

        doc = fitz.open(str(pdf))
        mat = fitz.Matrix(2, 2)
        for i, page in enumerate(doc):
            pix = page.get_pixmap(matrix=mat, alpha=False)
            pix.save(str(CIH_PAGES / f"p{i + 1:02d}.png"))
        doc.close()
    logo = CIH_ASSETS / "renca_logo.png"
    if not logo.is_file() and (CIH_PAGES / "p02.png").is_file():
        im = Image.open(CIH_PAGES / "p02.png")
        im.crop((1600, 8, 1910, 155)).save(logo)


def _portada_png(fecha_txt: str, linea1: str, linea2: str, linea3: str, dest: Path) -> Path:
    im = Image.open(CIH_PAGES / "p01.png").convert("RGB")
    dr = ImageDraw.Draw(im)
    # Tapa el título y la fecha del Nº 2 (coords del PDF × 2).
    dr.rectangle((600, 638, 1400, 830), fill="white")
    dr.rectangle((280, 948, 620, 1040), fill="white")
    f1 = _font(FONT_SANS_B, 50)
    f2 = _font(FONT_SANS_B, 38)
    f3 = _font(FONT_SANS, 22)
    fd = _font(FONT_SANS, 28)
    dr.text((622, 650), linea1, fill=(0x66, 0x66, 0x66), font=f1)
    dr.text((622, 712), linea2, fill=(0x1D, 0x53, 0x72), font=f2)
    dr.text((622, 770), linea3, fill=(0x66, 0x66, 0x66), font=f3)
    dr.text((305, 968), fecha_txt, fill=(0x7F, 0x7F, 0x7F), font=fd)
    dest.parent.mkdir(parents=True, exist_ok=True)
    im.save(dest, "PNG")
    return dest


def _find_run_dir(explicit: Optional[Path]) -> Path:
    if explicit is not None:
        p = explicit if explicit.is_absolute() else ROOT / explicit
        if not p.is_dir():
            raise SystemExit(f"No existe --run-dir: {p}")
        return p
    base = ROOT / "reports" / "reporte de auditoria"
    dirs = sorted(base.glob("auditoria_renca_semana_pasada_vs_esta_*"), key=lambda x: x.name, reverse=True)
    if not dirs:
        raise SystemExit("No hay carpeta auditoria_renca_semana_pasada_vs_esta_*")
    return dirs[0]


def _xlsx_de(run_dir: Path, rec: Recinto) -> Path:
    d = run_dir / rec.folder_glob
    xs = list(d.glob("consumo_consolidado_*.xlsx"))
    if not xs:
        raise FileNotFoundError(f"Sin Excel en {d}")
    return xs[0]


def _datos_recinto(xlsx: Path, hora_corte: int) -> dict:
    fechas, mats = leer_matriz_consolidado(xlsx)
    n = len(fechas)
    if n < 2 or n % 2:
        raise ValueError(f"Excel impar ({n} cols): {xlsx}")
    mid = n // 2
    con = mats[:mid]
    sin = mats[mid:]
    m3_con = float(sum(sum(col) for col in con))
    m3_sin = float(sum(sum(col) for col in sin))
    dias_con = [float(sum(col)) for col in con]
    dias_sin = [float(sum(col)) for col in sin]
    return {
        "fechas": fechas,
        "mats": mats,
        "mid": mid,
        "m3_con": m3_con,
        "m3_sin": m3_sin,
        "dias_con": dias_con,
        "dias_sin": dias_sin,
        "hora_corte": hora_corte,
        "n_dias": mid,
    }


def _hora_corte_run(run_dir: Path) -> Optional[int]:
    """Preferir hora_corte.txt / run_meta.json; si no, carpeta …_YYYYMMDD_HHMM → HH-1."""
    meta = run_dir / "run_meta.json"
    if meta.is_file():
        try:
            data = json.loads(meta.read_text(encoding="utf-8"))
            if "hora_corte" in data:
                return int(data["hora_corte"])
        except Exception:
            pass
    txt = run_dir / "hora_corte.txt"
    if txt.is_file():
        try:
            return int(txt.read_text(encoding="utf-8").strip())
        except Exception:
            pass
    parts = run_dir.name.rsplit("_", 2)
    if len(parts) >= 1 and parts[-1].isdigit() and len(parts[-1]) == 4:
        hh = int(parts[-1][:2])
        return max(0, min(23, hh - 1))
    return None


def _fecha_portada_run(run_dir: Path, fallback: datetime) -> datetime:
    meta = run_dir / "run_meta.json"
    if meta.is_file():
        try:
            data = json.loads(meta.read_text(encoding="utf-8"))
            if data.get("hasta"):
                d = date.fromisoformat(str(data["hasta"]))
                return datetime(d.year, d.month, d.day, 23, 59)
        except Exception:
            pass
    return fallback


def _hora_corte_desde_mats(mats: Sequence[Sequence[float]], mid: int) -> int:
    """Última hora con dato en el último día (ambos periodos cortados igual)."""
    last = np.array(mats[mid - 1], dtype=float)
    last_sin = np.array(mats[-1], dtype=float)
    for h in range(23, -1, -1):
        if last[h] > 1e-9 or last_sin[h] > 1e-9:
            return int(h)
    return 23


def _grafico_barras_dos(
    m3_con: float,
    m3_sin: float,
    dest: Path,
    lab_con: str = "Con WES",
    lab_sin: str = "Sin WES",
) -> Path:
    fig, ax = plt.subplots(figsize=(4.6, 3.15), dpi=140)
    fig.patch.set_facecolor("white")
    vals = [m3_con, m3_sin]
    bars = ax.bar([lab_con, lab_sin], vals, color=[COLOR_CON, COLOR_SIN], width=0.52)
    ax.set_ylabel("Σ (m³/h)", fontsize=10)
    ax.set_title("Total acumulado con WES V/S sin WES", fontsize=11, fontweight="bold")
    ax.set_ylim(0, max(vals + [1.0]) * 1.22)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    for b, v in zip(bars, vals):
        ax.text(
            b.get_x() + b.get_width() / 2,
            b.get_height(),
            f"{v:.1f}",
            ha="center",
            va="bottom",
            fontsize=11,
            fontweight="bold",
        )
    fig.tight_layout()
    dest.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(dest, dpi=140, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return dest


def _fmt_clp(clp: float) -> str:
    signo = "-" if clp < 0 else ""
    return f"{signo}${format_number_chilean(abs(clp), 0)}"


def _idx_martes(fechas: Sequence[date], mid: int) -> int:
    """Índice del martes homólogo en el periodo Con WES (0-based)."""
    for i, d in enumerate(fechas[:mid]):
        if d.weekday() == 1:
            return i
    return min(1, max(0, mid - 1))


def _grafico_dia(
    fechas,
    mats,
    j: int,
    dest: Path,
    lab_con: str,
    lab_sin: str,
    figsize: Tuple[float, float] = (7.35, 5.0),
) -> Path:
    mid = len(fechas) // 2
    d_con = fechas[j]
    d_sin = fechas[mid + j]
    y_con = np.array(mats[j], dtype=float)
    y_sin = np.array(mats[mid + j], dtype=float)
    horas = np.arange(24)
    wd = ("Lunes", "Martes", "Miércoles", "Jueves", "Viernes", "Sábado", "Domingo")[d_con.weekday()]
    fig, ax = plt.subplots(figsize=figsize, dpi=140)
    fig.patch.set_facecolor("white")
    _dibujar_comparativo_area_lineas(
        ax,
        horas,
        y_con,
        y_sin,
        f"{lab_con} ({d_con:%d-%m-%Y})",
        f"{lab_sin} ({d_sin:%d-%m-%Y})",
        f"Comparativo: {wd} — área + líneas (día homólogo)",
    )
    ax.title.set_fontsize(12)
    fig.tight_layout()
    dest.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(dest, dpi=140, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return dest


def _grafico_pares_barras(fechas, dias_con, dias_sin, dest: Path, lab_con: str, lab_sin: str) -> Path:
    mid = len(fechas) // 2
    labels = []
    for i, d in enumerate(fechas[:mid]):
        ds = fechas[mid + i]
        labels.append(f"{('Lun','Mar','Mié','Jue','Vie','Sáb','Dom')[d.weekday()]} {d:%d} vs {ds:%d}")
    x = np.arange(len(labels))
    w = 0.36
    fig, ax = plt.subplots(figsize=(7.4, 3.3), dpi=140)
    fig.patch.set_facecolor("white")
    ax.bar(x - w / 2, dias_con, width=w, color=COLOR_CON, label=lab_con, zorder=2)
    ax.bar(x + w / 2, dias_sin, width=w, color=COLOR_SIN, label=lab_sin, zorder=2)
    ax.set_xticks(list(x))
    ax.set_xticklabels(labels, fontsize=9)
    ax.set_ylabel("m³")
    ax.set_title("Pares homólogos (m³/día)", fontsize=12, fontweight="bold")
    ax.grid(axis="y", linestyle="--", alpha=0.45, zorder=0)
    ax.legend(fontsize=8, frameon=True)
    ax.set_ylim(bottom=0)
    fig.tight_layout()
    dest.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(dest, dpi=140, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return dest


def _grafico_5_puntos(nombres, m3_con, m3_sin, dest: Path) -> Path:
    x = np.arange(len(nombres))
    w = 0.36
    fig, ax = plt.subplots(figsize=(11.6, 4.6), dpi=140)
    fig.patch.set_facecolor("white")
    b1 = ax.bar(x - w / 2, m3_con, width=w, color=COLOR_CON, label="Con WES", zorder=2)
    b2 = ax.bar(x + w / 2, m3_sin, width=w, color=COLOR_SIN, label="Sin WES", zorder=2)
    ymax = max(list(m3_con) + list(m3_sin) + [1.0]) * 1.28
    ax.set_ylim(0, ymax)
    ax.set_xticks(list(x))
    ax.set_xticklabels(nombres, fontsize=10)
    ax.set_ylabel("m³/día")
    ax.set_title("Renca — 5 puntos (m³/día)  |  ICCP al final: 10–16 vs 17–23", fontsize=13, fontweight="bold", color="#1d5372")
    ax.grid(axis="y", linestyle="--", alpha=0.4, zorder=0)
    ax.legend(fontsize=9, loc="upper left")
    for bars in (b1, b2):
        for bar in bars:
            h = bar.get_height()
            ax.text(bar.get_x() + bar.get_width() / 2, h + ymax * 0.012, f"{h:.1f}", ha="center", va="bottom", fontsize=8)
    for i, (a, b) in enumerate(zip(m3_sin, m3_con)):
        pct = _pct(b, a)
        ax.text(x[i], ymax * 0.95, f"{pct:+.0f} %", ha="center", va="top", fontsize=10, fontweight="bold", color="#548235" if pct >= 0 else COLOR_SIN)
    fig.tight_layout()
    dest.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(dest, dpi=140, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return dest


def _add_picture_full(slide, png: Path) -> None:
    slide.shapes.add_picture(str(png), Inches(0), Inches(0), width=SLIDE_W, height=SLIDE_H)


def _add_picture(slide, png: Path, l, t, w, h) -> None:
    slide.shapes.add_picture(str(png), Inches(l), Inches(t), width=Inches(w), height=Inches(h))


def _fecha_fila(d: dict) -> str:
    fechas, mid = d["fechas"], d["mid"]
    a, b = fechas[0], fechas[mid - 1]
    c, e = fechas[mid], fechas[-1]
    return f"Con WES ({a:%d/%m} a {b:%d/%m}) vs Sin WES ({c:%d/%m} a {e:%d/%m})"


def _rango_txt(fechas: Sequence[date], mid: int, con: bool, hora_corte: int) -> str:
    xs = fechas[:mid] if con else fechas[mid:]
    a, b = xs[0], xs[-1]
    extra = "" if hora_corte >= 23 else f" 00:00–{hora_corte:02d}:59"
    if a == b:
        return f"{a:%d/%m/%Y}{extra}"
    return f"{a:%d/%m} al {b:%d/%m/%Y}{extra}"


def _slide_agregado(prs, filas: List[dict], n_horas: int, hora_corte: int, fechas_4, fechas_i) -> None:
    slide = _blank(prs)
    _banner(slide, "AGREGADO RESUMEN DE PUNTOS")
    rows = 1 + len(filas) + 1
    cols = 6
    table = slide.shapes.add_table(rows, cols, Inches(0.22), Inches(0.88), Inches(12.9), Inches(4.55)).table
    headers = [
        "Fecha auditada",
        "Recinto",
        "Eficiencia (%) vs línea base",
        "Ahorro mes estimado (m³)",
        "Ahorro mes estimado (CLP)",
        "Comentarios",
    ]
    for c, h in enumerate(headers):
        _cell(table, 0, c, h, size=9, bold=True, color=BLANCO, fill=AZUL)
    widths = [1.85, 2.05, 1.55, 1.55, 1.85, 4.05]
    for i, w in enumerate(widths):
        table.columns[i].width = Inches(w)

    tot4_con = tot4_sin = tot4_mes = 0.0
    for i, f in enumerate(filas, start=1):
        rec: Recinto = f["rec"]
        pct = f["pct"]
        mes = f["ahorro_mes"]
        clp = f["ahorro_clp"]
        if rec.iccp:
            fecha = _fecha_fila(f)
        else:
            fecha = _fecha_fila(f)
            tot4_con += f["m3_con"]
            tot4_sin += f["m3_sin"]
            tot4_mes += mes
        _cell(table, i, 0, fecha, size=8, align=PP_ALIGN.LEFT)
        _cell(table, i, 1, rec.corto, size=11, bold=True, align=PP_ALIGN.LEFT)
        _cell(table, i, 2, f"{format_number_chilean(pct, 1)} %", size=12, bold=True, color=_color_pct(pct))
        _cell(table, i, 3, format_number_chilean(mes, 0), size=12, bold=True, color=_color_pct(mes))
        signo = "-" if clp < 0 else ""
        _cell(table, i, 4, f"{signo}${format_number_chilean(abs(clp), 0)}", size=11, bold=True, color=_color_pct(clp))
        _cell(table, i, 5, rec.comentario, size=8, align=PP_ALIGN.LEFT)

    pct4 = _pct(tot4_con, tot4_sin)
    tot4_clp = tot4_mes * TARIFA
    r = rows - 1
    _cell(table, r, 0, "", size=9, fill=RGBColor(0xE8, 0xEE, 0xF3))
    _cell(table, r, 1, "TOTAL 4 puntos (sin ICCP)", size=10, bold=True, fill=RGBColor(0xE8, 0xEE, 0xF3), align=PP_ALIGN.LEFT)
    _cell(table, r, 2, "TOTAL", size=11, bold=True, fill=RGBColor(0xE8, 0xEE, 0xF3))
    _cell(table, r, 3, format_number_chilean(tot4_mes, 0), size=12, bold=True, fill=RGBColor(0xE8, 0xEE, 0xF3))
    _cell(table, r, 4, f"${format_number_chilean(tot4_clp, 0)}", size=11, bold=True, fill=RGBColor(0xE8, 0xEE, 0xF3))
    _cell(
        table,
        r,
        5,
        f"Eficiencia 4 puntos = {format_number_chilean(pct4, 1)} %  ·  "
        f"Ahorro {format_number_chilean(tot4_mes, 0)} m³ ({_fmt_clp(tot4_clp)}). ICCP no se mezcla.",
        size=9,
        bold=True,
        fill=RGBColor(0xE8, 0xEE, 0xF3),
        align=PP_ALIGN.LEFT,
    )

    notas = [
        f"●  Ventana 4 puntos: {_rango_txt(fechas_4, len(fechas_4)//2, True, hora_corte)} con WES vs "
        f"{_rango_txt(fechas_4, len(fechas_4)//2, False, hora_corte)} sin WES "
        f"(corte 00:00–{hora_corte:02d}:59, {n_horas} h). "
        f"ICCP al final: {_rango_txt(fechas_i, len(fechas_i)//2, True, 23)} vs "
        f"{_rango_txt(fechas_i, len(fechas_i)//2, False, 23)} (7 días, no se mezcla).",
        "●  Corresponde al 3er informe de auditoría (avance agosto 2026, formato CIH Nº 2).",
        "●  Ahorro mes = (m³/día sin − m³/día con) × 30, con m³/día prorateado a las horas de la ventana. "
        f"Tarifa 1.300 CLP/m³, sin sobreconsumo.",
    ]
    y = 5.55
    for ln in notas:
        _textbox(slide, 0.28, y, 12.7, 0.38, ln, 10, color=GRIS)
        y += 0.38


def _slide_ventanas(prs, png5: Path, hora_corte: int, filas: List[dict]) -> None:
    slide = _blank(prs)
    _banner(slide, "VENTANAS DE COMPARACIÓN — AGOSTO 2026")
    _shape_fill(slide, 0.28, 0.92, 6.2, 1.55, RGBColor(0xE8, 0xEE, 0xF3), "", size=10, tcolor=GRIS)
    _textbox(slide, 0.42, 0.98, 5.9, 0.32, "4 puntos con control (ICCO · Lo Velásquez · gimnasio · piscina)", 12, True, AZUL)
    f4 = next(f for f in filas if not f["rec"].iccp)
    _textbox(
        slide,
        0.42,
        1.32,
        5.9,
        0.95,
        f"Sin WES: {_rango_txt(f4['fechas'], f4['mid'], False, f4['hora_corte'])}\n"
        f"Con WES: {_rango_txt(f4['fechas'], f4['mid'], True, f4['hora_corte'])}\n"
        "Homólogo día a día, mismo corte de hora.",
        12,
        False,
        GRIS,
    )

    _shape_fill(slide, 6.85, 0.92, 6.2, 1.55, RGBColor(0xFF, 0xF2, 0xCC), "", size=10)
    _textbox(slide, 7.0, 0.98, 5.9, 0.32, "ICCP (al final) — punto sí o sí", 12, True, AZUL)
    _textbox(slide, 7.0, 1.32, 5.9, 0.95, "Con WES: lun 10 – dom 16 agosto\nSin WES: lun 17 – dom 23 agosto\nSemana completa. No entra al % de los 4 puntos.", 12, False, GRIS)

    _add_picture(slide, png5, 0.35, 2.62, 12.6, 4.55)


def _slide_comparativo(prs, rec: Recinto, d: dict, png_bar: Path, png_dia: Path, n_horas: int, lab_con: str, lab_sin: str) -> None:
    slide = _blank(prs)
    _banner(slide, rec.titulo)
    _add_picture(slide, png_bar, 0.22, 0.88, 4.85, 3.35)
    d_con, d_sin, mes, clp = _ahorro_mes(d["m3_con"], d["m3_sin"], n_horas)
    pct = d["pct"]
    fechas = d["fechas"]
    mid = d["mid"]
    hc = d["hora_corte"]
    y = 4.28
    _textbox(slide, 0.28, y, 4.8, 0.28, f"{lab_sin}:  {_rango_txt(fechas, mid, False, hc)}", 11, True, GRIS)
    _textbox(slide, 0.28, y + 0.30, 4.8, 0.28, f"{lab_con}:  {_rango_txt(fechas, mid, True, hc)}", 11, True, GRIS)
    _textbox(slide, 0.28, y + 0.68, 4.8, 0.42, f"Eficiencia:   {format_number_chilean(pct, 1)} %", 22, True, _color_pct(pct))
    _textbox(
        slide,
        0.28,
        y + 1.12,
        4.8,
        0.48,
        f"Ahorro estimado mes:   {format_number_chilean(mes, 0)} m³  ·  {_fmt_clp(clp)}",
        14,
        True,
        _color_pct(mes),
    )
    _textbox(
        slide,
        0.28,
        y + 1.58,
        4.9,
        0.45,
        f"*Considera: consumo promedio basal de {format_number_chilean(d_sin, 0)} m³/día "
        f"y post activación de {format_number_chilean(d_con, 0)} m³/día.",
        9,
        False,
        GRIS,
    )
    _textbox(slide, 5.35, 0.82, 7.6, 0.32, "Comparativo horario homologado — martes", 14, True, AZUL)
    _add_picture(slide, png_dia, 5.28, 1.14, 7.75, 5.53)


def _slide_comentarios(prs, rec: Recinto, d: dict, png_pares: Path, n_horas: int) -> None:
    slide = _blank(prs)
    _banner(slide, rec.titulo)
    _add_picture(slide, png_pares, 0.25, 0.90, 7.55, 3.55)
    d_con, d_sin, mes, clp = _ahorro_mes(d["m3_con"], d["m3_sin"], n_horas)
    pct = d["pct"]
    lab_sin = "Sin WES (línea base)"
    lab_con = "Con WES (control activo)"
    _textbox(slide, 8.05, 0.95, 4.9, 0.28, lab_sin, 11, True, GRIS)
    _textbox(slide, 8.05, 1.22, 4.9, 0.28, _rango_txt(d["fechas"], d["mid"], False, d["hora_corte"]), 12, False, GRIS)
    _textbox(slide, 8.05, 1.58, 4.9, 0.28, lab_con, 11, True, GRIS)
    _textbox(slide, 8.05, 1.85, 4.9, 0.28, _rango_txt(d["fechas"], d["mid"], True, d["hora_corte"]), 12, False, GRIS)
    _textbox(slide, 8.05, 2.28, 4.9, 0.42, f"Eficiencia:  {format_number_chilean(pct, 1)} %", 20, True, _color_pct(pct))
    _textbox(
        slide,
        8.05,
        2.75,
        4.9,
        0.48,
        f"Ahorro estimado mes:  {format_number_chilean(mes, 0)} m³  ·  {_fmt_clp(clp)}",
        13,
        True,
        _color_pct(mes),
    )
    _textbox(slide, 8.05, 3.28, 4.9, 0.32, f"{format_number_chilean(d['m3_con'], 1)} m³ con  vs  {format_number_chilean(d['m3_sin'], 1)} m³ sin", 11, False, GRIS)
    _textbox(slide, 8.05, 3.58, 4.9, 0.32, f"Tarifa 1.300 CLP/m³", 11, False, GRIS)

    _textbox(slide, 0.28, 4.55, 12.7, 0.32, "Comentarios:", 16, True, AZUL)
    bullets = [f">  {rec.comentario}"]
    # pares
    mid = d["mid"]
    fechas = d["fechas"]
    pares = []
    for i in range(mid):
        a, p = d["dias_sin"][i] - d["dias_con"][i], _pct(d["dias_con"][i], d["dias_sin"][i])
        pares.append(
            f"{('Lun','Mar','Mié','Jue','Vie','Sáb','Dom')[fechas[i].weekday()]} "
            f"{fechas[i]:%d} vs {fechas[mid+i]:%d}: "
            f"{format_number_chilean(p, 1)} % ({format_number_chilean(a, 1)} m³)"
        )
    if rec.iccp:
        bullets.append(">  " + "  ·  ".join(pares[:4]))
        bullets.append(">  " + "  ·  ".join(pares[4:]))
    else:
        bullets.append(">  " + "  ·  ".join(pares))
    y = 4.92
    for b in bullets:
        _textbox(slide, 0.28, y, 12.7, 0.55, b, 13, False, GRIS)
        y += 0.58


def _slide_anexo_vwb(prs) -> None:
    png = CIH_PAGES / "p20.png"
    slide = _blank(prs)
    if png.is_file():
        _add_picture_full(slide, png)
        return
    _banner(slide, "ANEXO: METODOLOGÍA - VWB")
    _textbox(
        slide,
        0.4,
        1.1,
        12.4,
        5.8,
        "VWB = [Extracción de referencia] − [Extracción con proyecto]\n\n"
        "Fuente: Reig, P. et al. 2019. Volumetric Water Benefit Accounting (VWBA). WRI.",
        16,
        False,
        GRIS,
    )


def _convertir_pdf(pptx: Path) -> Optional[Path]:
    soffice = shutil.which("soffice") or shutil.which("libreoffice") or "/usr/bin/soffice"
    if not Path(soffice).exists():
        return None
    profile = Path("/tmp/lo_profile_cih")
    profile.mkdir(parents=True, exist_ok=True)
    # Copiar a /tmp: LibreOffice Impress falla con espacios en la ruta.
    tmp_in = Path("/tmp") / pptx.name.replace(" ", "_")
    shutil.copy2(pptx, tmp_in)
    subprocess.run(
        [
            soffice,
            f"-env:UserInstallation=file://{profile}",
            "--headless",
            "--norestore",
            "--nolockcheck",
            "--nologo",
            "--nofirststartwizard",
            "--convert-to",
            "pdf",
            "--outdir",
            "/tmp",
            str(tmp_in),
        ],
        check=False,
        capture_output=True,
        text=True,
        timeout=180,
    )
    tmp_pdf = tmp_in.with_suffix(".pdf")
    pdf = pptx.with_suffix(".pdf")
    if tmp_pdf.is_file():
        shutil.copy2(tmp_pdf, pdf)
        return pdf
    return pdf if pdf.is_file() else None


def main() -> int:
    ap = argparse.ArgumentParser(description="PPT/PDF Informe de consumos Renca estilo CIH.")
    ap.add_argument("--run-dir", type=Path, default=None)
    ap.add_argument("--cih-pdf", type=Path, default=CIH_PDF_DEFAULT)
    ap.add_argument("-o", "--output", type=Path, default=None)
    args = ap.parse_args()

    run_dir = _find_run_dir(args.run_dir)
    pdf_cih = args.cih_pdf
    if not pdf_cih.is_file():
        raise SystemExit(f"Falta el PDF CIH de plantilla: {pdf_cih}")
    _ensure_cih_pages(pdf_cih)

    gxlsx.LABEL_P1 = "Con WES"
    gxlsx.LABEL_P2 = "Sin WES"

    datos: List[dict] = []
    hora_corte_run = _hora_corte_run(run_dir) or 17
    for rec in RECINTOS:
        xlsx = _xlsx_de(run_dir, rec)
        d = _datos_recinto(xlsx, hora_corte_run)
        d["rec"] = rec
        d["xlsx"] = xlsx
        d["pct"] = _pct(d["m3_con"], d["m3_sin"])
        if rec.iccp:
            d["hora_corte"] = 23
            d["n_horas"] = d["n_dias"] * 24
        else:
            d["hora_corte"] = hora_corte_run
            d["n_horas"] = _n_horas(d["n_dias"], hora_corte_run)
        _c, _s, mes, clp = _ahorro_mes(d["m3_con"], d["m3_sin"], d["n_horas"])
        d["ahorro_mes"] = mes
        d["ahorro_clp"] = clp
        datos.append(d)
    hora_corte = int(next(d["hora_corte"] for d in datos if not d["rec"].iccp))
    n_horas = next(d["n_horas"] for d in datos if not d["rec"].iccp)

    out_charts = run_dir / "graficos_cih"
    out_charts.mkdir(exist_ok=True)
    png5 = _grafico_5_puntos(
        ["ICCO", "Lo Velásquez", "Gimnasio", "Piscina", "ICCP"],
        [d["m3_con"] * 24.0 / d["n_horas"] for d in datos],
        [d["m3_sin"] * 24.0 / d["n_horas"] for d in datos],
        out_charts / "cinco_puntos.png",
    )

    hoy = datetime.now()
    try:
        from zoneinfo import ZoneInfo

        hoy = datetime.now(ZoneInfo("America/Santiago"))
    except Exception:
        pass
    hoy = _fecha_portada_run(run_dir, hoy)
    meses = (
        "ENERO", "FEBRERO", "MARZO", "ABRIL", "MAYO", "JUNIO",
        "JULIO", "AGOSTO", "SEPTIEMBRE", "OCTUBRE", "NOVIEMBRE", "DICIEMBRE",
    )
    fecha_portada = f"{hoy.day} {meses[hoy.month - 1]} {hoy.year}"
    cover = _portada_png(
        fecha_portada,
        "Informe de consumos Nº 3",
        "CENTRO DE INTELIGENCIA HÍDRICA",
        "Avance agosto 2026  ·  5 recintos Renca",
        out_charts / "portada_n3.png",
    )

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    s0 = _blank(prs)
    _add_picture_full(s0, cover)

    fechas_4 = next(d["fechas"] for d in datos if not d["rec"].iccp)
    fechas_i = next(d["fechas"] for d in datos if d["rec"].iccp)
    _slide_agregado(prs, datos, n_horas, hora_corte, fechas_4, fechas_i)
    _slide_ventanas(prs, png5, hora_corte, datos)

    for d in datos:
        rec: Recinto = d["rec"]
        intro = CIH_PAGES / f"p{INTRO_PDF[rec.nid]:02d}.png"
        if intro.is_file():
            s = _blank(prs)
            _add_picture_full(s, intro)
        lab_con = "Con WES"
        lab_sin = "Sin WES"
        gdir = out_charts / rec.nid
        bar = _grafico_barras_dos(
            d["m3_con"],
            d["m3_sin"],
            gdir / "barras.png",
            "Con WES",
            "Sin WES",
        )
        j_mar = _idx_martes(d["fechas"], d["mid"])
        p_mar = _grafico_dia(
            d["fechas"],
            d["mats"],
            j_mar,
            gdir / "martes.png",
            lab_con,
            lab_sin,
        )
        _slide_comparativo(prs, rec, d, bar, p_mar, d["n_horas"], lab_con, lab_sin)
        pares = _grafico_pares_barras(d["fechas"], d["dias_con"], d["dias_sin"], gdir / "pares.png", lab_con, lab_sin)
        _slide_comentarios(prs, rec, d, pares, d["n_horas"])

    _slide_anexo_vwb(prs)
    s_end = _blank(prs)
    _add_picture_full(s_end, cover)

    stamp = hoy.strftime("%Y%m%d_%H%M")
    out = (
        args.output.resolve()
        if args.output
        else (run_dir / f"Informe_consumos_CIH_Renca_N3_avance_{stamp}.pptx").resolve()
    )
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    pdf = _convertir_pdf(out)
    print(out)
    if pdf:
        print(pdf)
    else:
        print("(sin PDF: falta soffice)")
    # resumen consola
    print(f"DIR {run_dir}")
    print(f"Corte {hora_corte:02d}:59  |  {n_horas} h")
    tot_c = sum(d["m3_con"] for d in datos if not d["rec"].iccp)
    tot_s = sum(d["m3_sin"] for d in datos if not d["rec"].iccp)
    print(f"4 puntos {tot_c:.1f} vs {tot_s:.1f} = {_pct(tot_c, tot_s):.1f}%")
    for d in datos:
        print(f"  {d['rec'].corto}: {d['pct']:.1f}%  {d['m3_con']:.1f} vs {d['m3_sin']:.1f}  mes {d['ahorro_mes']:.0f} m³")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
