"""
Genera un PowerPoint alineado al informe «Informe de consumos» CIH:
una sola portada, AGREGADO RESUMEN DE PUNTOS, tabla, una diapositiva por establecimiento
(nombre + dirección + gráfico área+líneas + métricas), Pareto si hay PNG, anexo VWB.

Los gráficos deben generarse con ``generar_graficos_comparativos_desde_excel_consolidado.py``
(estilo área + líneas, ejes h / m³/h, leyenda con marco).

Réplica fiel de la portada CIH
------------------------------
Opción recomendada: en PowerPoint, abre el informe CIH original, deja solo la diapositiva 1,
«Guardar como» → un .pptx de una sola página (p. ej. ``portada_cih.pptx``) y genera así::

  python generar_ppt_informe_auditorias_renca.py --plantilla portada_cih.pptx

El script **abre ese archivo** y añade el resto de diapositivas después de la portada (fondos,
logos y tipografías quedan iguales al original). Sin ``--plantilla``, se usa la portada
generada por código (posiciones medidas desde el CIH).

Uso:
  python generar_ppt_informe_auditorias_renca.py
  python generar_ppt_informe_auditorias_renca.py --plantilla "C:\\ruta\\portada_cih.pptx"
  python generar_ppt_informe_auditorias_renca.py --regenerar-graficos
"""
from __future__ import annotations

import argparse
from datetime import datetime
from pathlib import Path
from typing import List, Optional, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from generar_graficos_comparativos_desde_excel_consolidado import (
    _limpiar_pngs_carpeta_graficos,
    generar_pngs,
    leer_matriz_consolidado,
    totales_rejilla_desde_excel_consolidado,
)
from generar_reporte_word import get_node_name

ROOT = Path(__file__).resolve().parent
BASE_AUDIT = (
    ROOT
    / "reports"
    / "reporte de auditoria"
    / "auditoria_puntos_renca_abril_2026"
)
XLSX_NOMBRE = "consumo_consolidado_parseo_filas_abr06-12_abr13-19_2026.xlsx"

AUDITORIAS: Tuple[Tuple[str, str], ...] = (
    ("000017-08", "Auditoria ICCO Renca 000017-08"),
    ("000017-04", "Auditoria Escuela Lo Velazquez 000017-04"),
    ("000017-06", "Auditoria Piscina Municipal 000017-06"),
    ("000017-05", "Auditoria Gimnasio 000017-05"),
    ("000017-07", "Auditoria Cumbre de condores 000017-07"),
)

# Direcciones (misma línea que informes municipales CIH donde aplica)
DIRECCION_POR_NODO: dict[str, str] = {
    "000017-08": "Colegio ICCO Renca — Renca (referencia auditoría WES)",
    "000017-04": "Av. José Miguel Infante 7401, Renca",
    "000017-06": "Av. José Miguel Infante 6502, Renca",
    "000017-05": "Av. Vicuña Mackenna 7836, Renca",
    "000017-07": "Av. Brasil 7965, Renca",
}

COLOR_TITULO = RGBColor(0, 51, 102)
COLOR_CUERPO = RGBColor(40, 40, 40)

TEXTO_ANEXO_VWB = """Diversas actividades pueden reducir el volumen de agua extraída de una fuente, incluyendo transacciones legales (por ejemplo, arrendamientos o compras de derechos de agua), medidas de eficiencia operativa, reparación de fugas, entubado de canales de riego, medidas de eficiencia y reutilización del agua.
El volumen de extracción reducido se calcula como la diferencia entre el volumen de extracción en la condición "con proyecto" en comparación con la condición de "referencia" (línea de base). La condición de "referencia" describe la extracción actual. La condición "con proyecto" representa la extracción después de la implementación de medidas de eficiencia, reutilización de agua, reparación de fugas o transacciones legales.
Para todas estas aplicaciones, el VWB (Volumetric Water Benefit - Beneficio Hídrico Volumétrico) se calcula como la disminución del volumen de extracción.

Fórmula de cálculo
VWB = [Extracción de referencia] - [Extracción con proyecto]

Fuente: Reig, P., W. Larson, S. Vionnet y J.B. Bayart. 2019. "Volumetric Water Benefit Accounting (VWBA): A Method for Implementing and Valuing Water Stewardship Activities." Working Paper. Washington, DC: World Resources Institute. Disponible en línea en: www.wri.org/publication/volumetric-water-benefit-accounting.
"""


def _fecha_portada_upper(d: datetime) -> str:
    meses = (
        "enero",
        "febrero",
        "marzo",
        "abril",
        "mayo",
        "junio",
        "julio",
        "agosto",
        "septiembre",
        "octubre",
        "noviembre",
        "diciembre",
    )
    return f"{d.day} {meses[d.month - 1].upper()} {d.year}"


def _pct(m3_con: float, m3_sin: float) -> Optional[float]:
    if m3_sin <= 1e-9:
        return None
    return (m3_sin - m3_con) / m3_sin * 100.0


def _cargar_datos() -> List[Tuple[str, str, Path, float, float, Optional[float]]]:
    out: List[Tuple[str, str, Path, float, float, Optional[float]]] = []
    for node_id, sub in AUDITORIAS:
        px = (BASE_AUDIT / sub / XLSX_NOMBRE).resolve()
        if not px.is_file():
            raise FileNotFoundError(f"Falta Excel de auditoría: {px}")
        tc, ts, _ = totales_rejilla_desde_excel_consolidado(px)
        nombre = (get_node_name(node_id) or "").strip() or node_id
        pct = _pct(float(tc), float(ts))
        out.append((node_id, nombre, px, float(tc), float(ts), pct))
    return out


def _regenerar_pngs_todas_las_auditorias() -> None:
    for _node_id, sub in AUDITORIAS:
        xlsx = (BASE_AUDIT / sub / XLSX_NOMBRE).resolve()
        out_dir = BASE_AUDIT / sub / "graficos_comparativos"
        if not xlsx.is_file():
            print(f"[skip] Sin Excel: {xlsx}")
            continue
        _limpiar_pngs_carpeta_graficos(out_dir)
        fechas, mats = leer_matriz_consolidado(xlsx)
        generar_pngs(fechas, mats, out_dir)
        print(f"[ok] PNG en {out_dir}")


def _ultimos_pareto_reporte_agregado() -> Tuple[Optional[Path], Optional[Path]]:
    pat = "Reporte_agregado_5_auditorias_*"
    dirs = sorted(BASE_AUDIT.glob(pat), key=lambda p: p.stat().st_mtime, reverse=True)
    if not dirs:
        return None, None
    d = dirs[0]
    p1 = d / "pareto_semana_CON_WES_13_19_abr_2026.png"
    p2 = d / "pareto_semana_SIN_WES_06_12_abr_2026.png"
    return (p1 if p1.is_file() else None, p2 if p2.is_file() else None)


def _png_perfil_informe(gdir: Path) -> Optional[Path]:
    """Prioriza Lunes homólogo (mismo criterio que ejemplo CIH), si no 01_promedio."""
    lunes = gdir / "04_area_Lunes.png"
    if lunes.is_file():
        return lunes
    p01 = gdir / "01_promedio_24h_dos_periodos.png"
    return p01 if p01.is_file() else None


def _add_portada_cih(
    prs: Presentation,
    fecha_inferior: str,
    linea1: str,
    linea2: str,
) -> None:
    """Portada tipo slide 1 / 15 / 20 del informe CIH (posiciones aproximadas)."""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    # Fecha o mes (banda inferior centrada)
    dbox = slide.shapes.add_textbox(Inches(2.0), Inches(6.22), Inches(9.3), Inches(0.55))
    dp = dbox.text_frame.paragraphs[0]
    dp.text = fecha_inferior
    dp.alignment = PP_ALIGN.CENTER
    dp.font.size = Pt(14)
    dp.font.bold = True
    dp.font.color.rgb = RGBColor(60, 60, 60)

    # Bloque central: dos líneas (Informe… / CENTRO DE INTELIGENCIA HÍDRICA)
    cbox = slide.shapes.add_textbox(Inches(3.75), Inches(4.25), Inches(6.2), Inches(2.0))
    tf = cbox.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = linea1
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(25)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)
    p2 = tf.add_paragraph()
    p2.text = linea2
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(19)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(0, 0, 0)
    p2.space_before = Pt(6)


def _add_agregado_resumen_cih(prs: Presentation, cuerpo: str) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    tbox = slide.shapes.add_textbox(Inches(0.25), Inches(0.2), Inches(8.2), Inches(0.9))
    tp = tbox.text_frame.paragraphs[0]
    tp.text = "AGREGADO RESUMEN DE PUNTOS"
    tp.font.size = Pt(30)
    tp.font.bold = True
    tp.font.color.rgb = COLOR_TITULO

    bbox = slide.shapes.add_textbox(Inches(0.25), Inches(1.05), Inches(12.5), Inches(5.8))
    btf = bbox.text_frame
    btf.word_wrap = True
    lines = cuerpo.strip().split("\n")
    btf.text = lines[0]
    btf.paragraphs[0].font.size = Pt(11)
    btf.paragraphs[0].font.color.rgb = COLOR_CUERPO
    for ln in lines[1:]:
        par = btf.add_paragraph()
        par.text = ln
        par.font.size = Pt(11)
        par.font.color.rgb = COLOR_CUERPO
        par.space_after = Pt(4)


def _add_tabla_resumen(prs: Presentation, filas: List[Tuple[str, str, float, float, Optional[float]]]) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    tbox = slide.shapes.add_textbox(Inches(0.25), Inches(0.35), Inches(12.0), Inches(0.75))
    tbox.text_frame.text = "AGREGADO — DETALLE POR ESTABLECIMIENTO"
    tbox.text_frame.paragraphs[0].font.size = Pt(24)
    tbox.text_frame.paragraphs[0].font.bold = True
    tbox.text_frame.paragraphs[0].font.color.rgb = COLOR_TITULO

    rows = 1 + len(filas)
    cols = 5
    left = Inches(0.45)
    top = Inches(1.1)
    width = Inches(12.35)
    height = Inches(0.38 * rows + 0.25)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    headers = ["Nodo", "Establecimiento (app)", "m³ sem. Con WES", "m³ sem. ref. Sin WES", "% rendimiento"]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(11)

    for r, (nid, nom, tc, ts, pct) in enumerate(filas, start=1):
        table.cell(r, 0).text = nid
        table.cell(r, 1).text = nom[:48] + ("…" if len(nom) > 48 else "")
        table.cell(r, 2).text = f"{tc:,.1f}".replace(",", ".")
        table.cell(r, 3).text = f"{ts:,.1f}".replace(",", ".")
        table.cell(r, 4).text = "—" if pct is None else f"{pct:.1f} %"
        for c in range(5):
            for p in table.cell(r, c).text_frame.paragraphs:
                p.font.size = Pt(10)


def _add_picture_slide(prs: Presentation, title: str, png: Path) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    tbox = slide.shapes.add_textbox(Inches(0.25), Inches(0.3), Inches(12.0), Inches(0.65))
    tbox.text_frame.text = title
    tbox.text_frame.paragraphs[0].font.size = Pt(20)
    tbox.text_frame.paragraphs[0].font.bold = True
    tbox.text_frame.paragraphs[0].font.color.rgb = COLOR_TITULO
    slide.shapes.add_picture(str(png), Inches(0.45), Inches(1.0), width=Inches(12.35), height=Inches(5.55))


def _add_diapositiva_establecimiento(
    prs: Presentation,
    nombre_mayus: str,
    direccion: str,
    node_id: str,
    m3_con: float,
    m3_sin: float,
    pct: Optional[float],
    png: Optional[Path],
) -> None:
    """Una sola diapositiva por establecimiento: título, dirección, gráfico y métricas (CIH)."""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    tbox = slide.shapes.add_textbox(Inches(0.25), Inches(0.18), Inches(11.5), Inches(0.55))
    tp = tbox.text_frame.paragraphs[0]
    tp.text = nombre_mayus
    tp.font.size = Pt(30)
    tp.font.bold = True
    tp.font.color.rgb = COLOR_TITULO

    abox = slide.shapes.add_textbox(Inches(0.25), Inches(0.72), Inches(12.0), Inches(0.38))
    ap = abox.text_frame.paragraphs[0]
    ap.text = direccion
    ap.font.size = Pt(10)
    ap.font.bold = True
    ap.font.color.rgb = COLOR_CUERPO

    y_chart = 1.12
    if png is not None and png.is_file():
        slide.shapes.add_picture(
            str(png),
            Inches(0.25),
            Inches(y_chart),
            width=Inches(7.25),
            height=Inches(3.25),
        )
    else:
        miss = slide.shapes.add_textbox(Inches(0.35), Inches(2.0), Inches(7.0), Inches(1.0))
        miss.text_frame.text = (
            "(Perfil comparativo no encontrado. Ejecutar generar_graficos_comparativos_desde_excel_consolidado.)"
        )
        miss.text_frame.paragraphs[0].font.size = Pt(11)
        miss.text_frame.paragraphs[0].font.italic = True

    s1 = slide.shapes.add_textbox(Inches(7.95), Inches(3.35), Inches(3.35), Inches(0.42))
    s1.text_frame.text = "Perfil semanal Con WES (13–19 abr 2026)"
    s1.text_frame.paragraphs[0].font.size = Pt(10)
    s1.text_frame.paragraphs[0].font.bold = True

    s2 = slide.shapes.add_textbox(Inches(7.95), Inches(6.45), Inches(3.35), Inches(0.42))
    s2.text_frame.text = "Referencia Sin WES (6–12 abr 2026)"
    s2.text_frame.paragraphs[0].font.size = Pt(10)
    s2.text_frame.paragraphs[0].font.bold = True

    bloque = (
        f"Nodo: {node_id}\n"
        f"Volumen semanal Con WES (Σ rejilla): {m3_con:,.1f} m³\n"
        f"Volumen semanal ref. Sin WES (Σ rejilla): {m3_sin:,.1f} m³\n"
    )
    if pct is not None:
        bloque += f"% rendimiento (ahorro vs ref.): {pct:.1f} %\n"
        bloque += "Definición: (Sin WES − Con WES) / Sin WES × 100."
    mbox = slide.shapes.add_textbox(Inches(0.65), Inches(4.45), Inches(6.6), Inches(2.55))
    mp = mbox.text_frame
    mp.word_wrap = True
    mp.text = bloque
    for p in mp.paragraphs:
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_CUERPO


def _add_anexo_vwb(prs: Presentation) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    tbox = slide.shapes.add_textbox(Inches(0.25), Inches(0.2), Inches(12.0), Inches(0.75))
    tbox.text_frame.text = "ANEXO: METODOLOGÍA - VWB"
    tbox.text_frame.paragraphs[0].font.size = Pt(24)
    tbox.text_frame.paragraphs[0].font.bold = True
    tbox.text_frame.paragraphs[0].font.color.rgb = COLOR_TITULO

    bbox = slide.shapes.add_textbox(Inches(0.35), Inches(1.0), Inches(12.2), Inches(6.0))
    btf = bbox.text_frame
    btf.word_wrap = True
    lines = TEXTO_ANEXO_VWB.split("\n")
    btf.text = lines[0]
    btf.paragraphs[0].font.size = Pt(11)
    btf.paragraphs[0].font.color.rgb = COLOR_CUERPO
    for ln in lines[1:]:
        par = btf.add_paragraph()
        par.text = ln
        par.font.size = Pt(11)
        par.font.color.rgb = COLOR_CUERPO
        par.space_after = Pt(3)


def main() -> int:
    ap = argparse.ArgumentParser(description="PPT informe auditorías Renca (estructura CIH).")
    ap.add_argument(
        "-o",
        "--output",
        type=Path,
        default=None,
        help="Ruta del .pptx de salida",
    )
    ap.add_argument(
        "--regenerar-graficos",
        action="store_true",
        help="Regenera PNG en graficos_comparativos/ de cada auditoría antes de armar la PPT",
    )
    ap.add_argument(
        "--plantilla",
        type=Path,
        default=None,
        help=".pptx de **una sola** diapositiva: portada copiada del informe CIH (fondos/logos idénticos)",
    )
    args = ap.parse_args()

    if args.regenerar_graficos:
        _regenerar_pngs_todas_las_auditorias()

    datos = _cargar_datos()
    filas_tab = [(d[0], d[1], d[3], d[4], d[5]) for d in datos]

    hoy = datetime.now()

    plantilla = args.plantilla.resolve() if args.plantilla else None
    if plantilla is not None and plantilla.is_file():
        prs = Presentation(str(plantilla))
        if len(prs.slides) != 1:
            raise SystemExit(
                f"La plantilla debe tener exactamente 1 diapositiva (solo portada). "
                f"Tiene {len(prs.slides)}. En PowerPoint: archivo nuevo → pegar solo la portada del CIH → guardar."
            )
    else:
        if args.plantilla:
            print(f"[WARN] Plantilla no encontrada: {args.plantilla} — usando portada generada por código.")
        prs = Presentation()
        prs.slide_width = int(Inches(13.333))
        prs.slide_height = int(Inches(7.5))
        _add_portada_cih(
            prs,
            _fecha_portada_upper(hoy),
            "Informe de auditorías WES ",
            "CENTRO DE INTELIGENCIA HÍDRICA",
        )

    cuerpo_agregado = (
        "El presente informe consolida las auditorías de consumo hídrico con monitoreo inteligente WES "
        "en cinco establecimientos de Renca (incluye Colegio ICCO Renca).\n\n"
        "Se comparan dos ventanas de siete días en abril 2026: operación con control WES (13–19 abr) "
        "y período de referencia sin control (6–12 abr), alineado a los informes individuales y a la "
        "matriz «Consolidado» de cada carpeta de auditoría.\n\n"
        "Corresponde a la consolidación de auditorías realizadas en el marco del programa de "
        "Centro de Inteligencia Hídrica. Los totales semanales provienen de la suma hora a hora de la rejilla. "
        "El porcentaje de rendimiento indica el ahorro de volumen de la semana con WES respecto a la semana de referencia."
    )
    _add_agregado_resumen_cih(prs, cuerpo_agregado)
    _add_tabla_resumen(prs, filas_tab)

    for node_id, nombre, _px, tc, ts, pct in datos:
        sub = next(s for n, s in AUDITORIAS if n == node_id)
        gdir = BASE_AUDIT / sub / "graficos_comparativos"
        png = _png_perfil_informe(gdir)
        nom_u = (nombre or node_id).upper()
        dir_txt = DIRECCION_POR_NODO.get(node_id, "Renca")
        _add_diapositiva_establecimiento(prs, nom_u, dir_txt, node_id, tc, ts, pct, png)

    pc, ps = _ultimos_pareto_reporte_agregado()
    if pc is not None:
        _add_picture_slide(prs, "Pareto — consumo semanal Con WES (13–19 abr 2026)", pc)
    if ps is not None:
        _add_picture_slide(prs, "Pareto — consumo semanal referencia Sin WES (6–12 abr 2026)", ps)

    _add_anexo_vwb(prs)

    stamp = hoy.strftime("%Y%m%d_%H%M")
    out = (
        args.output.resolve()
        if args.output
        else (BASE_AUDIT / f"Informe_auditorias_WES_Renca_{stamp}.pptx").resolve()
    )
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(out.resolve())
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
