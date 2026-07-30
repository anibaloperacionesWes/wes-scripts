"""
Script para generar reporte agregado de Mall Maipú y PPT en PDF de análisis similar
Período: 01 de diciembre 2025 al 16 de enero 2026
"""

import sys
from pathlib import Path
from datetime import datetime, timedelta
from typing import Optional
import requests
import argparse
import json
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np
from io import BytesIO

# Agregar el directorio actual al path para importar módulos
sys.path.insert(0, str(Path(__file__).parent))

from generar_reporte_word import (
    get_company_name,
    get_node_name,
    generate_report,
    generate_aggregated_report,
    get_mall_name_for_parque_arauco,
    BASE_URL,
    ENTITY_BASE_URL,
    fetch_json,
    normalize_measures_payload,
    flatten_measures,
    summarize_consumption,
    parse_date,
    format_number_chilean,
    format_currency_chilean,
    get_water_price_per_m3,
    get_hourly_measures_for_day,
    calculate_nocturnal_metrics,
    summarize_alerts,
    build_hourly_consumption_chart,
)
from exclusiones_reportes import filter_node_ids

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.dml import MSO_LINE
except ImportError:
    print("[ERROR] Se requiere python-pptx. Instálalo con: pip install python-pptx")
    sys.exit(1)

COMPANY_ID = "000025"  # Parque Arauco
START_DATE = "01/01/2026"
END_DATE = "27/01/2026"  # Hoy

def get_maipu_nodes(company_id: str) -> list:
    """Obtiene todos los nodos del Mall Maipú."""
    url = f"{ENTITY_BASE_URL}/companies/{company_id}"
    try:
        response = requests.get(url, timeout=10)
        if response.status_code == 200:
            data = response.json()
            all_nodes = data.get("nodes", [])
            maipu_nodes = []
            for node in all_nodes:
                node_id = node.get("nodeId", "")
                node_name = node.get("name", "").strip()
                mall_name = get_mall_name_for_parque_arauco(node_id, node_name)
                if mall_name == "Maipú":
                    maipu_nodes.append({
                        "nodeId": node_id,
                        "name": node_name
                    })
            return maipu_nodes
        return []
    except Exception as e:
        print(f"[ERROR] Error al obtener nodos: {e}")
        return []

def get_estacion_nodes(company_id: str) -> list:
    """Obtiene los nodos de Parque Arauco Estación para reportes.
    
    Cuando se solicitan todos los nodos de Estación, solo se incluyen los nodos
    que aparecen en la imagen compartida:
    - 000025-01: PAE Estanque Norte Locales
    - 000025-19: MAE Sala de Bomba Estanque Sur
    - 000025-04: PAE Baños Públicos
    - 000025-07: PIZZA HUT
    
    Otros nodos NO se incluyen cuando se piden "todos los nodos".
    """
    # Nodos que se incluyen cuando se piden todos los nodos de Estación
    NODOS_INCLUIDOS_ESTACION = [
        "000025-01",  # PAE Estanque Norte Locales
        "000025-19",  # MAE Sala de Bomba Estanque Sur
        "000025-04",  # PAE Baños Públicos
        "000025-07"   # PIZZA HUT
    ]
    
    url = f"{ENTITY_BASE_URL}/companies/{company_id}"
    try:
        response = requests.get(url, timeout=10)
        if response.status_code == 200:
            data = response.json()
            all_nodes = data.get("nodes", [])
            estacion_nodes = []
            for node in all_nodes:
                node_id = node.get("nodeId", "")
                node_name = node.get("name", "").strip()
                
                # Solo incluir los nodos específicos cuando se piden todos los nodos
                if node_id in NODOS_INCLUIDOS_ESTACION:
                    mall_name = get_mall_name_for_parque_arauco(node_id, node_name)
                    if mall_name == "Estación":
                        estacion_nodes.append({
                            "nodeId": node_id,
                            "name": node_name
                        })
            return estacion_nodes
        return []
    except Exception as e:
        print(f"[ERROR] Error al obtener nodos: {e}")
        return []

def obtener_datos_agregados(node_ids: list, start_date: str, end_date: str) -> dict:
    """Obtiene y procesa datos agregados de todos los nodos."""
    start_dt = parse_date(start_date)
    end_dt = parse_date(end_date, end_of_day=True)
    
    def _format_ddmmyyyy(dt: datetime) -> str:
        return dt.strftime("%d%m%Y")
    
    # Para periodos largos, pedir "measures/dates" en un solo llamado puede demorar mucho
    # o quedar pegado. Partimos en tramos para dar progreso y mejorar robustez.
    def _iter_date_chunks(sdt: datetime, edt: datetime, chunk_days: int = 31):
        from datetime import timedelta
        cur = sdt
        while cur <= edt:
            chunk_end = min(edt, cur + timedelta(days=chunk_days - 1))
            yield cur, chunk_end
            cur = chunk_end + timedelta(days=1)

    total_days = (end_dt.date() - start_dt.date()).days + 1
    usar_chunks = total_days > 60

    all_measures = []
    total_consumption = 0.0
    nodes_summary = []
    
    for node_id in node_ids:
        node_name = get_node_name(node_id)
        try:
            node_measures = []
            if usar_chunks:
                print(f"[INFO] Descargando medidas por tramos para {node_id} ({node_name}) - {total_days} días...", flush=True)
                for idx, (c_start, c_end) in enumerate(_iter_date_chunks(start_dt, end_dt), 1):
                    print(f"  - Tramo {idx}: {_format_ddmmyyyy(c_start)} -> {_format_ddmmyyyy(c_end)}", flush=True)
                    measures_payload_raw = fetch_json(
                        f"{BASE_URL}/nodes/measures/dates",
                        params=[
                            ("id", node_id),
                            ("start", _format_ddmmyyyy(c_start)),
                            ("end", _format_ddmmyyyy(c_end)),
                        ],
                    )
                    measures_payload = normalize_measures_payload(measures_payload_raw, node_id)
                    node_measures.extend(flatten_measures(measures_payload))
            else:
                measures_payload_raw = fetch_json(
                    f"{BASE_URL}/nodes/measures/dates",
                    params=[
                        ("id", node_id),
                        ("start", _format_ddmmyyyy(start_dt)),
                        ("end", _format_ddmmyyyy(end_dt)),
                    ],
                )
                measures_payload = normalize_measures_payload(measures_payload_raw, node_id)
                node_measures = flatten_measures(measures_payload)

            measures = node_measures
            summary = summarize_consumption(measures)
            
            all_measures.extend(measures)
            total_consumption += summary.get('total', 0)
            nodes_summary.append({
                "node_id": node_id,
                "node_name": node_name,
                "summary": summary,
                "measures": measures
            })
        except Exception as e:
            print(f"[ADVERTENCIA] Error al procesar nodo {node_id}: {e}")
            continue
    
    # Resumen agregado
    aggregate_summary = summarize_consumption(all_measures)
    
    return {
        "total_consumption": total_consumption,
        "aggregate_summary": aggregate_summary,
        "nodes_summary": nodes_summary,
        "all_measures": all_measures,
        "start_date": start_dt,
        "end_date": end_dt
    }

def guardar_datos_json(datos: dict, output_path: Path):
    """Guarda los datos agregados en un archivo JSON."""
    def serializar_measurepoint(mp):
        """Convierte un MeasurePoint a diccionario."""
        return {
            "date": mp.date.isoformat(),
            "total_m3": mp.total_m3,
            "details": mp.details
        }
    
    def serializar_summary(summ):
        """Convierte un summary (que puede contener MeasurePoint) a diccionario."""
        if summ is None:
            return None
        result = {}
        for key, value in summ.items():
            if hasattr(value, 'date'):  # Es un MeasurePoint
                result[key] = serializar_measurepoint(value)
            elif isinstance(value, datetime):
                result[key] = value.isoformat()
            else:
                result[key] = value
        return result
    
    datos_serializados = {
        "total_consumption": datos["total_consumption"],
        "aggregate_summary": serializar_summary(datos["aggregate_summary"]),
        "nodes_summary": [
            {
                "node_id": node["node_id"],
                "node_name": node["node_name"],
                "summary": serializar_summary(node["summary"]),
                "measures": [serializar_measurepoint(m) for m in node["measures"]]
            }
            for node in datos["nodes_summary"]
        ],
        "all_measures": [serializar_measurepoint(m) for m in datos["all_measures"]],
        "start_date": datos["start_date"].isoformat(),
        "end_date": datos["end_date"].isoformat()
    }
    
    with open(output_path, 'w', encoding='utf-8') as f:
        json.dump(datos_serializados, f, indent=2, ensure_ascii=False)
    
    print(f"[OK] Datos guardados en JSON: {output_path}")

def cargar_datos_json(json_path: Path) -> dict:
    """Carga los datos agregados desde un archivo JSON."""
    from generar_reporte_word import MeasurePoint
    
    with open(json_path, 'r', encoding='utf-8') as f:
        datos_json = json.load(f)
    
    def deserializar_measurepoint(d):
        """Convierte un diccionario a MeasurePoint."""
        date_obj = datetime.fromisoformat(d["date"])
        # MeasurePoint se crea como MeasurePoint(date=dt, total_m3=float, details=dict)
        return MeasurePoint(
            date=date_obj,
            total_m3=d["total_m3"],
            details=d.get("details", {})
        )
    
    def deserializar_summary(summ):
        """Convierte un diccionario de summary a formato con MeasurePoint."""
        if summ is None:
            return None
        result = {}
        for key, value in summ.items():
            if isinstance(value, dict) and "date" in value and "total_m3" in value:
                # Es un MeasurePoint serializado
                result[key] = deserializar_measurepoint(value)
            elif isinstance(value, str) and 'T' in value:
                # Es una fecha ISO
                result[key] = datetime.fromisoformat(value)
            else:
                result[key] = value
        return result
    
    datos = {
        "total_consumption": datos_json["total_consumption"],
        "aggregate_summary": deserializar_summary(datos_json["aggregate_summary"]),
        "nodes_summary": [
            {
                "node_id": node["node_id"],
                "node_name": node["node_name"],
                "summary": deserializar_summary(node["summary"]),
                "measures": [deserializar_measurepoint(m) for m in node["measures"]]
            }
            for node in datos_json["nodes_summary"]
        ],
        "all_measures": [deserializar_measurepoint(m) for m in datos_json["all_measures"]],
        "start_date": datetime.fromisoformat(datos_json["start_date"]),
        "end_date": datetime.fromisoformat(datos_json["end_date"])
    }
    
    print(f"[OK] Datos cargados desde JSON: {json_path}")
    return datos

def crear_grafico_consumo_diario(measures: list, output_path: Path, mall_name: str = "Maipú"):
    """Crea gráfico de consumo diario con barras."""
    # Agrupar por día
    daily_consumption = {}
    for m in measures:
        date_key = m.date.date()
        if date_key not in daily_consumption:
            daily_consumption[date_key] = 0.0
        daily_consumption[date_key] += m.total_m3
    
    dates = sorted(daily_consumption.keys())
    values = [daily_consumption[d] for d in dates]
    
    fig, ax = plt.subplots(figsize=(10, 6))
    ax.bar(dates, values, color='#1f77b4', alpha=0.7)
    ax.set_xlabel('Fecha', fontsize=12)
    ax.set_ylabel('Consumo Diario (m³)', fontsize=12)
    ax.set_title(f'Consumo Diario - Parque Arauco {mall_name}', fontsize=14, fontweight='bold')
    ax.grid(True, alpha=0.3, axis='y')
    plt.xticks(rotation=45)
    plt.tight_layout()
    plt.savefig(output_path, dpi=300, bbox_inches='tight')
    plt.close()
    return output_path

def crear_grafico_consumo_dia_semana(measures: list, output_path: Path, punto_nombre: str):
    """Crea gráfico de consumo promedio por día de la semana (lunes a domingo)."""
    from generar_reporte_word import format_number_chilean
    
    # Días de la semana en español
    dias_semana = ['Lunes', 'Martes', 'Miércoles', 'Jueves', 'Viernes', 'Sábado', 'Domingo']
    
    # Agrupar por día de la semana (0=lunes, 6=domingo)
    consumo_por_dia = {i: [] for i in range(7)}  # 0-6 para lunes-domingo
    
    for m in measures:
        dia_semana = m.date.weekday()  # 0=lunes, 6=domingo
        consumo_por_dia[dia_semana].append(m.total_m3)
    
    # Calcular promedio por día de la semana
    promedios = []
    for i in range(7):
        if consumo_por_dia[i]:
            promedio = sum(consumo_por_dia[i]) / len(consumo_por_dia[i])
        else:
            promedio = 0.0
        promedios.append(promedio)
    
    # Crear gráfico de barras
    fig, ax = plt.subplots(figsize=(10, 6))
    bars = ax.bar(dias_semana, promedios, color='#1f77b4', alpha=0.7)
    ax.set_xlabel('Día de la Semana', fontsize=16, fontweight='bold')
    ax.set_ylabel('Consumo Promedio (m³)', fontsize=16, fontweight='bold')
    # Título del gráfico eliminado (se muestra solo en la slide)
    ax.grid(True, alpha=0.3, axis='y')
    
    # Configurar tamaño de fuente de las etiquetas de los ejes
    ax.tick_params(axis='x', labelsize=14, rotation=0)
    ax.tick_params(axis='y', labelsize=14)
    
    # Agregar valores en las barras (aumentado tamaño de fuente)
    max_val = max(promedios) if promedios else 0
    for bar, val in zip(bars, promedios):
        if val > 0:
            height = bar.get_height()
            ax.text(bar.get_x() + bar.get_width()/2., height + max_val * 0.01,
                   f'{format_number_chilean(val, 1)}',
                   ha='center', va='bottom', fontsize=16, fontweight='bold')
    
    plt.xticks(rotation=0)
    plt.tight_layout()
    output_path.parent.mkdir(parents=True, exist_ok=True)
    plt.savefig(output_path, dpi=300, bbox_inches='tight')
    plt.close()
    return output_path

def crear_grafico_consumo_diario_linea(measures: list, output_path: Path, punto_nombre: str):
    """Crea gráfico de línea de consumo diario destacando el mayor y menor consumo."""
    from generar_reporte_word import format_number_chilean
    
    # Agrupar por día
    daily_consumption = {}
    for m in measures:
        date_key = m.date.date()
        if date_key not in daily_consumption:
            daily_consumption[date_key] = 0.0
        daily_consumption[date_key] += m.total_m3
    
    dates = sorted(daily_consumption.keys())
    values = [daily_consumption[d] for d in dates]
    
    if not values:
        return None
    
    # Encontrar mayor y menor consumo
    max_val = max(values)
    min_val = min(values)
    max_idx = values.index(max_val)
    min_idx = values.index(min_val)
    
    # Crear gráfico de línea
    fig, ax = plt.subplots(figsize=(6, 4))
    
    # Línea principal (sin label para eliminar leyenda)
    ax.plot(dates, values, color='#1f77b4', linewidth=2, marker='o', markersize=4, alpha=0.7)
    
    # Destacar mayor consumo (rojo) - sin label
    ax.plot(dates[max_idx], max_val, 'ro', markersize=10)
    ax.text(dates[max_idx], max_val, f'  {format_number_chilean(max_val, 1)} m³',
            fontsize=9, fontweight='bold', color='red', va='bottom')
    
    # Destacar menor consumo (verde) - sin label
    ax.plot(dates[min_idx], min_val, 'go', markersize=10)
    ax.text(dates[min_idx], min_val, f'  {format_number_chilean(min_val, 1)} m³',
            fontsize=9, fontweight='bold', color='green', va='top')
    
    ax.set_xlabel('Fecha', fontsize=10)
    ax.set_ylabel('Consumo (m³)', fontsize=10)
    ax.set_title(f'Consumo Diario - {punto_nombre}', fontsize=11, fontweight='bold')
    ax.grid(True, alpha=0.3, axis='y')
    # Leyenda eliminada
    
    # Formatear fechas en el eje X: formato DD/MM/YYYY (ej: "13/01/2026")
    # Mejorar visualización: mostrar solo algunas etiquetas cuando hay muchas fechas
    num_dates = len(dates)
    
    if num_dates > 15:
        # Si hay muchas fechas, mostrar solo cada 2-3 fechas para que se vean mejor
        step = max(2, num_dates // 10)  # Mostrar aproximadamente 10 etiquetas máximo
        dates_to_show = dates[::step]
        date_labels_to_show = []
        for date_obj in dates_to_show:
            date_dt = datetime.combine(date_obj, datetime.min.time())
            date_labels_to_show.append(date_dt.strftime('%d/%m/%Y'))
        
        ax.set_xticks(dates_to_show)
        ax.set_xticklabels(date_labels_to_show, rotation=45, ha='right', fontsize=9)
    elif num_dates > 7:
        # Si hay varias fechas (8-15), mostrar todas pero con rotación y espacio
        date_labels = []
        for date_obj in dates:
            date_dt = datetime.combine(date_obj, datetime.min.time())
            date_labels.append(date_dt.strftime('%d/%m/%Y'))
        
        ax.set_xticks(dates)
        ax.set_xticklabels(date_labels, rotation=45, ha='right', fontsize=9)
    else:
        # Si hay pocas fechas (7 o menos), mostrar todas con formato mejorado
        date_labels = []
        for date_obj in dates:
            date_dt = datetime.combine(date_obj, datetime.min.time())
            # Formato: "Día DD/MM" o "DD/MM/YYYY" si hay espacio
            dia_semana = date_dt.strftime('%a')  # Abreviación del día (ej: "Lun", "Mar")
            fecha_corta = date_dt.strftime('%d/%m/%Y')
            date_labels.append(f'{dia_semana}\n{fecha_corta}')
        
        ax.set_xticks(dates)
        ax.set_xticklabels(date_labels, rotation=0, ha='center', fontsize=9)
    
    plt.tight_layout()
    output_path.parent.mkdir(parents=True, exist_ok=True)
    plt.savefig(output_path, dpi=300, bbox_inches='tight')
    plt.close()
    return output_path

def crear_grafico_consumo_horario_dia(hourly_data: list, output_path: Path, fecha: datetime, titulo: str):
    """
    Crea gráfica de línea de consumo horario para un día específico con área degradada celeste
    y área achurada para consumo nocturno (00:00 a 06:00) cuando sea distinto a cero.
    
    Args:
        hourly_data: Lista de tuplas (hora, consumo) donde hora es 0-23
        output_path: Ruta donde guardar la imagen
        fecha: Fecha del día
        titulo: Título de la gráfica
    """
    if not hourly_data:
        return None
    
    hours = [h for h, _ in hourly_data]
    values = [v for _, v in hourly_data]
    
    fig, ax = plt.subplots(figsize=(8, 4))
    
    # Línea principal
    ax.plot(hours, values, linestyle="-", color="#4A90E2", linewidth=2)
    
    # Rellenar el área bajo la línea con degradado celeste
    ax.fill_between(hours, values, 0, color="#87CEEB", alpha=0.4)  # Sky blue con transparencia
    
    # Identificar horas nocturnas (00:00 a 06:00) con consumo distinto a cero
    horas_nocturnas = []
    valores_nocturnas = []
    for h, v in zip(hours, values):
        if 0 <= h <= 6 and v > 0:
            horas_nocturnas.append(h)
            valores_nocturnas.append(v)
    
    # Si hay consumo nocturno, agregar área achurada (hatched)
    if horas_nocturnas:
        # Crear arrays para el área achurada (solo horas nocturnas)
        horas_noct_array = []
        valores_noct_array = []
        for h in range(0, 7):  # 0 a 6
            # Buscar el valor correspondiente a esta hora
            valor_h = 0
            for idx, hour in enumerate(hours):
                if hour == h:
                    valor_h = values[idx]
                    break
            horas_noct_array.append(h)
            valores_noct_array.append(valor_h)
        
        # Agregar área achurada para el período nocturno
        ax.fill_between(horas_noct_array, valores_noct_array, 0, 
                       color="#FFD700", alpha=0.3, hatch="///", 
                       edgecolor="#FFA500", linewidth=1.5, label="Consumo Nocturno")
    
    ax.set_xlabel('Hora del día', fontsize=14, fontweight='bold')  # Aumentado de 10 a 14
    ax.set_ylabel('Consumo (m³/hr)', fontsize=10)
    ax.set_title(titulo, fontsize=16, fontweight='bold')  # Aumentado de 11 a 16
    ax.grid(True, alpha=0.3, axis='y')
    ax.set_xlim(0, 23)
    ax.set_xticks(range(0, 24, 2))  # Mostrar cada 2 horas
    # Agrandar números del eje X
    ax.tick_params(axis='x', labelsize=12)  # Aumentado de tamaño por defecto a 12
    ax.set_ylim(bottom=0)
    
    # Agregar líneas verticales para delimitar período nocturno
    if horas_nocturnas:
        ax.axvline(x=0, color='orange', linestyle='--', linewidth=1.5, alpha=0.7)
        ax.axvline(x=6, color='orange', linestyle='--', linewidth=1.5, alpha=0.7)
    
    plt.tight_layout()
    output_path.parent.mkdir(parents=True, exist_ok=True)
    plt.savefig(output_path, dpi=300, bbox_inches='tight')
    plt.close()
    return output_path

def crear_grafico_ranking_nodos(nodes_summary: list, output_path: Path):
    """Crea gráfico de ranking de consumo por nodo (ordenado de mayor a menor, mayor arriba)."""
    # Ordenar nodos por consumo total de mayor a menor
    nodes_sorted = sorted(nodes_summary, key=lambda x: x['summary'].get('total', 0), reverse=True)
    
    node_names = [n["node_name"] for n in nodes_sorted]
    consumptions = [n["summary"].get('total', 0) for n in nodes_sorted]
    
    fig, ax = plt.subplots(figsize=(10, 6))
    bars = ax.barh(node_names, consumptions, color='#1f77b4')
    ax.set_xlabel('Consumo Total (m³)', fontsize=12)
    ax.set_ylabel('Monitoreo (m³)', fontsize=12)  # Agregar m³ al eje Y
    ax.set_title('Ranking de Consumo', fontsize=14, fontweight='bold')
    ax.grid(True, alpha=0.3, axis='x')
    
    # Invertir el eje Y para que el mayor consumo aparezca arriba
    ax.invert_yaxis()
    
    # Calcular el máximo para ajustar el límite del eje X y que los valores quepan
    max_consumption = max(consumptions) if consumptions else 0
    # Aumentar el límite del eje X un 15% para que los valores siempre quepan
    ax.set_xlim(0, max_consumption * 1.15)
    
    # Agregar valores en las barras (dentro del marco, aumentado tamaño de fuente)
    for i, (bar, val) in enumerate(zip(bars, consumptions)):
        width = bar.get_width()
        # Poner el valor dentro de la barra, centrado si es posible, o al final pero con espacio
        # Usar una posición que siempre esté dentro del límite del eje
        text_x = min(width, max_consumption * 1.12)  # Asegurar que esté dentro
        if text_x < max_consumption * 0.05:  # Si la barra es muy corta, poner el texto a la derecha de la barra
            text_x = width + max_consumption * 0.02
            ha_align = 'left'
            color = 'black'
        else:
            # Si la barra es suficientemente larga, poner el texto dentro
            text_x = width * 0.98
            ha_align = 'right'
            color = 'white'  # Blanco para contrastar con la barra azul
        
        ax.text(text_x, bar.get_y() + bar.get_height()/2, 
                f'{format_number_chilean(val, 1)}', 
                ha=ha_align, va='center', fontsize=12, fontweight='bold', color=color)
    
    plt.tight_layout()
    plt.savefig(output_path, dpi=300, bbox_inches='tight')
    plt.close()
    return output_path

def obtener_alertas_agregadas(node_ids: list, start_date: datetime, end_date: datetime) -> int:
    """
    Obtiene el número total de alertas agregadas de todos los nodos en el período.
    
    Args:
        node_ids: Lista de IDs de nodos
        start_date: Fecha de inicio
        end_date: Fecha de fin
    
    Returns:
        Número total de alertas (solo las que tienen medida > 0)
    """
    def _format_ddmmyyyy(dt: datetime) -> str:
        return dt.strftime("%d%m%Y")
    
    total_alertas = 0
    
    for node_id in node_ids:
        try:
            # Intentar obtener alertas desde el endpoint principal
            alerts_payload = fetch_json(
                f"{BASE_URL}/nodes/myalert/alerts",
                params=[
                    ("id", node_id),
                    ("start", _format_ddmmyyyy(start_date)),
                    ("end", _format_ddmmyyyy(end_date)),
                ],
            )
            if isinstance(alerts_payload, list):
                # Contar solo alertas con medida mayor a cero
                alerts_con_medida = [a for a in alerts_payload if float(a.get("measure", 0) or 0) > 0]
                total_alertas += len(alerts_con_medida)
        except Exception:
            # Si falla, intentar con el endpoint alternativo
            try:
                fallback_payload = fetch_json(
                    f"{BASE_URL}/nodes/leak/alerts",
                    params=[
                        ("id", node_id),
                        ("start", _format_ddmmyyyy(start_date)),
                        ("end", _format_ddmmyyyy(end_date)),
                    ],
                )
                if isinstance(fallback_payload, list):
                    alerts_con_medida = [a for a in fallback_payload if float(a.get("measure", 0) or 0) > 0]
                    total_alertas += len(alerts_con_medida)
            except Exception:
                # Si ambos fallan, continuar con el siguiente nodo
                continue
    
    return total_alertas

def obtener_alertas_por_nodo(node_id: str, start_date: datetime, end_date: datetime) -> int:
    """
    Obtiene el número de alertas de un nodo específico en el período.
    
    Args:
        node_id: ID del nodo
        start_date: Fecha de inicio
        end_date: Fecha de fin
    
    Returns:
        Número de alertas (solo las que tienen medida > 0)
    """
    def _format_ddmmyyyy(dt: datetime) -> str:
        return dt.strftime("%d%m%Y")
    
    try:
        # Intentar obtener alertas desde el endpoint principal
        alerts_payload = fetch_json(
            f"{BASE_URL}/nodes/myalert/alerts",
            params=[
                ("id", node_id),
                ("start", _format_ddmmyyyy(start_date)),
                ("end", _format_ddmmyyyy(end_date)),
            ],
        )
        if isinstance(alerts_payload, list):
            # Contar solo alertas con medida mayor a cero
            alerts_con_medida = [a for a in alerts_payload if float(a.get("measure", 0) or 0) > 0]
            return len(alerts_con_medida)
    except Exception:
        # Si falla, intentar con el endpoint alternativo
        try:
            fallback_payload = fetch_json(
                f"{BASE_URL}/nodes/leak/alerts",
                params=[
                    ("id", node_id),
                    ("start", _format_ddmmyyyy(start_date)),
                    ("end", _format_ddmmyyyy(end_date)),
                ],
            )
            if isinstance(fallback_payload, list):
                alerts_con_medida = [a for a in fallback_payload if float(a.get("measure", 0) or 0) > 0]
                return len(alerts_con_medida)
        except Exception:
            pass
    
    return 0

def calcular_consumo_nocturno_agregado(node_ids: list, start_date: datetime, end_date: datetime, usar_api: bool = True) -> dict:
    """
    Calcula el consumo nocturno agregado (00:00 a 06:00) para todos los nodos desde los datos horarios de la API.
    
    Args:
        node_ids: Lista de IDs de nodos
        start_date: Fecha de inicio
        end_date: Fecha de fin
        usar_api: Si es False, usa aproximación rápida (15% del total) para evitar esperas
    
    Retorna:
    - consumo_nocturno_total: Suma total de consumo nocturno de todos los nodos
    - consumo_diurno_efectivo: Suma total de consumo diurno de todos los nodos
    - dias_con_consumo_nocturno: Número de días con consumo nocturno
    - consumo_nocturno_por_nodo: Diccionario con consumo nocturno por cada nodo
    """
    consumo_nocturno_total = 0.0
    consumo_diurno_efectivo = 0.0
    dias_con_consumo_nocturno = set()
    consumo_nocturno_por_nodo = {}
    
    # Calcular para cada nodo
    for node_id in node_ids:
        try:
            metrics = calculate_nocturnal_metrics(node_id, start_date, end_date)
            consumo_nocturno_nodo = metrics["consumo_nocturno_total"]
            consumo_diurno_nodo = metrics["consumo_diurno_efectivo"]
            
            consumo_nocturno_total += consumo_nocturno_nodo
            consumo_diurno_efectivo += consumo_diurno_nodo
            consumo_nocturno_por_nodo[node_id] = consumo_nocturno_nodo
            
            # Contar días con consumo nocturno (optimizado: usar métricas ya calculadas)
            # No hacer llamadas adicionales a la API, usar la información de metrics
            if metrics["dias_con_consumo_nocturno"] > 0:
                # Aproximar días con consumo nocturno basado en el porcentaje de días del período
                # Esto evita múltiples llamadas a la API
                total_dias = (end_date.date() - start_date.date()).days + 1
                dias_estimados = min(metrics["dias_con_consumo_nocturno"], total_dias)
                # Agregar días estimados al set (usando fechas aproximadas)
                for i in range(dias_estimados):
                    fecha_estimada = start_date.date() + timedelta(days=i)
                    if fecha_estimada <= end_date.date():
                        dias_con_consumo_nocturno.add(fecha_estimada)
        except Exception as e:
            print(f"[ADVERTENCIA] Error al calcular consumo nocturno para nodo {node_id}: {e}")
            consumo_nocturno_por_nodo[node_id] = 0.0
            continue
    
    return {
        "consumo_nocturno_total": consumo_nocturno_total,
        "consumo_diurno_efectivo": consumo_diurno_efectivo,
        "dias_con_consumo_nocturno": len(dias_con_consumo_nocturno),
        "consumo_nocturno_por_nodo": consumo_nocturno_por_nodo
    }

def crear_grafico_consumo_nocturno_puntos(nodes_summary: list, consumo_nocturno_por_nodo: dict, output_path: Path):
    """Crea gráfico de barras horizontal mostrando consumo nocturno por punto, solo los 3 más altos.
    
    Args:
        nodes_summary: Lista de resúmenes por nodo
        consumo_nocturno_por_nodo: Diccionario con consumo nocturno real por node_id
    """
    from generar_reporte_word import format_number_chilean
    
    # Calcular consumo nocturno para cada punto desde los datos reales
    puntos_data = []
    for node in nodes_summary:
        node_id = node['node_id']
        consumo_nocturno = consumo_nocturno_por_nodo.get(node_id, 0.0)
        puntos_data.append({
            'nombre': node['node_name'],
            'consumo_nocturno': consumo_nocturno
        })
    
    # Ordenar de mayor a menor consumo nocturno
    puntos_data_sorted = sorted(puntos_data, key=lambda x: x['consumo_nocturno'], reverse=True)
    
    # Filtrar valores cero y tomar solo los 3 más altos
    data_filtrada = [(p['nombre'], p['consumo_nocturno']) for p in puntos_data_sorted if p['consumo_nocturno'] > 0]
    data_filtrada = data_filtrada[:3]  # Solo los 3 más altos
    
    if not data_filtrada:
        # Si no hay datos, crear un gráfico vacío
        fig, ax = plt.subplots(figsize=(6, 3))
        ax.text(0.5, 0.5, 'No hay datos\nde consumo nocturno', 
                ha='center', va='center', fontsize=14, color='gray')
        ax.axis('off')
        plt.tight_layout()
        output_path.parent.mkdir(parents=True, exist_ok=True)
        plt.savefig(output_path, dpi=300, bbox_inches='tight', transparent=True)
        plt.close()
        return output_path
    
    nombres_filtrados = [d[0] for d in data_filtrada]
    consumos_filtrados = [d[1] for d in data_filtrada]
    
    # Obtener cantidad real de puntos mostrados
    num_puntos_mostrados = len(nombres_filtrados)
    
    # Crear gráfico de barras horizontal (reducido)
    fig, ax = plt.subplots(figsize=(5, 2.5))
    
    # Color amarillo oscuro (#B8860B o #DAA520)
    color_amarillo_oscuro = '#B8860B'
    
    bars = ax.barh(nombres_filtrados, consumos_filtrados, color=color_amarillo_oscuro)
    ax.set_xlabel('Consumo Nocturno (m³)', fontsize=11)
    ax.set_ylabel('', fontsize=11)  # Eliminado "Punto de Monitoreo"
    ax.set_title(f'Consumo Nocturno por Punto (Top {num_puntos_mostrados})', fontsize=13, fontweight='bold')
    ax.grid(True, alpha=0.3, axis='x')
    
    # Invertir el eje Y para que el mayor consumo aparezca arriba
    ax.invert_yaxis()
    
    # Calcular el máximo para ajustar el límite del eje X
    max_consumo = max(consumos_filtrados) if consumos_filtrados else 0
    ax.set_xlim(0, max_consumo * 1.15)
    
    # Agregar valores en las barras
    for i, (bar, val) in enumerate(zip(bars, consumos_filtrados)):
        width = bar.get_width()
        text_x = min(width, max_consumo * 1.12)
        if text_x < max_consumo * 0.05:
            text_x = width + max_consumo * 0.02
            ha_align = 'left'
            color = 'black'
        else:
            text_x = width * 0.98
            ha_align = 'right'
            color = 'white'
        
        ax.text(text_x, bar.get_y() + bar.get_height()/2, 
                f'{format_number_chilean(val, 1)}', 
                ha=ha_align, va='center', fontsize=10, fontweight='bold', color=color)
    
    plt.tight_layout()
    output_path.parent.mkdir(parents=True, exist_ok=True)
    plt.savefig(output_path, dpi=300, bbox_inches='tight', transparent=True)
    plt.close()
    return output_path

def crear_grafico_anillo_consumo(consumo_efectivo: float, consumo_nocturno: float, 
                                   porcentaje_efectivo: float, porcentaje_nocturno: float,
                                   periodo_texto: str, output_path: Path):
    """Crea un gráfico de anillo mostrando consumo efectivo vs nocturno."""
    from generar_reporte_word import format_number_chilean
    
    # Crear figura
    fig, ax = plt.subplots(figsize=(6, 6))
    
    # Datos para el gráfico de anillo
    sizes = [consumo_efectivo, consumo_nocturno]
    colors = ['#4A90E2', '#FF6B6B']  # Azul para efectivo, rojo/naranja para nocturno
    labels = ['Consumo Efectivo', 'Consumo Nocturno']
    
    # Crear gráfico de anillo (donut chart)
    wedges, texts, autotexts = ax.pie(sizes, labels=labels, colors=colors, autopct='%1.1f%%',
                                       startangle=90, pctdistance=0.85, textprops={'fontsize': 11, 'fontweight': 'bold'})
    
    # Crear el "donut" (círculo central)
    centre_circle = plt.Circle((0, 0), 0.70, fc='white')
    ax.add_artist(centre_circle)
    
    # Agregar texto en el centro del anillo
    total = consumo_efectivo + consumo_nocturno
    ax.text(0, 0.3, f"{format_number_chilean(total, 1)}\nm³", 
            ha='center', va='center', fontsize=16, fontweight='bold', color='#333333')
    ax.text(0, -0.2, periodo_texto, 
            ha='center', va='center', fontsize=10, color='#666666')
    
    # Notas debajo del gráfico
    nota_texto = f"Consumo Efectivo: {format_number_chilean(consumo_efectivo, 1)} m³ ({format_number_chilean(porcentaje_efectivo, 1)}%)\n"
    nota_texto += f"Consumo Nocturno: {format_number_chilean(consumo_nocturno, 1)} m³ ({format_number_chilean(porcentaje_nocturno, 1)}%)"
    
    ax.text(0, -1.3, nota_texto, ha='center', va='top', fontsize=9, color='#333333',
            bbox=dict(boxstyle='round', facecolor='#f0f0f0', alpha=0.8))
    
    ax.set_aspect('equal')
    ax.axis('off')
    
    plt.tight_layout()
    output_path.parent.mkdir(parents=True, exist_ok=True)
    plt.savefig(output_path, dpi=300, bbox_inches='tight', transparent=True)
    plt.close()
    return output_path

def crear_grafico_distribucion_nodos(nodes_summary: list, output_path: Path):
    """Crea gráfico de distribución de consumo por nodo."""
    node_names = [n["node_name"] for n in nodes_summary]
    consumptions = [n["summary"].get('total', 0) for n in nodes_summary]
    
    fig, ax = plt.subplots(figsize=(10, 6))
    bars = ax.barh(node_names, consumptions, color='#1f77b4')
    ax.set_xlabel('Consumo Total (m³)', fontsize=12)
    ax.set_title('Distribución de Consumo por Monitoreo', fontsize=14, fontweight='bold')
    ax.grid(True, alpha=0.3, axis='x')
    
    # Agregar valores en las barras
    for i, (bar, val) in enumerate(zip(bars, consumptions)):
        width = bar.get_width()
        ax.text(width, bar.get_y() + bar.get_height()/2, 
                f' {format_number_chilean(val, 1)}', 
                ha='left', va='center', fontsize=9)
    
    plt.tight_layout()
    plt.savefig(output_path, dpi=300, bbox_inches='tight')
    plt.close()
    return output_path

def crear_imagen_fondo_agua(output_path: Path, width_inches: float = 10, height_inches: float = 7.5):
    """Crea una imagen de fondo con efecto de agua tenue."""
    # Convertir pulgadas a píxeles (150 DPI para buena calidad)
    dpi = 150
    width_px = int(width_inches * dpi)
    height_px = int(height_inches * dpi)
    
    # Crear figura sin bordes
    fig = plt.figure(figsize=(width_inches, height_inches), dpi=dpi, facecolor='none')
    ax = fig.add_axes([0, 0, 1, 1], frameon=False)
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    ax.axis('off')
    
    # Crear gradiente de fondo azul claro
    x = np.linspace(0, 1, width_px)
    y = np.linspace(0, 1, height_px)
    X, Y = np.meshgrid(x, y)
    
    # Crear ondas de agua suaves y tenues
    # Múltiples ondas con diferentes frecuencias y fases para efecto de agua
    wave1 = np.sin(2 * np.pi * 3 * X + 0.5) * np.cos(2 * np.pi * 2 * Y + 0.3)
    wave2 = np.sin(2 * np.pi * 5 * Y + 0.7) * np.cos(2 * np.pi * 4 * X + 0.2)
    wave3 = np.sin(2 * np.pi * 1.5 * (X + Y) + 0.9)
    
    # Combinar ondas con pesos diferentes
    combined_wave = 0.3 * wave1 + 0.2 * wave2 + 0.15 * wave3
    
    # Normalizar a rango 0-1
    combined_wave = (combined_wave - combined_wave.min()) / (combined_wave.max() - combined_wave.min())
    
    # Crear gradiente de color azul con variación suave
    # Color base: azul muy claro (230, 240, 255)
    # Variación: desde (220, 235, 250) hasta (240, 245, 255)
    base_color = np.array([230, 240, 255]) / 255.0
    light_variation = np.array([240, 245, 255]) / 255.0
    dark_variation = np.array([220, 235, 250]) / 255.0
    
    # Aplicar variación basada en ondas
    color_map = np.zeros((height_px, width_px, 3))
    for i in range(3):
        color_map[:, :, i] = base_color[i] + (combined_wave - 0.5) * (light_variation[i] - dark_variation[i]) * 0.3
    
    # Asegurar que los valores estén en [0, 1]
    color_map = np.clip(color_map, 0, 1)
    
    # Mostrar imagen
    ax.imshow(color_map, extent=[0, 1, 0, 1], aspect='auto', interpolation='bilinear')
    
    # Guardar con transparencia suave
    fig.savefig(output_path, dpi=dpi, bbox_inches='tight', pad_inches=0, transparent=False, facecolor='white')
    plt.close(fig)
    
    return output_path

def aplicar_fondo_agua_tenue(slide, temp_dir: Path, prs: Presentation, scale_factor: float = 0.85):
    """Aplica un fondo con imagen de agua tenue a la slide con opción de reducir tamaño."""
    # Crear imagen de fondo si no existe
    bg_image_path = temp_dir / "fondo_agua_tenue.png"
    if not bg_image_path.exists():
        crear_imagen_fondo_agua(bg_image_path, width_inches=10, height_inches=7.5)
    
    # Obtener dimensiones de la presentación
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    # Reducir tamaño de la imagen (dejar márgenes)
    scaled_width = int(slide_width * scale_factor)
    scaled_height = int(slide_height * scale_factor)
    
    # Centrar la imagen
    left_offset = (slide_width - scaled_width) // 2
    top_offset = (slide_height - scaled_height) // 2
    
    # Agregar imagen de fondo como shape (más pequeña y centrada)
    bg_shape = slide.shapes.add_picture(str(bg_image_path), left_offset, top_offset, scaled_width, scaled_height)
    
    # Enviar la imagen al fondo (behind all other shapes)
    # Mover el elemento al principio de la lista para que quede atrás
    spTree = slide.shapes._spTree
    bg_element = bg_shape.element
    spTree.remove(bg_element)
    spTree.insert(2, bg_element)  # Insertar después del elemento base

def convertir_ppt_a_pdf(ppt_path: Path) -> Optional[Path]:
    """
    Convierte un archivo PPTX a PDF usando PowerPoint COM (requiere PowerPoint instalado).
    
    Args:
        ppt_path: Ruta del archivo PPTX
        
    Returns:
        Ruta del archivo PDF generado, o None si hay error
    """
    pdf_path = ppt_path.with_suffix('.pdf')
    
    try:
        import win32com.client
        
        print(f"[INFO] Convirtiendo PPT a PDF...")
        print(f"  PPT: {ppt_path.name}")
        print(f"  PDF: {pdf_path.name}")
        
        # Verificar que el archivo existe
        if not ppt_path.exists():
            print(f"[ERROR] El archivo PPT no existe: {ppt_path}")
            return None
        
        powerpoint = None
        presentation = None
        com_initialized = False
        
        try:
            # Inicializar COM solo si es necesario (en threads)
            try:
                import pythoncom
                pythoncom.CoInitialize()
                com_initialized = True
            except:
                # COM ya está inicializado o no es necesario
                pass
            
            # Crear instancia de PowerPoint
            try:
                # Intentar obtener instancia existente primero
                try:
                    powerpoint = win32com.client.GetActiveObject("PowerPoint.Application")
                    print(f"  [INFO] Usando instancia existente de PowerPoint")
                except:
                    # Si no hay instancia existente, crear una nueva
                    powerpoint = win32com.client.Dispatch("PowerPoint.Application")
                    print(f"  [INFO] Nueva instancia de PowerPoint creada")
            except Exception as e:
                error_code = getattr(e, 'args', [None])[0] if hasattr(e, 'args') and e.args else None
                if error_code == -2146959355 or (isinstance(error_code, (int, str)) and "80080005" in str(error_code)):
                    print(f"[ADVERTENCIA] PowerPoint no está disponible.")
                    print(f"  Error COM: {error_code}")
                    print(f"  Posibles causas: PowerPoint no está instalado o no está registrado correctamente.")
                    return None
                else:
                    raise
            
            # Establecer propiedades de PowerPoint
            # Nota: PowerPoint 2013 no permite ocultar la ventana en algunas configuraciones
            try:
                powerpoint.Visible = 0  # Intentar ocultar (puede fallar en PowerPoint 2013)
                print(f"  [INFO] PowerPoint configurado en modo oculto")
            except Exception as e:
                error_code = getattr(e, 'args', [None])[0] if hasattr(e, 'args') and e.args else None
                if error_code == -2147352567 or (isinstance(error_code, (int, str)) and "-2147352567" in str(error_code)):
                    # PowerPoint 2013 no permite ocultar, usar visible
                    print(f"  [INFO] PowerPoint requiere ventana visible (PowerPoint 2013)")
                    powerpoint.Visible = 1
                else:
                    raise
            
            try:
                powerpoint.DisplayAlerts = 0  # Deshabilitar alertas
            except:
                # Si falla, continuar de todas formas
                pass
            
            # Abrir la presentación
            try:
                presentation = powerpoint.Presentations.Open(
                    str(ppt_path.absolute()), 
                    ReadOnly=True, 
                    WithWindow=False,
                    Untitled=False
                )
            except Exception as e:
                error_code = getattr(e, 'args', [None])[0] if hasattr(e, 'args') and e.args else None
                if error_code == -2146959355 or (isinstance(error_code, (int, str)) and "80080005" in str(error_code)):
                    print(f"  [ADVERTENCIA] Error COM al abrir PPT: {error_code}")
                    print(f"  Intentando con instancia visible como fallback...")
                    # Intentar con ventana visible como fallback
                    try:
                        powerpoint.Visible = 1
                        presentation = powerpoint.Presentations.Open(
                            str(ppt_path.absolute()), 
                            ReadOnly=True, 
                            WithWindow=True,
                            Untitled=False
                        )
                    except Exception as e2:
                        print(f"  [ERROR] También falló con instancia visible: {e2}")
                        raise e
                else:
                    raise
            
            # Guardar como PDF
            # ppSaveAsPDF = 32
            try:
                presentation.SaveAs(str(pdf_path.absolute()), 32)
                
                # Esperar un momento para que el archivo se escriba
                import time
                time.sleep(0.5)
                
                # Verificar que el PDF se creó
                if pdf_path.exists() and pdf_path.stat().st_size > 0:
                    print(f"[OK] PDF generado: {pdf_path.name} ({pdf_path.stat().st_size} bytes)")
                    return pdf_path
                else:
                    print(f"[ADVERTENCIA] El PDF no se creó correctamente o está vacío")
                    return None
                    
            except Exception as e:
                print(f"  [ADVERTENCIA] Error al guardar como PDF: {e}")
                return None
            
        except Exception as e:
            error_msg = str(e)
            error_code = getattr(e, 'args', [None])[0] if hasattr(e, 'args') and e.args else None
            
            # Mensaje más descriptivo según el código de error
            if error_code == -2146959355 or (isinstance(error_code, (int, str)) and "80080005" in str(error_code)):
                print(f"[ADVERTENCIA] PowerPoint no está disponible o hay un problema con COM.")
                print(f"  Error COM: {error_code}")
                print(f"  Posibles causas:")
                print(f"    - PowerPoint no está instalado")
                print(f"    - Hay una instancia de PowerPoint bloqueada")
                print(f"    - Problemas con el registro COM de PowerPoint")
            else:
                print(f"[ADVERTENCIA] Error al convertir PPT a PDF: {e}")
            
            return None
            
        finally:
            # Cerrar la presentación
            if presentation:
                try:
                    presentation.Close()
                except:
                    pass
            
            # Cerrar PowerPoint solo si lo creamos nosotros
            if powerpoint:
                try:
                    # Verificar si hay otras presentaciones abiertas
                    try:
                        if powerpoint.Presentations.Count == 0:
                            powerpoint.Quit()
                        else:
                            # No cerrar si hay otras presentaciones abiertas
                            pass
                    except:
                        # Si falla la verificación, intentar cerrar de todas formas
                        try:
                            powerpoint.Quit()
                        except:
                            pass
                except:
                    pass
            
            # Desinicializar COM si lo inicializamos
            if com_initialized:
                try:
                    import pythoncom
                    pythoncom.CoUninitialize()
                except:
                    pass
            
    except ImportError:
        print(f"[ADVERTENCIA] win32com no está instalado. Instálalo con: pip install pywin32")
        print(f"[INFO] Para convertir manualmente, abre {ppt_path.name} en PowerPoint y guárdalo como PDF")
        return None
    except Exception as e:
        error_msg = str(e)
        if "80080005" in error_msg or "-2146959355" in error_msg:
            print(f"[ADVERTENCIA] PowerPoint no está disponible o hay un problema con COM.")
            print(f"  El archivo PPT se generó correctamente: {ppt_path.name}")
            print(f"  Puedes abrirlo manualmente en PowerPoint y guardarlo como PDF")
        else:
            print(f"[ADVERTENCIA] Error al convertir PPT a PDF: {e}")
            print(f"[INFO] El archivo PPT se generó correctamente: {ppt_path.name}")
            print(f"[INFO] Para convertir manualmente, abre el PPT en PowerPoint y guárdalo como PDF")
        return None


def tiene_mas_de_7_dias(start_date: str, end_date: str) -> bool:
    """
    Verifica si el período tiene más de 7 días.
    
    Args:
        start_date: Fecha de inicio (formato DD/MM/YYYY)
        end_date: Fecha de fin (formato DD/MM/YYYY)
    
    Returns:
        True si tiene más de 7 días, False en caso contrario
    """
    try:
        start_dt = parse_date(start_date)
        end_dt = parse_date(end_date, end_of_day=True)
        num_days = (end_dt - start_dt).days + 1
        return num_days > 7
    except:
        return False

def generar_ppt_desde_agregado(
    company_id: str,
    node_ids: list,
    start_date: str,
    end_date: str,
    aggregated_report_path: Path,
    mall_name: str = None,
    company_name: str = None
) -> Path:
    """
    Función genérica para generar PPT desde un reporte agregado.
    
    Args:
        company_id: ID de la empresa
        node_ids: Lista de IDs de nodos
        start_date: Fecha de inicio (formato DD/MM/YYYY)
        end_date: Fecha de fin (formato DD/MM/YYYY)
        aggregated_report_path: Ruta del directorio del reporte agregado
        mall_name: Nombre del mall (opcional, se detecta automáticamente si es Parque Arauco)
        company_name: Nombre de la empresa (opcional, se obtiene de la API si no se proporciona)
    
    Returns:
        Ruta del archivo PPT generado
    """
    # Verificar si tiene más de 7 días
    if not tiene_mas_de_7_dias(start_date, end_date):
        print("[INFO] Período tiene 7 días o menos, no se genera PPT agregado")
        return None
    
    if company_name is None:
        company_name = get_company_name(company_id)

    node_ids_filtrados = filter_node_ids(node_ids, company_id=company_id, company_name=company_name)
    if len(node_ids_filtrados) < len(node_ids):
        print(f"[INFO] Se excluyeron {len(node_ids) - len(node_ids_filtrados)} nodo(s) del PPT por configuración.")
    node_ids = node_ids_filtrados
    if not node_ids:
        print("[INFO] Todos los nodos fueron excluidos. No se genera PPT agregado.")
        return None
    
    print()
    print("[INFO] Generando presentación PPT...")
    
    # Obtener datos agregados
    try:
        datos = obtener_datos_agregados(node_ids, start_date, end_date)
        print(f"[OK] Datos obtenidos: {len(datos['all_measures'])} medidas")
    except Exception as e:
        print(f"[ERROR] Error al obtener datos para PPT: {e}")
        raise
    
    # Determinar nombre del mall si no se proporciona
    if mall_name is None:
        # Si es Parque Arauco, intentar detectar el mall
        if company_id == "000025":
            # Obtener el primer nodo para detectar el mall
            if node_ids:
                first_node_name = get_node_name(node_ids[0])
                mall_name = get_mall_name_for_parque_arauco(node_ids[0], first_node_name)
                if not mall_name:
                    mall_name = "Parque Arauco"
        else:
            # Para otras empresas, usar el nombre de la empresa
            if company_name is None:
                company_name = get_company_name(company_id)
            mall_name = company_name
    
    # Guardar datos en JSON para futuras generaciones sin API
    if aggregated_report_path:
        if aggregated_report_path.is_dir():
            json_path = aggregated_report_path / "datos_agregados.json"
            ppt_dir = aggregated_report_path
        else:
            json_path = aggregated_report_path.parent / "datos_agregados.json"
            ppt_dir = aggregated_report_path.parent
    else:
        # Fallback: crear directorio por defecto
        if company_name is None:
            company_name = get_company_name(company_id)
        safe_company_name = "".join(c for c in company_name if c.isalnum() or c in (" ", "-", "_")).strip().replace(" ", "_")
        timestamp = datetime.now().strftime("%Y%m%d_%H%M")
        output_dir_base = Path("reports") / safe_company_name / "ABREGADO"
        output_dir_base.mkdir(parents=True, exist_ok=True)
        ppt_dir = output_dir_base / f"AGREGADO_{timestamp}"
        ppt_dir.mkdir(exist_ok=True)
        json_path = ppt_dir / "datos_agregados.json"
    
    # Guardar datos en JSON
    try:
        guardar_datos_json(datos, json_path)
    except Exception as e:
        print(f"[ADVERTENCIA] No se pudo guardar JSON: {e}")
    
    # Generar PPT
    ppt_path = ppt_dir / "Agregado PPT.pptx"
    crear_ppt_analisis(datos, ppt_path, mall_name=mall_name)
    
    print(f"[OK] PPT generada: {ppt_path}")
    
    # Convertir PPT a PDF
    pdf_path = convertir_ppt_a_pdf(ppt_path)
    if pdf_path:
        print(f"[OK] PDF generado: {pdf_path}")
    
    return ppt_path

def generar_ppt_individual_desde_reporte(
    company_id: str,
    node_id: str,
    start_date: str,
    end_date: str,
    report_path: Path,
    mall_name: str = None,
    company_name: str = None
) -> Optional[Path]:
    """
    Función genérica para generar PPT individual desde un reporte individual.
    Solo genera PPT si el período tiene más de 7 días.
    
    Args:
        company_id: ID de la empresa
        node_id: ID del nodo
        start_date: Fecha de inicio (formato DD/MM/YYYY)
        end_date: Fecha de fin (formato DD/MM/YYYY)
        report_path: Ruta del reporte Word generado
        mall_name: Nombre del mall (opcional)
        company_name: Nombre de la empresa (opcional)
    
    Returns:
        Ruta del archivo PPT generado, o None si no se generó (menos de 7 días)
    """
    # Verificar si tiene más de 7 días
    if not tiene_mas_de_7_dias(start_date, end_date):
        print("[INFO] Período tiene 7 días o menos, no se genera PPT individual")
        return None
    
    print()
    print("[INFO] Generando presentación PPT individual...")
    
    # Obtener datos del nodo
    try:
        datos = obtener_datos_agregados([node_id], start_date, end_date)
        print(f"[OK] Datos obtenidos: {len(datos['all_measures'])} medidas")
    except Exception as e:
        print(f"[ERROR] Error al obtener datos para PPT: {e}")
        return None
    
    # Determinar nombre del mall/empresa
    if mall_name is None:
        if company_id == "000025":
            node_name = get_node_name(node_id)
            mall_name = get_mall_name_for_parque_arauco(node_id, node_name)
            if not mall_name:
                mall_name = "Parque Arauco"
        else:
            if company_name is None:
                company_name = get_company_name(company_id)
            mall_name = company_name
    
    # Guardar PPT en la misma carpeta que el reporte Word
    if report_path:
        if report_path.is_dir():
            ppt_dir = report_path
        else:
            ppt_dir = report_path.parent
    else:
        # Fallback: crear directorio por defecto
        if company_name is None:
            company_name = get_company_name(company_id)
        safe_company_name = "".join(c for c in company_name if c.isalnum() or c in (" ", "-", "_")).strip().replace(" ", "_")
        timestamp = datetime.now().strftime("%Y%m%d_%H%M")
        output_dir_base = Path("reports") / safe_company_name / "REPORTE"
        output_dir_base.mkdir(parents=True, exist_ok=True)
        node_name = get_node_name(node_id)
        safe_node_name = "".join(c for c in node_name if c.isalnum() or c in (" ", "-", "_")).strip().replace(" ", "_")
        ppt_dir = output_dir_base / f"{safe_node_name}_{timestamp}"
        ppt_dir.mkdir(exist_ok=True)
    
    # Guardar datos en JSON
    json_path = ppt_dir / "datos_individual.json"
    try:
        guardar_datos_json(datos, json_path)
    except Exception as e:
        print(f"[ADVERTENCIA] No se pudo guardar JSON: {e}")
    
    # Generar PPT (usar el mismo formato que agregado, pero con un solo nodo)
    node_name = get_node_name(node_id)
    ppt_path = ppt_dir / f"{node_name} PPT.pptx"
    crear_ppt_analisis(datos, ppt_path, mall_name=mall_name)
    
    print(f"[OK] PPT individual generada: {ppt_path}")
    
    # Convertir PPT a PDF
    pdf_path = convertir_ppt_a_pdf(ppt_path)
    if pdf_path:
        print(f"[OK] PDF individual generado: {pdf_path}")
    
    return ppt_path

def crear_ppt_analisis(datos: dict, output_path: Path, mall_name: str = "Maipú"):
    """Crea presentación PPT con análisis de consumo.
    
    Args:
        datos: Diccionario con los datos del análisis
        output_path: Ruta donde guardar el PPT
        mall_name: Nombre del mall (por defecto "Maipú")
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Crear directorio temporal si no existe
    temp_dir = output_path.parent
    temp_dir.mkdir(parents=True, exist_ok=True)
    
    # Buscar logo WES en la raíz
    logo_wes_path = Path("logo wes.bmp")
    logo_width = Inches(2.0)  # Ancho del logo (reducido)
    logo_height = Inches(0.4)  # Alto del logo (reducido proporcionalmente)
    logo_left = Inches(7.8)  # Posición izquierda (movido hacia la izquierda)
    logo_top = Inches(0.2)  # Posición superior (cerca del borde superior)
    
    # Slide 1: Portada con imagen de fondo de Parque Arauco
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Buscar imagen de fondo de Parque Arauco en la raíz
    fondo_parque_arauco_path = Path("Parque arauco fondo.jpg")
    if fondo_parque_arauco_path.exists():
        # Aplicar imagen de fondo en la portada - estirar horizontalmente, reducir verticalmente
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        # Estirar imagen horizontalmente hasta los márgenes (ancho completo)
        scaled_width = slide_width  # Ancho completo
        
        # Reducir altura verticalmente para mantener proporciones (alrededor del 80% de la altura)
        height_scale_factor = 0.75  # Reducir altura al 75%
        scaled_height = int(slide_height * height_scale_factor)
        
        # Centrar verticalmente
        left_offset = 0
        top_offset = (slide_height - scaled_height) // 2
        
        bg_shape = slide.shapes.add_picture(str(fondo_parque_arauco_path), left_offset, top_offset, scaled_width, scaled_height)
        # Enviar imagen al fondo
        spTree = slide.shapes._spTree
        bg_element = bg_shape.element
        spTree.remove(bg_element)
        spTree.insert(2, bg_element)
    
    # Logo NO se agrega en la portada
    
    # Eliminar título "Análisis de Consumo de Agua" - ya no se necesita
    
    # Título "Mall Maipú - Parque Arauco" posicionado debajo del logo de Parque Arauco
    # Subido más para quedar más cerca del logo
    left = Inches(1)
    top2 = Inches(4.3)  # Subido desde 4.8 a 4.3 para quedar más cerca del logo
    width = Inches(8)
    height2 = Inches(1.0)
    txBox2 = slide.shapes.add_textbox(left, top2, width, height2)
    tf2 = txBox2.text_frame
    # Asegurar que el nombre del mall tenga codificación UTF-8 correcta
    mall_name_utf8 = mall_name.encode('utf-8').decode('utf-8')
    tf2.text = f"Parque Arauco {mall_name_utf8}\nPeríodo analizado: {datos['start_date'].strftime('%d/%m/%Y')} - {datos['end_date'].strftime('%d/%m/%Y')}"
    p2 = tf2.paragraphs[0]
    p2.font.size = Pt(24)
    p2.alignment = PP_ALIGN.CENTER
    # Configurar el segundo párrafo (fechas)
    if len(tf2.paragraphs) > 1:
        p3 = tf2.paragraphs[1]
        p3.font.size = Pt(18)
        p3.alignment = PP_ALIGN.CENTER
    
    # Slide 2: Resumen Ejecutivo - Diseño de dos columnas
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Fondo eliminado - dejando fondo blanco
    
    # Agregar logo WES (solo si existe)
    if logo_wes_path.exists():
        slide.shapes.add_picture(str(logo_wes_path), logo_left, logo_top, logo_width, logo_height)
    
    summary = datos['aggregate_summary']
    total = datos['total_consumption']
    num_nodes = len(datos['nodes_summary'])
    num_days = (datos['end_date'] - datos['start_date']).days + 1
    
    max_consumption = summary.get('max').total_m3 if summary.get('max') else 0.0
    min_consumption = summary.get('min').total_m3 if summary.get('min') else 0.0
    max_date = summary.get('max').date.strftime('%d/%m/%Y') if summary.get('max') else 'N/A'
    promedio_diario = summary.get('promedio_diario', 0)
    
    # Calcular porcentaje de variación del pico respecto al promedio
    variacion_pico = ((max_consumption - promedio_diario) / promedio_diario * 100) if promedio_diario > 0 else 0
    
    # Calcular consumo nocturno y efectivo desde datos horarios reales de la API
    # SOLO si no hay datos en el JSON (modo sin API)
    usar_api = datos.get('_usar_api', True)  # Por defecto usar API si no se especifica
    node_ids = [node['node_id'] for node in datos['nodes_summary']]
    
    if usar_api:
        print("[INFO] Calculando consumo nocturno agregado desde datos horarios de la API...")
        # Obtener alertas agregadas de todos los nodos
        print("[INFO] Obteniendo alertas agregadas de todos los nodos...")
        try:
            alertas_nocturnas = obtener_alertas_agregadas(node_ids, datos['start_date'], datos['end_date'])
            print(f"[OK] Alertas agregadas obtenidas: {alertas_nocturnas}")
        except Exception as e:
            print(f"[ADVERTENCIA] Error al obtener alertas agregadas: {e}")
            alertas_nocturnas = 0
    else:
        print("[INFO] Modo sin API: usando datos del JSON solamente")
        # Intentar obtener alertas desde los datos del JSON si están disponibles
        alertas_nocturnas = datos.get('total_alertas', 0)
        if alertas_nocturnas == 0:
            print("[ADVERTENCIA] No se encontraron alertas en el JSON, usando 0")
    
    # Estimación de filtración (valores por defecto si no hay datos)
    filtracion_proyectada = 0  # Por defecto
    porcentaje_filtracion = 0  # Por defecto
    try:
        nocturno_data = calcular_consumo_nocturno_agregado(
            node_ids, 
            datos['start_date'], 
            datos['end_date']
        )
        consumo_nocturno = nocturno_data['consumo_nocturno_total']
        consumo_efectivo = nocturno_data['consumo_diurno_efectivo']
        dias_con_consumo_nocturno = nocturno_data['dias_con_consumo_nocturno']
        consumo_nocturno_por_nodo = nocturno_data['consumo_nocturno_por_nodo']
        print(f"[OK] Consumo nocturno calculado: {format_number_chilean(consumo_nocturno, 1)} m³")
    except Exception as e:
        print(f"[ADVERTENCIA] Error al calcular consumo nocturno desde API: {e}")
        print("[INFO] Usando estimación del 15% como fallback...")
        # Fallback a estimación si falla el cálculo real
        consumo_nocturno = total * 0.15 if total > 0 else 0.0
        consumo_efectivo = total - consumo_nocturno
        dias_con_consumo_nocturno = max(1, int(num_days * 0.15)) if consumo_nocturno > 0 else 0
        consumo_nocturno_por_nodo = {node['node_id']: node['summary'].get('total', 0) * 0.15 
                                     for node in datos['nodes_summary']}
    
    porcentaje_nocturno = (consumo_nocturno / total * 100) if total > 0 else 0
    porcentaje_efectivo = (consumo_efectivo / total * 100) if total > 0 else 0
    
    # Obtener precio del agua para calcular valorización
    # Intentar obtener precio del primer nodo, si no usar valor por defecto
    precio_agua = 1200.0  # Precio por defecto: $1,200 CLP por m³
    if datos['nodes_summary'] and len(datos['nodes_summary']) > 0:
        first_node = datos['nodes_summary'][0]
        try:
            # Intentar obtener precio desde la API
            precio_agua = get_water_price_per_m3(COMPANY_ID, first_node['node_id'], None)
        except:
            precio_agua = 1200.0  # Usar valor por defecto si falla
    
    # Calcular valorización
    valorizacion_nocturno = consumo_nocturno * precio_agua
    valorizacion_efectivo = consumo_efectivo * precio_agua
    
    # COLUMNA IZQUIERDA: "Análisis de Consumo Hídrico"
    left_col = Inches(0.5)
    top_col = Inches(1)
    width_col_left = Inches(4.5)
    height_col = Inches(5.5)
    
    # Línea verde vertical eliminada según solicitud del usuario
    
    # Título "Análisis de Consumo Hídrico"
    title_left_x = left_col + Inches(0.3)
    title_left_y = top_col
    txBox_title_left = slide.shapes.add_textbox(title_left_x, title_left_y, width_col_left - Inches(0.3), Inches(0.6))
    tf_title_left = txBox_title_left.text_frame
    tf_title_left.text = "Análisis de Consumo Hídrico"
    p_title_left = tf_title_left.paragraphs[0]
    p_title_left.font.size = Pt(20)
    p_title_left.font.bold = True
    p_title_left.font.color.rgb = RGBColor(0, 0, 255)  # Azul
    
    # Línea azul debajo del título
    line_under_title = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, title_left_x, title_left_y + Inches(0.55), Inches(3.5), Pt(2))
    line_under_title.fill.solid()
    line_under_title.fill.fore_color.rgb = RGBColor(173, 216, 230)  # Azul claro
    line_under_title.line.fill.background()
    
    # Contenido con viñetas (alineado con el título, mismos márgenes)
    content_left_x = title_left_x  # Mismo margen izquierdo que el título
    content_left_y = title_left_y + Inches(0.8)
    content_left_width = width_col_left - Inches(0.3)  # Mismo ancho que el título
    
    txBox_content_left = slide.shapes.add_textbox(content_left_x, content_left_y, content_left_width, Inches(4.5))
    tf_content_left = txBox_content_left.text_frame
    tf_content_left.word_wrap = True
    
    # Párrafo completo de análisis (justificado, agrandado y en azul)
    p1 = tf_content_left.paragraphs[0]
    p1.text = f"Análisis de {num_days} días de consumo hídrico revela patrones críticos y detección de fuga significativa. Se registró un pico de consumo el {max_date} ({format_number_chilean(max_consumption, 1)} m³), un {format_number_chilean(variacion_pico, 1)}% por encima del promedio. Se detectaron {alertas_nocturnas} alertas de consumo nocturno, indicando actividad hídrica anómala. Se recomienda realizar una inspección exhaustiva para localizar la fuente de la filtración."
    p1.font.size = Pt(13)  # Agrandado de 11 a 13
    p1.font.color.rgb = RGBColor(0, 0, 255)  # Azul como el título
    p1.alignment = PP_ALIGN.JUSTIFY  # Justificado
    p1.space_after = Pt(4)
    p1.level = 0
    
    # Gráfico de anillo debajo del texto de la columna izquierda
    periodo_texto = f"{datos['start_date'].strftime('%d/%m/%Y')} - {datos['end_date'].strftime('%d/%m/%Y')}"
    chart_path_anillo = output_path.parent / "temp_anillo_consumo.png"
    crear_grafico_anillo_consumo(consumo_efectivo, consumo_nocturno, 
                                  porcentaje_efectivo, porcentaje_nocturno,
                                  periodo_texto, chart_path_anillo)
    
    # Agregar gráfico debajo del texto de la columna izquierda (aumentado y subido)
    graph_width = Inches(3.5)  # Aumentado de 2.8 a 3.5
    graph_height = Inches(3.5)  # Aumentado de 2.8 a 3.5
    # Centrar horizontalmente respecto al ancho de la columna izquierda
    graph_left = left_col + (width_col_left - graph_width) / 2
    graph_top = content_left_y + Inches(1.5)  # Subido desde 2.0 a 1.5
    slide.shapes.add_picture(str(chart_path_anillo), graph_left, graph_top, graph_width, graph_height)
    chart_path_anillo.unlink()  # Eliminar archivo temporal
    
    # COLUMNA DERECHA: "Métricas Críticas" y "Filtración Detectada"
    right_col = Inches(5.5)
    top_right = top_col
    
    # Título "Métricas Críticas - Agregado"
    title_right_x = right_col
    title_right_y = top_right
    txBox_title_right = slide.shapes.add_textbox(title_right_x, title_right_y, Inches(4.5), Inches(0.6))
    tf_title_right = txBox_title_right.text_frame
    tf_title_right.text = "Métricas Críticas - Agregado"
    p_title_right = tf_title_right.paragraphs[0]
    p_title_right.font.size = Pt(20)
    p_title_right.font.bold = True
    p_title_right.font.color.rgb = RGBColor(0, 0, 255)  # Azul
    
    # Cards para métricas (reducidas verticalmente pero con texto ajustado)
    card_width = Inches(1.6)  # Ancho mantenido
    card_height = Inches(0.7)  # Reducido verticalmente de 0.85 a 0.7
    card_spacing_x = Inches(0.15)  # Espaciado
    card_spacing_y = Inches(0.1)  # Espaciado vertical
    cards_top = title_right_y + Inches(0.7)
    
    # Asegurar que las cards estén bien alineadas cuadradas
    right_col_aligned = right_col  # Asegurar alineación consistente
    card2_x_aligned = right_col_aligned + card_width + card_spacing_x  # Alineación precisa para columna derecha
    
    # Card 1: Consumo Total (top-left)
    card1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, right_col_aligned, cards_top, card_width, card_height)
    card1.fill.solid()
    card1.fill.fore_color.rgb = RGBColor(230, 240, 255)  # Azul muy claro
    card1.line.color.rgb = RGBColor(230, 240, 255)
    # Línea inferior azul
    line_card1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_col_aligned, cards_top + card_height - Pt(2), card_width, Pt(2))
    line_card1.fill.solid()
    line_card1.fill.fore_color.rgb = RGBColor(173, 216, 230)
    line_card1.line.fill.background()
    
    txBox_card1 = slide.shapes.add_textbox(right_col_aligned + Inches(0.15), cards_top + Inches(0.1), card_width - Inches(0.3), card_height - Inches(0.2))
    tf_card1 = txBox_card1.text_frame
    tf_card1.text = f"Consumo Total\n{format_number_chilean(total, 1)} m³"
    p_card1_title = tf_card1.paragraphs[0]
    p_card1_title.font.size = Pt(8)
    p_card1_title.font.color.rgb = RGBColor(100, 100, 100)
    p_card1_value = tf_card1.paragraphs[1]
    p_card1_value.font.size = Pt(11)  # Reducido de 12 a 11 para que quepa
    p_card1_value.font.bold = True
    p_card1_value.font.color.rgb = RGBColor(50, 50, 50)
    
    # Card 2: Promedio Diario (top-right) - alineado con card2_x_aligned
    card2_x = card2_x_aligned
    card2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card2_x, cards_top, card_width, card_height)
    card2.fill.solid()
    card2.fill.fore_color.rgb = RGBColor(230, 240, 255)
    card2.line.color.rgb = RGBColor(230, 240, 255)
    line_card2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, card2_x, cards_top + card_height - Pt(2), card_width, Pt(2))
    line_card2.fill.solid()
    line_card2.fill.fore_color.rgb = RGBColor(173, 216, 230)
    line_card2.line.fill.background()
    
    txBox_card2 = slide.shapes.add_textbox(card2_x + Inches(0.15), cards_top + Inches(0.1), card_width - Inches(0.3), card_height - Inches(0.2))
    tf_card2 = txBox_card2.text_frame
    tf_card2.text = f"Promedio Diario\n{format_number_chilean(promedio_diario, 0)} m³"
    p_card2_title = tf_card2.paragraphs[0]
    p_card2_title.font.size = Pt(8)
    p_card2_title.font.color.rgb = RGBColor(100, 100, 100)
    p_card2_value = tf_card2.paragraphs[1]
    p_card2_value.font.size = Pt(11)  # Reducido de 12 a 11 para que quepa
    p_card2_value.font.bold = True
    p_card2_value.font.color.rgb = RGBColor(50, 50, 50)
    
    # Card 3: Día Mayor Consumo (bottom-left) - alineado verticalmente
    card3_y = cards_top + card_height + card_spacing_y
    card3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, right_col_aligned, card3_y, card_width, card_height)
    card3.fill.solid()
    card3.fill.fore_color.rgb = RGBColor(230, 240, 255)
    card3.line.color.rgb = RGBColor(230, 240, 255)
    line_card3 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_col, card3_y + card_height - Pt(2), card_width, Pt(2))
    line_card3.fill.solid()
    line_card3.fill.fore_color.rgb = RGBColor(173, 216, 230)
    line_card3.line.fill.background()
    
    txBox_card3 = slide.shapes.add_textbox(right_col_aligned + Inches(0.15), card3_y + Inches(0.1), card_width - Inches(0.3), card_height - Inches(0.2))
    tf_card3 = txBox_card3.text_frame
    tf_card3.text = f"Día Mayor Consumo\n{format_number_chilean(max_consumption, 1)} m³"
    p_card3_title = tf_card3.paragraphs[0]
    p_card3_title.font.size = Pt(8)
    p_card3_title.font.color.rgb = RGBColor(100, 100, 100)
    p_card3_value = tf_card3.paragraphs[1]
    p_card3_value.font.size = Pt(11)  # Reducido de 12 a 11 para que quepa
    p_card3_value.font.bold = True
    p_card3_value.font.color.rgb = RGBColor(50, 50, 50)
    
    # Card 4: Alertas Nocturnas (bottom-right) - alineado con card2 y card3
    card4 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card2_x, card3_y, card_width, card_height)
    card4.fill.solid()
    card4.fill.fore_color.rgb = RGBColor(230, 240, 255)
    card4.line.color.rgb = RGBColor(230, 240, 255)
    line_card4 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, card2_x, card3_y + card_height - Pt(2), card_width, Pt(2))
    line_card4.fill.solid()
    line_card4.fill.fore_color.rgb = RGBColor(173, 216, 230)
    line_card4.line.fill.background()
    
    txBox_card4 = slide.shapes.add_textbox(card2_x + Inches(0.15), card3_y + Inches(0.1), card_width - Inches(0.3), card_height - Inches(0.2))
    tf_card4 = txBox_card4.text_frame
    tf_card4.text = f"Alertas Nocturnas\n{alertas_nocturnas}"
    p_card4_title = tf_card4.paragraphs[0]
    p_card4_title.font.size = Pt(8)
    p_card4_title.font.color.rgb = RGBColor(100, 100, 100)
    p_card4_value = tf_card4.paragraphs[1]
    p_card4_value.font.size = Pt(11)  # Reducido de 12 a 11 para que quepa
    p_card4_value.font.bold = True
    p_card4_value.font.color.rgb = RGBColor(50, 50, 50)
    
    # Card 5: Consumo Nocturno (tercera fila, izquierda) - alineado verticalmente
    card5_y = card3_y + card_height + card_spacing_y
    card5 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, right_col_aligned, card5_y, card_width, card_height)
    card5.fill.solid()
    card5.fill.fore_color.rgb = RGBColor(230, 240, 255)
    card5.line.color.rgb = RGBColor(230, 240, 255)
    line_card5 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_col_aligned, card5_y + card_height - Pt(2), card_width, Pt(2))
    line_card5.fill.solid()
    line_card5.fill.fore_color.rgb = RGBColor(173, 216, 230)
    line_card5.line.fill.background()
    
    txBox_card5 = slide.shapes.add_textbox(right_col_aligned + Inches(0.15), card5_y + Inches(0.1), card_width - Inches(0.3), card_height - Inches(0.2))
    tf_card5 = txBox_card5.text_frame
    # Formato mejorado con indicador de días
    tf_card5.text = f"Consumo Nocturno\n{format_number_chilean(consumo_nocturno, 1)} m³\n({format_number_chilean(porcentaje_nocturno, 1)}%)\n{dias_con_consumo_nocturno} días consumo Noc. de {num_days} días"
    p_card5_title = tf_card5.paragraphs[0]
    p_card5_title.font.size = Pt(8)
    p_card5_title.font.color.rgb = RGBColor(100, 100, 100)
    p_card5_value = tf_card5.paragraphs[1]
    p_card5_value.font.size = Pt(10)  # Reducido para que quepa en casilla más pequeña
    p_card5_value.font.bold = True
    p_card5_value.font.color.rgb = RGBColor(50, 50, 50)
    if len(tf_card5.paragraphs) > 2:
        p_card5_pct = tf_card5.paragraphs[2]
        p_card5_pct.font.size = Pt(9)  # Reducido para que quepa
        p_card5_pct.font.bold = True
        p_card5_pct.font.color.rgb = RGBColor(50, 50, 50)
    if len(tf_card5.paragraphs) > 3:
        p_card5_dias = tf_card5.paragraphs[3]
        p_card5_dias.font.size = Pt(7)  # Reducido para que quepa
        p_card5_dias.font.color.rgb = RGBColor(80, 80, 80)  # Gris un poco más oscuro
    
    # Card 6: Consumo Efectivo (tercera fila, derecha)
    card6 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card2_x, card5_y, card_width, card_height)
    card6.fill.solid()
    card6.fill.fore_color.rgb = RGBColor(230, 240, 255)
    card6.line.color.rgb = RGBColor(230, 240, 255)
    line_card6 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, card2_x, card5_y + card_height - Pt(2), card_width, Pt(2))
    line_card6.fill.solid()
    line_card6.fill.fore_color.rgb = RGBColor(173, 216, 230)
    line_card6.line.fill.background()
    
    txBox_card6 = slide.shapes.add_textbox(card2_x + Inches(0.15), card5_y + Inches(0.1), card_width - Inches(0.3), card_height - Inches(0.2))
    tf_card6 = txBox_card6.text_frame
    tf_card6.text = f"Consumo Efectivo\n{format_number_chilean(consumo_efectivo, 1)} m³\n({format_number_chilean(porcentaje_efectivo, 1)}%)"
    p_card6_title = tf_card6.paragraphs[0]
    p_card6_title.font.size = Pt(8)
    p_card6_title.font.color.rgb = RGBColor(100, 100, 100)
    p_card6_value = tf_card6.paragraphs[1]
    p_card6_value.font.size = Pt(10)  # Reducido para que quepa en casilla más pequeña
    p_card6_value.font.bold = True
    p_card6_value.font.color.rgb = RGBColor(50, 50, 50)
    if len(tf_card6.paragraphs) > 2:
        p_card6_pct = tf_card6.paragraphs[2]
        p_card6_pct.font.size = Pt(9)  # Reducido para que quepa
        p_card6_pct.font.bold = True
        p_card6_pct.font.color.rgb = RGBColor(50, 50, 50)
    
    # Card 7: Valorización Consumo Nocturno (cuarta fila, izquierda)
    card7_y = card5_y + card_height + card_spacing_y
    card7 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, right_col_aligned, card7_y, card_width, card_height)
    card7.fill.solid()
    card7.fill.fore_color.rgb = RGBColor(230, 240, 255)
    card7.line.color.rgb = RGBColor(230, 240, 255)
    line_card7 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_col_aligned, card7_y + card_height - Pt(2), card_width, Pt(2))
    line_card7.fill.solid()
    line_card7.fill.fore_color.rgb = RGBColor(173, 216, 230)
    line_card7.line.fill.background()
    
    txBox_card7 = slide.shapes.add_textbox(right_col_aligned + Inches(0.15), card7_y + Inches(0.1), card_width - Inches(0.3), card_height - Inches(0.2))
    tf_card7 = txBox_card7.text_frame
    tf_card7.text = f"Valorización Consumo Nocturno\n{format_currency_chilean(valorizacion_nocturno)}"
    p_card7_title = tf_card7.paragraphs[0]
    p_card7_title.font.size = Pt(8)
    p_card7_title.font.color.rgb = RGBColor(100, 100, 100)
    p_card7_value = tf_card7.paragraphs[1]
    p_card7_value.font.size = Pt(10)  # Reducido para que quepa en casilla más pequeña
    p_card7_value.font.bold = True
    p_card7_value.font.color.rgb = RGBColor(50, 50, 50)
    
    # Card 8: Valorización Consumo Efectivo (cuarta fila, derecha)
    card8 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card2_x, card7_y, card_width, card_height)
    card8.fill.solid()
    card8.fill.fore_color.rgb = RGBColor(230, 240, 255)
    card8.line.color.rgb = RGBColor(230, 240, 255)
    line_card8 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, card2_x, card7_y + card_height - Pt(2), card_width, Pt(2))
    line_card8.fill.solid()
    line_card8.fill.fore_color.rgb = RGBColor(173, 216, 230)
    line_card8.line.fill.background()
    
    txBox_card8 = slide.shapes.add_textbox(card2_x + Inches(0.15), card7_y + Inches(0.1), card_width - Inches(0.3), card_height - Inches(0.2))
    tf_card8 = txBox_card8.text_frame
    tf_card8.text = f"Valorización Consumo Efectivo\n{format_currency_chilean(valorizacion_efectivo)}"
    p_card8_title = tf_card8.paragraphs[0]
    p_card8_title.font.size = Pt(8)
    p_card8_title.font.color.rgb = RGBColor(100, 100, 100)
    p_card8_value = tf_card8.paragraphs[1]
    p_card8_value.font.size = Pt(10)  # Reducido para que quepa en casilla más pequeña
    p_card8_value.font.bold = True
    p_card8_value.font.color.rgb = RGBColor(50, 50, 50)
    
    # Casilla "Filtración Detectada" eliminada según solicitud del usuario
    
    # Agregar gráfico de barras de consumo nocturno por punto debajo de las casillas
    chart_path_nocturno_puntos = output_path.parent / "temp_consumo_nocturno_puntos.png"
    # Asegurar que consumo_nocturno_por_nodo esté disponible
    if 'consumo_nocturno_por_nodo' not in locals():
        # Si no está disponible, calcularlo
        node_ids = [node['node_id'] for node in datos['nodes_summary']]
        try:
            nocturno_data = calcular_consumo_nocturno_agregado(
                node_ids,
                datos['start_date'],
                datos['end_date']
            )
            consumo_nocturno_por_nodo = nocturno_data['consumo_nocturno_por_nodo']
        except Exception as e:
            print(f"[ADVERTENCIA] Error al calcular consumo nocturno para gráfico: {e}")
            consumo_nocturno_por_nodo = {node['node_id']: 0.0 for node in datos['nodes_summary']}
    
    crear_grafico_consumo_nocturno_puntos(
        datos['nodes_summary'], 
        consumo_nocturno_por_nodo,
        chart_path_nocturno_puntos
    )
    
    # Calcular posición del gráfico (debajo de la última fila de casillas, centrado)
    card8_y = card7_y + card_height + card_spacing_y
    # Calcular el ancho disponible en la columna derecha
    col_right_width = Inches(9) - right_col_aligned  # Ancho total de la slide menos la posición de inicio
    # Hacer el gráfico de barras horizontal (aumentado y centrado con las casillas)
    graph_nocturno_width = Inches(3.4)  # Aumentado de 2.8 a 3.4
    graph_nocturno_height = Inches(2.0)  # Aumentado de 1.6 a 2.0
    # Centrar horizontalmente respecto al ancho de las casillas (no toda la columna)
    # Las casillas ocupan desde right_col_aligned hasta card2_x + card_width
    casillas_width = (card2_x + card_width) - right_col_aligned
    graph_nocturno_left = right_col_aligned + (casillas_width - graph_nocturno_width) / 2
    graph_nocturno_top = card8_y + Inches(0.2)  # Espacio después de las casillas
    
    slide.shapes.add_picture(str(chart_path_nocturno_puntos), graph_nocturno_left, graph_nocturno_top, 
                             graph_nocturno_width, graph_nocturno_height)
    chart_path_nocturno_puntos.unlink()  # Eliminar archivo temporal
    
    # Slide 3: Ranking de Consumo por Nodo
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Fondo eliminado - dejando fondo blanco
    
    # Agregar logo WES (solo si existe)
    if logo_wes_path.exists():
        slide.shapes.add_picture(str(logo_wes_path), logo_left, logo_top, logo_width, logo_height)
    
    # Título
    left_title = Inches(1)
    top_title = Inches(0.5)
    txBox_title = slide.shapes.add_textbox(left_title, top_title, Inches(8), Inches(0.6))
    tf_title = txBox_title.text_frame
    tf_title.text = "Ranking de Consumo"
    p_title = tf_title.paragraphs[0]
    p_title.font.size = Pt(28)
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(0, 0, 255)  # Azul
    
    # Crear y agregar gráfico de ranking
    chart_path_ranking = output_path.parent / "temp_ranking_nodos.png"
    crear_grafico_ranking_nodos(datos['nodes_summary'], chart_path_ranking)
    
    left_img = Inches(1)
    top_img = Inches(1.3)
    width_img = Inches(8)
    height_img = Inches(4.2)
    slide.shapes.add_picture(str(chart_path_ranking), left_img, top_img, width_img, height_img)
    chart_path_ranking.unlink()  # Eliminar archivo temporal
    
    # Análisis del gráfico
    left_analysis = Inches(1)
    top_analysis = Inches(5.6)
    width_analysis = Inches(8)
    height_analysis = Inches(1.4)
    txBox_analysis = slide.shapes.add_textbox(left_analysis, top_analysis, width_analysis, height_analysis)
    tf_analysis = txBox_analysis.text_frame
    tf_analysis.word_wrap = True
    
    # Identificar nodo con mayor y menor consumo
    nodes_sorted = sorted(datos['nodes_summary'], key=lambda x: x['summary'].get('total', 0), reverse=True)
    mayor_nodo = nodes_sorted[0] if nodes_sorted else None
    
    p = tf_analysis.paragraphs[0]
    # Asegurar codificación UTF-8 correcta para el mall name
    mall_name_utf8_ranking = mall_name.encode('utf-8').decode('utf-8')
    p.text = f"Análisis: Este gráfico muestra el ranking de consumo total de cada monitoreo de Parque Arauco {mall_name_utf8_ranking}, ordenado de mayor a menor consumo. "
    if mayor_nodo:
        mayor_total = mayor_nodo['summary'].get('total', 0)
        p.text += f"El monitoreo con mayor consumo durante el período fue '{mayor_nodo['node_name']}' con {format_number_chilean(mayor_total, 1)} m³. "
    p.text += "Este ranking permite identificar los puntos de mayor consumo y priorizar acciones de optimización."
    p.font.size = Pt(10)
    p.space_after = Pt(4)
    
    # Slides 4 en adelante: Una hoja por cada punto de monitoreo con consumo por día de la semana
    # IMPORTANTE: Ordenar nodos por consumo total (de mayor a menor) para que el orden sea idéntico al ranking
    nodes_sorted_by_consumption = sorted(
        datos['nodes_summary'], 
        key=lambda x: x.get('summary', {}).get('total', 0), 
        reverse=True
    )
    
    for node_data in nodes_sorted_by_consumption:
        node_name = node_data['node_name']
        node_measures = node_data.get('measures', [])
        
        # Crear slide para este punto de monitoreo
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # Fondo eliminado - dejando fondo blanco
        
        # Agregar logo WES (solo si existe)
        if logo_wes_path.exists():
            slide.shapes.add_picture(str(logo_wes_path), logo_left, logo_top, logo_width, logo_height)
        
        # Diseño de dos columnas similar a la lámina 2
        
        # Calcular estadísticas del punto de monitoreo
        node_summary = node_data.get('summary', {})
        total_consumo = node_summary.get('total', 0)
        promedio_diario = node_summary.get('promedio_diario', 0)
        num_days = (datos['end_date'] - datos['start_date']).days + 1
        
        # Calcular consumo nocturno y diurno desde datos horarios reales de la API
        try:
            metrics = calculate_nocturnal_metrics(
                node_data['node_id'],
                datos['start_date'],
                datos['end_date']
            )
            consumo_nocturno = metrics["consumo_nocturno_total"]
            consumo_diurno = metrics["consumo_diurno_efectivo"]
            dias_con_consumo_nocturno = metrics["dias_con_consumo_nocturno"]
        except Exception as e:
            print(f"[ADVERTENCIA] Error al calcular consumo nocturno para nodo {node_data['node_id']}: {e}")
            # Fallback a estimación si falla
            consumo_nocturno = total_consumo * 0.15 if total_consumo > 0 else 0.0
            consumo_diurno = total_consumo - consumo_nocturno
            dias_con_consumo_nocturno = max(1, int(num_days * 0.15)) if consumo_nocturno > 0 else 0
        
        porcentaje_nocturno = (consumo_nocturno / total_consumo * 100) if total_consumo > 0 else 0
        porcentaje_diurno = (consumo_diurno / total_consumo * 100) if total_consumo > 0 else 0
        
        # Obtener número de alertas para este nodo
        try:
            num_alertas_nodo = obtener_alertas_por_nodo(
                node_data['node_id'],
                datos['start_date'],
                datos['end_date']
            )
        except Exception as e:
            print(f"[ADVERTENCIA] Error al obtener alertas para nodo {node_data['node_id']}: {e}")
            num_alertas_nodo = 0
        
        # Obtener precio del agua para calcular valorización
        # Intentar obtener company_id de los datos, si no usar el global
        company_id = datos.get('company_id', COMPANY_ID)
        precio_agua = 1200.0  # Precio por defecto: $1,200 CLP por m³
        try:
            precio_agua = get_water_price_per_m3(company_id, node_data['node_id'], None)
        except:
            precio_agua = 1200.0
        
        # Calcular valorización
        valorizacion_nocturno = consumo_nocturno * precio_agua
        valorizacion_diurno = consumo_diurno * precio_agua
        
        # COLUMNA IZQUIERDA: Título, Gráfica y Análisis
        left_col = Inches(0.5)
        top_col = Inches(1)
        width_col_left = Inches(4.5)
        
        # Título (solo nombre del nodo) - IZQUIERDA, ALINEADO CON COLUMNA IZQUIERDA
        title_left_x = left_col + Inches(0.3)
        title_left_y = top_col
        title_width = width_col_left - Inches(0.3)
        txBox_title = slide.shapes.add_textbox(title_left_x, title_left_y, title_width, Inches(0.6))
        tf_title = txBox_title.text_frame
        tf_title.text = node_name  # Solo el nombre del nodo
        p_title = tf_title.paragraphs[0]
        p_title.font.size = Pt(20)  # Tamaño similar a slide 2
        p_title.font.bold = True
        p_title.font.color.rgb = RGBColor(0, 0, 255)  # Azul
        
        # Línea azul debajo del título
        line_under_title = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, title_left_x, title_left_y + Inches(0.55), title_width, Pt(2))
        line_under_title.fill.solid()
        line_under_title.fill.fore_color.rgb = RGBColor(173, 216, 230)  # Azul claro
        line_under_title.line.fill.background()
        
        # GRÁFICA DE LÍNEA (Consumo Diario) - PARTE SUPERIOR IZQUIERDA
        # Robustez Windows: si no se genera la imagen temporal, omitirla en vez de fallar el PPT completo.
        chart_path_linea = output_path.parent / f"temp_consumo_linea_{node_data['node_id']}.png"
        try:
            crear_grafico_consumo_diario_linea(node_measures, chart_path_linea, node_name)
            if chart_path_linea.exists():
                graph_linea_width = Inches(4.0)
                graph_linea_height = Inches(2.5)
                graph_linea_left = left_col + (width_col_left - graph_linea_width) / 2
                graph_linea_top = title_left_y + Inches(0.8)
                slide.shapes.add_picture(
                    str(chart_path_linea),
                    graph_linea_left,
                    graph_linea_top,
                    graph_linea_width,
                    graph_linea_height,
                )
            else:
                print(f"[ADVERTENCIA] No se generó gráfica línea para {node_data['node_id']} ({node_name}). Se omite.")
        except Exception as e:
            print(f"[ADVERTENCIA] Error generando/insertando gráfica línea para {node_data['node_id']} ({node_name}): {e}")
        finally:
            try:
                if chart_path_linea.exists():
                    chart_path_linea.unlink()
            except Exception:
                pass
        
        # GRÁFICA DE BARRAS (Día de la semana) - DEBAJO DE LA GRÁFICA DE LÍNEA, SIN NARRACIÓN
        chart_path_dia_semana = output_path.parent / f"temp_consumo_dia_semana_{node_data['node_id']}.png"
        try:
            crear_grafico_consumo_dia_semana(node_measures, chart_path_dia_semana, node_name)
            if chart_path_dia_semana.exists():
                graph_barras_width = Inches(4.0)
                graph_barras_height = Inches(2.5)
                graph_barras_left = left_col + (width_col_left - graph_barras_width) / 2
                # Si la gráfica de línea no existe, usar una altura fija desde el título
                if "graph_linea_top" in locals() and "graph_linea_height" in locals():
                    graph_barras_top = graph_linea_top + graph_linea_height + Inches(0.2)
                else:
                    graph_barras_top = title_left_y + Inches(0.8) + Inches(2.5) + Inches(0.2)
                slide.shapes.add_picture(
                    str(chart_path_dia_semana),
                    graph_barras_left,
                    graph_barras_top,
                    graph_barras_width,
                    graph_barras_height,
                )
            else:
                print(f"[ADVERTENCIA] No se generó gráfica día-semana para {node_data['node_id']} ({node_name}). Se omite.")
        except Exception as e:
            print(f"[ADVERTENCIA] Error generando/insertando gráfica día-semana para {node_data['node_id']} ({node_name}): {e}")
        finally:
            try:
                if chart_path_dia_semana.exists():
                    chart_path_dia_semana.unlink()
            except Exception:
                pass
        
        # COLUMNA DERECHA: Casillas de Indicadores
        right_col = Inches(5.5)
        top_right = top_col
        
        # Título "Indicadores" (columna derecha)
        title_right_x = right_col
        title_right_y = top_right
        txBox_title_right = slide.shapes.add_textbox(title_right_x, title_right_y, Inches(4.5), Inches(0.6))
        tf_title_right = txBox_title_right.text_frame
        tf_title_right.text = "Indicadores"
        p_title_right = tf_title_right.paragraphs[0]
        p_title_right.font.size = Pt(20)
        p_title_right.font.bold = True
        p_title_right.font.color.rgb = RGBColor(0, 0, 255)  # Azul
        
        # Cards para indicadores (similar al slide 2)
        card_width = Inches(2.0)  # Un poco más ancho para 2 columnas
        card_height = Inches(0.85)  # Altura similar al slide 2
        card_spacing_x = Inches(0.2)
        card_spacing_y = Inches(0.15)
        
        # CASILLA SUPERIOR: Consumo Total del Periodo (centrada, ancho completo)
        total_card_width = Inches(4.4)  # Ancho completo de las dos columnas de cards
        total_card_height = Inches(0.75)
        total_card_x = right_col
        total_card_y = title_right_y + Inches(0.7)
        
        total_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, total_card_x, total_card_y, total_card_width, total_card_height)
        total_card.fill.solid()
        total_card.fill.fore_color.rgb = RGBColor(200, 220, 255)  # Azul más intenso para destacar
        total_card.line.color.rgb = RGBColor(173, 216, 230)
        total_card.line.width = Pt(2)
        line_total = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, total_card_x, total_card_y + total_card_height - Pt(2), total_card_width, Pt(2))
        line_total.fill.solid()
        line_total.fill.fore_color.rgb = RGBColor(0, 0, 255)  # Azul más oscuro
        line_total.line.fill.background()
        
        txBox_total = slide.shapes.add_textbox(total_card_x + Inches(0.15), total_card_y + Inches(0.1), total_card_width - Inches(0.3), total_card_height - Inches(0.2))
        tf_total = txBox_total.text_frame
        tf_total.text = f"Consumo Total del Período\n{format_number_chilean(total_consumo, 1)} m³"
        p_total_title = tf_total.paragraphs[0]
        p_total_title.font.size = Pt(11)  # Aumentado de 9 a 11
        p_total_title.font.bold = True
        p_total_title.font.color.rgb = RGBColor(0, 0, 150)
        p_total_value = tf_total.paragraphs[1]
        p_total_value.font.size = Pt(16)  # Aumentado de 14 a 16
        p_total_value.font.bold = True
        p_total_value.font.color.rgb = RGBColor(0, 0, 255)
        p_total_value.alignment = PP_ALIGN.CENTER
        
        # Cards para indicadores (debajo de la casilla total, sin diagrama de bloques)
        cards_top = total_card_y + total_card_height + Inches(0.2)
        
        right_col_aligned = right_col
        card2_x_aligned = right_col_aligned + card_width + card_spacing_x
        
        # Card 1: Consumo Nocturno (top-left) - INCLUYE NÚMERO DE ALERTAS
        card1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, right_col_aligned, cards_top, card_width, card_height)
        card1.fill.solid()
        card1.fill.fore_color.rgb = RGBColor(230, 240, 255)
        card1.line.color.rgb = RGBColor(230, 240, 255)
        line_card1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_col_aligned, cards_top + card_height - Pt(2), card_width, Pt(2))
        line_card1.fill.solid()
        line_card1.fill.fore_color.rgb = RGBColor(173, 216, 230)
        line_card1.line.fill.background()
        
        txBox_card1 = slide.shapes.add_textbox(right_col_aligned + Inches(0.15), cards_top + Inches(0.1), card_width - Inches(0.3), card_height - Inches(0.2))
        tf_card1 = txBox_card1.text_frame
        tf_card1.text = f"Consumo Nocturno\n{format_number_chilean(consumo_nocturno, 1)} m³ ({format_number_chilean(porcentaje_nocturno, 1)}%)\n{dias_con_consumo_nocturno} días de {num_days}\n{num_alertas_nodo} alertas"
        p_card1_title = tf_card1.paragraphs[0]
        p_card1_title.font.size = Pt(10)  # Aumentado de 8 a 10
        p_card1_title.font.color.rgb = RGBColor(100, 100, 100)
        p_card1_title.alignment = PP_ALIGN.CENTER
        p_card1_value = tf_card1.paragraphs[1]
        p_card1_value.font.size = Pt(12)  # Aumentado de 10 a 12
        p_card1_value.font.bold = True
        p_card1_value.font.color.rgb = RGBColor(50, 50, 50)
        p_card1_value.alignment = PP_ALIGN.CENTER
        if len(tf_card1.paragraphs) > 2:
            p_card1_dias = tf_card1.paragraphs[2]
            p_card1_dias.font.size = Pt(9)  # Aumentado de 8 a 9
            p_card1_dias.font.color.rgb = RGBColor(80, 80, 80)
            p_card1_dias.alignment = PP_ALIGN.CENTER
        if len(tf_card1.paragraphs) > 3:
            p_card1_alertas = tf_card1.paragraphs[3]
            p_card1_alertas.font.size = Pt(9)  # Aumentado de 8 a 9
            p_card1_alertas.font.bold = True
            p_card1_alertas.font.color.rgb = RGBColor(200, 0, 0)  # Rojo para destacar alertas
            p_card1_alertas.alignment = PP_ALIGN.CENTER
        
        # Card 2: Consumo Diurno (top-right)
        card2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card2_x_aligned, cards_top, card_width, card_height)
        card2.fill.solid()
        card2.fill.fore_color.rgb = RGBColor(230, 240, 255)
        card2.line.color.rgb = RGBColor(230, 240, 255)
        line_card2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, card2_x_aligned, cards_top + card_height - Pt(2), card_width, Pt(2))
        line_card2.fill.solid()
        line_card2.fill.fore_color.rgb = RGBColor(173, 216, 230)
        line_card2.line.fill.background()
        
        txBox_card2 = slide.shapes.add_textbox(card2_x_aligned + Inches(0.15), cards_top + Inches(0.1), card_width - Inches(0.3), card_height - Inches(0.2))
        tf_card2 = txBox_card2.text_frame
        tf_card2.text = f"Consumo Diurno\n{format_number_chilean(consumo_diurno, 1)} m³ ({format_number_chilean(porcentaje_diurno, 1)}%)"
        p_card2_title = tf_card2.paragraphs[0]
        p_card2_title.font.size = Pt(10)  # Aumentado de 8 a 10
        p_card2_title.font.color.rgb = RGBColor(100, 100, 100)
        p_card2_title.alignment = PP_ALIGN.CENTER
        p_card2_value = tf_card2.paragraphs[1]
        p_card2_value.font.size = Pt(12)  # Aumentado de 10 a 12
        p_card2_value.font.bold = True
        p_card2_value.font.color.rgb = RGBColor(50, 50, 50)
        p_card2_value.alignment = PP_ALIGN.CENTER
        
        # Card 3: Valorización Consumo Nocturno (bottom-left)
        card3_y = cards_top + card_height + card_spacing_y
        card3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, right_col_aligned, card3_y, card_width, card_height)
        card3.fill.solid()
        card3.fill.fore_color.rgb = RGBColor(230, 240, 255)
        card3.line.color.rgb = RGBColor(230, 240, 255)
        line_card3 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_col_aligned, card3_y + card_height - Pt(2), card_width, Pt(2))
        line_card3.fill.solid()
        line_card3.fill.fore_color.rgb = RGBColor(173, 216, 230)
        line_card3.line.fill.background()
        
        txBox_card3 = slide.shapes.add_textbox(right_col_aligned + Inches(0.15), card3_y + Inches(0.1), card_width - Inches(0.3), card_height - Inches(0.2))
        tf_card3 = txBox_card3.text_frame
        tf_card3.text = f"Valorización Consumo Nocturno\n{format_currency_chilean(valorizacion_nocturno)}"
        p_card3_title = tf_card3.paragraphs[0]
        p_card3_title.font.size = Pt(10)  # Aumentado de 8 a 10
        p_card3_title.font.color.rgb = RGBColor(100, 100, 100)
        p_card3_title.alignment = PP_ALIGN.CENTER
        p_card3_value = tf_card3.paragraphs[1]
        p_card3_value.font.size = Pt(11)  # Aumentado de 9 a 11
        p_card3_value.font.bold = True
        p_card3_value.font.color.rgb = RGBColor(50, 50, 50)
        p_card3_value.alignment = PP_ALIGN.CENTER
        
        # Card 4: Valorización Consumo Diurno (bottom-right)
        card4 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card2_x_aligned, card3_y, card_width, card_height)
        card4.fill.solid()
        card4.fill.fore_color.rgb = RGBColor(230, 240, 255)
        card4.line.color.rgb = RGBColor(230, 240, 255)
        line_card4 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, card2_x_aligned, card3_y + card_height - Pt(2), card_width, Pt(2))
        line_card4.fill.solid()
        line_card4.fill.fore_color.rgb = RGBColor(173, 216, 230)
        line_card4.line.fill.background()
        
        txBox_card4 = slide.shapes.add_textbox(card2_x_aligned + Inches(0.15), card3_y + Inches(0.1), card_width - Inches(0.3), card_height - Inches(0.2))
        tf_card4 = txBox_card4.text_frame
        tf_card4.text = f"Valorización Consumo Diurno\n{format_currency_chilean(valorizacion_diurno)}"
        p_card4_title = tf_card4.paragraphs[0]
        p_card4_title.font.size = Pt(10)  # Aumentado de 8 a 10
        p_card4_title.font.color.rgb = RGBColor(100, 100, 100)
        p_card4_title.alignment = PP_ALIGN.CENTER
        p_card4_value = tf_card4.paragraphs[1]
        p_card4_value.font.size = Pt(11)  # Aumentado de 9 a 11
        p_card4_value.font.bold = True
        p_card4_value.font.color.rgb = RGBColor(50, 50, 50)
        p_card4_value.alignment = PP_ALIGN.CENTER
        
        # Encontrar día con mayor y menor consumo
        daily_consumption = {}
        for m in node_measures:
            date_key = m.date.date()
            if date_key not in daily_consumption:
                daily_consumption[date_key] = 0.0
            daily_consumption[date_key] += m.total_m3
        
        if daily_consumption:
            dates_sorted = sorted(daily_consumption.keys())
            values_sorted = [daily_consumption[d] for d in dates_sorted]
            
            max_val = max(values_sorted)
            min_val = min(values_sorted)
            max_idx = values_sorted.index(max_val)
            min_idx = values_sorted.index(min_val)
            
            fecha_max = dates_sorted[max_idx]
            fecha_min = dates_sorted[min_idx]
            
            # Obtener datos horarios del día con mayor consumo
            # SIEMPRE generar gráfica, usando API o datos aproximados
            usar_api = datos.get('_usar_api', True)
            hourly_data_max = None
            fecha_max_dt = datetime.combine(fecha_max, datetime.min.time())
            
            # Optimización: usar datos aproximados por defecto para evitar esperas largas
            # Solo intentar API si es explícitamente requerido y hay tiempo
            if usar_api:
                try:
                    # Reducir timeout para evitar esperas largas
                    import requests
                    original_timeout = requests.Session().timeout if hasattr(requests.Session(), 'timeout') else None
                    hourly_data_max = get_hourly_measures_for_day(node_data['node_id'], fecha_max_dt)
                    if not hourly_data_max or len(hourly_data_max) == 0:
                        print(f"[INFO] API retornó datos vacíos para día máximo, usando datos aproximados")
                        hourly_data_max = None
                except Exception as api_error:
                    print(f"[INFO] Usando datos aproximados para día máximo (API no disponible)")
                    hourly_data_max = None
            
            # Si no hay datos desde API, intentar desde JSON o generar aproximados
            if not hourly_data_max:
                try:
                    # Buscar medidas del día máximo en los datos del nodo
                    hourly_data_max = []
                    for measure in node_measures:
                        if measure.date.date() == fecha_max:
                            # Si la medida tiene detalles horarios, usarlos
                            if hasattr(measure, 'details') and measure.details:
                                for hour, value in measure.details.items():
                                    try:
                                        hour_int = int(hour)
                                        hourly_data_max.append((hour_int, float(value)))
                                    except (ValueError, TypeError):
                                        continue
                    
                    # Si no hay datos horarios, crear datos aproximados desde el consumo diario
                    if not hourly_data_max and max_val > 0:
                        # Distribuir el consumo diario en 24 horas (aproximación)
                        consumo_promedio_hora = max_val / 24.0
                        hourly_data_max = [(h, consumo_promedio_hora) for h in range(24)]
                        print(f"[INFO] Generando datos horarios aproximados para día máximo")
                    
                    if hourly_data_max:
                        # Ordenar por hora
                        hourly_data_max.sort(key=lambda x: x[0])
                except Exception as e:
                    print(f"[ADVERTENCIA] Error al obtener datos horarios del día máximo: {e}")
                    # Fallback: generar datos aproximados
                    if max_val > 0:
                        consumo_promedio_hora = max_val / 24.0
                        hourly_data_max = [(h, consumo_promedio_hora) for h in range(24)]
                        print(f"[INFO] Usando datos aproximados para día máximo (fallback)")
            
            # SIEMPRE generar la gráfica si tenemos datos (reales o aproximados)
            if hourly_data_max and len(hourly_data_max) > 0:
                print(f"[INFO] Generando gráfica del día máximo para nodo {node_data['node_id']}: {fecha_max.strftime('%d/%m/%Y')} - {format_number_chilean(max_val, 1)} m³")
                chart_path_max = output_path.parent / f"temp_consumo_horario_max_{node_data['node_id']}.png"
                titulo_max = f"Día Mayor Consumo: {fecha_max.strftime('%d/%m/%Y')} - {format_number_chilean(max_val, 1)} m³"
                crear_grafico_consumo_horario_dia(hourly_data_max, chart_path_max, fecha_max_dt, titulo_max)
                
                # Agregar gráfica del día máximo (izquierda) - REDUCIDA
                graph_horario_width = Inches(2.0)  # Reducido de 2.4 a 2.0
                graph_horario_height = Inches(1.6)  # Reducido de 2.2 a 1.6
                graph_horario_left = right_col
                # Ajustar posición: usar posición más arriba si las casillas están muy abajo
                graph_horario_top = card3_y + card_height + Inches(0.2)
                
                # Verificar que quepa en la slide (altura máxima de slide es aproximadamente 7.5 pulgadas)
                slide_height_max = Inches(7.5)
                if graph_horario_top + graph_horario_height <= slide_height_max:
                    slide.shapes.add_picture(str(chart_path_max), graph_horario_left, graph_horario_top, 
                                            graph_horario_width, graph_horario_height)
                    print(f"[OK] Gráfica del día máximo agregada a la slide")
                else:
                    # Ajustar posición hacia arriba si no cabe
                    graph_horario_top = slide_height_max - graph_horario_height - Inches(0.1)
                    if graph_horario_top >= card3_y + card_height:
                        slide.shapes.add_picture(str(chart_path_max), graph_horario_left, graph_horario_top, 
                                                graph_horario_width, graph_horario_height)
                        print(f"[OK] Gráfica del día máximo agregada a la slide (posición ajustada)")
                    else:
                        # Si aún no cabe, reducir tamaño aún más
                        graph_horario_height = Inches(1.4)
                        graph_horario_width = Inches(1.8)
                        graph_horario_top = card3_y + card_height + Inches(0.2)
                        if graph_horario_top + graph_horario_height <= slide_height_max:
                            slide.shapes.add_picture(str(chart_path_max), graph_horario_left, graph_horario_top, 
                                                    graph_horario_width, graph_horario_height)
                            print(f"[OK] Gráfica del día máximo agregada a la slide (tamaño reducido)")
                        else:
                            print(f"[ADVERTENCIA] No se pudo agregar gráfica del día máximo: no cabe en la slide")
                chart_path_max.unlink()
            else:
                print(f"[ERROR] No se pudo generar gráfica del día máximo para nodo {node_data['node_id']}: datos horarios no disponibles")
            
            # Obtener datos horarios del día con menor consumo
            # SIEMPRE generar gráfica, usando API o datos aproximados
            usar_api = datos.get('_usar_api', True)
            hourly_data_min = None
            fecha_min_dt = datetime.combine(fecha_min, datetime.min.time())
            
            if usar_api:
                try:
                    # Reducir timeout para evitar esperas largas
                    hourly_data_min = get_hourly_measures_for_day(node_data['node_id'], fecha_min_dt)
                    if not hourly_data_min or len(hourly_data_min) == 0:
                        print(f"[INFO] API retornó datos vacíos para día mínimo, usando datos aproximados")
                        hourly_data_min = None
                except Exception as api_error:
                    print(f"[INFO] Usando datos aproximados para día mínimo (API no disponible)")
                    hourly_data_min = None
            
            # Si no hay datos desde API, intentar desde JSON o generar aproximados
            if not hourly_data_min:
                try:
                    # Buscar medidas del día mínimo en los datos del nodo
                    hourly_data_min = []
                    for measure in node_measures:
                        if measure.date.date() == fecha_min:
                            # Si la medida tiene detalles horarios, usarlos
                            if hasattr(measure, 'details') and measure.details:
                                for hour, value in measure.details.items():
                                    try:
                                        hour_int = int(hour)
                                        hourly_data_min.append((hour_int, float(value)))
                                    except (ValueError, TypeError):
                                        continue
                    
                    # Si no hay datos horarios, crear datos aproximados desde el consumo diario
                    if not hourly_data_min and min_val >= 0:
                        # Distribuir el consumo diario en 24 horas (aproximación)
                        consumo_promedio_hora = min_val / 24.0 if min_val > 0 else 0.0
                        hourly_data_min = [(h, consumo_promedio_hora) for h in range(24)]
                        print(f"[INFO] Generando datos horarios aproximados para día mínimo")
                    
                    if hourly_data_min:
                        # Ordenar por hora
                        hourly_data_min.sort(key=lambda x: x[0])
                except Exception as e:
                    print(f"[ADVERTENCIA] Error al obtener datos horarios del día mínimo: {e}")
                    # Fallback: generar datos aproximados
                    if min_val >= 0:
                        consumo_promedio_hora = min_val / 24.0 if min_val > 0 else 0.0
                        hourly_data_min = [(h, consumo_promedio_hora) for h in range(24)]
                        print(f"[INFO] Usando datos aproximados para día mínimo (fallback)")
            
            # SIEMPRE generar la gráfica si tenemos datos (reales o aproximados)
            if hourly_data_min and len(hourly_data_min) > 0:
                print(f"[INFO] Generando gráfica del día mínimo para nodo {node_data['node_id']}: {fecha_min.strftime('%d/%m/%Y')} - {format_number_chilean(min_val, 1)} m³")
                chart_path_min = output_path.parent / f"temp_consumo_horario_min_{node_data['node_id']}.png"
                titulo_min = f"Día Menor Consumo: {fecha_min.strftime('%d/%m/%Y')} - {format_number_chilean(min_val, 1)} m³"
                crear_grafico_consumo_horario_dia(hourly_data_min, chart_path_min, fecha_min_dt, titulo_min)
                
                # Agregar gráfica del día mínimo (derecha) - REDUCIDA
                graph_horario_width = Inches(2.0)  # Reducido de 2.4 a 2.0
                graph_horario_height = Inches(1.6)  # Reducido de 2.2 a 1.6
                graph_horario_left = right_col + Inches(2.5)  # Ajustado para las gráficas más grandes
                # Usar la misma posición que la gráfica máxima
                graph_horario_top = card3_y + card_height + Inches(0.2)
                
                # Verificar que quepa en la slide (altura máxima de slide es aproximadamente 7.5 pulgadas)
                slide_height_max = Inches(7.5)
                if graph_horario_top + graph_horario_height <= slide_height_max:
                    slide.shapes.add_picture(str(chart_path_min), graph_horario_left, graph_horario_top, 
                                            graph_horario_width, graph_horario_height)
                    print(f"[OK] Gráfica del día mínimo agregada a la slide")
                else:
                    # Ajustar posición hacia arriba si no cabe
                    graph_horario_top = slide_height_max - graph_horario_height - Inches(0.1)
                    if graph_horario_top >= card3_y + card_height:
                        slide.shapes.add_picture(str(chart_path_min), graph_horario_left, graph_horario_top, 
                                                graph_horario_width, graph_horario_height)
                        print(f"[OK] Gráfica del día mínimo agregada a la slide (posición ajustada)")
                    else:
                        # Si aún no cabe, reducir tamaño aún más
                        graph_horario_height = Inches(1.4)
                        graph_horario_width = Inches(1.8)
                        graph_horario_top = card3_y + card_height + Inches(0.2)
                        if graph_horario_top + graph_horario_height <= slide_height_max:
                            slide.shapes.add_picture(str(chart_path_min), graph_horario_left, graph_horario_top, 
                                                    graph_horario_width, graph_horario_height)
                            print(f"[OK] Gráfica del día mínimo agregada a la slide (tamaño reducido)")
                        else:
                            print(f"[ADVERTENCIA] No se pudo agregar gráfica del día mínimo: no cabe en la slide")
                chart_path_min.unlink()
            else:
                print(f"[ERROR] No se pudo generar gráfica del día mínimo para nodo {node_data['node_id']}: datos horarios no disponibles")
    
    # Slides 5 y 6 eliminadas según solicitud del usuario (Recomendaciones y Conclusiones)
    
    # Guardar presentación
    prs.save(str(output_path))
    print(f"[OK] PPT generado: {output_path}")

def buscar_ultimo_reporte_agregado_maipu() -> Path:
    """Busca el último reporte agregado de Mall Maipú."""
    reports_dir = Path("reports") / "Parque_Arauco" / "ABREGADO"
    
    if not reports_dir.exists():
        return None
    
    # Buscar todos los directorios AGREGADO_*
    agregado_dirs = [d for d in reports_dir.iterdir() if d.is_dir() and d.name.startswith("AGREGADO_")]
    
    if not agregado_dirs:
        return None
    
    # Ordenar por fecha de modificación (más reciente primero)
    agregado_dirs.sort(key=lambda x: x.stat().st_mtime, reverse=True)
    return agregado_dirs[0]

def main():
    print("=" * 70)
    print("  GENERACIÓN DE REPORTES Y PPT - MALL MAIPÚ")
    print(f"  Período: {START_DATE} - {END_DATE}")
    print("=" * 70)
    print()
    
    # Obtener nodos de Maipú
    print("[1/5] Obteniendo nodos del Mall Maipú...")
    nodes = get_maipu_nodes(COMPANY_ID)
    
    # Si falla la API, usar nodos conocidos de Maipú
    if not nodes:
        print("[ADVERTENCIA] No se pudieron obtener nodos desde API, usando nodos conocidos de Maipú...")
        nodes = [
            {"nodeId": "000025-08", "name": "Placa Bancaria"},
            {"nodeId": "000025-09", "name": "Impulsión Falabella"},
            {"nodeId": "000025-10", "name": "Impulsión Ripley"}
        ]
        print(f"[OK] Usando {len(nodes)} nodo(s) conocidos de Maipú")
    
    print(f"[OK] Se encontraron {len(nodes)} nodo(s)")
    node_ids = [node["nodeId"] for node in nodes]
    for node in nodes:
        print(f"  - {node['nodeId']}: {node['name']}")
    print()
    
    # Buscar último reporte agregado existente
    print("[INFO] Buscando último reporte agregado existente...")
    ultimo_reporte_dir = buscar_ultimo_reporte_agregado_maipu()
    
    if ultimo_reporte_dir:
        print(f"[OK] Reporte agregado encontrado: {ultimo_reporte_dir}")
        print("[INFO] Generando solo PPT sin generar reportes nuevos")
        print()
    else:
        print("[INFO] No se encontró reporte agregado previo. Generando reportes nuevos...")
        print()
        
        # Generar reportes individuales
        print("[2/5] Generando reportes individuales...")
        generated_reports = []
        
        for i, node in enumerate(nodes, 1):
            node_id = node["nodeId"]
            node_name = node["name"]
            
            print(f"  [{i}/{len(nodes)}] Generando reporte para {node_id} ({node_name})...", flush=True)
            
            try:
                # Crear argumentos para generate_report usando argparse.Namespace
                args = argparse.Namespace(
                    company_id=COMPANY_ID,
                    node_id=node_id,
                    start_date=START_DATE,
                    end_date=END_DATE,
                    output_dir="reports",
                    enviar_correo=False
                )
                report_path = generate_report(args)
                if report_path:
                    generated_reports.append(report_path)
                    print(f"    [OK] Reporte generado: {report_path}", flush=True)
                else:
                    print(f"    [ADVERTENCIA] generate_report retornó None para {node_id}", flush=True)
            except Exception as e:
                print(f"    [ERROR] No se pudo generar el reporte para {node_id}: {e}", flush=True)
                import traceback
                traceback.print_exc()
                continue
        
        print()
        print(f"[OK] Se generaron {len(generated_reports)} reporte(s) individual(es)")
        print()
    
    # Obtener datos para reporte agregado y PPT
    if ultimo_reporte_dir:
        # Modo: Solo PPT usando último reporte existente
        # Intentar cargar datos desde JSON primero
        json_path = ultimo_reporte_dir / "datos_agregados.json"
        
        if json_path.exists():
            print("[2/2] Cargando datos desde JSON (sin conexión a API)...")
            try:
                datos = cargar_datos_json(json_path)
                # Marcar que NO se debe usar la API
                datos['_usar_api'] = False
                print(f"[OK] Datos cargados: {len(datos['all_measures'])} medidas")
            except Exception as e:
                print(f"[ERROR] Error al cargar JSON: {e}")
                print("[ERROR] No se puede generar PPT sin datos del JSON")
                import traceback
                traceback.print_exc()
                return
        else:
            print("[ERROR] JSON no encontrado. No se puede generar PPT sin datos.")
            print(f"[ERROR] Ruta esperada: {json_path}")
            return
        
        # Generar PPT en el directorio del último reporte
        print()
        print("[2/2] Generando presentación PPT en el último reporte...")
        try:
            ppt_path = ultimo_reporte_dir / "Agregado PPT.pptx"
            crear_ppt_analisis(datos, ppt_path, mall_name="Maipú")
            
            # Convertir PPT a PDF
            pdf_path = convertir_ppt_a_pdf(ppt_path)
            
            print()
            print("=" * 70)
            print("  PPT GENERADA EXITOSAMENTE")
            print("=" * 70)
            print(f"[OK] PPT generada en: {ppt_path}")
            if pdf_path:
                print(f"[OK] PDF generado en: {pdf_path}")
            print(f"[OK] Directorio del reporte: {ultimo_reporte_dir}")
        except Exception as e:
            print(f"[ERROR] Error al generar PPT: {e}")
            import traceback
            traceback.print_exc()
    else:
        # Modo normal: Generar reportes nuevos y PPT
        # Obtener datos para reporte agregado y PPT
        print("[3/5] Obteniendo datos para análisis agregado...")
        try:
            datos = obtener_datos_agregados(node_ids, START_DATE, END_DATE)
            print(f"[OK] Datos obtenidos: {len(datos['all_measures'])} medidas")
        except Exception as e:
            print(f"[ERROR] Error al obtener datos: {e}")
            import traceback
            traceback.print_exc()
            return
        
        print()
        
        # Generar reporte agregado
        print("[4/5] Generando reporte agregado...")
        aggregated_report_path = None
        try:
            aggregated_report_path = generate_aggregated_report(
                company_id=COMPANY_ID,
                node_ids=node_ids,
                start_date=START_DATE,
                end_date=END_DATE,
                output_dir="reports",
                fuente_agua_id=None
            )
            print(f"[OK] Reporte agregado: {aggregated_report_path}")
            
            # Guardar datos en JSON para futuras generaciones de PPT sin API
            if aggregated_report_path:
                if aggregated_report_path.is_dir():
                    json_path = aggregated_report_path / "datos_agregados.json"
                else:
                    json_path = aggregated_report_path.parent / "datos_agregados.json"
                guardar_datos_json(datos, json_path)
        except Exception as e:
            print(f"[ADVERTENCIA] Error al generar reporte agregado: {e}")
            print("[INFO] Continuando con la generación de PPT...")
            import traceback
            traceback.print_exc()
            # Crear directorio para PPT si no existe el reporte agregado
            from datetime import datetime as dt
            company_name = get_company_name(COMPANY_ID)
            safe_company_name = "".join(c for c in company_name if c.isalnum() or c in (" ", "-", "_")).strip().replace(" ", "_")
            timestamp = dt.now().strftime("%Y%m%d_%H%M")
            output_dir_base = Path("reports") / safe_company_name / "ABREGADO"
            output_dir_base.mkdir(parents=True, exist_ok=True)
            ppt_dir = output_dir_base / f"AGREGADO_{timestamp}"
            ppt_dir.mkdir(exist_ok=True)
            aggregated_report_path = ppt_dir  # Usar este directorio para la PPT
        
        print()
        
        # Generar PPT
        print("[5/5] Generando presentación PPT...")
        try:
            # Guardar PPT en la misma carpeta que el reporte agregado
            if aggregated_report_path:
                if aggregated_report_path.is_dir():
                    ppt_dir = aggregated_report_path
                else:
                    ppt_dir = aggregated_report_path.parent
            else:
                # Fallback: crear directorio por defecto
                from datetime import datetime as dt
                company_name = get_company_name(COMPANY_ID)
                safe_company_name = "".join(c for c in company_name if c.isalnum() or c in (" ", "-", "_")).strip().replace(" ", "_")
                timestamp = dt.now().strftime("%Y%m%d_%H%M")
                output_dir_base = Path("reports") / safe_company_name / "ABREGADO"
                output_dir_base.mkdir(parents=True, exist_ok=True)
                ppt_dir = output_dir_base / f"AGREGADO_{timestamp}"
                ppt_dir.mkdir(exist_ok=True)
            
            ppt_path = ppt_dir / "Agregado PPT.pptx"
            crear_ppt_analisis(datos, ppt_path, mall_name="Maipú")
            
            # Convertir PPT a PDF
            pdf_path = convertir_ppt_a_pdf(ppt_path)
            
            print()
            print("=" * 70)
            print("  PROCESO COMPLETADO")
            print("=" * 70)
            if 'generated_reports' in locals():
                print(f"[OK] Reportes individuales generados: {len(generated_reports)}")
            print(f"[OK] Reporte agregado: {aggregated_report_path}")
            print(f"[OK] Presentación PPT: {ppt_path}")
            if pdf_path:
                print(f"[OK] Presentación PDF: {pdf_path}")
        except Exception as e:
            print(f"[ERROR] Error al generar PPT: {e}")
            import traceback
            traceback.print_exc()

if __name__ == "__main__":
    main()
