import sys
from pptx import Presentation
from pptx.util import Emu
sys.stdout.reconfigure(encoding="utf-8", errors="replace")

path = r"G:\Mi unidad\Agente WES\wes-scripts\reports\_tmp_pa_7malls_charts\Informe WES __ Parque Arauco 7 Malls (07.07.2026).pptx"
p = Presentation(path)
sl = p.slides[4]
print("=== rId por imagen (orden T) ===")
pics = [sh for sh in sl.shapes if sh.shape_type == 13 and Emu(sh.width).inches > 0.6]
pics.sort(key=lambda s: (Emu(s.left).inches, Emu(s.top).inches))
rids = {}
for sh in pics:
    rid = sh._element.blip_rId
    pos = f"L={Emu(sh.left).inches:.2f} T={Emu(sh.top).inches:.2f}"
    sha = sh.image.sha1[:12]
    dup = " *** COMPARTE rId con otra ***" if rid in rids else ""
    print(f"  {pos} rId={rid} sha={sha}{dup}")
    rids[rid] = pos
