# -*- coding: utf-8 -*-
"""
Fichas ejecutivas Parque Arauco — una por recinto (mall), 5 variables:

  1) Equipos instalados (puntos activos WES)
  2) Consumo mensualizado (junio / julio / agosto a la fecha + proyección)
  3) Hallazgos / conclusiones  (noche con control: no se lee como fuga)
  4) Solicitudes / mensajes al recinto
  5) Noche con control — estado (fuera del análisis de fugas)

Período: 01/06/2026 – 16/08/2026.
Fuente de datos: API WES. Narrativa alineada al PPT 7 Malls
(reports/.../entrega_diego_anibal) y a la gestión operativa de controles.

Uso:
  python generar_fichas_ejecutivas_pa.py
"""

from __future__ import annotations

import json
import sys
from datetime import date, datetime
from pathlib import Path
from typing import Any, Dict, List, Tuple

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
from docx import Document
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_LINE_SPACING
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Cm, Inches, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGB
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn as ppt_qn
from pptx.util import Emu, Inches as PptInches, Pt as PptPt

from generar_reporte_word import format_number_chilean

ROOT = Path(__file__).resolve().parent
OUT_DIR = ROOT / "reports" / "Parque_Arauco" / "TMP_7MALLS" / "entrega_diego_anibal"
CHARTS = OUT_DIR / "charts_fichas"
JSON_DATOS = OUT_DIR / "datos_fichas_jun_ago.json"
JSON_NOCHES = OUT_DIR / "cortes_control_fichas.json"
LOGO = ROOT / "logo wes.bmp"
FONDO = ROOT / "Parque arauco fondo.jpg"

NODOS_FICHAS = [
    "000025-01", "000025-02", "000025-04", "000025-07", "000025-08",
    "000025-09", "000025-10", "000025-11", "000025-12", "000025-13",
    "000025-17", "000025-18", "000025-19", "000025-20", "000025-21",
    "000025-22", "000025-23", "000025-24", "000025-27", "000025-28",
    "000025-29", "000025-32", "000025-33", "000025-34", "000025-35",
    "000025-36", "000025-37", "000025-38",
]
NOCHES_CLAVE = [
    ("000025-18", "2026-07-15", "SI500 víspera control"),
    ("000025-18", "2026-07-20", "SI500 post control"),
    ("000025-18", "2026-08-10", "SI500 agosto"),
    ("000025-18", "2026-08-16", "SI500 16/08"),
    ("000025-07", "2026-08-10", "Pizza ago"),
    ("000025-07", "2026-08-16", "Pizza 16/08"),
    ("000025-19", "2026-08-10", "Sur ago corte"),
    ("000025-01", "2026-08-04", "Norte pre 5/8"),
    ("000025-01", "2026-08-08", "Norte post 5/8"),
    ("000025-01", "2026-08-16", "Norte 16/08"),
    ("000025-08", "2026-08-10", "Placa"),
    ("000025-09", "2026-08-12", "Falabella 12/08"),
    ("000025-09", "2026-08-13", "Falabella 13/08"),
    ("000025-09", "2026-08-15", "Falabella 15/08"),
    ("000025-09", "2026-08-16", "Falabella 16/08"),
    ("000025-13", "2026-08-10", "MAQ matriz"),
    ("000025-13", "2026-08-16", "MAQ matriz 16/08"),
    ("000025-17", "2026-08-10", "SI300"),
    ("000025-11", "2026-08-10", "AEB matriz"),
    ("000025-12", "2026-08-10", "AEB anillo"),
    ("000025-27", "2026-08-10", "PAK DL"),
    ("000025-04", "2026-08-10", "Baños MAE"),
]

DESDE = date(2026, 6, 1)
HASTA = date(2026, 8, 16)
AGO_DIAS = 16
AGO_MES = 31
PERIODO = "01/06/2026 – 16/08/2026"
FECHA_EMISION = "17 agosto 2026"
AGO_ETQ = f"1–{AGO_DIAS}"

NAVY = (13, 59, 102)
GOLD = (201, 162, 39)
TEAL = (31, 119, 180)
GREEN = (39, 124, 91)
RED = (153, 45, 48)
GRAY = (90, 90, 90)
LIGHT = (245, 247, 250)
WHITE = (255, 255, 255)

# Puntos activos por mall (WES + PPT 7 Malls). Inactivos/relocalizados fuera.
MALLS: List[Dict[str, Any]] = [
    {
        "code": "MAE",
        "nombre": "Estación",
        "titulo": "Mall Arauco Estación",
        "nodes": ["000025-01", "000025-04", "000025-07", "000025-19"],
        "extra_nodes": ["000025-02"],
        "pendiente": [],
        "recepcion": "20/10/2025",
        "capacitacion": "18/02/2025",
        "usuarios": [
            "Equipo medioambiente.dcl@parauco.com",
            "Sala de monitores MAE (salademonitoresmae@gmail.com)",
            "Sergio Fuenzalida — Analista Gestión Ambiental",
        ],
    },
    {
        "code": "MAM",
        "nombre": "Maipú",
        "titulo": "Mall Arauco Maipú",
        "nodes": ["000025-08", "000025-09", "000025-10", "000025-32", "000025-33"],
        "extra_nodes": [],
        "pendiente": [],
        "notas_nodo": {
            "000025-09": "activo desde 11/08/2026 (OC / cambio de equipo)",
        },
        "recepcion": "06/11/2025",
        "capacitacion": "14/11/2025",
        "usuarios": [
            "Miguel Rupayan — Encargado de Operaciones",
            "Constanza Vilches — Analista Ambiental",
            "Equipo Mantención: C. Bustamante, O. Cuevas y Supervisor Eléctrico",
        ],
    },
    {
        "code": "MAQ",
        "nombre": "Quilicura",
        "titulo": "Mall Arauco Quilicura",
        "nodes": ["000025-13", "000025-34"],
        "extra_nodes": [],
        "pendiente": [],
        "recepcion": "06/11/2025 (original) / 17/02/2026 (relocalizado)",
        "capacitacion": "14/11/2025",
        "usuarios": [
            "Mario Freitez — Jefe de Operaciones",
            "Tomás Saba — Center Manager",
            "Sebastián Araneda — Analista Ambiental",
            "Equipo Mantención: Iván Dustan, Katihuska Varas, Lucas Méndez, Carlos Leyto",
        ],
    },
    {
        "code": "BOM",
        "nombre": "Buenaventura",
        "titulo": "Buenaventura (San Ignacio)",
        "nodes": ["000025-17", "000025-18"],
        "extra_nodes": [],
        "pendiente": [],
        "recepcion": "18/11/2025",
        "capacitacion": "11/12/2025",
        "usuarios": [
            "Aliro Cortés — Jefe de Operaciones",
            "Tomás Saba — Center Manager",
            "Sebastián Araneda — Analista Ambiental",
        ],
    },
    {
        "code": "AEB",
        "nombre": "El Bosque",
        "titulo": "Arauco El Bosque",
        "nodes": ["000025-11", "000025-12"],
        "extra_nodes": [],
        "pendiente": [],
        "recepcion": "29/10/2025 (original) / 16/01/2026 (relocalizado)",
        "capacitacion": "20/11/2025",
        "usuarios": [
            "Tamara Martínez — Jefa de Operaciones",
            "Tomás Saba — Center Manager",
            "Sebastián Araneda — Analista Ambiental",
        ],
    },
    {
        "code": "CUR",
        "nombre": "Curauma",
        "titulo": "Arauco Curauma",
        "nodes": ["000025-37", "000025-38"],
        "extra_nodes": [],
        "pendiente": [],
        "recepcion": "29/04/2026 (Anillo Norte / Anillo Sur)",
        "capacitacion": "12/12/2025",
        "usuarios": [
            "Joceline Lazo — Jefe de Operaciones",
            "Constanza Vilches — Analista Ambiental",
        ],
    },
    {
        "code": "PAK",
        "nombre": "Kennedy",
        "titulo": "Parque Arauco Kennedy",
        "nodes": [
            "000025-20",
            "000025-21",
            "000025-22",
            "000025-23",
            "000025-24",
            "000025-27",
            "000025-28",
            "000025-29",
            "000025-35",
            "000025-36",
        ],
        "extra_nodes": [],
        "pendiente": [],
        "cabecera": [
            "000025-20",
            "000025-21",
            "000025-22",
            "000025-23",
            "000025-24",
            "000025-28",
            "000025-29",
        ],
        "recepcion": "12/12/2025",
        "capacitacion": "17/12/2025",
        "usuarios": [
            "Francisco Jeldres — Jefe de Operaciones",
            "Paula Azolas — Analista Gestión Ambiental",
            "Equipo Mantención (C. Naranjo, M. Jara, R. Moreno, R. Díaz, J. Gutiérrez, H. Fierro)",
        ],
    },
]


def fn(v: float, dec: int = 1) -> str:
    return format_number_chilean(float(v), dec)


def _avg(daily: Dict[str, float], d0: str, d1: str) -> float:
    vals = [float(v) for d, v in daily.items() if d0 <= d <= d1]
    return (sum(vals) / len(vals)) if vals else 0.0


def _n06(hourly: Dict[str, Any], nid: str, dia: str) -> float:
    rec = hourly.get(f"{nid}_{dia}") or {}
    return float(rec.get("n06") or 0.0)


def _dia(daily: Dict[str, float], dia: str) -> float:
    return float(daily.get(dia) or 0.0)


def refrescar_datos() -> None:
    from generar_reportes_y_ppt_mall_maipu import guardar_datos_json, obtener_datos_agregados

    print(f"[INFO] Descargando medidas WES {PERIODO}…", flush=True)
    datos = obtener_datos_agregados(
        NODOS_FICHAS,
        DESDE.strftime("%d/%m/%Y"),
        HASTA.strftime("%d/%m/%Y"),
    )
    datos["all_measures"] = []
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    guardar_datos_json(datos, JSON_DATOS)


def refrescar_noches() -> Dict[str, Any]:
    from generar_reporte_word import get_hourly_measures_for_day

    hourly: Dict[str, Any] = {}
    if JSON_NOCHES.is_file():
        try:
            hourly = json.loads(JSON_NOCHES.read_text(encoding="utf-8")).get("hourly") or {}
        except Exception:
            hourly = {}
    for nid, dia, label in NOCHES_CLAVE:
        key = f"{nid}_{dia}"
        print(f"[INFO] Noche 0–6 h {nid} {dia}…", flush=True)
        try:
            serie = get_hourly_measures_for_day(nid, datetime.fromisoformat(dia)) or []
        except Exception as exc:
            print(f"[AVISO] {key}: {exc}", flush=True)
            continue
        by_h = {str(int(h)): float(v) for h, v in serie}
        n06 = sum(float(v) for h, v in serie if int(h) < 6)
        rest = sum(float(v) for h, v in serie if int(h) >= 6)
        hourly[key] = {
            "label": label,
            "n06": round(n06, 2),
            "rest": round(rest, 2),
            "by_h": by_h,
        }
    payload = {"hourly": hourly}
    JSON_NOCHES.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")
    return hourly


def cargar_datos() -> Tuple[Dict[str, str], Dict[str, Dict[str, float]]]:
    raw = json.loads(JSON_DATOS.read_text(encoding="utf-8"))
    names: Dict[str, str] = {}
    by: Dict[str, Dict[str, float]] = {}
    for ns in raw["nodes_summary"]:
        nid = ns["node_id"]
        names[nid] = ns["node_name"]
        row: Dict[str, Any] = {
            "jun": 0.0,
            "jul": 0.0,
            "ago": 0.0,
            "total": 0.0,
            "dias": 0.0,
            "daily": {},
        }
        days = set()
        for m in ns["measures"]:
            d = m["date"][:10]
            month = d[5:7]
            v = float(m["total_m3"])
            row["daily"][d] = row["daily"].get(d, 0.0) + v
            if month == "06":
                row["jun"] += v
            elif month == "07":
                row["jul"] += v
            elif month == "08":
                row["ago"] += v
            row["total"] += v
            days.add(d)
        row["dias"] = float(len(days))
        by[nid] = row
    return names, by


def cargar_noches() -> Dict[str, Any]:
    if not JSON_NOCHES.is_file():
        return {}
    return json.loads(JSON_NOCHES.read_text(encoding="utf-8")).get("hourly") or {}


def sum_mes(by: Dict[str, Dict[str, float]], nids: List[str]) -> Dict[str, float]:
    out = {"jun": 0.0, "jul": 0.0, "ago": 0.0, "total": 0.0}
    for nid in nids:
        r = by.get(nid) or {}
        for k in out:
            out[k] += float(r.get(k) or 0)
    out["ago_d"] = out["ago"] / AGO_DIAS if AGO_DIAS else 0.0
    out["jun_d"] = out["jun"] / 30.0
    out["jul_d"] = out["jul"] / 31.0
    out["ago_proy"] = out["ago_d"] * AGO_MES
    return out


def chart_mensual(path: Path, jun: float, jul: float, ago: float, proy: float) -> None:
    fig, ax = plt.subplots(figsize=(5.4, 2.15), dpi=140)
    labels = ["Junio", "Julio", f"Agosto\n({AGO_ETQ})", "Ago. proy.\n(31 d)"]
    vals = [jun, jul, ago, proy]
    colors = ["#1F77B4", "#0D3B66", "#C9A227", "#8FA4B8"]
    bars = ax.bar(labels, vals, color=colors, width=0.62, zorder=3)
    ax.set_ylabel("m³", fontsize=8, color="#0D3B66")
    ax.tick_params(labelsize=7.5, colors="#0D3B66")
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.yaxis.grid(True, linestyle=":", alpha=0.5, zorder=0)
    ax.set_axisbelow(True)
    ymax = max(vals) * 1.22 if max(vals) > 0 else 1
    ax.set_ylim(0, ymax)
    for b, v in zip(bars, vals):
        ax.text(
            b.get_x() + b.get_width() / 2,
            b.get_height() + ymax * 0.02,
            fn(v, 0),
            ha="center",
            va="bottom",
            fontsize=7,
            color="#0D3B66",
            fontweight="bold",
        )
    fig.tight_layout(pad=0.25)
    path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(path, bbox_inches="tight", facecolor="white")
    plt.close(fig)


def _equipos_lineas(mall: Dict[str, Any], names: Dict[str, str], by: Dict[str, Dict[str, float]]) -> List[str]:
    lines: List[str] = []
    notas = mall.get("notas_nodo") or {}
    for nid in mall["nodes"]:
        nm = names.get(nid, nid)
        extra = notas.get(nid)
        if extra:
            lines.append(f"• {nm} ({nid}) — {extra}")
        else:
            lines.append(f"• {nm} ({nid})")
    for nid in mall.get("extra_nodes") or []:
        nm = names.get(nid, nid)
        lines.append(f"• {nm} ({nid}) — activo, fuera del deck 7 Malls")
    for nid in mall.get("pendiente") or []:
        nm = names.get(nid, nid)
        dias = int((by.get(nid) or {}).get("dias") or 0)
        if dias:
            lines.append(f"• {nm} ({nid}) — pendiente OC; {dias} día(s) con dato en agosto")
        else:
            lines.append(f"• {nm} ({nid}) — a la espera de OC / cambio de equipo")
    return lines


def _consumo_lineas(mall: Dict[str, Any], tot: Dict[str, float], by: Dict[str, Dict[str, float]], names: Dict[str, str]) -> List[str]:
    lines = [
        f"Junio: {fn(tot['jun'])} m³  ({fn(tot['jun_d'])} m³/día)",
        f"Julio: {fn(tot['jul'])} m³  ({fn(tot['jul_d'])} m³/día)",
        f"Agosto {AGO_ETQ}: {fn(tot['ago'])} m³  ({fn(tot['ago_d'])} m³/día)",
        f"Proyección agosto: {fn(tot['ago_proy'])} m³",
    ]
    ranked = sorted(mall["nodes"], key=lambda n: -float((by.get(n) or {}).get("total") or 0))
    if mall["code"] == "PAK":
        lines.append("Cifras de cabecera (20+21+22+23+24+28+29). DL/Bazar/Kennedy no se suman.")
        top = "000025-22"
        t = by[top]
        lines.append(
            f"Mayor cabecera: {names.get(top, top)} — "
            f"jun {fn(t['jun'])} / jul {fn(t['jul'])} / ago {fn(t['ago'])} m³"
        )
    else:
        top = ranked[0]
        t = by[top]
        lines.append(
            f"Mayor aporte: {names.get(top, top)} — "
            f"jun {fn(t['jun'])} / jul {fn(t['jul'])} / ago {fn(t['ago'])} m³"
        )
    var_jj = ((tot["jul_d"] / tot["jun_d"]) - 1) * 100 if tot["jun_d"] else 0
    var_aj = ((tot["ago_d"] / tot["jul_d"]) - 1) * 100 if tot["jul_d"] else 0
    lines.append(f"Variación m³/día: jul vs jun {fn(var_jj, 0)}%  ·  ago vs jul {fn(var_aj, 0)}%")
    if mall["code"] == "MAE":
        lines.append("Suma de los 4 puntos del deck (sin 000025-02).")
    return lines


def contenidos(
    names: Dict[str, str],
    by: Dict[str, Dict[str, float]],
    hourly: Dict[str, Any] | None = None,
) -> List[Dict[str, Any]]:
    """Arma las 5 variables por mall con cifras reales. Noche con control ≠ fuga."""
    hourly = hourly or {}
    ago_hasta = HASTA.isoformat()
    out: List[Dict[str, Any]] = []
    for mall in MALLS:
        nids_consumo = mall.get("cabecera") or mall["nodes"]
        tot = sum_mes(by, nids_consumo)
        extra = sum_mes(by, mall["nodes"] + (mall.get("extra_nodes") or []))
        equipos = _equipos_lineas(mall, names, by)
        consumo = _consumo_lineas(mall, tot, by, names)
        if mall["code"] == "MAE":
            d_sur = by["000025-19"]["daily"]
            d_nor = by["000025-01"]["daily"]
            sur_pre = _avg(d_sur, "2026-06-01", "2026-06-09")
            sur_post = _avg(d_sur, "2026-06-11", "2026-06-30")
            sur_jul = float(by["000025-19"]["jul"]) / 31.0
            sur_ago = float(by["000025-19"]["ago"]) / AGO_DIAS
            piz_jun = float(by["000025-07"]["jun"]) / 30.0
            piz_jul = float(by["000025-07"]["jul"]) / 31.0
            piz_ago = float(by["000025-07"]["ago"]) / AGO_DIAS
            piz_n10 = _n06(hourly, "000025-07", "2026-08-10")
            nor_pre = _avg(d_nor, "2026-08-01", "2026-08-04")
            nor_post = _avg(d_nor, "2026-08-05", ago_hasta)
            ban_jun = float(by["000025-04"]["jun"]) / 30.0
            ban_jul = float(by["000025-04"]["jul"]) / 31.0
            ban_ago = float(by["000025-04"]["ago"]) / AGO_DIAS
            hallazgos = [
                f"Reparación 10/06 en red Sur (validada con mantención): Estanque Sur pasó de {fn(sur_pre)} a {fn(sur_post)} m³/día y se mantiene ({fn(sur_jul)} jul / {fn(sur_ago)} ago m³/día). No es fuga residual.",
                f"Pizza Hut: control nocturno desde 01/07 funciona (noche {fn(piz_n10, 1)} m³ el 10/08). El alza de julio ({fn(piz_jun)} → {fn(piz_jul)} m³/día) es diurna; agosto baja a {fn(piz_ago)} m³/día. Esa noche no se lee como fuga.",
                f"Estanque Norte: control desde 05/08. La noche no explica el alza (ago 1–4: {fn(nor_pre)} m³/día; 5–{HASTA.day:02d}: {fn(nor_post)} m³/día). Revisar locales mall en horario hábil.",
                f"Baños Públicos: {fn(ban_jun)} → {fn(ban_jul)} → {fn(ban_ago)} m³/día. Noches ya ~0. Control instalado sin funcionamiento: no es prioridad.",
            ]
            solicitudes = [
                "Confirmar que el corte on/off de Estanque Sur queda a cargo permanente de mantención nocturna.",
                "No abrir orden de fuga por noches de Norte / Pizza Hut / Sur: controles activos o corte operativo.",
                "Revisar alza diurna de Estanque Norte en agosto (locales mall).",
                "Incorporar Abastecimiento Sur Terminal (000025-02) al tablero del recinto si corresponde a la cuenta del mall.",
            ]
            controles = [
                f"Estanque Sur (000025-19): corte on/off de mantención nocturna. 10/08 noche {fn(_n06(hourly,'000025-19','2026-08-10'), 2)} m³. La madrugada NO se lee como fuga.",
                f"Pizza Hut (000025-07): control 00:00–06:00 desde 01/07. Noche OK (10/08: {fn(piz_n10, 1)} m³). No se lee como fuga.",
                "Estanque Norte (000025-01): control desde 05/08. La madrugada NO se lee como fuga (el alza de agosto es diurna).",
                "Baños Públicos: control instalado sin uso; noches ya ~0. Tampoco se interpreta como fuga.",
            ]
        elif mall["code"] == "MAM":
            d_pla = by["000025-08"]["daily"]
            placa_ago_d = float(by["000025-08"]["ago"]) / AGO_DIAS if AGO_DIAS else 0.0
            placa_17 = _dia(d_pla, "2026-06-17")
            placa_18 = _dia(d_pla, "2026-06-18")
            placa_jul_d = float(by["000025-08"]["jul"]) / 31.0
            placa_n10 = _n06(hourly, "000025-08", "2026-08-10")
            rip_jul = float(by["000025-10"]["jul"])
            rip_ago_d = float(by["000025-10"]["ago"]) / AGO_DIAS
            pas_tot = float(by["000025-32"]["total"])
            arr_tot = float(by["000025-33"]["total"])
            daily_f = by["000025-09"].get("daily") or {}
            dias_f = sorted((d, v) for d, v in daily_f.items() if d >= "2026-08-11")
            if dias_f:
                d0 = f"{dias_f[0][0][8:10]}/{dias_f[0][0][5:7]}"
                d1 = f"{dias_f[-1][0][8:10]}/{dias_f[-1][0][5:7]}"
                serie_f = f"{d0}–{d1}: " + " / ".join(fn(v) for _, v in dias_f) + " m³"
            else:
                serie_f = "sin serie aún"
            tot_f = sum(v for _, v in dias_f)
            n12 = _n06(hourly, "000025-09", "2026-08-12")
            n13 = _n06(hourly, "000025-09", "2026-08-13")
            hallazgos = [
                f"Placa Bancaria concentra ~{fn(by['000025-08']['total']/tot['total']*100,0)}% del volumen monitoreado (jun–ago).",
                f"Auditoría 18/06: Placa subió de {fn(placa_17)} a {fn(placa_18)} m³/día. Julio {fn(placa_jul_d)} m³/día. Agosto {AGO_ETQ} vuelve a {fn(placa_ago_d)} m³/día (bajo el nivel pre-auditoría). Noche 10/08 = {fn(placa_n10, 1)} m³.",
                f"Ripley: noches tendiendo a cero (lámina 9 del deck); volumen jul {fn(rip_jul, 0)} m³, estable en agosto (~{fn(rip_ago_d)} m³/día).",
                f"Pasillo Técnico Boulevard y salida ARROW: consumo residual ({fn(pas_tot)} y {fn(arr_tot)} m³ en el período).",
                "Impulsión Falabella (000025-09) ACTIVA desde el 11/08/2026. "
                f"Jun–jul = 0 m³ (equipo fuera). Serie {serie_f} "
                f"(acum. {fn(tot_f)} m³). Noche 12/08 {fn(n12, 1)} m³ y 13/08 {fn(n13, 1)} m³ (0–6 h). "
                "Sin control nocturno; baseline 2–3 semanas.",
            ]
            solicitudes = [
                "Incorporar Impulsión Falabella al tablero diario: ya no está en espera de OC; activo desde el 11/08.",
                "No comparar Falabella con junio–julio (serie en cero). Seguir 2–3 semanas para fijar baseline diurno y nocturno.",
                "Mantener seguimiento de Placa Bancaria: el alza del 18/06 se revirtió en agosto; no reabrir como fuga nocturna.",
                "Pasillo y ARROW: dejar como referencia de red (no priorizar).",
            ]
            controles = [
                "Sin control nocturno WES declarado en este recinto.",
                "Ripley y Placa: patrón nocturno del deck ya era ~0; esa madrugada no se lee como fuga.",
                f"Falabella: las noches desde el 12/08 ({fn(n12, 0)}–{fn(n13, 0)} m³ en 0–6 h) SÍ entran al análisis; el punto recién entra en medición.",
            ]
        elif mall["code"] == "MAQ":
            mat = by["000025-13"]
            ban = by["000025-34"]
            mat_jun_d = float(mat["jun"]) / 30.0
            mat_jul_d = float(mat["jul"]) / 31.0
            mat_ago_d = float(mat["ago"]) / AGO_DIAS
            n10 = _n06(hourly, "000025-13", "2026-08-10")
            dia_10 = _dia(mat["daily"], "2026-08-10")
            hallazgos = [
                f"Matriz Principal = {fn(mat['total']/tot['total']*100,1)}% del recinto. Baños = {fn(ban['total']/tot['total']*100,1)}%.",
                f"Alza clara: {fn(mat_jun_d, 0)} m³/día en junio → {fn(mat_jul_d, 0)} m³/día en julio y {fn(mat_ago_d)} m³/día en agosto ({AGO_ETQ}).",
                f"10/08 Matriz: {fn(n10, 1)} m³ en 0–6 h (de {fn(dia_10, 0)} m³ del día). Sigue el patrón del deck (sin noches en cero).",
                f"Alimentación Baños: bajo y estable (jun {fn(ban['jun'], 0)} / jul {fn(ban['jul'], 0)} / ago {fn(ban['ago'], 0)} m³); uso hábil.",
                "Red de Incendio (000025-14) relocalizada: 0 m³. No forma parte del activo.",
            ]
            solicitudes = [
                "Implementar control on/off 00:00–08:00 en Matriz Principal (igual que estanques MAE).",
                f"Oportunidad de orden de magnitud: ~{fn(n10, 0)} m³/noche × 30 ≈ {fn(n10 * 30, 0)} m³/mes si se corta el caudal inhábil.",
                "No hay control nocturno activo: este es el hallazgo principal del recinto.",
            ]
            controles = [
                "No hay control nocturno activo en Quilicura.",
                "El consumo de madrugada de Matriz Principal SÍ entra al análisis: es la variable a gestionar.",
            ]
        elif mall["code"] == "BOM":
            d500 = by["000025-18"]["daily"]
            d300 = by["000025-17"]["daily"]
            s500_pre = _avg(d500, "2026-06-01", "2026-06-25")
            s500_26 = _avg(d500, "2026-06-26", "2026-06-30")
            s500_1_15 = _avg(d500, "2026-07-01", "2026-07-15")
            s500_16_31 = _avg(d500, "2026-07-16", "2026-07-31")
            s500_ago = float(by["000025-18"]["ago"]) / AGO_DIAS
            n15 = _n06(hourly, "000025-18", "2026-07-15")
            n20 = _n06(hourly, "000025-18", "2026-07-20")
            n10s = _n06(hourly, "000025-18", "2026-08-10")
            s300_jun = float(by["000025-17"]["jun"]) / 30.0
            s300_jul = float(by["000025-17"]["jul"]) / 31.0
            s300_ago = float(by["000025-17"]["ago"]) / AGO_DIAS
            n300 = _n06(hourly, "000025-17", "2026-08-10")
            d300_10 = _dia(d300, "2026-08-10")
            pct300 = (n300 / d300_10 * 100) if d300_10 else 0.0
            hallazgos = [
                f"San Ignacio 500 = ~{fn(by['000025-18']['total']/tot['total']*100,0)}% del recinto. Alza desde el 26/06 ({fn(s500_pre)} → {fn(s500_26)} m³/día) y 1–15/07 en {fn(s500_1_15)} m³/día.",
                f"Control nocturno 500 desde 16/07 FUNCIONA: 15/07 noche {fn(n15, 1)} m³ → 20/07 {fn(n20, 1)} m³ → 10/08 {fn(n10s, 1)} m³. Esa noche no se lee como fuga.",
                f"El volumen diurno no volvió a la base de junio: 16–31/07 {fn(s500_16_31)} m³/día y {AGO_ETQ}/08 {fn(s500_ago)} m³/día vs {fn(s500_pre)} m³/día (1–25/06). Queda alza operacional de día.",
                f"San Ignacio 300 (solo monitoreo): {fn(s300_jun)} → {fn(s300_jul)} → {fn(s300_ago)} m³/día. El 10/08, {fn(n300, 1)} m³ de {fn(d300_10, 1)} m³ fueron en 0–6 h (~{fn(pct300, 0)}% nocturno). Sin control.",
            ]
            solicitudes = [
                "500: no reabrir fuga nocturna. Pedir a operaciones la causa del caudal diurno/vespertino que quedó alto desde el 26/06.",
                "300: revisar consumo de madrugada sostenido e implementar control on/off (no está cubierto).",
                "Mensaje al JO (Aliro Cortés): el control de 500 ya aporta; el ahorro grande ahora está en el día y en el 300.",
            ]
            controles = [
                "San Ignacio 500 (000025-18): control desde 16/07. Noche 42 → 2 m³. Esa madrugada NO se lee como fuga.",
                "San Ignacio 300: sin control. La noche SÍ entra al análisis (el 10/08 fue ~51% nocturna).",
            ]
        elif mall["code"] == "AEB":
            n11 = _n06(hourly, "000025-11", "2026-08-10")
            n12a = _n06(hourly, "000025-12", "2026-08-10")
            d11_10 = _dia(by["000025-11"]["daily"], "2026-08-10")
            d12_10 = _dia(by["000025-12"]["daily"], "2026-08-10")
            mat_d = float(by["000025-11"]["total"]) / max(float(by["000025-11"]["dias"]) or 1.0, 1.0)
            hallazgos = [
                "Puntos activos: Matriz principal 1° piso (000025-11) y Anillo Plaza (000025-12). Matriz A.A. (000025-30) = 0 m³ en el período (no operativo).",
                f"Matriz 11: jun {fn(by['000025-11']['jun'])} / jul {fn(by['000025-11']['jul'])} / ago {fn(by['000025-11']['ago'])} m³ (~{fn(mat_d)} m³/día, estable).",
                f"10/08 Matriz: {fn(n11, 1)} m³ en 0–6 h (día {fn(d11_10, 1)} m³). Sin control.",
                f"Anillo Plaza: jul {fn(by['000025-12']['jul'])} m³ (alza vs jun {fn(by['000025-12']['jun'])}); 10/08 noche {fn(n12a, 1)} de {fn(d12_10, 1)} m³.",
                "La portada del deck 7 Malls aún nombra Matriz A.A.; el dato activo es 000025-11.",
            ]
            solicitudes = [
                f"Replicar control on/off 00:00–08:00 en Matriz (como estanques MAE). Oportunidad ~{fn(n11, 0)} m³/noche.",
                "Revisar llaves / equipos del Anillo Plaza para bajar la base nocturna.",
                "Actualizar el listado del recinto: activo 11+12; 30 en cero.",
            ]
            controles = [
                "No hay control nocturno activo en El Bosque.",
                "La noche de Matriz y Anillo SÍ entra al análisis: es la oportunidad de ahorro del recinto.",
            ]
        elif mall["code"] == "CUR":
            hallazgos = [
                "Red reconfigurada: Anillo Norte (000025-38) y Anillo Sur (000025-37). Matriz/Baños 15–16 = 0 m³ (no comparables).",
                f"Volumen conjunto estable: jun {fn(tot['jun'])} / jul {fn(tot['jul'])} / ago {fn(tot['ago'])} m³ (~20–22 m³/día).",
                "Ambos anillos aportan en partes similares (Sur levemente mayor).",
                "En el deck se cuadró WES vs boleta Esval N° 2368516 (18/05–16/06); la diferencia se explicó por desfase de lecturas, no por pérdida.",
            ]
            solicitudes = [
                "Enviar la(s) boleta(s) Esval de junio–agosto para repetir el cuadre WES vs cuenta.",
                "No hay control nocturno que activar como prioridad: el recinto está estable y acotado.",
                "Mantener a Joceline Lazo / medioambiente.dcl en el tablero de los dos anillos.",
            ]
            controles = [
                "Sin control nocturno WES en Curauma.",
                "No se declara fuga nocturna en este recinto para el período.",
            ]
        else:  # PAK
            n_dl = _n06(hourly, "000025-27", "2026-08-10")
            d_dl = _dia(by["000025-27"]["daily"], "2026-08-10")
            hallazgos = [
                "10 puntos activos. Baños 5/6 (25–26) relocalizados a Bazar Gourmet (35) y DL Kennedy (36).",
                "Cadena: Sandía Antigua (22) y Sandía Nueva (28) alimentan DL (27), que se reparte en Bazar (35) y DL Kennedy (36). No sumar la cadena para facturar.",
                f"Consumo de cabecera (sin doble conteo): jun {fn(tot['jun'])} / jul {fn(tot['jul'])} / ago {fn(tot['ago'])} m³.",
                f"10/08 DL: {fn(n_dl, 1)} m³ en 0–6 h (de {fn(d_dl, 0)} m³ del día). El deck ya marcaba al DL como el mayor nocturno de la cadena.",
                f"Andén Locales Gastronómicos (21) sube: jun {fn(by['000025-21']['jun'])} → jul {fn(by['000025-21']['jul'])} m³. Piletas: volumen menor.",
            ]
            solicitudes = [
                "Evaluar control on/off 00:00–08:00 en Distrito de Lujo y salas de bomba Sandía (mismo criterio MAE).",
                "No usar la suma de los 10 puntos como consumo del mall: hay doble conteo de la cadena DL.",
                "Mensaje a JO (Francisco Jeldres) / mantención: la noche del DL es la variable de mayor impacto.",
            ]
            controles = [
                "Sin control nocturno WES activo en Kennedy.",
                "El patrón 0–8 h del DL SÍ entra al análisis: es el hallazgo a gestionar.",
            ]

        out.append(
            {
                "mall": mall,
                "tot": tot,
                "extra": extra,
                "equipos": equipos,
                "consumo": consumo,
                "hallazgos": hallazgos,
                "solicitudes": solicitudes,
                "controles": controles,
            }
        )
    return out


# ---------------------------------------------------------------------------
# PPT
# ---------------------------------------------------------------------------
def _rgb(t: Tuple[int, int, int]) -> PptRGB:
    return PptRGB(*t)


def _set_run(run, text: str, size: int, bold: bool = False, color=NAVY) -> None:
    run.text = text
    run.font.size = PptPt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)
    run.font.name = "Calibri"


def _caja(slide, l, t, w, h, fill=LIGHT, line=None):
    sh = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        PptInches(l),
        PptInches(t),
        PptInches(w),
        PptInches(h),
    )
    sh.fill.solid()
    sh.fill.fore_color.rgb = _rgb(fill)
    sh.line.fill.background()
    if line:
        sh.line.fill.solid()
        sh.line.color.rgb = _rgb(line)
        sh.line.width = Emu(6350)
    try:
        sh.adjustments[0] = 0.08
    except Exception:
        pass
    return sh


def _tb(slide, l, t, w, h, lines: List[Tuple[str, int, bool, Tuple[int, int, int]]], align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(PptInches(l), PptInches(t), PptInches(w), PptInches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, (text, size, bold, color) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = PptPt(2)
        p.level = 0
        run = p.add_run()
        _set_run(run, text, size, bold, color)
    return box


def _header_bar(slide, prs, titulo: str, sub: str) -> None:
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, PptInches(0.92)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb(NAVY)
    bar.line.fill.background()
    gold = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, PptInches(0.92), prs.slide_width, PptInches(0.06)
    )
    gold.fill.solid()
    gold.fill.fore_color.rgb = _rgb(GOLD)
    gold.line.fill.background()
    if LOGO.is_file():
        slide.shapes.add_picture(str(LOGO), PptInches(12.15), PptInches(0.18), height=PptInches(0.58))
    _tb(slide, 0.28, 0.12, 11.6, 0.42, [(titulo, 20, True, WHITE)])
    _tb(slide, 0.28, 0.50, 11.6, 0.36, [(sub, 11, False, (220, 230, 240))])


def _card_title(slide, l, t, w, titulo: str, color=TEAL) -> None:
    _tb(slide, l, t, w, 0.28, [(titulo.upper(), 11, True, color)])


def build_ppt(fichas: List[Dict[str, Any]], names: Dict[str, str], by: Dict[str, Dict[str, float]]) -> Path:
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)

    # Portada
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    if FONDO.is_file():
        pic = sl.shapes.add_picture(str(FONDO), 0, 0, width=prs.slide_width, height=prs.slide_height)
        spTree = sl.shapes._spTree
        spTree.remove(pic.element)
        spTree.insert(2, pic.element)
    veil = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, PptInches(3.55), prs.slide_width, PptInches(3.95))
    veil.fill.solid()
    veil.fill.fore_color.rgb = _rgb((13, 59, 102))
    try:
        veil.fill.fore_color.brightness = 0  # type: ignore
    except Exception:
        pass
    # transparencia aproximada no crítica
    veil.line.fill.background()
    _tb(sl, 0.6, 3.75, 12, 0.5, [("WES  ·  Parque Arauco", 16, True, GOLD)])
    _tb(sl, 0.6, 4.20, 12, 0.7, [("Fichas ejecutivas por recinto", 32, True, WHITE)])
    _tb(
        sl,
        0.6,
        4.95,
        12,
        1.2,
        [
            (f"Período {PERIODO}   |   Emisión {FECHA_EMISION}", 16, False, WHITE),
            ("5 variables por mall  ·  si hay control nocturno, esa madrugada no se lee como fuga", 14, False, (220, 230, 240)),
            ("MAE · MAM · MAQ · BOM · AEB · CUR · PAK   |   Puntos activos WES", 14, False, GOLD),
        ],
    )

    # Metodología
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _header_bar(
        sl, prs,
        "Cómo leer las fichas",
        "Misma lógica del PPT 7 Malls (entrega Diego / Aníbal), actualizado a agosto",
    )
    bloques = [
        (
            "1. Equipos instalados",
            "Puntos WES activos del recinto. Se excluyen relocalizados en cero (Red de Incendio, Matriz/Baños CUR antiguos, Baños 5/6 PAK, KFC/Poniente 7/Locales comida). Impulsión Falabella (MAM) quedó activa el 11/08/2026.",
        ),
        (
            "2. Consumo mensualizado",
            f"Suma de puntos WES del recinto: junio, julio y agosto {AGO_ETQ}, más proyección lineal a 31 días. En Kennedy no se suma la cadena DL (doble conteo).",
        ),
        (
            "3. Hallazgos",
            "Tendencia jun–ago y puntos de mayor aporte. La noche de un punto CON control no se lee como fuga.",
        ),
        (
            "4. Solicitudes / mensajes",
            "Lo que queremos pasar al JO / mantención / medioambiente del mall.",
        ),
        (
            "5. Noche con control: no se lee como fuga",
            "Si el punto tiene control on/off (o corte de mantención), el caudal de madrugada no se interpreta como fuga. San Ignacio 500 desde 16/07 · Pizza Hut desde 01/07 · Estanque Norte desde 05/08 · Estanque Sur: corte de mantención nocturna.",
        ),
    ]
    y = 1.15
    for tit, txt in bloques:
        _caja(sl, 0.35, y, 12.6, 1.05)
        _tb(sl, 0.5, y + 0.08, 12.3, 0.28, [(tit, 13, True, NAVY)])
        _tb(sl, 0.5, y + 0.38, 12.3, 0.58, [(txt, 12, False, GRAY)])
        y += 1.12

    # Una ficha por mall
    for item in fichas:
        mall = item["mall"]
        tot = item["tot"]
        sl = prs.slides.add_slide(prs.slide_layouts[6])
        n_eq = len(mall["nodes"])
        extra_n = len(mall.get("extra_nodes") or [])
        pend_n = len(mall.get("pendiente") or [])
        pts = f"{n_eq} puntos activos"
        if extra_n:
            pts += f" + {extra_n} adicional"
        if pend_n:
            pts += f" + {pend_n} pendiente OC"
        _header_bar(
            sl, prs,
            f"{mall['code']}  ·  {mall['titulo']}",
            f"{pts}   |   {PERIODO}   |   Recepción {mall['recepcion']}",
        )

        chart_path = CHARTS / f"mensual_{mall['code']}.png"
        chart_mensual(chart_path, tot["jun"], tot["jul"], tot["ago"], tot["ago_proy"])

        # 5 cards
        _caja(sl, 0.25, 1.12, 4.35, 2.55)
        _card_title(sl, 0.38, 1.18, 4.1, "1. Equipos instalados")
        eq_lines = [(ln, 10 if len(item["equipos"]) > 8 else 11, False, NAVY) for ln in item["equipos"]]
        _tb(sl, 0.38, 1.48, 4.1, 2.1, eq_lines)

        _caja(sl, 4.70, 1.12, 8.35, 2.55)
        _card_title(sl, 4.85, 1.18, 4.0, "2. Consumo mensualizado (puntos WES)")
        cons_lines = [(ln, 11, False, NAVY) for ln in item["consumo"]]
        _tb(sl, 4.85, 1.48, 3.95, 2.1, cons_lines)
        sl.shapes.add_picture(str(chart_path), PptInches(8.90), PptInches(1.45), width=PptInches(3.95))

        _caja(sl, 0.25, 3.78, 6.40, 2.05)
        _card_title(sl, 0.38, 3.84, 6.1, "3. Hallazgos / conclusiones", RED)
        hall = [(f"• {h}" if not h.startswith("•") else h, 10, False, NAVY) for h in item["hallazgos"][:5]]
        _tb(sl, 0.38, 4.14, 6.12, 1.62, hall)

        _caja(sl, 6.80, 3.78, 6.25, 2.05)
        _card_title(sl, 6.93, 3.84, 6.0, "4. Solicitudes / mensajes", GREEN)
        sol = [(f"• {h}", 10, False, NAVY) for h in item["solicitudes"][:4]]
        _tb(sl, 6.93, 4.14, 5.98, 1.62, sol)

        _caja(sl, 0.25, 5.95, 12.80, 1.38, fill=(255, 249, 235), line=GOLD)
        _card_title(sl, 0.38, 6.00, 12.5, "5. Noche con control: no se lee como fuga", GOLD)
        ctrl = [(f"• {h}", 11, False, NAVY) for h in item["controles"]]
        _tb(sl, 0.38, 6.28, 12.5, 0.98, ctrl)

    # Consolidado
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _header_bar(sl, prs, "Consolidado  ·  mensaje por recinto", f"Puntos WES  |  {PERIODO}")
    rows = [["Mall", "Puntos", "Junio m³", "Julio m³", f"Ago {AGO_ETQ} m³", "m³/día ago", "Mensaje"]]
    mensajes_corto = {
        "MAE": "Norte/Pizza/Sur: noche con control, no es fuga. Norte: alza diurna ago.",
        "MAM": "Placa se revirtió en ago. Falabella activa desde 11/08.",
        "MAQ": "Matriz ~24 m³/noche — pedir control on/off.",
        "BOM": "500 noche OK desde 16/07; queda alza diurna + noche del 300.",
        "AEB": "Activo 11+12. Pedir control inhábil en Matriz.",
        "CUR": "Estable. Pedir boletas Esval jun–ago.",
        "PAK": "Cifras de cabecera (sin DL). Noche DL ~47 m³ (10/08) — control.",
    }
    for item in fichas:
        m = item["mall"]
        t = item["tot"]
        n = len(m["nodes"]) + len(m.get("extra_nodes") or [])
        rows.append(
            [
                f"{m['code']} {m['nombre']}",
                str(n),
                fn(t["jun"], 0),
                fn(t["jul"], 0),
                fn(t["ago"], 0),
                fn(t["ago_d"], 1),
                mensajes_corto[m["code"]],
            ]
        )

    table_shape = sl.shapes.add_table(len(rows), 7, PptInches(0.28), PptInches(1.20), PptInches(12.75), PptInches(5.6))
    table = table_shape.table
    widths = [1.9, 0.85, 1.35, 1.35, 1.5, 1.2, 4.6]
    for i, w in enumerate(widths):
        table.columns[i].width = PptInches(w)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = PptPt(10 if r else 11)
                p.font.bold = r == 0 or c == 0
                p.font.name = "Calibri"
                p.font.color.rgb = _rgb(WHITE if r == 0 else NAVY)
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = _rgb(NAVY if r == 0 else (LIGHT if r % 2 else WHITE))

    _tb(
        sl,
        0.28,
        6.95,
        12.7,
        0.4,
        [("Fuente: API WES. Controles informados por operación (16/07 SI500, 01/07 Pizza Hut, 05/08 Estanque Norte, corte Sur por mantención nocturna).", 10, False, GRAY)],
    )

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    path = OUT_DIR / f"Fichas_ejecutivas_Parque_Arauco_{HASTA.strftime('%Y%m%d')}.pptx"
    prs.save(str(path))
    print(f"[OK] PPT {path}")
    return path


# ---------------------------------------------------------------------------
# Word
# ---------------------------------------------------------------------------
def _shade(cell, hex_color: str) -> None:
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:fill"), hex_color)
    shd.set(qn("w:val"), "clear")
    tcPr.append(shd)


def _set_cell(cell, text: str, *, bold=False, color=NAVY, size=10, fill=None, center=False) -> None:
    if fill:
        _shade(cell, fill)
    cell.text = ""
    p = cell.paragraphs[0]
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER if center else WD_ALIGN_PARAGRAPH.LEFT
    run = p.add_run(text)
    run.bold = bold
    run.font.size = Pt(size)
    run.font.color.rgb = RGBColor(*color)
    run.font.name = "Calibri"


def _add_bullets(doc: Document, items: List[str]) -> None:
    for it in items:
        p = doc.add_paragraph(style="List Bullet")
        p.paragraph_format.space_after = Pt(4)
        p.paragraph_format.space_before = Pt(0)
        run = p.add_run(it)
        run.font.size = Pt(11)
        run.font.name = "Calibri"
        run.font.color.rgb = RGBColor(*NAVY)


def build_word(fichas: List[Dict[str, Any]]) -> Path:
    doc = Document()
    for s in doc.sections:
        s.top_margin = Cm(1.6)
        s.bottom_margin = Cm(1.6)
        s.left_margin = Cm(1.8)
        s.right_margin = Cm(1.8)

    t = doc.add_paragraph()
    r = t.add_run("WES  ·  Parque Arauco")
    r.bold = True
    r.font.size = Pt(12)
    r.font.color.rgb = RGBColor(*GOLD)
    r.font.name = "Calibri"

    h = doc.add_paragraph()
    r = h.add_run("Fichas ejecutivas por recinto")
    r.bold = True
    r.font.size = Pt(22)
    r.font.color.rgb = RGBColor(*NAVY)
    r.font.name = "Calibri"

    p = doc.add_paragraph()
    r = p.add_run(
        f"Período {PERIODO}. Emisión {FECHA_EMISION}. "
        "Cinco variables por mall: equipos instalados, consumo mensualizado de los puntos WES, "
        "hallazgos, solicitudes, y noche con control (no se lee como fuga). "
        "Alineado al PPT 7 Malls de la carpeta entrega_diego_anibal, con datos junio–agosto."
    )
    r.font.size = Pt(11)
    r.font.name = "Calibri"

    p = doc.add_paragraph()
    r = p.add_run("Noche con control: no se lee como fuga. ")
    r.bold = True
    r.font.size = Pt(11)
    r.font.color.rgb = RGBColor(*NAVY)
    r = p.add_run(
        "San Ignacio 500 desde el 16/07/2026; Pizza Hut desde el 01/07/2026; "
        "Estanque Norte desde el 05/08/2026; Estanque Sur con corte on/off a cargo de "
        "personal de mantención nocturno. Esas madrugadas no se interpretan como fuga."
    )
    r.font.size = Pt(11)

    for item in fichas:
        mall = item["mall"]
        tot = item["tot"]
        doc.add_page_break()
        h = doc.add_paragraph()
        r = h.add_run(f"{mall['code']}  —  {mall['titulo']}")
        r.bold = True
        r.font.size = Pt(18)
        r.font.color.rgb = RGBColor(*NAVY)

        meta = doc.add_paragraph()
        r = meta.add_run(
            f"Recepción {mall['recepcion']}  ·  Capacitación {mall['capacitacion']}  ·  "
            f"Usuarios: {'; '.join(mall['usuarios'])}"
        )
        r.font.size = Pt(10)
        r.font.color.rgb = RGBColor(*GRAY)
        r.italic = True

        for titulo, lineas in (
            ("1. Equipos instalados", item["equipos"]),
            ("2. Consumo mensualizado del recinto (puntos WES)", item["consumo"]),
        ):
            hh = doc.add_paragraph()
            rr = hh.add_run(titulo)
            rr.bold = True
            rr.font.size = Pt(13)
            rr.font.color.rgb = RGBColor(*TEAL)
            _add_bullets(doc, lineas)

        # tabla mensual
        table = doc.add_table(rows=2, cols=5)
        table.style = "Table Grid"
        table.alignment = WD_TABLE_ALIGNMENT.CENTER
        hdr = ["Junio", "Julio", f"Agosto {AGO_ETQ}", "m³/día ago", "Proyección agosto"]
        vals = [
            f"{fn(tot['jun'])} m³",
            f"{fn(tot['jul'])} m³",
            f"{fn(tot['ago'])} m³",
            f"{fn(tot['ago_d'])} m³/d",
            f"{fn(tot['ago_proy'])} m³",
        ]
        for i, x in enumerate(hdr):
            _set_cell(table.rows[0].cells[i], x, bold=True, color=WHITE, fill="0D3B66", center=True, size=10)
        for i, x in enumerate(vals):
            _set_cell(table.rows[1].cells[i], x, bold=True, fill="F5F7FA", center=True, size=11)

        chart = CHARTS / f"mensual_{mall['code']}.png"
        if chart.is_file():
            doc.add_paragraph().add_run().add_picture(str(chart), width=Inches(5.8))

        for titulo, lineas, col in (
            ("3. Hallazgos / conclusiones", item["hallazgos"], TEAL),
            ("4. Solicitudes / mensajes a pasar", item["solicitudes"], GREEN),
            ("5. Noche con control: no se lee como fuga", item["controles"], GOLD),
        ):
            hh = doc.add_paragraph()
            rr = hh.add_run(titulo)
            rr.bold = True
            rr.font.size = Pt(13)
            rr.font.color.rgb = RGBColor(*col)
            _add_bullets(doc, lineas)

    # consolidado
    doc.add_page_break()
    h = doc.add_paragraph()
    r = h.add_run("Consolidado")
    r.bold = True
    r.font.size = Pt(18)
    r.font.color.rgb = RGBColor(*NAVY)
    table = doc.add_table(rows=1 + len(fichas), cols=6)
    table.style = "Table Grid"
    for i, x in enumerate(["Mall", "Junio m³", "Julio m³", f"Ago {AGO_ETQ} m³", "m³/día ago", "Mensaje clave"]):
        _set_cell(table.rows[0].cells[i], x, bold=True, color=WHITE, fill="0D3B66", center=True, size=9)
    msgs = {
        "MAE": "Norte/Pizza/Sur: noche con control, no es fuga. Norte: alza diurna agosto.",
        "MAM": "Placa se revirtió. Falabella activa desde el 11/08.",
        "MAQ": "Pedir control nocturno en Matriz (~24 m³/noche).",
        "BOM": "500 noche OK; queda alza diurna y noche del 300.",
        "AEB": "Activos 11 y 12. Pedir control inhábil en Matriz.",
        "CUR": "Estable. Pedir boletas Esval junio–agosto.",
        "PAK": "Cifras de cabecera (sin DL). Control nocturno en DL.",
    }
    for r_i, item in enumerate(fichas, 1):
        m = item["mall"]
        t = item["tot"]
        vals = [
            f"{m['code']} {m['nombre']}",
            fn(t["jun"], 0),
            fn(t["jul"], 0),
            fn(t["ago"], 0),
            fn(t["ago_d"], 1),
            msgs[m["code"]],
        ]
        fill = "F5F7FA" if r_i % 2 else "FFFFFF"
        for c, x in enumerate(vals):
            _set_cell(table.rows[r_i].cells[c], x, fill=fill, size=9, bold=c == 0)

    p = doc.add_paragraph()
    r = p.add_run(
        f"Fuente: API de medidas WES, {PERIODO}. "
        f"La proyección de agosto es lineal ({AGO_DIAS} días → 31). "
        "Si el punto tiene control nocturno, esa madrugada no se interpreta como fuga."
    )
    r.font.size = Pt(9)
    r.italic = True
    r.font.color.rgb = RGBColor(*GRAY)

    path = OUT_DIR / f"Fichas_ejecutivas_Parque_Arauco_{HASTA.strftime('%Y%m%d')}.docx"
    doc.save(str(path))
    print(f"[OK] Word {path}")
    return path


def convertir_pdf(docx_path: Path) -> Path | None:
    import shutil
    import subprocess

    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        print("[AVISO] LibreOffice no está disponible; se omite PDF.")
        return None
    cmd = [
        soffice,
        "--headless",
        "--norestore",
        "--convert-to",
        "pdf",
        "--outdir",
        str(docx_path.parent),
        str(docx_path),
    ]
    print(f"[INFO] PDF via {soffice}", flush=True)
    subprocess.run(cmd, check=True, timeout=180)
    pdf = docx_path.with_suffix(".pdf")
    if pdf.is_file():
        print(f"[OK] PDF {pdf}")
        return pdf
    return None


def main() -> int:
    skip = "--skip-refresh" in sys.argv
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    CHARTS.mkdir(parents=True, exist_ok=True)
    if not skip or not JSON_DATOS.is_file():
        refrescar_datos()
    if not skip or not JSON_NOCHES.is_file():
        refrescar_noches()
    names, by = cargar_datos()
    hourly = cargar_noches()
    fichas = contenidos(names, by, hourly)
    ppt = build_ppt(fichas, names, by)
    docx = build_word(fichas)
    pdf = convertir_pdf(docx)
    print("\n=== SALIDA ===")
    print(ppt)
    print(docx)
    if pdf:
        print(pdf)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
