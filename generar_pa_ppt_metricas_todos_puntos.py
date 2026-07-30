"""
Genera PPT de métricas Parque Arauco (000025) con TODOS los puntos, sin aplicar exclusiones,
para el período indicado.

Incluye una diapositiva final con narrativa (hallazgos) + notas sobre instalación/habilitados.
"""

from __future__ import annotations

import argparse
import os
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List

import requests

ENTITY_BASE = "http://104.248.53.141:7001/wes/api/acl-entities/v1"

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from generar_pa_agregado_todos_puntos import _kpi_fechas, _parrafo_hallazgos_q3
from generar_reporte_word import format_number_chilean
from generar_reportes_y_ppt_mall_maipu import obtener_datos_agregados, crear_ppt_analisis


NODOS_NO_OPERATIVOS = {
    "000025-03": "No operativo desde 14-10-2025",
}


def _cargar_nodos_pa() -> Dict[str, Any]:
    r = requests.get(f"{ENTITY_BASE}/companies/000025", timeout=60)
    r.raise_for_status()
    return r.json()


def _agregar_slide_notas(prs: Presentation, payload: dict, datos: dict, periodo: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Título
    box_t = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(12.2), Inches(0.8))
    tf_t = box_t.text_frame
    tf_t.text = "Notas y hallazgos (PA - todos los puntos)"
    p = tf_t.paragraphs[0]
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 255)

    # Cuerpo
    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(12.0), Inches(5.6))
    tf = box.text_frame
    tf.word_wrap = True

    nodes = sorted(payload.get("nodes") or [], key=lambda x: str(x.get("nodeId", "")))
    total = float((datos.get("aggregate_summary") or {}).get("total") or 0)

    lineas: List[str] = []
    lineas.append(f"Período: {periodo}")
    lineas.append(f"Consumo total agregado: {format_number_chilean(total, 1)} m³")
    lineas.append("")
    lineas.append("Instalado (referencia API KPI):")
    for n in nodes[:8]:
        nid = str(n.get("nodeId", "")).strip()
        nombre = str(n.get("name", "")).strip()
        cre, _exp = _kpi_fechas((n.get("configuration") or {}))
        lineas.append(f"- {nid}: {nombre} (KPI creationDate: {cre})")
    if len(nodes) > 8:
        lineas.append(f"- … y {len(nodes) - 8} punto(s) más (ver Word agregado para el listado completo)")
    lineas.append("")
    lineas.append("Habilitados:")
    lineas.append("La API usada en scripts no expone usuarios habilitados; confirmar en backoffice WES.")
    lineas.append("")
    lineas.append("Hallazgos (resumen):")
    # Tomar primeras líneas de la narrativa larga
    hall = _parrafo_hallazgos_q3(datos, periodo).splitlines()
    for ln in hall[3:18]:
        if ln.strip():
            lineas.append(ln.strip())

    if NODOS_NO_OPERATIVOS:
        lineas.append("")
        lineas.append("Nodos no operativos (excluidos del PPT):")
        for nid, nota in sorted(NODOS_NO_OPERATIVOS.items()):
            lineas.append(f"- {nid}: {nota}")

    tf.text = lineas[0]
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    for ln in lineas[1:]:
        par = tf.add_paragraph()
        par.text = ln
        par.font.size = Pt(12)
        par.font.color.rgb = RGBColor(0, 0, 0)


def main() -> int:
    parser = argparse.ArgumentParser(description="PPT métricas PA todos los puntos")
    parser.add_argument("--desde", default="01/03/2026", help="DD/MM/YYYY")
    parser.add_argument("--hasta", default="27/03/2026", help="DD/MM/YYYY")
    args = parser.parse_args()

    payload = _cargar_nodos_pa()
    nodes = payload.get("nodes") or []
    node_ids = sorted(
        {
            str(n.get("nodeId", "")).strip()
            for n in nodes
            if n.get("nodeId") and str(n.get("nodeId", "")).strip() not in NODOS_NO_OPERATIVOS
        }
    )
    if not node_ids:
        print("[ERROR] Sin nodos en empresa 000025")
        return 1

    periodo = f"{args.desde} a {args.hasta}"
    print(f"[INFO] Obteniendo datos agregados PA ({len(node_ids)} nodos) período {periodo} ...")
    datos = obtener_datos_agregados(node_ids, args.desde, args.hasta)

    ts = datetime.now().strftime("%Y%m%d_%H%M")
    out_dir = Path("reports") / "Parque_Arauco" / "PA_Todos_Los_Puntos" / "ABREGADO" / f"AGREGADO_PPT_{ts}"
    out_dir.mkdir(parents=True, exist_ok=True)

    base_ppt = out_dir / "Agregado PPT.pptx"
    crear_ppt_analisis(datos, base_ppt, mall_name="PA - Todos los puntos")

    prs = Presentation(str(base_ppt))
    _agregar_slide_notas(prs, payload, datos, periodo)

    final_ppt = out_dir / f"Agregado PPT PA Todos los puntos {args.desde.replace('/','')}-{args.hasta.replace('/','')}.pptx"
    prs.save(str(final_ppt))
    print(f"[OK] PPT: {final_ppt}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

