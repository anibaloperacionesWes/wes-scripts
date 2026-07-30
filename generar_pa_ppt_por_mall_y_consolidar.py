"""
Genera PPT de métricas de Parque Arauco (000025) por mall (uno a uno) y un consolidado final.

Requerimiento:
- Dejar fuera nodos no operativos / excluidos:
  - 000025-03 (no operativo desde 14-10-2025)
  - 000025-05
  - 000025-06

Salida:
- reports/Parque_Arauco/<Mall>/ABREGADO/AGREGADO_PPT_<ts>/Agregado PPT <Mall> <desde>-<hasta>.pptx
- reports/Parque_Arauco/CONSOLIDADO/Agregado PPT Parque Arauco CONSOLIDADO <desde>-<hasta>.pptx
"""

from __future__ import annotations

import argparse
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Tuple

import requests
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from generar_reporte_word import format_number_chilean, get_mall_name_for_parque_arauco
from generar_reportes_y_ppt_mall_maipu import obtener_datos_agregados, crear_ppt_analisis


ENTITY_BASE = "http://104.248.53.141:7001/wes/api/acl-entities/v1"

NODOS_EXCLUIDOS: Dict[str, str] = {
    "000025-03": "No operativo desde 14-10-2025",
    "000025-05": "Excluido (no solicitar graficar)",
    "000025-06": "Excluido (no solicitar graficar)",
}


def _cargar_empresa_pa() -> dict:
    r = requests.get(f"{ENTITY_BASE}/companies/000025", timeout=60)
    r.raise_for_status()
    return r.json()


def _slug(s: str) -> str:
    s2 = "".join(c for c in s if c.isalnum() or c in (" ", "-", "_")).strip()
    return s2.replace(" ", "_") or "SIN_NOMBRE"


def _titulo_slide(tf, text: str) -> None:
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 255)


def _agregar_slide_notas(prs: Presentation, mall: str, desde: str, hasta: str, excluidos: Dict[str, str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    t = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(12.2), Inches(0.8)).text_frame
    _titulo_slide(t, f"Notas operativas ({mall})")

    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(12.0), Inches(5.6))
    tf = box.text_frame
    tf.word_wrap = True

    lineas: List[str] = []
    lineas.append(f"Período: {desde} a {hasta}")
    lineas.append("")
    if excluidos:
        lineas.append("Puntos excluidos / no operativos (no graficados):")
        for nid, nota in sorted(excluidos.items()):
            lineas.append(f"- {nid}: {nota}")
    else:
        lineas.append("Sin exclusiones específicas para este mall.")

    tf.text = lineas[0]
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    for ln in lineas[1:]:
        par = tf.add_paragraph()
        par.text = ln
        par.font.size = Pt(12)
        par.font.color.rgb = RGBColor(0, 0, 0)


def _resumen_mall(datos: dict) -> Dict[str, Any]:
    agg = datos.get("aggregate_summary") or {}
    total = float(agg.get("total") or 0)
    prom = float(agg.get("promedio_diario") or 0)
    dias = int(agg.get("dias") or 0)

    ranked: List[Tuple[float, str, str]] = []
    for n in (datos.get("nodes_summary") or []):
        t = float((n.get("summary") or {}).get("total") or 0)
        ranked.append((t, str(n.get("node_id") or ""), str(n.get("node_name") or "")))
    ranked.sort(key=lambda x: -x[0])
    top = ranked[0] if ranked else (0.0, "—", "—")

    return {
        "total": total,
        "promedio_diario": prom,
        "dias": dias,
        "top_total": float(top[0]),
        "top_id": top[1],
        "top_name": top[2],
        "num_puntos": len(ranked),
    }


def _crear_consolidado(
    out_path: Path,
    desde: str,
    hasta: str,
    resumenes: Dict[str, Dict[str, Any]],
    rutas_ppt: Dict[str, Path],
) -> Path:
    prs = Presentation()
    # limpiar slides por defecto (tiene 1)
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    t = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(12.2), Inches(0.9)).text_frame
    _titulo_slide(t, "Parque Arauco — Consolidado de métricas (por mall)")
    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(12.0), Inches(5.3))
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = f"Período: {desde} a {hasta}"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True

    for mall in sorted(resumenes.keys()):
        r = resumenes[mall]
        p = tf.add_paragraph()
        p.text = (
            f"{mall}: total {format_number_chilean(r['total'], 1)} m³ | "
            f"prom. {format_number_chilean(r['promedio_diario'], 1)} m³/día | "
            f"top {r['top_name']} ({r['top_id']}) {format_number_chilean(r['top_total'], 1)} m³ | "
            f"puntos {r['num_puntos']}"
        )
        p.font.size = Pt(12)

    # Slide por mall con ruta del PPT
    for mall in sorted(resumenes.keys()):
        r = resumenes[mall]
        ppt_path = rutas_ppt.get(mall)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        t = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(12.2), Inches(0.8)).text_frame
        _titulo_slide(t, f"{mall} — Resumen")
        box = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(12.0), Inches(5.6))
        tf = box.text_frame
        tf.word_wrap = True
        tf.text = f"Período: {desde} a {hasta}"
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.bold = True

        lines = [
            f"Consumo total: {format_number_chilean(r['total'], 1)} m³",
            f"Promedio diario agregado: {format_number_chilean(r['promedio_diario'], 1)} m³/día",
            f"Días con medición (agregado): {r['dias']}",
            f"Top punto por consumo: {r['top_name']} ({r['top_id']}) — {format_number_chilean(r['top_total'], 1)} m³",
            f"N° puntos incluidos: {r['num_puntos']}",
            "",
            f"PPT detallado del mall: {str(ppt_path) if ppt_path else '—'}",
        ]
        for ln in lines:
            p = tf.add_paragraph()
            p.text = ln
            p.font.size = Pt(12)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return out_path


def main() -> int:
    parser = argparse.ArgumentParser(description="PPT por mall (PA) + consolidado")
    parser.add_argument("--desde", default="01/03/2026", help="DD/MM/YYYY")
    parser.add_argument("--hasta", default="27/03/2026", help="DD/MM/YYYY")
    args = parser.parse_args()

    empresa = _cargar_empresa_pa()
    nodes = empresa.get("nodes") or []

    mall_to_nodes: Dict[str, List[Tuple[str, str]]] = {}
    for n in nodes:
        nid = str(n.get("nodeId") or "").strip()
        name = str(n.get("name") or "").strip()
        if not nid:
            continue
        mall = get_mall_name_for_parque_arauco(nid, name) or "Parque Arauco"
        mall_to_nodes.setdefault(mall, []).append((nid, name))

    ts = datetime.now().strftime("%Y%m%d_%H%M")
    rutas_ppt: Dict[str, Path] = {}
    resumenes: Dict[str, Dict[str, Any]] = {}

    for mall, pairs in sorted(mall_to_nodes.items()):
        node_ids = [nid for nid, _ in pairs if nid not in NODOS_EXCLUIDOS]
        if not node_ids:
            print(f"[INFO] {mall}: sin nodos tras exclusiones, se omite.")
            continue

        print(f"[INFO] Generando PPT mall {mall}: {len(node_ids)} punto(s)...")
        datos = obtener_datos_agregados(node_ids, args.desde, args.hasta)
        resumenes[mall] = _resumen_mall(datos)

        out_dir = (
            Path("reports")
            / "Parque_Arauco"
            / _slug(mall)
            / "ABREGADO"
            / f"AGREGADO_PPT_{ts}"
        )
        out_dir.mkdir(parents=True, exist_ok=True)

        base_ppt = out_dir / "Agregado PPT.pptx"
        crear_ppt_analisis(datos, base_ppt, mall_name=mall)

        prs = Presentation(str(base_ppt))
        excl_mall = {nid: nota for nid, nota in NODOS_EXCLUIDOS.items() if nid in {p[0] for p in pairs}}
        _agregar_slide_notas(prs, mall, args.desde, args.hasta, excl_mall)

        final_ppt = out_dir / f"Agregado PPT { _slug(mall) } {args.desde.replace('/','')}-{args.hasta.replace('/','')}.pptx"
        prs.save(str(final_ppt))
        rutas_ppt[mall] = final_ppt
        print(f"[OK] PPT {mall}: {final_ppt}")

    consolidado_dir = Path("reports") / "Parque_Arauco" / "CONSOLIDADO"
    out_cons = consolidado_dir / f"Agregado PPT Parque Arauco CONSOLIDADO {args.desde.replace('/','')}-{args.hasta.replace('/','')}.pptx"
    _crear_consolidado(out_cons, args.desde, args.hasta, resumenes, rutas_ppt)
    print(f"[OK] CONSOLIDADO: {out_cons}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

